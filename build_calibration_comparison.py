#!/usr/bin/env python3
"""
Build a calibration comparison presentation (NTE200 vs 730E ECM parameters).
Data source: ECM_对比_NTE200_730E_中文.xlsx + ECM Excel exports.
"""
import io
import os
import sys
import textwrap
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import gridspec
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

# ─── Brand palette ───────────────────────────────────────────────────────────
C_DARK  = '#293136'
C_DARK2 = '#1e2529'
C_ACC   = '#3EF0AF'
C_RED   = '#E05252'
C_ORG   = '#E08C2E'
C_YLW   = '#D4C94A'
C_GRY   = '#707880'
C_WHT   = '#FFFFFF'
C_BLU   = '#4EC9F0'

def hex2rgb(h): r,g,b=int(h[1:3],16),int(h[3:5],16),int(h[5:7],16); return r,g,b
def rgb(h):     return RGBColor(*hex2rgb(h))

W = Inches(13.33)
H = Inches(7.5)
SLIDE_W = 13.33
SLIDE_H = 7.5

# ─── Calibration parameters data ─────────────────────────────────────────────
# Translated from ECM_对比_NTE200_730E_中文.xlsx
PARAMS = [
    {
        'param':  'Fuel Code\n(карта впрыска)',
        'code':   'Fuel Code',
        'val_730': 'FC0BU52',
        'val_nte': 'FC0WH95',
        'unit':   '—',
        'risk':   'WARN',
        'effect': 'Разные давление CR, тайминг, дымоограничитель и карта\nнаддува — принципиально разный режим сгорания',
    },
    {
        'param':  'Флаг ограничения об/мин\nпри перегреве масла',
        'code':   'C_EPD_OT_RPM_Drt_En',
        'val_730': '✓ ВКЛ (1)',
        'val_nte': '✗ ОТКЛ (0)',
        'unit':   'флаг',
        'risk':   'CRITICAL',
        'effect': 'NTE200: при перегреве масла об/мин НЕ снижаются.\n730E автоматически уходит в дерейт.',
    },
    {
        'param':  'Об/мин дерейта\n(перегрев масла)',
        'code':   'C_EPD_OT_Spd_Derate',
        'val_730': '2200',
        'val_nte': '1675',
        'unit':   'об/мин',
        'risk':   'NOTE',
        'effect': 'На NTE200 флаг отключён — значение не работает.',
    },
    {
        'param':  'Время нарастания до\nмакс. момента (масло)',
        'code':   'C_EPD_CP_TB_Time_To_Max_Trq_Brt',
        'val_730': '0 с (мгновенно)',
        'val_nte': '51 с',
        'unit':   'с',
        'risk':   'HIGH',
        'effect': 'NTE200 работает на полной мощности 51 с\nпри критически низком давлении масла.',
    },
    {
        'param':  'Время нарастания до\nмакс. момента (ОЖ)',
        'code':   'C_EPD_CT_TB_Time_To_Max_Trq_Drt',
        'val_730': '0 с (мгновенно)',
        'val_nte': '43 с',
        'unit':   'с',
        'risk':   'HIGH',
        'effect': 'NTE200 не снижает нагрузку 43 с при перегреве ОЖ.',
    },
    {
        'param':  'Макс. рабочие об/мин',
        'code':   'C_ABS_Max_Acctr_Speed',
        'val_730': '1990',
        'val_nte': '2100',
        'unit':   'об/мин',
        'risk':   'HIGH',
        'effect': '+110 об/мин = +25–35 °C EGT, выше ударная нагрузка\nна клапаны (~5,5% больше циклов)',
    },
    {
        'param':  'Максимальный дроуп\nрегулятора',
        'code':   'C_ABS_DroopMax',
        'val_730': '0% (изохрон)',
        'val_nte': '20%',
        'unit':   '%',
        'risk':   'WARN',
        'effect': '730E держит строго постоянные об/мин.\nNTE200 допускает ±20% колебания.',
    },
    {
        'param':  'Порог тормозного\nотключения (макс.)',
        'code':   'C_HSST_Net_Torque_High_Thd',
        'val_730': '5509',
        'val_nte': '3650',
        'unit':   'Н·м',
        'risk':   'HIGH',
        'effect': 'NTE200 срабатывает на 1859 Н·м раньше — агрессивные\nпиковые нагрузки перед отключением.',
    },
    {
        'param':  'Порог тормозного\nотключения (мин.)',
        'code':   'C_HSST_Net_Torque_Low_Thd',
        'val_730': '4722',
        'val_nte': '1850',
        'unit':   'Н·м',
        'risk':   'CRITICAL',
        'effect': 'Гистерезис в 2,5 раза ниже — нестабильное управление\nнагрузкой при тепловом останове.',
    },
    {
        'param':  'Лимит коррекции TIB\n(топливо)',
        'code':   'C_TIB_Fuel_Adjust_Upper_Limit',
        'val_730': '100%',
        'val_nte': '200%',
        'unit':   '%',
        'risk':   'CRITICAL',
        'effect': 'NTE200 удваивает дозу впрыска —\nриск переобогащения смеси, перегрев EGT.',
    },
    {
        'param':  'Порог включения\nвентилятора наддува',
        'code':   'C_FCC_Charge_Temp_Y',
        'val_730': '50 °C',
        'val_nte': '46 °C',
        'unit':   '°C',
        'risk':   'NOTE',
        'effect': 'NTE200 включает охлаждение раньше на 4 °C —\nкомпенсирующая мера, не защита.',
    },
    {
        'param':  'Цифровой выход: порог\nтемпературы наддува',
        'code':   'C_DO_01_A_High_Threshold',
        'val_730': '54 °C',
        'val_nte': '✗ ОТКЛ',
        'unit':   '°C',
        'risk':   'CRITICAL',
        'effect': 'Только 730E сигнализирует внешнему контроллеру\nо перегреве воздуха наддува.',
    },
    {
        'param':  'Цифровой выход: порог\nтемпературы ОЖ',
        'code':   'C_DO_07_A_High_Threshold',
        'val_730': '85 °C',
        'val_nte': '✗ ОТКЛ',
        'unit':   '°C',
        'risk':   'CRITICAL',
        'effect': 'Только 730E передаёт сигнал перегрева ОЖ\nна внешний контроллер машины.',
    },
    {
        'param':  'Цифровой выход: порог\nтемпературы топлива',
        'code':   'C_DO_08_A_High_Threshold',
        'val_730': '68 °C',
        'val_nte': '✗ ОТКЛ',
        'unit':   '°C',
        'risk':   'CRITICAL',
        'effect': 'Только 730E подаёт сигнал перегрева топлива.',
    },
    {
        'param':  'Макс. об/мин\n(позиция 2/3)',
        'code':   'T_ABS_MaxRef_2/3',
        'val_730': '1990',
        'val_nte': '2225,8',
        'unit':   'об/мин',
        'risk':   'HIGH',
        'effect': 'NTE200 может разогнаться до 2225 об/мин\nв отдельных режимах работы.',
    },
]

RISK_COLOR = {
    'CRITICAL': C_RED,
    'HIGH':     C_ORG,
    'WARN':     C_YLW,
    'NOTE':     C_GRY,
}
RISK_LABEL = {
    'CRITICAL': 'КРИТИЧНО',
    'HIGH':     'ВЫСОКИЙ',
    'WARN':     'ПРЕДУПРЕЖДЕНИЕ',
    'NOTE':     'ПРИМЕЧАНИЕ',
}

# ─── Helpers ─────────────────────────────────────────────────────────────────
def add_slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def clear_placeholders(slide):
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = rgb(line); shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=11, bold=False, color=C_WHT,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return tb

def add_title_bar(slide, title, subtitle=''):
    add_rect(slide, 0, 0, 33.87, 1.9, fill=C_DARK2)
    add_rect(slide, 0, 0, 0.5, 1.9, fill=C_ACC)
    add_text(slide, title, 0.8, 0.15, 30, 1.0, size=20, bold=True, color=C_WHT)
    if subtitle:
        add_text(slide, subtitle, 0.8, 1.2, 30, 0.75, size=10, color=C_GRY)

def add_footer(slide, text='АО Развитие  •  QSK50 MCRS  •  Полюс Магадан  •  2026'):
    add_rect(slide, 0, 18.5, 33.87, 0.7, fill=C_DARK2)
    add_text(slide, text, 0.5, 18.55, 32, 0.6, size=8, color=C_GRY)

def slide_background(slide, color=C_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)

def img_to_pptx(fig, slide, x, y, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    slide.shapes.add_picture(buf, Cm(x), Cm(y), Cm(w), Cm(h))
    plt.close(fig)

# ─── Slide builders ──────────────────────────────────────────────────────────

def build_title_slide(prs):
    slide = add_slide(prs)
    clear_placeholders(slide)
    slide_background(slide, C_DARK)

    # Big accent bar left
    add_rect(slide, 0, 0, 1.2, 19.05, fill=C_ACC)

    # Title block
    add_text(slide, 'Анализ калибровок ECM', 2.0, 3.0, 28, 2.0,
             size=36, bold=True, color=C_WHT)
    add_text(slide, 'QSK50 MCRS — NTE200 vs 730E-DC', 2.0, 5.2, 28, 1.2,
             size=22, bold=False, color=C_ACC)
    add_text(slide,
             'Сравнение параметров защиты двигателя по данным\n'
             'INSITE ECM Snapshots (19–20 мая 2026 г.)',
             2.0, 6.8, 28, 1.8, size=14, color=C_GRY)

    # Divider line
    add_rect(slide, 2.0, 8.9, 20, 0.06, fill=C_ACC)

    # Source
    add_text(slide,
             'Источник: INSITE ECM Image Archive «Данные сняты под нагрузкой»\n'
             'Полюс Магадан — АО Развитие, технический анализ',
             2.0, 9.2, 28, 1.4, size=10, color=C_GRY, italic=True)

    add_footer(slide)


def _risk_badge_fig(risk):
    """Small horizontal badge image for risk level."""
    color = RISK_COLOR[risk]
    label = RISK_LABEL[risk]
    fig, ax = plt.subplots(figsize=(2.2, 0.35))
    fig.patch.set_facecolor(color)
    ax.set_facecolor(color)
    ax.text(0.5, 0.5, label, color='black', fontsize=8, fontweight='bold',
            ha='center', va='center', transform=ax.transAxes)
    ax.axis('off')
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=120, bbox_inches='tight', facecolor=color)
    buf.seek(0)
    plt.close(fig)
    return buf


def build_comparison_table_slide(prs, params_subset, slide_title, slide_sub=''):
    """Build a table slide for a subset of parameters."""
    slide = add_slide(prs)
    clear_placeholders(slide)
    slide_background(slide, C_DARK)
    add_title_bar(slide, slide_title, slide_sub)
    add_footer(slide)

    fig = plt.figure(figsize=(16.5, 8.8), facecolor=C_DARK)
    ax = fig.add_subplot(111)
    ax.set_facecolor(C_DARK)
    ax.axis('off')

    # Column layout
    col_x  = [0.01, 0.21, 0.36, 0.50, 0.62, 0.65]  # normalized x starts
    col_w  = [0.19, 0.14, 0.13, 0.12, 0.03, 0.34]   # normalized widths
    headers = ['Параметр', '730E-DC', 'NTE200', 'Ед.изм.', 'Риск', 'Влияние на двигатель']
    hdr_colors = [C_DARK2, C_DARK2, C_DARK2, C_DARK2, C_DARK2, C_DARK2]

    n = len(params_subset)
    row_h = 0.85 / (n + 1)   # fractional height per row
    y0    = 0.97              # start y (top)

    def draw_row(ax, y, cells, bg, text_colors=None, is_header=False):
        if text_colors is None:
            text_colors = [C_WHT] * len(cells)
        for i, (text, col_xs, cw) in enumerate(zip(cells, col_x, col_w)):
            rect = mpatches.FancyBboxPatch(
                (col_xs, y - row_h + 0.005),
                cw - 0.005, row_h - 0.01,
                boxstyle='round,pad=0.005',
                linewidth=0,
                facecolor=bg[i] if isinstance(bg, list) else bg,
                transform=ax.transAxes, clip_on=False, zorder=2
            )
            ax.add_patch(rect)
            fs = 7.5 if not is_header else 8.5
            fw = 'bold' if is_header else 'normal'
            ax.text(col_xs + 0.005, y - row_h / 2, str(text),
                    color=text_colors[i], fontsize=fs, fontweight=fw,
                    va='center', ha='left', transform=ax.transAxes,
                    clip_on=False, zorder=3,
                    wrap=True,
                    multialignment='left')

    # Header
    draw_row(ax, y0, headers, [C_DARK2]*6,
             text_colors=[C_ACC]*6, is_header=True)
    y = y0 - row_h

    for idx, p in enumerate(params_subset):
        risk_c = RISK_COLOR[p['risk']]
        risk_l = RISK_LABEL[p['risk']]
        row_bg_even = '#222930'
        row_bg_odd  = '#1e2529'
        bg = row_bg_even if idx % 2 == 0 else row_bg_odd

        # Value colors: red if NTE worse, green if better
        val_nte_color = C_WHT
        if p['risk'] in ('CRITICAL', 'HIGH'):
            val_nte_color = '#FF9999'

        draw_row(ax, y,
                 [p['param'], p['val_730'], p['val_nte'], p['unit'],
                  '', p['effect']],
                 bg=[bg]*6,
                 text_colors=[C_WHT, '#90EE90', val_nte_color,
                               C_GRY, C_WHT, '#C0C8D0'])

        # Risk color stripe
        stripe = mpatches.FancyBboxPatch(
            (col_x[4], y - row_h + 0.005),
            col_w[4] - 0.005, row_h - 0.01,
            boxstyle='round,pad=0.005',
            linewidth=0, facecolor=risk_c,
            transform=ax.transAxes, clip_on=False, zorder=2
        )
        ax.add_patch(stripe)

        y -= row_h

    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    fig.tight_layout(pad=0.1)

    img_to_pptx(fig, slide, 0.3, 2.0, 33.2, 16.5)


def build_bar_comparison_slide(prs):
    """Bar chart comparing numeric parameters."""
    slide = add_slide(prs)
    clear_placeholders(slide)
    slide_background(slide, C_DARK)
    add_title_bar(slide, 'Числовые параметры: 730E-DC vs NTE200',
                  'Сравнение ключевых установок ECM QSK50 MCRS')
    add_footer(slide)

    # Numeric params only
    numeric = [
        ('Макс. об/мин', 1990, 2100, 'об/мин'),
        ('Об/мин дерейта (масло)', 2200, 1675, 'об/мин'),
        ('Время нарастания\n(давл. масло, с)', 0, 51, 'с'),
        ('Время нарастания\n(темп. ОЖ, с)', 0, 43, 'с'),
        ('TIB лимит впрыска', 100, 200, '%'),
        ('HSST момент HIGH', 5509, 3650, 'Н·м'),
        ('HSST момент LOW', 4722, 1850, 'Н·м'),
        ('Дроуп регулятора', 0, 20, '%'),
    ]

    fig = plt.figure(figsize=(15, 7.5), facecolor=C_DARK)
    fig.patch.set_facecolor(C_DARK)
    axes_bg = C_DARK2

    n = len(numeric)
    cols = 4
    rows_n = (n + cols - 1) // cols
    gs = gridspec.GridSpec(rows_n, cols, figure=fig,
                           hspace=0.6, wspace=0.4,
                           left=0.04, right=0.98, top=0.92, bottom=0.08)

    for i, (name, v730, vnte, unit) in enumerate(numeric):
        ax = fig.add_subplot(gs[i // cols, i % cols])
        ax.set_facecolor(axes_bg)
        ax.tick_params(colors=C_GRY, labelsize=7)
        for sp in ax.spines.values(): sp.set_edgecolor(C_GRY); sp.set_linewidth(0.5)

        vals  = [v730, vnte]
        colors = ['#4EC9F0', C_RED if vnte > v730 * 1.05 else ('#90EE90' if vnte < v730 * 0.95 else C_GRY)]
        bars = ax.bar(['730E', 'NTE200'], vals, color=colors, width=0.55, zorder=3)

        for bar, val in zip(bars, vals):
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(vals)*0.04,
                    str(val), ha='center', va='bottom', color=C_WHT, fontsize=8, fontweight='bold')

        ax.set_title(name, color=C_WHT, fontsize=7.5, pad=4)
        ax.set_ylabel(unit, color=C_GRY, fontsize=7)
        ax.set_ylim(0, max(vals) * 1.25)
        ax.yaxis.label.set_color(C_GRY)
        ax.tick_params(axis='y', colors=C_GRY)
        ax.tick_params(axis='x', colors=C_WHT)
        ax.grid(axis='y', color=C_GRY, alpha=0.2, zorder=0)

    img_to_pptx(fig, slide, 0.3, 2.0, 33.2, 16.0)


def build_disabled_features_slide(prs):
    """Visual showing disabled protections on NTE200."""
    slide = add_slide(prs)
    clear_placeholders(slide)
    slide_background(slide, C_DARK)
    add_title_bar(slide, 'Отключённые функции защиты на NTE200',
                  'Параметры, активные на 730E и деактивированные на NTE200')
    add_footer(slide)

    disabled = [
        ('C_EPD_OT_RPM_Drt_En', 'Ограничение об/мин\nпри перегреве масла', 'флаг'),
        ('C_DO_01_A_High_Threshold', 'Сигнал перегрева\nвоздуха наддува', '54 °C → ОТКЛ'),
        ('C_DO_07_A_High_Threshold', 'Сигнал перегрева\nохлаждающей жидкости', '85 °C → ОТКЛ'),
        ('C_DO_08_A_High_Threshold', 'Сигнал перегрева\nтоплива', '68 °C → ОТКЛ'),
    ]
    elevated = [
        ('C_TIB_Fuel_Adjust_Upper_Limit', 'Лимит коррекции\nвпрыска TIB', '100% → 200%'),
        ('C_ABS_Max_Acctr_Speed', 'Максимальные\nоб/мин', '1990 → 2100'),
        ('C_ABS_DroopMax', 'Дроуп регулятора', '0% → 20%'),
        ('C_EPD_CT_TB_Time_To_Max_Trq_Drt', 'Задержка дерейта\n(темп. ОЖ)', '0 с → 43 с'),
        ('C_EPD_CP_TB_Time_To_Max_Trq_Brt', 'Задержка дерейта\n(давл. масло)', '0 с → 51 с'),
    ]

    fig = plt.figure(figsize=(15.5, 7.8), facecolor=C_DARK)
    fig.patch.set_facecolor(C_DARK)

    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_facecolor(C_DARK)
    ax.axis('off')
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)

    # Left panel — disabled
    ax.text(0.25, 0.97, '✗  Отключено на NTE200', color=C_RED,
            fontsize=13, fontweight='bold', ha='center', va='top')
    ax.text(0.25, 0.91, '(активно на 730E)', color=C_GRY,
            fontsize=9, ha='center', va='top')

    for i, (code, name, change) in enumerate(disabled):
        y_top = 0.83 - i * 0.18
        # Card
        card = mpatches.FancyBboxPatch((0.02, y_top - 0.14), 0.45, 0.14,
                                        boxstyle='round,pad=0.01',
                                        facecolor='#3a1a1a', edgecolor=C_RED,
                                        linewidth=1.2, transform=ax.transAxes,
                                        clip_on=False)
        ax.add_patch(card)
        ax.text(0.04, y_top - 0.03, name, color=C_WHT,
                fontsize=9, fontweight='bold', va='top')
        ax.text(0.04, y_top - 0.08, code, color=C_GRY,
                fontsize=7, va='top', style='italic')
        ax.text(0.44, y_top - 0.07, change, color=C_RED,
                fontsize=9, fontweight='bold', ha='right', va='center')

    # Right panel — elevated risks
    ax.text(0.75, 0.97, '⚠  Повышенные риски на NTE200',
            color=C_ORG, fontsize=13, fontweight='bold', ha='center', va='top')
    ax.text(0.75, 0.91, '(значения опаснее, чем на 730E)',
            color=C_GRY, fontsize=9, ha='center', va='top')

    for i, (code, name, change) in enumerate(elevated):
        y_top = 0.83 - i * 0.145
        card = mpatches.FancyBboxPatch((0.52, y_top - 0.125), 0.46, 0.125,
                                        boxstyle='round,pad=0.01',
                                        facecolor='#2e2000', edgecolor=C_ORG,
                                        linewidth=1.2, transform=ax.transAxes,
                                        clip_on=False)
        ax.add_patch(card)
        ax.text(0.54, y_top - 0.025, name, color=C_WHT,
                fontsize=8.5, fontweight='bold', va='top')
        ax.text(0.54, y_top - 0.075, code, color=C_GRY,
                fontsize=6.5, va='top', style='italic')
        ax.text(0.95, y_top - 0.06, change, color=C_ORG,
                fontsize=8.5, fontweight='bold', ha='right', va='center')

    # Vertical divider
    ax.axvline(0.505, color=C_GRY, lw=1, alpha=0.4, ymin=0.05, ymax=0.97)

    img_to_pptx(fig, slide, 0.3, 2.0, 33.2, 16.5)


def build_fuel_code_slide(prs):
    """Fuel Code comparison with explanation."""
    slide = add_slide(prs)
    clear_placeholders(slide)
    slide_background(slide, C_DARK)
    add_title_bar(slide, 'Fuel Code: FC0BU52 (730E) vs FC0WH95 (NTE200)',
                  'Различия в картах программирования впрыска')
    add_footer(slide)

    fig = plt.figure(figsize=(15.5, 7.8), facecolor=C_DARK)
    fig.patch.set_facecolor(C_DARK)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_facecolor(C_DARK)
    ax.axis('off')
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)

    items_730 = [
        ('Карта давления CR',    'Штатная для\nгорных условий'),
        ('Тайминг впрыска',      'Оптимизирован для\n1990 об/мин'),
        ('Дымоограничитель',     'Активный, ограничивает\nсмесеобразование'),
        ('Карта наддува',        'Соответствует 1990 об/мин\nи высотной коррекции'),
        ('EGT выхлопа',          '430–480 °C при номинале'),
    ]
    items_nte = [
        ('Карта давления CR',    'Модифицированная —\nдругой паттерн пилот/основной'),
        ('Тайминг впрыска',      'Смещён под 2100 об/мин\n(выше нагрев ГБЦ)'),
        ('Дымоограничитель',     'Ослаблен — больше\nтопливо при TIB=200%'),
        ('Карта наддува',        'Под 2100 об/мин — выше\nтемпература воздуха'),
        ('EGT выхлопа',          '470–560 °C при тех же\nусловиях нагружения'),
    ]

    ax.text(0.25, 0.97, 'FC0BU52  —  730E-DC', color=C_BLU,
            fontsize=14, fontweight='bold', ha='center', va='top')
    ax.text(0.75, 0.97, 'FC0WH95  —  NTE200', color=C_RED,
            fontsize=14, fontweight='bold', ha='center', va='top')
    ax.axvline(0.505, color=C_GRY, lw=1, alpha=0.4, ymin=0.05, ymax=0.94)

    for i, ((lbl, v730), (_, vnte)) in enumerate(zip(items_730, items_nte)):
        y = 0.86 - i * 0.155

        for side, val, xc, col, edge in [
            ('left',  v730, 0.02, '#1a2d3a', C_BLU),
            ('right', vnte, 0.53, '#2d1a1a', C_RED),
        ]:
            card = mpatches.FancyBboxPatch((xc, y - 0.12), 0.46, 0.12,
                                           boxstyle='round,pad=0.01',
                                           facecolor=col, edgecolor=edge,
                                           linewidth=1, transform=ax.transAxes,
                                           clip_on=False)
            ax.add_patch(card)

        # Label in center
        ax.text(0.505, y - 0.055, lbl, color=C_ACC, fontsize=9,
                fontweight='bold', ha='center', va='center')

        # Values
        ax.text(0.46, y - 0.055, v730, color=C_WHT,
                fontsize=8, ha='right', va='center', multialignment='right')
        ax.text(0.555, y - 0.055, vnte, color='#FF9999',
                fontsize=8, ha='left', va='center', multialignment='left')

    img_to_pptx(fig, slide, 0.3, 2.0, 33.2, 16.5)


def build_conclusion_slide(prs):
    """Final conclusions slide."""
    slide = add_slide(prs)
    clear_placeholders(slide)
    slide_background(slide, C_DARK)
    add_title_bar(slide, 'Выводы: ECM-конфигурация NTE200 vs 730E',
                  'Ключевые различия и связь с отказами клапанов')
    add_footer(slide)

    conclusions = [
        (C_RED,  'КРИТИЧНО',
         'На NTE200 отключены ВСЕ 4 цифровых выхода защиты (DO01/07/08, EPD_OT_RPM_Drt).\n'
         'Машина работает без внешней сигнализации перегрева и без автодерейта.'),
        (C_ORG,  'ВЫСОКИЙ РИСК',
         'TIB=200% → двойная доза коррекции впрыска при балансировке.\n'
         'Max. об/мин = 2100 vs 1990 на 730E → +5.5% ударных нагрузок на клапаны, +25–35 °C EGT.'),
        (C_ORG,  'ВЫСОКИЙ РИСК',
         'Задержки дерейта 51 с (масло) и 43 с (ОЖ) означают, что двигатель продолжает\n'
         'работать на полной мощности при критических условиях.'),
        (C_YLW,  'ПРЕДУПРЕЖДЕНИЕ',
         'Fuel Code FC0WH95 имеет другие карты давления CR и таймингов — режим сгорания\n'
         'отличается от штатного 730E. Возможно избыточное тепловыделение в конце расширения.'),
        (C_YLW,  'ПРЕДУПРЕЖДЕНИЕ',
         'Дроуп регулятора 20% vs 0% — нестабильность об/мин приводит к циклической\n'
         'нагрузке на клапанный механизм (переменная температура выхлопа).'),
        (C_ACC,  'РЕКОМЕНДАЦИЯ',
         'Привести ECM NTE200 к конфигурации 730E: включить DO01/07/08, EPD_OT_RPM_Drt_En,\n'
         'снизить Max RPM до 1990, ограничить TIB до 100%, установить задержки дерейта = 0 с.'),
    ]

    fig = plt.figure(figsize=(15.5, 7.8), facecolor=C_DARK)
    fig.patch.set_facecolor(C_DARK)
    ax = fig.add_axes([0.01, 0.01, 0.98, 0.98])
    ax.set_facecolor(C_DARK)
    ax.axis('off')
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)

    n = len(conclusions)
    h = 0.14
    gap = 0.01
    total = n * h + (n - 1) * gap
    y0 = (1 - total) / 2 + total

    for i, (color, label, text) in enumerate(conclusions):
        y = y0 - i * (h + gap)
        # Left color bar
        bar = mpatches.FancyBboxPatch((0.0, y - h + 0.005), 0.006, h - 0.01,
                                      boxstyle='square,pad=0',
                                      facecolor=color, linewidth=0,
                                      transform=ax.transAxes, clip_on=False)
        ax.add_patch(bar)
        # Card background
        card = mpatches.FancyBboxPatch((0.008, y - h + 0.005), 0.99, h - 0.01,
                                       boxstyle='round,pad=0.005',
                                       facecolor=C_DARK2, linewidth=0,
                                       transform=ax.transAxes, clip_on=False)
        ax.add_patch(card)
        ax.text(0.018, y - 0.025, label, color=color,
                fontsize=8.5, fontweight='bold', va='top')
        ax.text(0.018, y - 0.065, text, color=C_WHT,
                fontsize=8, va='top', multialignment='left')

    img_to_pptx(fig, slide, 0.3, 2.0, 33.2, 16.5)


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    out = '/home/user/NTE200/Калибровки_ECM_NTE200_vs_730E.pptx'

    # Load template
    tpl_path = '/home/user/NTE200/Развитие_шаблон.pptx'
    prs = Presentation(tpl_path)
    prs.slide_width  = Emu(int(13.33 * 914400))
    prs.slide_height = Emu(int(7.50  * 914400))

    build_title_slide(prs)
    print('Slide 1: title done')

    # Table slides (split 15 params into two pages)
    build_comparison_table_slide(
        prs, PARAMS[:8],
        'ECM-параметры защиты: сравнение 730E-DC vs NTE200  [1/2]',
        'Управление нагрузкой, Fuel Code, защита по температуре'
    )
    print('Slide 2: table 1 done')

    build_comparison_table_slide(
        prs, PARAMS[8:],
        'ECM-параметры защиты: сравнение 730E-DC vs NTE200  [2/2]',
        'Тепловые защиты, цифровые выходы, референсные об/мин'
    )
    print('Slide 3: table 2 done')

    build_bar_comparison_slide(prs)
    print('Slide 4: bar chart done')

    build_disabled_features_slide(prs)
    print('Slide 5: disabled features done')

    build_fuel_code_slide(prs)
    print('Slide 6: fuel code done')

    build_conclusion_slide(prs)
    print('Slide 7: conclusions done')

    prs.save(out)
    print(f'\nSaved: {out}')


if __name__ == '__main__':
    main()
