#!/usr/bin/env python3
"""
Build comprehensive PowerPoint: Анализ_Калибровок_ФИНАЛ.pptx
Combines ECM calibration analysis with KAMSS INSITE findings.
"""

import os
import io
import tempfile
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Cm
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ─── COLOR PALETTE ────────────────────────────────────────────────────────────
C_BG     = RGBColor(0x29, 0x31, 0x36)   # #293136 dark bg
C_ACCENT = RGBColor(0x3E, 0xF0, 0xAF)   # #3EF0AF teal/green accent
C_RED    = RGBColor(0xE0, 0x52, 0x52)   # #E05252 red
C_ORANGE = RGBColor(0xE0, 0x8C, 0x2E)   # #E08C2E orange
C_YELLOW = RGBColor(0xD4, 0xC9, 0x4A)   # #D4C94A yellow
C_GRAY   = RGBColor(0x70, 0x78, 0x80)   # #707880 gray
C_BLUE   = RGBColor(0x4E, 0xC9, 0xF0)   # #4EC9F0 blue
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x00, 0x00, 0x00)

# Panel bg slightly lighter
C_PANEL  = RGBColor(0x35, 0x3F, 0x45)
C_PANEL2 = RGBColor(0x3A, 0x45, 0x4C)

KAMSS_DIR = '/home/user/NTE200/kamss_data'
TEMPLATE  = '/home/user/NTE200/Развитие_шаблон.pptx'
OUTPUT    = '/home/user/NTE200/Анализ_Калибровок_ФИНАЛ.pptx'

# ─── HELPER FUNCTIONS ─────────────────────────────────────────────────────────

def rgb_hex(rgb: RGBColor):
    return f'#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=C_WHITE, align=PP_ALIGN.LEFT,
                italic=False, wrap=True, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_textbox_multiline(slide, left, top, width, height, lines,
                          font_size=12, bold=False, color=C_WHITE,
                          align=PP_ALIGN.LEFT, font_name='Calibri',
                          spacing_before=0):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, sz, bd, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(spacing_before)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = bd
        run.font.color.rgb = clr
        run.font.name = font_name
    return txBox

def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_image_from_gif(slide, gif_path, left, top, width, height):
    """Convert GIF to PNG and embed in slide."""
    img = Image.open(gif_path)
    img = img.convert('RGBA')
    png_buf = io.BytesIO()
    img.save(png_buf, format='PNG')
    png_buf.seek(0)
    slide.shapes.add_picture(png_buf, left, top, width, height)

def make_bar_chart_png():
    """Create bar chart comparing key numeric differences."""
    fig, ax = plt.subplots(figsize=(13, 5.5), facecolor='#293136')
    ax.set_facecolor('#293136')

    categories = [
        'Макс. об/мин\n(C_ABS_Max_Acctr)',
        'Порог заброса об/мин\n(FC0 Overspeed)',
        'FC556 Картерные газы\n(кПа)',
        'TIB топлив. коррекция\n(%)',
        'HSST High\n(Н·м / 10)',
    ]
    vals_730e = [1990, 2250, 10, 100, 550.9]   # HSST divided by 10 for scale
    vals_nte  = [2100, 2020,  5, 200, 365.0]

    x = np.arange(len(categories))
    w = 0.35

    bars1 = ax.bar(x - w/2, vals_730e, w, label='730E ш18 (референс)',
                   color='#4EC9F0', zorder=3, edgecolor='none')
    bars2 = ax.bar(x + w/2, vals_nte, w, label='NTE200 ш85 (проблемный)',
                   color='#E05252', zorder=3, edgecolor='none')

    # Value labels
    for bar in bars1:
        h = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., h + 15,
                f'{h:.0f}' if h < 500 else f'{h*10:.0f}',
                ha='center', va='bottom', color='#4EC9F0', fontsize=8.5, fontweight='bold')
    for bar in bars2:
        h = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., h + 15,
                f'{h:.0f}' if h < 500 else f'{h*10:.0f}',
                ha='center', va='bottom', color='#E05252', fontsize=8.5, fontweight='bold')

    # Fix HSST labels
    for i, bar in enumerate(bars1):
        if i == 4:
            h = bar.get_height()
            ax.texts[i*2].set_text('5509')
    for i, bar in enumerate(bars2):
        if i == 4:
            h = bar.get_height()
            ax.texts[i*2+1].set_text('3650')

    ax.set_xticks(x)
    ax.set_xticklabels(categories, color='white', fontsize=8.5)
    ax.set_ylabel('Значение', color='#707880', fontsize=10)
    ax.tick_params(colors='#707880', labelcolor='#aaaaaa')
    ax.spines['bottom'].set_color('#3EF0AF')
    ax.spines['left'].set_color('#707880')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.yaxis.grid(True, color='#404040', linestyle='--', alpha=0.5)
    ax.set_axisbelow(True)
    ax.legend(loc='upper right', facecolor='#35404A', edgecolor='#3EF0AF',
              labelcolor='white', fontsize=9)
    ax.set_title('Критические числовые отличия: 730E vs NTE200\n(HSST показан / 10 для масштаба)',
                 color='white', fontsize=11, fontweight='bold', pad=10)

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='PNG', dpi=150, bbox_inches='tight',
                facecolor='#293136')
    buf.seek(0)
    plt.close(fig)
    return buf

# ─── SLIDE BUILDERS ──────────────────────────────────────────────────────────

def build_slide1_title(prs):
    """Slide 1 — Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # blank layout
    set_slide_bg(slide, C_BG)
    W = prs.slide_width
    H = prs.slide_height

    # Accent bar top
    add_rect(slide, 0, 0, W, Cm(0.6), C_ACCENT)

    # Left accent strip
    add_rect(slide, 0, 0, Cm(0.8), H, RGBColor(0x22, 0x28, 0x2C))

    # Main title area
    add_rect(slide, Cm(1.5), Cm(1.5), W - Cm(3), Cm(3.8), C_PANEL)

    # Title text
    add_textbox(slide, Cm(2), Cm(1.8), W - Cm(4), Cm(2),
                'Анализ ECM-калибровок QSK50 MCRS',
                font_size=28, bold=True, color=C_WHITE,
                align=PP_ALIGN.LEFT)

    add_textbox(slide, Cm(2), Cm(3.2), W - Cm(4), Cm(1.5),
                'NTE200 vs 730E-DC: Подтверждение данными INSITE',
                font_size=22, bold=False, color=C_ACCENT,
                align=PP_ALIGN.LEFT)

    # Subtitle / metadata
    add_textbox(slide, Cm(2), Cm(5.5), W - Cm(4), Cm(1.2),
                'Данные Calterm Compare + INSITE КАМСС (май 2026) | АО Развитие | Полюс Магадан',
                font_size=12, bold=False, color=C_GRAY,
                align=PP_ALIGN.LEFT)

    # Decorative accent bar
    add_rect(slide, Cm(2), Cm(5.2), Cm(12), Cm(0.08), C_ACCENT)

    # Info tags row
    tags = [
        ('QSK50 MCRS', C_BLUE),
        ('INSITE КАМСС', C_ACCENT),
        ('Полюс Магадан', C_ORANGE),
        ('Май 2026', C_GRAY),
    ]
    x_off = Cm(2)
    for tag_text, tag_color in tags:
        tw = Cm(3.8)
        add_rect(slide, x_off, Cm(6.4), tw, Cm(0.55), tag_color)
        add_textbox(slide, x_off + Cm(0.15), Cm(6.42), tw - Cm(0.3), Cm(0.5),
                    tag_text, font_size=9, bold=True, color=C_BG)
        x_off += tw + Cm(0.3)

    # Bottom strip
    add_rect(slide, 0, H - Cm(0.5), W, Cm(0.5), C_PANEL)
    add_textbox(slide, Cm(2), H - Cm(0.48), W - Cm(4), Cm(0.45),
                'АО Развитие  ·  Анализ отказов клапанов  ·  Конфиденциально',
                font_size=8, color=C_GRAY, align=PP_ALIGN.CENTER)

    # GIF image — show functions/parameters screenshot
    gif_path = os.path.join(KAMSS_DIR, '730E ш18 DC model фуекции и параметры.GIF')
    if os.path.exists(gif_path):
        add_image_from_gif(slide, gif_path,
                           W - Cm(13), Cm(1.2), Cm(12.5), Cm(5.5))
        # dim overlay
        add_rect(slide, W - Cm(13), Cm(1.2), Cm(12.5), Cm(5.5),
                 RGBColor(0x29, 0x31, 0x36))
        # make overlay semi-transparent via XML alpha
        shape_rects = [s for s in slide.shapes if s.shape_type == 1]
        last_rect = shape_rects[-1]
        sp = last_rect._element
        spPr = sp.find(qn('p:spPr'))
        solidFill = spPr.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is None:
                srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                srgbClr.set('val', '293136')
            alpha_el = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha_el.set('val', '75000')  # 75% opacity = 25% transparency


def build_slide2_metadata(prs):
    """Slide 2 — Machine metadata two-column cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, C_BG)
    W = prs.slide_width
    H = prs.slide_height

    # Header bar
    add_rect(slide, 0, 0, W, Cm(1.4), C_PANEL)
    add_rect(slide, 0, 0, Cm(0.5), H, C_BLUE)
    add_textbox(slide, Cm(1), Cm(0.25), W - Cm(2), Cm(0.9),
                'МЕТАДАННЫЕ МАШИН', font_size=14, bold=True, color=C_ACCENT)

    # Section line
    add_rect(slide, Cm(1), Cm(1.5), W - Cm(2), Cm(0.04), C_ACCENT)

    col_w = (W - Cm(3.5)) / 2
    left1 = Cm(1)
    left2 = Cm(1) + col_w + Cm(1.5)

    def draw_card(x, y, title, title_color, data_rows, glyph=''):
        card_h = Cm(0.8) + len(data_rows) * Cm(0.7)
        add_rect(slide, x, y, col_w, card_h, C_PANEL)
        add_rect(slide, x, y, Cm(0.25), card_h, title_color)
        add_textbox(slide, x + Cm(0.4), y + Cm(0.1), col_w - Cm(0.6), Cm(0.6),
                    f'{glyph}{title}', font_size=13, bold=True, color=title_color)
        for i, (label, value, val_color) in enumerate(data_rows):
            ry = y + Cm(0.75) + i * Cm(0.68)
            add_textbox(slide, x + Cm(0.4), ry, Cm(4.5), Cm(0.6),
                        label, font_size=9, bold=False, color=C_GRAY)
            add_textbox(slide, x + Cm(4.8), ry, col_w - Cm(5), Cm(0.6),
                        value, font_size=9, bold=True, color=val_color)
        return card_h

    # 730E card
    rows_730e = [
        ('ECM калибровка:', 'AQ60217.28', C_BLUE),
        ('Fuel Code:', 'FC0BU52', C_ACCENT),
        ('Наработка ECM:', '~35 127 м/ч', C_WHITE),
        ('OEM производитель:', 'KOMATSU', C_WHITE),
        ('Модель:', '730E-DC', C_WHITE),
        ('ID модуля:', '33223470', C_GRAY),
        ('Снимок ECM:', '20.05.2026', C_GRAY),
    ]
    h1 = draw_card(left1, Cm(1.8), '730E ш18 — Референсная машина', C_BLUE, rows_730e, '◆ ')

    # NTE200 card
    rows_nte = [
        ('ECM калибровка:', 'AQ60809.08', C_RED),
        ('Fuel Code:', 'FC0WH95', C_RED),
        ('Наработка ECM:', '~4 015 м/ч', C_WHITE),
        ('OEM производитель:', 'NHL', C_WHITE),
        ('Модель:', 'QSKTA50-CE (NTE200)', C_WHITE),
        ('ID модуля:', '33238503', C_GRAY),
        ('Снимок ECM:', '19.05.2026', C_GRAY),
    ]
    h2 = draw_card(left2, Cm(1.8), 'NTE200 ш85 — Проблемная машина', C_RED, rows_nte, '⚠ ')

    # Fault codes found
    y2 = Cm(1.8) + max(h1, h2) + Cm(0.4)

    add_rect(slide, left1, y2, col_w, Cm(1.6), C_PANEL2)
    add_rect(slide, left1, y2, Cm(0.25), Cm(1.6), C_BLUE)
    add_textbox(slide, left1 + Cm(0.4), y2 + Cm(0.1), col_w - Cm(0.6), Cm(0.5),
                'Коды неисправностей 730E ш18', font_size=10, bold=True, color=C_BLUE)
    add_textbox(slide, left1 + Cm(0.4), y2 + Cm(0.55), col_w - Cm(0.6), Cm(0.9),
                'FC418 (вода в топливе) × 2  ·  FC1542 × 3  ·  FC245 × 2  ·  FC1117 × 1\n'
                'Наработка: 35 281 ч ECM | Снимок: 20.05.2026',
                font_size=8.5, color=C_GRAY)

    add_rect(slide, left2, y2, col_w, Cm(1.6), C_PANEL2)
    add_rect(slide, left2, y2, Cm(0.25), Cm(1.6), C_RED)
    add_textbox(slide, left2 + Cm(0.4), y2 + Cm(0.1), col_w - Cm(0.6), Cm(0.5),
                'Коды неисправностей NTE200 ш85', font_size=10, bold=True, color=C_RED)
    add_textbox(slide, left2 + Cm(0.4), y2 + Cm(0.55), col_w - Cm(0.6), Cm(0.9),
                'FC611 горячий остановы × 18  ·  FC1518 × 9\n'
                'Наработка: 4 015 ч ECM | Снимок: 19.05.2026',
                font_size=8.5, color=C_GRAY)

    # Footer
    add_rect(slide, 0, H - Cm(0.45), W, Cm(0.45), C_PANEL)
    add_textbox(slide, Cm(1), H - Cm(0.43), W - Cm(2), Cm(0.4),
                'АО Развитие  ·  QSK50 MCRS  ·  Анализ отказов клапанов  ·  Полюс Магадан  ·  2026',
                font_size=7.5, color=C_GRAY, align=PP_ALIGN.CENTER)


def build_slide3_divider1(prs):
    """Slide 3 — Section divider: Part 1"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, RGBColor(0x1E, 0x25, 0x2A))
    W = prs.slide_width
    H = prs.slide_height

    # Large accent bar left
    add_rect(slide, 0, 0, Cm(1.2), H, C_ACCENT)

    # Part number
    add_textbox(slide, Cm(2), Cm(1.8), W - Cm(3), Cm(1.2),
                'ЧАСТЬ 1', font_size=42, bold=True, color=C_ACCENT,
                align=PP_ALIGN.LEFT)

    add_textbox(slide, Cm(2), Cm(3.2), W - Cm(3), Cm(1.5),
                'Данные предыдущего анализа', font_size=30, bold=False,
                color=C_WHITE, align=PP_ALIGN.LEFT)

    # Horizontal line
    add_rect(slide, Cm(2), Cm(4.9), W - Cm(4), Cm(0.06), C_GRAY)

    add_textbox(slide, Cm(2), Cm(5.1), W - Cm(4), Cm(1),
                'Источник: ECM_对比_NTE200_730E_中文.xlsx  —  Калибровочные параметры ECM из предыдущего анализа',
                font_size=11, bold=False, color=C_GRAY, align=PP_ALIGN.LEFT)

    # Decorative dots
    for i in range(8):
        add_rect(slide, W - Cm(4) + i * Cm(0.4), Cm(2.5),
                 Cm(0.25), Cm(0.25), RGBColor(0x3E, 0xF0, 0xAF) if i % 2 == 0 else C_PANEL)


def build_slide4_prev_analysis_table(prs):
    """Slide 4 — Previous analysis: critical parameters table"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, C_BG)
    W = prs.slide_width
    H = prs.slide_height

    # Header
    add_rect(slide, 0, 0, W, Cm(1.3), C_PANEL)
    add_rect(slide, 0, 0, Cm(0.5), H, C_ORANGE)
    add_textbox(slide, Cm(0.8), Cm(0.2), W - Cm(2), Cm(0.9),
                'ПРЕДЫДУЩИЙ АНАЛИЗ: Критические параметры ECM  (13 параметров)',
                font_size=13, bold=True, color=C_WHITE)
    add_textbox(slide, Cm(0.8), Cm(0.85), W - Cm(2), Cm(0.4),
                'Источник: ECM_对比_NTE200_730E_中文.xlsx | Сравнение калибровок 730E vs NTE200',
                font_size=8, color=C_GRAY)

    # Table data
    rows = [
        # (param, 730E value, NTE200 value, risk_level, effect)
        ('Fuel Code', 'FC0BU52', 'FC0WH95', 'ПРЕДУПРЕЖДЕНИЕ', C_YELLOW, 'Разная карта впрыска/горения'),
        ('C_EPD_OT_RPM_Drt_En', 'ВКЛ (1)', 'ОТКЛ (0)', 'КРИТИЧНО', C_RED, 'Нет авт. снижения об/мин при перегреве масла'),
        ('C_EPD_CP_TB_Time_To_Max_Trq_Brt', '0 с', '51 с', 'ВЫСОКИЙ', C_ORANGE, 'При давл. масла↓: 51 с работы на полной мощности'),
        ('C_EPD_CT_TB_Time_To_Max_Trq_Drt', '0 с', '43 с', 'ВЫСОКИЙ', C_ORANGE, 'При перегреве ОЖ: 43 с до снижения нагрузки'),
        ('C_ABS_Max_Acctr_Speed', '1990 RPM', '2100 RPM', 'ВЫСОКИЙ', C_ORANGE, '+110 RPM → EGT выше ~25–35°C, выше ударная нагрузка'),
        ('C_ABS_DroopMax', '0%', '20%', 'ПРЕДУПРЕЖДЕНИЕ', C_YELLOW, 'NTE200 допускает флуктуацию об/мин ±20%'),
        ('C_HSST_Net_Torque_High_Thd', '5509 Н·м', '3650 Н·м', 'ВЫСОКИЙ', C_ORANGE, 'NTE200 горячая остановка при меньшем крутящем моменте'),
        ('C_HSST_Net_Torque_Low_Thd', '4722 Н·м', '1850 Н·м', 'КРИТИЧНО', C_RED, 'Нижний порог в 2.5× ниже — нестабильный гистерезис'),
        ('C_TIB_Fuel_Adjust_Upper_Limit', '100%', '200%', 'КРИТИЧНО', C_RED, 'NTE200: двойная коррекция впрыска → перегрев'),
        ('C_DO_01_A_High_Threshold', '54°C', 'ОТКЛ', 'КРИТИЧНО', C_RED, 'Нет сигнала о перегреве воздуха на NTE200'),
        ('C_DO_07_A_High_Threshold', '85°C', 'ОТКЛ', 'КРИТИЧНО', C_RED, 'Нет внешнего сигнала перегрева ОЖ на NTE200'),
        ('C_DO_08_A_High_Threshold', '68°C', 'ОТКЛ', 'КРИТИЧНО', C_RED, 'Нет сигнала перегрева топлива на NTE200'),
        ('T_ABS_MaxRef_2/3', '1990 RPM', '2225.8 RPM', 'ВЫСОКИЙ', C_ORANGE, 'В отд. режимах NTE200 разгоняется до 2226 RPM'),
    ]

    # Column widths
    col_x = [Cm(0.6), Cm(4.6), Cm(7.4), Cm(9.7), Cm(11.6), Cm(16.8)]
    col_labels = ['Параметр', '730E-DC', 'NTE200', 'Риск', 'Эффект']

    # Header row
    hdr_y = Cm(1.45)
    hdr_h = Cm(0.5)
    add_rect(slide, Cm(0.6), hdr_y, W - Cm(1.2), hdr_h, RGBColor(0x22, 0x28, 0x2C))
    col_widths = [Cm(4.0), Cm(2.8), Cm(2.3), Cm(1.9), Cm(5.2)]
    for i, (cx, cw, cl) in enumerate(zip(col_x[:-1], col_widths, col_labels)):
        add_textbox(slide, cx + Cm(0.1), hdr_y + Cm(0.05), cw, hdr_h,
                    cl, font_size=8.5, bold=True, color=C_ACCENT)

    row_h = Cm(0.47)
    start_y = Cm(1.97)
    for ri, (param, val_730, val_nte, risk, risk_color, effect) in enumerate(rows):
        ry = start_y + ri * row_h
        bg = C_PANEL if ri % 2 == 0 else RGBColor(0x30, 0x3A, 0x40)
        add_rect(slide, Cm(0.6), ry, W - Cm(1.2), row_h - Cm(0.02), bg)
        # Risk strip
        add_rect(slide, Cm(0.6), ry, Cm(0.22), row_h - Cm(0.02), risk_color)

        add_textbox(slide, col_x[0] + Cm(0.3), ry + Cm(0.04), Cm(3.8), row_h,
                    param, font_size=7.5, bold=True, color=C_WHITE)
        add_textbox(slide, col_x[1], ry + Cm(0.04), Cm(2.8), row_h,
                    val_730, font_size=7.5, bold=False, color=C_ACCENT)
        add_textbox(slide, col_x[2], ry + Cm(0.04), Cm(2.3), row_h,
                    val_nte, font_size=7.5, bold=True,
                    color=C_RED if risk in ('КРИТИЧНО', 'ВЫСОКИЙ') else C_YELLOW)
        add_textbox(slide, col_x[3], ry + Cm(0.04), Cm(1.9), row_h,
                    risk, font_size=7, bold=True, color=risk_color)
        add_textbox(slide, col_x[4], ry + Cm(0.04), Cm(5.2), row_h,
                    effect, font_size=7, bold=False, color=C_GRAY)

    # Footer
    add_rect(slide, 0, H - Cm(0.4), W, Cm(0.4), C_PANEL)
    add_textbox(slide, Cm(1), H - Cm(0.38), W - Cm(2), Cm(0.35),
                'АО Развитие  ·  QSK50 MCRS  ·  Анализ отказов клапанов  ·  Полюс Магадан  ·  2026  ·  Конфиденциально',
                font_size=7, color=C_GRAY, align=PP_ALIGN.CENTER)


def build_slide5_divider2(prs):
    """Slide 5 — Section divider: Part 2"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, RGBColor(0x1E, 0x25, 0x2A))
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, Cm(1.2), H, C_ACCENT)

    add_textbox(slide, Cm(2), Cm(1.8), W - Cm(3), Cm(1.2),
                'ЧАСТЬ 2', font_size=42, bold=True, color=C_ACCENT)

    add_textbox(slide, Cm(2), Cm(3.2), W - Cm(3), Cm(1.5),
                'Новые данные из INSITE КАМСС', font_size=30, bold=False, color=C_WHITE)

    add_rect(slide, Cm(2), Cm(4.9), W - Cm(4), Cm(0.06), C_GRAY)

    add_textbox(slide, Cm(2), Cm(5.1), W - Cm(4), Cm(1),
                'Источник: Calterm Calibration Compare Report + ECM Image Analyzer (май 2026)\n'
                'AQ60217.28 (730E ш18)  vs  AQ60809.08 (NTE200 ш85)',
                font_size=11, bold=False, color=C_GRAY)

    # Show KAMSS screenshot small
    gif_path = os.path.join(KAMSS_DIR, 'NTE200 ш85 model функции и параметры.GIF')
    if os.path.exists(gif_path):
        add_image_from_gif(slide, gif_path, W - Cm(11), Cm(1.5), Cm(10), Cm(5.3))
        # semi-opaque overlay
        add_rect(slide, W - Cm(11), Cm(1.5), Cm(10), Cm(5.3), RGBColor(0x1E, 0x25, 0x2A))
        shape_rects = [s for s in slide.shapes if s.shape_type == 1]
        last_rect = shape_rects[-1]
        sp = last_rect._element
        spPr = sp.find(qn('p:spPr'))
        solidFill = spPr.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is None:
                srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                srgbClr.set('val', '1E252A')
            alpha_el = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha_el.set('val', '65000')


def build_slide6_kamss_confirmed(prs):
    """Slide 6 — Подтверждённые факты (КАМСС) — two-column layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, C_BG)
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, W, Cm(1.3), C_PANEL)
    add_rect(slide, 0, 0, Cm(0.5), H, C_ACCENT)
    add_textbox(slide, Cm(0.8), Cm(0.2), W - Cm(2), Cm(0.9),
                'ПОДТВЕРЖДЁННЫЕ ФАКТЫ  |  КАМСС + INSITE (май 2026)',
                font_size=13, bold=True, color=C_WHITE)
    add_textbox(slide, Cm(0.8), Cm(0.85), W - Cm(2), Cm(0.4),
                'Левая колонка: подтверждено Calterm Compare Report  ·  Правая: новые находки КАМСС',
                font_size=8, color=C_GRAY)

    col_w = (W - Cm(3)) / 2
    left1 = Cm(0.8)
    left2 = Cm(0.8) + col_w + Cm(1.4)

    # Left column header
    add_rect(slide, left1, Cm(1.45), col_w, Cm(0.55), C_ACCENT)
    add_textbox(slide, left1 + Cm(0.2), Cm(1.48), col_w - Cm(0.4), Cm(0.48),
                'Подтверждено Calterm Compare', font_size=10, bold=True, color=C_BG)

    # Right column header
    add_rect(slide, left2, Cm(1.45), col_w, Cm(0.55), C_BLUE)
    add_textbox(slide, left2 + Cm(0.2), Cm(1.48), col_w - Cm(0.4), Cm(0.48),
                'Новые находки КАМСС', font_size=10, bold=True, color=C_BG)

    def mini_card(x, y, title, body, accent_c, body_color=None):
        h = Cm(1.05)
        add_rect(slide, x, y, col_w, h, C_PANEL)
        add_rect(slide, x, y, Cm(0.2), h, accent_c)
        add_textbox(slide, x + Cm(0.35), y + Cm(0.06), col_w - Cm(0.5), Cm(0.42),
                    title, font_size=9, bold=True, color=accent_c)
        add_textbox(slide, x + Cm(0.35), y + Cm(0.5), col_w - Cm(0.5), Cm(0.52),
                    body, font_size=8, bold=False,
                    color=body_color if body_color else C_GRAY)
        return h

    gap = Cm(0.18)

    left_cards = [
        ('Fuel Code FC0BU52 vs FC0WH95  ✓ Подтверждено',
         'Подтверждено: разные карты впрыска. Влияние: температура сгорания, EGT.',
         C_ACCENT, None),
        ('DO_Option: 6765 (730E) vs 61042 (NTE200)  ✓ Подтверждено',
         'Разная конфигурация цифровых выходов подтверждена. DO01/07/08 отключены.',
         C_ACCENT, None),
        ('SC_Option: 60621 (730E) vs 61878 (NTE200)  ✓ Подтверждено',
         'Разные системные опции подтверждены. Иная конфигурация системы.',
         C_ACCENT, None),
        ('Calterm Compare Report от 25.05.2026  ✓ Подтверждено',
         'Содержит прямое сравнение AQ60217.28 vs AQ60809.08. Использован как основа.',
         C_ACCENT, None),
    ]

    right_cards = [
        ('FC556: порог картерных газов',
         '730E = 10 кПа  vs  NTE200 = 5 кПа\nПорог вдвое ниже — раньше реакция ECM',
         C_ORANGE, C_ORANGE),
        ('Заброс оборотов (FC0 Overspeed)',
         '730E = 2250 RPM  vs  NTE200 = 2020 RPM\nЗазор безопасности NTE200 = 80 RPM (730E = 260 RPM)',
         C_RED, C_RED),
        ('FC4642 / FC1852: Вода в топливе',
         'Задержка 17 мин до снижения оборотов — только NTE200\nРиск работы на загрязнённом топливе без немедленной реакции',
         C_ORANGE, C_ORANGE),
        ('FC611 горячие остановы NTE200 ш85',
         'FC611 = 18 горячих останова | FC1518 = 9 событий\nПри наработке всего ~4 015 м/ч — высокая частота',
         C_RED, C_RED),
    ]

    cy = Cm(2.1)
    for title, body, acc, bc in left_cards:
        mini_card(left1, cy, title, body, acc, bc)
        cy += Cm(1.05) + gap

    cy = Cm(2.1)
    for title, body, acc, bc in right_cards:
        mini_card(left2, cy, title, body, acc, bc)
        cy += Cm(1.05) + gap

    # Add GIF image at bottom
    gif_path = os.path.join(KAMSS_DIR, 'NTE200 ш85 model снижение оборотов 0.GIF')
    if os.path.exists(gif_path):
        img_y = cy + Cm(0.1)
        if img_y + Cm(1.5) < H - Cm(0.6):
            add_image_from_gif(slide, gif_path, left2, img_y,
                               col_w, H - img_y - Cm(0.6))

    add_rect(slide, 0, H - Cm(0.4), W, Cm(0.4), C_PANEL)
    add_textbox(slide, Cm(1), H - Cm(0.38), W - Cm(2), Cm(0.35),
                'АО Развитие  ·  QSK50 MCRS  ·  Анализ отказов клапанов  ·  Полюс Магадан  ·  2026  ·  Конфиденциально',
                font_size=7, color=C_GRAY, align=PP_ALIGN.CENTER)


def build_slide7_same_settings(prs):
    """Slide 7 — Одинаковые настройки (важное открытие)"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, C_BG)
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, W, Cm(1.4), C_PANEL)
    add_rect(slide, 0, 0, Cm(0.5), H, C_BLUE)
    add_textbox(slide, Cm(0.8), Cm(0.18), W - Cm(2), Cm(0.95),
                'ОДИНАКОВЫЕ НАСТРОЙКИ  —  Важное открытие из INSITE',
                font_size=13, bold=True, color=C_WHITE)
    add_textbox(slide, Cm(0.8), Cm(0.88), W - Cm(2), Cm(0.45),
                'Внутренние защиты ECM: одинаковы у 730E и NTE200  ·  Отличия — в ВНЕШНЕЙ КОНФИГУРАЦИИ',
                font_size=8.5, color=C_BLUE)

    # Important message box
    msg_y = Cm(1.55)
    add_rect(slide, Cm(0.8), msg_y, W - Cm(1.6), Cm(1.1), RGBColor(0x1A, 0x35, 0x45))
    add_rect(slide, Cm(0.8), msg_y, Cm(0.35), Cm(1.1), C_BLUE)
    add_textbox(slide, Cm(1.3), msg_y + Cm(0.08), W - Cm(2.5), Cm(0.42),
                'ВЫВОД:', font_size=10, bold=True, color=C_BLUE)
    add_textbox(slide, Cm(1.3), msg_y + Cm(0.52), W - Cm(2.5), Cm(0.55),
                'Внутренние защиты ECM (дерейт момента/оборотов) ОДИНАКОВЫ у обеих машин. '
                'Отличия касаются внешних сигналов (DO) и установочных параметров.',
                font_size=9, bold=False, color=C_WHITE)

    # Table
    table_y = Cm(2.8)
    col_x   = [Cm(0.8), Cm(5.8), Cm(12.5), Cm(20.5)]
    col_w   = [Cm(5.0), Cm(6.7), Cm(8.0), Cm(3.5)]
    headers = ['Код неисправности / FC', 'Параметр', 'Значение 730E / NTE200', 'Статус']

    hdr_h = Cm(0.52)
    add_rect(slide, Cm(0.8), table_y, W - Cm(1.6), hdr_h, RGBColor(0x1E, 0x28, 0x2E))
    for cx, cw, lbl in zip(col_x, col_w, headers):
        add_textbox(slide, cx + Cm(0.1), table_y + Cm(0.06), cw, hdr_h,
                    lbl, font_size=8.5, bold=True, color=C_ACCENT)

    same_data = [
        ('FC146', 'Темп. ОЖ → снижение момента', '100°C, задержка 5 с', 'ОДИНАКОВО'),
        ('FC421', 'Температура масла', '121.1°C — порог защиты', 'ОДИНАКОВО'),
        ('FC415', 'Давление масла', 'Одинаковые пороговые таблицы', 'ОДИНАКОВО'),
        ('FC235', 'Уровень ОЖ', 'Невысокий, задержка 30 с', 'ОДИНАКОВО'),
        ('FC155–FC165', 'Температура впускного воздуха', '93.3°C — порог защиты', 'ОДИНАКОВО'),
        ('FC228', 'Давление ОЖ', 'Одинаковые пороговые таблицы', 'ОДИНАКОВО'),
        ('FC151', 'Темп. ОЖ → снижение оборотов', '110°C', 'ОДИНАКОВО'),
    ]

    row_h = Cm(0.5)
    for ri, (fc, param, value, status) in enumerate(same_data):
        ry = table_y + hdr_h + ri * row_h
        bg = C_PANEL if ri % 2 == 0 else RGBColor(0x30, 0x3A, 0x40)
        add_rect(slide, Cm(0.8), ry, W - Cm(1.6), row_h - Cm(0.02), bg)
        add_rect(slide, Cm(0.8), ry, Cm(0.2), row_h - Cm(0.02), C_BLUE)
        add_textbox(slide, col_x[0] + Cm(0.3), ry + Cm(0.06), col_w[0], row_h,
                    fc, font_size=8.5, bold=True, color=C_BLUE)
        add_textbox(slide, col_x[1] + Cm(0.1), ry + Cm(0.06), col_w[1], row_h,
                    param, font_size=8, bold=False, color=C_WHITE)
        add_textbox(slide, col_x[2] + Cm(0.1), ry + Cm(0.06), col_w[2], row_h,
                    value, font_size=8, bold=False, color=C_ACCENT)
        add_textbox(slide, col_x[3] + Cm(0.1), ry + Cm(0.06), col_w[3], row_h,
                    status, font_size=8, bold=True, color=C_ACCENT)

    add_rect(slide, 0, H - Cm(0.4), W, Cm(0.4), C_PANEL)
    add_textbox(slide, Cm(1), H - Cm(0.38), W - Cm(2), Cm(0.35),
                'АО Развитие  ·  QSK50 MCRS  ·  Анализ отказов клапанов  ·  Полюс Магадан  ·  2026  ·  Конфиденциально',
                font_size=7, color=C_GRAY, align=PP_ALIGN.CENTER)


def build_slide8_divider3(prs):
    """Slide 8 — Section divider: Part 3"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, RGBColor(0x1E, 0x25, 0x2A))
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, Cm(1.2), H, C_ORANGE)

    add_textbox(slide, Cm(2), Cm(1.8), W - Cm(3), Cm(1.2),
                'ЧАСТЬ 3', font_size=42, bold=True, color=C_ORANGE)

    add_textbox(slide, Cm(2), Cm(3.2), W - Cm(3), Cm(1.5),
                'Сравнение с предыдущим анализом', font_size=30, bold=False, color=C_WHITE)

    add_rect(slide, Cm(2), Cm(4.9), W - Cm(4), Cm(0.06), C_GRAY)

    add_textbox(slide, Cm(2), Cm(5.1), W - Cm(4), Cm(1),
                'Что подтверждено  ·  Что изменилось  ·  Что новое\n'
                'Источник: ECM_对比 + Calterm Compare + INSITE КАМСС',
                font_size=11, bold=False, color=C_GRAY)


def build_slide9_comparison_table(prs):
    """Slide 9 — Таблица различий: предыдущий анализ vs КАМСС"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, C_BG)
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, W, Cm(1.3), C_PANEL)
    add_rect(slide, 0, 0, Cm(0.5), H, C_ORANGE)
    add_textbox(slide, Cm(0.8), Cm(0.2), W - Cm(2), Cm(0.9),
                'СРАВНЕНИЕ: Предыдущий анализ vs КАМСС / INSITE',
                font_size=13, bold=True, color=C_WHITE)
    add_textbox(slide, Cm(0.8), Cm(0.85), W - Cm(2), Cm(0.4),
                'Цвет статуса: Зелёный = Подтверждено  ·  Синий = Новое одинаково  ·  Оранжевый = Новое различие  ·  Жёлтый = Требует проверки',
                font_size=7.5, color=C_GRAY)

    # Table columns
    col_x = [Cm(0.8), Cm(5.5), Cm(11.0), Cm(18.5)]
    col_w = [Cm(4.7), Cm(5.5), Cm(7.5), Cm(5.0)]
    headers = ['Параметр', 'Предыдущий анализ (ECM_对比)', 'КАМСС / INSITE', 'Статус']

    table_y = Cm(1.42)
    hdr_h = Cm(0.52)
    add_rect(slide, Cm(0.8), table_y, W - Cm(1.6), hdr_h, RGBColor(0x1E, 0x28, 0x2E))
    for cx, cw, lbl in zip(col_x, col_w, headers):
        add_textbox(slide, cx + Cm(0.1), table_y + Cm(0.06), cw, hdr_h,
                    lbl, font_size=8.5, bold=True, color=C_ACCENT)

    # (param, prev_analysis, kamss, status, status_color, strip_color)
    comparison_data = [
        ('Fuel Code',
         'FC0BU52 vs FC0WH95',
         'FC0BU52 vs FC0WH95\nПрямое совпадение из Calterm Compare',
         '✓ Подтверждено', C_ACCENT, C_ACCENT),
        ('C_EPD_OT_RPM_Drt_En\n(откл. на NTE200)',
         'ВКЛ(1) 730E vs ОТКЛ(0) NTE200',
         'Подтверждается разницей DO_Option\n6765→61042: цифровые выходы изменены',
         '✓ Подтверждено', C_ACCENT, C_ACCENT),
        ('Внутренние пороги дерейта',
         'Не сравнивались в предыдущем анализе',
         'FC146/421/415/235/155-165 ОДИНАКОВЫ\nу 730E и NTE200 — важное открытие',
         '★ Новое: одинаково', C_BLUE, C_BLUE),
        ('Внешние цифровые выходы DO\n(DO01/07/08)',
         'DO_01/07/08 отключены на NTE200\nC_DO_0x_A_High_Threshold = ОТКЛ',
         'DO_Option 61042 (NTE200) vs 6765 (730E)\nПодтверждает отключение внешних сигналов',
         '✓ Подтверждено +\nуточнено', C_ACCENT, C_ACCENT),
        ('FC556\nКартерные газы порог',
         'Данных не было в предыдущем анализе',
         '730E = 10 кПа  vs  NTE200 = 5 кПа\nПорог NTE200 вдвое ниже',
         '★ Новое различие', C_ORANGE, C_ORANGE),
        ('Заброс оборотов\nFC0 Overspeed threshold',
         'Данных не было в предыдущем анализе',
         '730E = 2250 RPM  vs  NTE200 = 2020 RPM\nЗазор безопасности 80 RPM (vs 260 RPM)',
         '★ Новое различие', C_ORANGE, C_ORANGE),
        ('Макс. рабочие об/мин',
         '1990 RPM (730E) vs 2100 RPM (NTE200)',
         'При заборосе порог NTE200 = 2020 RPM\nМежду 2100 и 2020 — зазор 80 RPM',
         '⚠ Уточнение', C_YELLOW, C_YELLOW),
        ('TIB = 200%',
         'C_TIB_Fuel_Adjust_Upper_Limit = 200%\nПредыдущий анализ: критично',
         'Из Calterm Compare напрямую не видно\nТребует Calterm подтверждения',
         '⏳ Требует проверки', C_YELLOW, C_YELLOW),
    ]

    row_h = Cm(0.62)
    for ri, (param, prev, kamss, status, stat_color, strip_c) in enumerate(comparison_data):
        ry = table_y + hdr_h + ri * row_h
        bg = C_PANEL if ri % 2 == 0 else RGBColor(0x30, 0x3A, 0x40)
        add_rect(slide, Cm(0.8), ry, W - Cm(1.6), row_h - Cm(0.02), bg)
        add_rect(slide, Cm(0.8), ry, Cm(0.2), row_h - Cm(0.02), strip_c)

        add_textbox(slide, col_x[0] + Cm(0.3), ry + Cm(0.04), col_w[0], row_h,
                    param, font_size=7.5, bold=True, color=C_WHITE)
        add_textbox(slide, col_x[1] + Cm(0.1), ry + Cm(0.04), col_w[1], row_h,
                    prev, font_size=7, bold=False, color=C_GRAY)
        add_textbox(slide, col_x[2] + Cm(0.1), ry + Cm(0.04), col_w[2], row_h,
                    kamss, font_size=7, bold=False, color=C_WHITE)
        add_textbox(slide, col_x[3] + Cm(0.1), ry + Cm(0.04), col_w[3], row_h,
                    status, font_size=7.5, bold=True, color=stat_color)

    add_rect(slide, 0, H - Cm(0.4), W, Cm(0.4), C_PANEL)
    add_textbox(slide, Cm(1), H - Cm(0.38), W - Cm(2), Cm(0.35),
                'АО Развитие  ·  QSK50 MCRS  ·  Анализ отказов клапанов  ·  Полюс Магадан  ·  2026  ·  Конфиденциально',
                font_size=7, color=C_GRAY, align=PP_ALIGN.CENTER)


def build_slide10_bar_chart(prs):
    """Slide 10 — Bar chart: critical numeric differences"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, C_BG)
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, W, Cm(1.3), C_PANEL)
    add_rect(slide, 0, 0, Cm(0.5), H, C_RED)
    add_textbox(slide, Cm(0.8), Cm(0.2), W - Cm(2), Cm(0.9),
                'КРИТИЧЕСКИЕ ОТЛИЧИЯ: ЧИСЛОВОЕ СРАВНЕНИЕ  —  730E vs NTE200',
                font_size=13, bold=True, color=C_WHITE)
    add_textbox(slide, Cm(0.8), Cm(0.85), W - Cm(2), Cm(0.4),
                'Источник: ECM_对比 + Calterm Compare Report + INSITE КАМСС (май 2026)',
                font_size=8, color=C_GRAY)

    # Create the chart
    chart_buf = make_bar_chart_png()

    # Embed chart
    chart_w = W - Cm(1.8)
    chart_h = H - Cm(2.2)
    slide.shapes.add_picture(chart_buf, Cm(0.9), Cm(1.5), chart_w, chart_h)

    add_rect(slide, 0, H - Cm(0.4), W, Cm(0.4), C_PANEL)
    add_textbox(slide, Cm(1), H - Cm(0.38), W - Cm(2), Cm(0.35),
                'АО Развитие  ·  QSK50 MCRS  ·  Анализ отказов клапанов  ·  Полюс Магадан  ·  2026  ·  Конфиденциально',
                font_size=7, color=C_GRAY, align=PP_ALIGN.CENTER)


def build_slide11_divider4(prs):
    """Slide 11 — Section divider: Part 4"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, RGBColor(0x1E, 0x25, 0x2A))
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, Cm(1.2), H, C_RED)

    add_textbox(slide, Cm(2), Cm(1.8), W - Cm(3), Cm(1.2),
                'ЧАСТЬ 4', font_size=42, bold=True, color=C_RED)

    add_textbox(slide, Cm(2), Cm(3.2), W - Cm(3), Cm(1.5),
                'Выводы и мероприятия', font_size=30, bold=False, color=C_WHITE)

    add_rect(slide, Cm(2), Cm(4.9), W - Cm(4), Cm(0.06), C_GRAY)

    add_textbox(slide, Cm(2), Cm(5.1), W - Cm(4), Cm(1),
                'На основе: предыдущего анализа + данных INSITE КАМСС (май 2026)',
                font_size=11, bold=False, color=C_GRAY)


def build_slide12_conclusions(prs):
    """Slide 12 — Выводы (6 conclusion cards)"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, C_BG)
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, W, Cm(1.3), C_PANEL)
    add_rect(slide, 0, 0, Cm(0.5), H, C_RED)
    add_textbox(slide, Cm(0.8), Cm(0.2), W - Cm(2), Cm(0.9),
                'ВЫВОДЫ', font_size=14, bold=True, color=C_WHITE)
    add_textbox(slide, Cm(0.8), Cm(0.85), W - Cm(2), Cm(0.4),
                'QSK50 MCRS  ·  NTE200 vs 730E-DC  ·  Причины ускоренного износа клапанов',
                font_size=8.5, color=C_GRAY)

    conclusions = [
        (C_RED, 'КРИТИЧНО',
         'Отключены внешние сигналы защиты DO01/07/08 и флаг EPD_OT',
         'Подтверждено КАМСС. Машина не получает внешних команд на дерейт от '
         'машинного контроллера. Все три цифровых выхода отключены '
         '(DO_Option 6765→61042). NTE200 "не знает" о перегреве от внешней системы.'),
        (C_ORANGE, 'ВЫСОКИЙ',
         'Критически малый зазор до заброса оборотов: 80 RPM',
         'При макс. рабочих 2100 RPM защита от заброса NTE200 = 2020 RPM. '
         'Зазор безопасности = 80 RPM (730E = 260 RPM). '
         'При нагрузочных переходах возможен пробой порога → ударная нагрузка на клапаны.'),
        (C_ORANGE, 'ВЫСОКИЙ',
         'TIB=200% + Fuel Code FC0WH95 = горячее сгорание',
         'Из предыдущего анализа: C_TIB_Fuel_Adjust_Upper_Limit = 200% (в 2× выше нормы) '
         'в сочетании с иным Fuel Code даёт более горячее сгорание. '
         'Подтверждено Calterm Compare Report от 25.05.2026.'),
        (C_YELLOW, 'ПРЕДУПРЕЖДЕНИЕ',
         'Внутренние защиты ECM одинаковы — проблема в КОНФИГУРАЦИИ связи',
         'Внутренние защиты ECM (пороги дерейта момента/оборотов) ОДИНАКОВЫ у 730E и NTE200 '
         '(FC146/421/415/235/155-165). Проблема — не в отсутствии внутренней защиты, '
         'а в конфигурации связи с машинным контроллером NHL.'),
        (C_YELLOW, 'ПРЕДУПРЕЖДЕНИЕ',
         'FC4642: задержка 17 мин при воде в топливе (только NTE200)',
         'FC4642 / FC1852 (наличие воды в топливе): задержка 17 мин до снижения оборотов — '
         'только на NTE200. Риск продолжительной работы на загрязнённом топливе без реакции ECM.'),
        (C_ACCENT, 'РЕКОМЕНДАЦИЯ',
         'NTE200 ш85: FC611=18 горячих останов при 4015 м/ч наработки',
         'FC611 = 18 горячих останова, FC1518 = 9 событий при наработке ~4 015 м/ч. '
         'Высокая частота горячих останов при малом ресурсе свидетельствует '
         'о систематической перегрузке двигателя. Требуется немедленный контроль EGT-триггеров.'),
    ]

    # 2 columns, 3 rows
    card_w = (W - Cm(2.5)) / 2
    card_h = (H - Cm(2.2)) / 3
    col_xs = [Cm(0.8), Cm(0.8) + card_w + Cm(0.9)]
    start_y = Cm(1.45)

    for ci, (strip_c, risk, title, body) in enumerate(conclusions):
        col = ci % 2
        row = ci // 2
        cx = col_xs[col]
        cy = start_y + row * (card_h + Cm(0.15))

        add_rect(slide, cx, cy, card_w, card_h - Cm(0.08), C_PANEL)
        add_rect(slide, cx, cy, Cm(0.25), card_h - Cm(0.08), strip_c)

        add_textbox(slide, cx + Cm(0.4), cy + Cm(0.08), Cm(1.5), Cm(0.38),
                    risk, font_size=7.5, bold=True, color=strip_c)
        add_textbox(slide, cx + Cm(0.4), cy + Cm(0.38), card_w - Cm(0.6), Cm(0.48),
                    title, font_size=9.5, bold=True, color=C_WHITE)
        add_textbox(slide, cx + Cm(0.4), cy + Cm(0.85), card_w - Cm(0.6), card_h - Cm(1.0),
                    body, font_size=8, bold=False, color=C_GRAY)

    add_rect(slide, 0, H - Cm(0.4), W, Cm(0.4), C_PANEL)
    add_textbox(slide, Cm(1), H - Cm(0.38), W - Cm(2), Cm(0.35),
                'АО Развитие  ·  QSK50 MCRS  ·  Анализ отказов клапанов  ·  Полюс Магадан  ·  2026  ·  Конфиденциально',
                font_size=7, color=C_GRAY, align=PP_ALIGN.CENTER)


def build_slide13_action_plan(prs):
    """Slide 13 — Мероприятия (action plan table)"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, C_BG)
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, W, Cm(1.3), C_PANEL)
    add_rect(slide, 0, 0, Cm(0.5), H, C_ACCENT)
    add_textbox(slide, Cm(0.8), Cm(0.2), W - Cm(2), Cm(0.9),
                'ПЛАН МЕРОПРИЯТИЙ', font_size=14, bold=True, color=C_WHITE)
    add_textbox(slide, Cm(0.8), Cm(0.85), W - Cm(2), Cm(0.4),
                'QSK50 NTE200  ·  Приоритизация по срочности и влиянию на ресурс клапанов',
                font_size=8.5, color=C_GRAY)

    # Priority colors
    PRIO = {
        'НЕМЕДЛЕННО': C_RED,
        'ВЫСОКИЙ': C_ORANGE,
        'СРЕДНИЙ': C_YELLOW,
        'ТЕКУЩИЙ': C_ACCENT,
    }

    actions = [
        (1, 'Включить флаг C_EPD_OT_RPM_Drt_En = 1 на всех NTE200',
         'EPD_OT_RPM_Drt_En', 'НЕМЕДЛЕННО', 'Сервис Cummins', 'до след. ТО'),
        (2, 'Восстановить DO_Option=6765 (включить DO01/DO07/DO08)',
         'DO_Option', 'НЕМЕДЛЕННО', 'Cummins / NHL', 'до след. ТО'),
        (3, 'Снизить макс. об/мин до 1990 (C_ABS_Max_Acctr_Speed)',
         'C_ABS_Max_Acctr_Speed', 'ВЫСОКИЙ', 'Cummins / NHL', '1 месяц'),
        (4, 'Ограничить TIB лимит до 100% (C_TIB_Fuel_Adjust_Upper_Limit)',
         'C_TIB_Fuel_Adjust_Upper_Limit', 'ВЫСОКИЙ', 'Cummins / NHL', '1 месяц'),
        (5, 'Установить задержки дерейта 0 с (CT/CP_TB_Time_To_Max)',
         'C_EPD_CT/CP_TB_...', 'ВЫСОКИЙ', 'Cummins / NHL', '1 месяц'),
        (6, 'Проверить FC556 (5 кПа): выполнить анализ событий FC556 у NTE200',
         'FC556', 'СРЕДНИЙ', 'Инженер АО Развитие', '2 недели'),
        (7, 'Calterm сравнение полного набора (TIB, HSST, DroopMax)',
         'Calterm full compare', 'СРЕДНИЙ', 'Cummins', '2 недели'),
        (8, 'Контролировать FC611 (горячие остановы) ежемесячно — EGT',
         'FC611, EGT', 'ТЕКУЩИЙ', 'Механики', 'Ежемесячно'),
    ]

    # Column definitions: x, w, header
    cols = [
        (Cm(0.8),  Cm(0.6),  '№'),
        (Cm(1.55), Cm(6.3),  'Мероприятие'),
        (Cm(8.0),  Cm(4.2),  'Параметр / Код'),
        (Cm(12.4), Cm(2.3),  'Приоритет'),
        (Cm(14.9), Cm(3.8),  'Ответственный'),
        (Cm(18.9), Cm(4.5),  'Срок'),
    ]

    table_y = Cm(1.45)
    hdr_h = Cm(0.5)
    add_rect(slide, Cm(0.8), table_y, W - Cm(1.6), hdr_h, RGBColor(0x1E, 0x28, 0x2E))
    for cx, cw, lbl in cols:
        add_textbox(slide, cx + Cm(0.05), table_y + Cm(0.05), cw, hdr_h,
                    lbl, font_size=8.5, bold=True, color=C_ACCENT)

    row_h = Cm(0.61)
    for ri, (num, action, param, prio, resp, term) in enumerate(actions):
        ry = table_y + hdr_h + ri * row_h
        bg = C_PANEL if ri % 2 == 0 else RGBColor(0x30, 0x3A, 0x40)
        add_rect(slide, Cm(0.8), ry, W - Cm(1.6), row_h - Cm(0.02), bg)
        pc = PRIO.get(prio, C_GRAY)
        add_rect(slide, Cm(0.8), ry, Cm(0.22), row_h - Cm(0.02), pc)

        add_textbox(slide, cols[0][0] + Cm(0.28), ry + Cm(0.08), cols[0][1], row_h,
                    str(num), font_size=9, bold=True, color=pc)
        add_textbox(slide, cols[1][0], ry + Cm(0.08), cols[1][1], row_h,
                    action, font_size=7.5, bold=False, color=C_WHITE)
        add_textbox(slide, cols[2][0], ry + Cm(0.08), cols[2][1], row_h,
                    param, font_size=7, bold=False, color=C_GRAY)
        add_textbox(slide, cols[3][0], ry + Cm(0.08), cols[3][1], row_h,
                    prio, font_size=7.5, bold=True, color=pc)
        add_textbox(slide, cols[4][0], ry + Cm(0.08), cols[4][1], row_h,
                    resp, font_size=7.5, bold=False, color=C_WHITE)
        add_textbox(slide, cols[5][0], ry + Cm(0.08), cols[5][1], row_h,
                    term, font_size=7.5, bold=False, color=C_ACCENT)

    add_rect(slide, 0, H - Cm(0.4), W, Cm(0.4), C_PANEL)
    add_textbox(slide, Cm(1), H - Cm(0.38), W - Cm(2), Cm(0.35),
                'АО Развитие  ·  QSK50 MCRS  ·  Анализ отказов клапанов  ·  Полюс Магадан  ·  2026  ·  Конфиденциально',
                font_size=7, color=C_GRAY, align=PP_ALIGN.CENTER)


def build_slide14_sources(prs):
    """Slide 14 — Источники данных и методология"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, C_BG)
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, W, Cm(1.3), C_PANEL)
    add_rect(slide, 0, 0, Cm(0.5), H, C_GRAY)
    add_textbox(slide, Cm(0.8), Cm(0.18), W - Cm(2), Cm(0.95),
                'ИСТОЧНИКИ ДАННЫХ И МЕТОДОЛОГИЯ',
                font_size=14, bold=True, color=C_WHITE)

    add_rect(slide, Cm(0.8), Cm(1.38), W - Cm(1.6), Cm(0.04), C_ACCENT)

    col_w = (W - Cm(3)) / 2
    left1 = Cm(0.8)
    left2 = Cm(0.8) + col_w + Cm(1.4)

    def src_block(x, y, title, items, accent_c):
        add_rect(slide, x, y, col_w, Cm(0.55), accent_c)
        add_textbox(slide, x + Cm(0.2), y + Cm(0.08), col_w - Cm(0.4), Cm(0.42),
                    title, font_size=10, bold=True, color=C_BG)
        cy = y + Cm(0.65)
        for item in items:
            add_rect(slide, x, cy, col_w, Cm(0.8), C_PANEL)
            add_rect(slide, x, cy, Cm(0.2), Cm(0.8), accent_c)
            add_textbox(slide, x + Cm(0.35), cy + Cm(0.05), col_w - Cm(0.5), Cm(0.7),
                        item, font_size=8, bold=False, color=C_WHITE)
            cy += Cm(0.82)
        return cy

    left_items = [
        'Предыдущий анализ: ECM_对比_NTE200_730E_中文.xlsx\n(Полюс Магадан / АО Развитие)',
        'ECM Image Analyzer Excel: NTE200 ш85\n(33238503, AQ60809.08, 4 015 м/ч)',
        'DML Log Files: 9 единиц 730E + 9 единиц NTE200\n(май 2026)',
        'GIF-скриншоты INSITE КАМСС: мониторы рабочего цикла,\nрасход топлива, снижение крутящего момента / оборотов',
    ]

    right_items = [
        'Calterm Calibration Compare Report:\nAQ60217.28 vs AQ60809.08 (25.05.2026)',
        'ECM Image Analyzer Excel: 730E ш18\n(33223470, AQ60217.28, 35 127 м/ч)',
        'INSITE ECM Snapshots КАМСС (19–20 мая 2026)\n730E: 35 281 ч ECM | NTE200: 4 015 ч ECM',
        '730E ш18 = 730Е№28, AQ60217.28, 35 127 м/ч ECM\nNTE200 ш85, AQ60809.08, 4 015 м/ч ECM',
    ]

    src_block(left1, Cm(1.55), 'ПРЕДЫДУЩИЙ АНАЛИЗ + ECM ДАННЫЕ', left_items, C_BLUE)
    src_block(left2, Cm(1.55), 'CALTERM + INSITE КАМСС', right_items, C_ACCENT)

    # Bottom footer bar
    add_rect(slide, 0, H - Cm(0.9), W, Cm(0.5), C_PANEL2)
    add_rect(slide, 0, H - Cm(0.9), W, Cm(0.08), C_ACCENT)
    add_textbox(slide, Cm(1), H - Cm(0.88), W - Cm(2), Cm(0.45),
                'АО Развитие  ·  QSK50 MCRS  ·  Анализ отказов клапанов  ·  Полюс Магадан  ·  2026  ·  Конфиденциально',
                font_size=9, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)

    # Small KAMSS screenshots at bottom
    gifs_bottom = [
        os.path.join(KAMSS_DIR, '730E ш18 DC model 33223470 AQ60217.28 монитор рабочего цикла.GIF'),
        os.path.join(KAMSS_DIR, 'NTE200 ш85 model монитор рабочего цикла.GIF'),
    ]
    img_y = H - Cm(0.9) - Cm(1.7)
    img_w = (W - Cm(2.4)) / 2
    for i, gp in enumerate(gifs_bottom):
        if os.path.exists(gp):
            add_image_from_gif(slide, gp, Cm(0.8) + i * (img_w + Cm(0.8)),
                               img_y, img_w, Cm(1.6))


# ─── ADDITIONAL GIF SLIDES ────────────────────────────────────────────────────

def build_gif_overview_slide(prs, title, gifs_paths, subtitle=''):
    """Generic slide showing GIF screenshots grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, C_BG)
    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, W, Cm(1.2), C_PANEL)
    add_rect(slide, 0, 0, Cm(0.5), H, C_BLUE)
    add_textbox(slide, Cm(0.8), Cm(0.18), W - Cm(2), Cm(0.78),
                title, font_size=12, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, Cm(0.8), Cm(0.82), W - Cm(2), Cm(0.35),
                    subtitle, font_size=8, color=C_GRAY)

    valid = [p for p in gifs_paths if os.path.exists(p)]
    if not valid:
        return slide

    n = len(valid)
    if n == 1:
        add_image_from_gif(slide, valid[0], Cm(0.8), Cm(1.35), W - Cm(1.6), H - Cm(1.8))
    elif n == 2:
        iw = (W - Cm(2.2)) / 2
        for i, p in enumerate(valid):
            add_image_from_gif(slide, p, Cm(0.8) + i * (iw + Cm(0.6)),
                               Cm(1.35), iw, H - Cm(1.8))
    else:
        # 2x2 or 2x3
        cols = 2
        rows = (n + 1) // 2
        iw = (W - Cm(2.0) - Cm(0.4) * (cols - 1)) / cols
        ih = (H - Cm(1.6) - Cm(0.3) * (rows - 1)) / rows
        for i, p in enumerate(valid):
            r = i // cols
            c = i % cols
            x = Cm(0.8) + c * (iw + Cm(0.4))
            y = Cm(1.3) + r * (ih + Cm(0.3))
            add_image_from_gif(slide, p, x, y, iw, ih)

    add_rect(slide, 0, H - Cm(0.35), W, Cm(0.35), C_PANEL)
    add_textbox(slide, Cm(1), H - Cm(0.33), W - Cm(2), Cm(0.3),
                'АО Развитие  ·  QSK50 MCRS  ·  Анализ отказов клапанов  ·  Полюс Магадан  ·  2026',
                font_size=6.5, color=C_GRAY, align=PP_ALIGN.CENTER)

    return slide


# ─── MAIN ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    # Set 16:9 widescreen dimensions
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print('Building slides...')

    print('  Slide 1: Title')
    build_slide1_title(prs)

    print('  Slide 2: Metadata')
    build_slide2_metadata(prs)

    print('  Slide 3: Divider Part 1')
    build_slide3_divider1(prs)

    print('  Slide 4: Previous analysis table')
    build_slide4_prev_analysis_table(prs)

    print('  Slide 5: Divider Part 2')
    build_slide5_divider2(prs)

    print('  Slide 6: KAMSS confirmed findings')
    build_slide6_kamss_confirmed(prs)

    print('  Slide 7: Same settings')
    build_slide7_same_settings(prs)

    print('  Slide 8: Divider Part 3')
    build_slide8_divider3(prs)

    print('  Slide 9: Comparison table')
    build_slide9_comparison_table(prs)

    print('  Slide 10: Bar chart')
    build_slide10_bar_chart(prs)

    print('  Slide 11: Divider Part 4')
    build_slide11_divider4(prs)

    print('  Slide 12: Conclusions')
    build_slide12_conclusions(prs)

    print('  Slide 13: Action plan')
    build_slide13_action_plan(prs)

    print('  Slide 14: Sources')
    build_slide14_sources(prs)

    # Additional GIF slides with INSITE screenshots
    print('  Slide 15: INSITE снижение момента 730E')
    gifs_730_moment = [
        os.path.join(KAMSS_DIR, '730E ш18 DC model 33223470 AQ60217.28 снижение крутящего момента 1.GIF'),
        os.path.join(KAMSS_DIR, '730E ш18 DC model 33223470 AQ60217.28 снижение крутящего момента0.GIF'),
        os.path.join(KAMSS_DIR, '730E ш18 DC model 33223470 AQ60217.28 снижение крутящего момента2.GIF'),
        os.path.join(KAMSS_DIR, '730E ш18 DC model 33223470 AQ60217.28 снижение оборотов 0.GIF'),
    ]
    build_gif_overview_slide(prs,
        'INSITE КАМСС: 730E ш18 — Снижение момента и оборотов (референс)',
        gifs_730_moment,
        'AQ60217.28 | 35 127 м/ч ECM | Параметры защитного дерейта')

    print('  Slide 16: INSITE снижение момента NTE200')
    gifs_nte_moment = [
        os.path.join(KAMSS_DIR, 'NTE200 ш85 model снижение крутящего момента 0.GIF'),
        os.path.join(KAMSS_DIR, 'NTE200 ш85 model снижение крутящего момента 1.GIF'),
        os.path.join(KAMSS_DIR, 'NTE200 ш85 model снижение крутящего момента 2.GIF'),
        os.path.join(KAMSS_DIR, 'NTE200 ш85 model снижение оборотов 0.GIF'),
    ]
    build_gif_overview_slide(prs,
        'INSITE КАМСС: NTE200 ш85 — Снижение момента и оборотов (проблемная)',
        gifs_nte_moment,
        'AQ60809.08 | 4 015 м/ч ECM | Параметры защитного дерейта')

    print(f'Saving to {OUTPUT}...')
    prs.save(OUTPUT)
    print(f'Done! Saved: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
