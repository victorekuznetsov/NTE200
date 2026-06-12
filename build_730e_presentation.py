#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DML Log Analysis Presentation Builder — 730E Fleet
Generates: DML_Анализ_LogFiles_730E.pptx
Structure mirrors DML_Анализ_LogFiles_NTE200.pptx
"""

import os
import sys
import tempfile
import warnings
import numpy as np
import pandas as pd

warnings.filterwarnings('ignore')

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.gridspec as gridspec

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Brand colors ─────────────────────────────────────────────────────────────
C_BG    = '#293136'
C_AXES  = '#1e2529'
C_GREEN = '#3EF0AF'
C_WHITE = '#FFFFFF'
C_RED   = '#FF4444'
C_ORANGE= '#FF8C42'
C_YELLOW= '#FFD166'
C_BLUE  = '#4EC9F0'
C_GRAY  = '#667788'
C_DARK2 = '#20272b'

def hex2rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def rgb_color(h):
    r, g, b = hex2rgb(h)
    return RGBColor(r, g, b)

plt.rcParams.update({
    'figure.facecolor': C_BG,
    'axes.facecolor':   C_AXES,
    'axes.edgecolor':   C_GRAY,
    'axes.labelcolor':  C_WHITE,
    'xtick.color':      C_GRAY,
    'ytick.color':      C_GRAY,
    'text.color':       C_WHITE,
    'grid.color':       C_GRAY,
    'grid.alpha':       0.3,
    'legend.facecolor': C_AXES,
    'legend.edgecolor': C_GRAY,
    'font.family':      'DejaVu Sans',
    'font.size':        9,
    'axes.titlesize':   11,
    'axes.labelsize':   9,
})

SLIDE_W = Cm(33.86)
SLIDE_H = Cm(19.05)

LOG_DIR = '/home/user/NTE200/log_files_730e'

FILE_CFG = {
    '24':  (f'{LOG_DIR}/730E №24 21.05.2026.csv',  'utf-8',  'ru', '21.05.2026'),
    '19':  (f'{LOG_DIR}/730Е №19 16.05.2026.csv',  'cp1251', 'en', '16.05.2026'),
    '23a': (f'{LOG_DIR}/730Е №23 19.05.2026.csv',  'cp1251', 'ru', '19.05.2026'),
    '23b': (f'{LOG_DIR}/730Е №23 20.05.2026.csv',  'cp1251', 'ru', '20.05.2026'),
    '27':  (f'{LOG_DIR}/730Е №27 18.05.2026.csv',  'cp1251', 'ru', '18.05.2026'),
    '28':  (f'{LOG_DIR}/730Е №28 20.05.2026.csv',  'cp1251', 'ru', '20.05.2026'),
    '39a': (f'{LOG_DIR}/730Е №39 16.05.2026.csv',  'cp1251', 'en', '16.05.2026'),
    '39b': (f'{LOG_DIR}/730Е №39 21.05.2026.csv',  'utf-8',  'ru', '21.05.2026'),
    '40':  (f'{LOG_DIR}/730Е №40 20.05.2026.csv',  'cp1251', 'ru', '20.05.2026'),
}

UNIT_LABELS = {
    '24':  '#24',
    '19':  '#19',
    '23a': '#23 (19.05)',
    '23b': '#23 (20.05)',
    '27':  '#27',
    '28':  '#28',
    '39a': '#39 (16.05)',
    '39b': '#39 (21.05)',
    '40':  '#40',
}

RISK_COLOR = {
    'КРИТИЧНО': C_RED,
    'ВНИМАНИЕ': C_ORANGE,
    'НОРМА':    C_GREEN,
}

# ─── Data loading ─────────────────────────────────────────────────────────────

def get_col(df, *keywords, exclude=None):
    for col in df.columns:
        cl = col.lower()
        if all(k.lower() in cl for k in keywords):
            if exclude and any(e.lower() in cl for e in exclude):
                continue
            return col
    return None

def to_sec(s):
    try:
        p = str(s).split(':')
        return float(p[0])*3600 + float(p[1])*60 + float(p[2])
    except Exception:
        return np.nan

def clean_cyl(v):
    v = np.array(v, dtype=float)
    return np.where((v < 50) | (v > 750), np.nan, v)

def load_unit(uid):
    path, enc, lang, date = FILE_CFG[uid]
    df = pd.read_csv(path, skiprows=26, encoding=enc, decimal=',',
                     on_bad_lines='skip', low_memory=False)

    time_col = 'Time' if lang == 'en' else 'Время'
    if time_col not in df.columns:
        time_col = df.columns[1]
    t = df[time_col].apply(to_sec)
    sec = (t - t.dropna().iloc[0]).values

    if lang == 'en':
        rpm_col  = get_col(df, 'Engine Speed (RPM)', 'Identifier0')
        load_col = get_col(df, 'Percent Load', 'Identifier0')
        oil_col  = get_col(df, 'Engine Oil Temperature', 'Identifier0', exclude=['Voltage'])
        boost_col= get_col(df, 'Intake Manifold Pressure', 'Identifier0')
        oil_p_col= get_col(df, 'Engine Oil Pressure', 'Identifier0')
        egt_L_col= get_col(df, 'Average Exhaust Temperature - Left Bank', 'Identifier144')
        egt_R_col= get_col(df, 'Average Exhaust Temperature - Right Bank', 'Identifier1')
        cyl_col_fn = lambda n, sfx: get_col(df, f'Exhaust Temperature Sensor Cylinder {n} (', sfx, exclude=['Voltage'])
        sfx_fn = lambda n: 'Identifier144' if n % 2 == 1 else 'Identifier1'
    else:
        rpm_col  = get_col(df, 'Частота вращения двигателя', 'Идентификатор0')
        if not rpm_col:
            rpm_col = get_col(df, 'Частота вращения двигателя', 'Идентификатор1')
        load_col = get_col(df, 'Относительная нагрузка', 'Идентификатор0')
        oil_col  = get_col(df, 'Температура масла', 'Идентификатор0')
        boost_col= get_col(df, 'Давление во впускном коллекторе', 'Идентификатор0')
        oil_p_col= get_col(df, 'Давление масла (кПа)', 'Идентификатор0')
        egt_L_col= get_col(df, 'левый ряд', 'Идентификатор144')
        egt_R_col= get_col(df, 'правый ряд', 'Идентификатор1')
        cyl_col_fn = lambda n, sfx: get_col(df, f'цилиндра {n} (°C)', sfx, exclude=['Напряжение', 'Voltage'])
        sfx_fn = lambda n: 'Идентификатор144' if n % 2 == 1 else 'Идентификатор1'

    def gs(col, lo=0, hi=9999):
        if not col or col not in df.columns:
            return np.full(len(df), np.nan)
        v = df[col].values
        if v.dtype == object:
            v = pd.to_numeric(pd.Series(v).str.replace(',', '.'), errors='coerce').values
        v = v.astype(float)
        return np.where((v < lo) | (v > hi), np.nan, v)

    cyl_egt = {}
    for n in range(1, 17):
        c = cyl_col_fn(n, sfx_fn(n))
        if c:
            cyl_egt[n] = clean_cyl(gs(c, 0, 9999))

    cyl_avgs, cyl_maxes = {}, {}
    for n, v in cyl_egt.items():
        finite = v[np.isfinite(v)]
        if len(finite) > 10:
            cyl_avgs[n] = np.mean(finite)
            cyl_maxes[n] = np.max(finite)

    rpm_ts   = gs(rpm_col, 100, 3000)
    load_ts  = gs(load_col, 0, 110)
    oil_ts   = gs(oil_col, 30, 200)
    boost_ts = gs(boost_col, 50, 500)
    oil_p_ts = gs(oil_p_col, 50, 2000)
    egt_L_ts = gs(egt_L_col, 80, 750)
    egt_R_ts = gs(egt_R_col, 80, 750)

    if cyl_avgs:
        egt_avg = np.mean(list(cyl_avgs.values()))
        egt_max = np.max(list(cyl_maxes.values()))
        hot_cyl = max(cyl_maxes, key=cyl_maxes.get)
    else:
        egt_avg = egt_max = np.nan
        hot_cyl = '?'

    nm = lambda v: np.nanmean(v) if np.any(np.isfinite(v)) else np.nan
    nx = lambda v: np.nanmax(v)  if np.any(np.isfinite(v)) else np.nan

    egt_fault = len(cyl_avgs) < 12 or (not np.isnan(egt_avg) and egt_avg < 150)

    if not np.isnan(egt_max) and egt_max >= 500:
        risk = 'КРИТИЧНО'
    elif not np.isnan(egt_max) and egt_max >= 450:
        risk = 'ВНИМАНИЕ'
    elif egt_fault or (not np.isnan(nx(oil_ts)) and nx(oil_ts) >= 105):
        risk = 'ВНИМАНИЕ'
    else:
        risk = 'НОРМА'

    return {
        'date': date, 'dur': (t.dropna().max() - t.dropna().min()) / 60,
        'rows': len(df), 'lang': lang,
        'sec': sec, 'rpm_ts': rpm_ts, 'load_ts': load_ts,
        'oil_ts': oil_ts, 'boost_ts': boost_ts, 'oil_p_ts': oil_p_ts,
        'egt_L_ts': egt_L_ts, 'egt_R_ts': egt_R_ts,
        'cyl_egt': cyl_egt, 'cyl_avgs': cyl_avgs, 'cyl_maxes': cyl_maxes,
        'rpm_avg': nm(rpm_ts), 'rpm_max': nx(rpm_ts), 'load': nm(load_ts),
        'egt_avg': egt_avg, 'egt_max': egt_max, 'oil_max': nx(oil_ts),
        'boost_avg': nm(boost_ts), 'oil_p_avg': nm(oil_p_ts),
        'hot_cyl': hot_cyl, 'risk': risk, 'egt_fault': egt_fault,
    }


print("Loading 730E DML files...")
DATA = {}
for uid in FILE_CFG:
    print(f"  Loading unit {uid}...")
    try:
        DATA[uid] = load_unit(uid)
    except Exception as e:
        print(f"  WARNING {uid}: {e}")
        DATA[uid] = None
print("All loaded.\n")

# ─── Temp file management ─────────────────────────────────────────────────────
_tmp = []

def save_fig(fig):
    fd, path = tempfile.mkstemp(suffix='.png')
    os.close(fd)
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor=C_BG, edgecolor='none')
    plt.close(fig)
    _tmp.append(path)
    return path

def cleanup():
    for f in _tmp:
        try: os.remove(f)
        except: pass

# ─── PPTX helpers ─────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color(C_BG)

def add_rect(slide, x, y, w, h, fill=None, line=None, lw=0):
    sh = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb_color(fill)
    else:
        sh.fill.background()
    if line and lw > 0:
        sh.line.color.rgb = rgb_color(line)
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def add_text(slide, text, x, y, w, h,
             fs=12, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb_color(color)
    return tb

def add_img(slide, path, x, y, w, h):
    slide.shapes.add_picture(path, Cm(x), Cm(y), Cm(w), Cm(h))

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 33.86, 2.0, fill=C_DARK2)
    add_rect(slide, 0, 1.88, 33.86, 0.12, fill=C_GREEN)
    add_text(slide, title, 0.5, 0.1, 32, 1.1, fs=18, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.5, 1.2, 32, 0.7, fs=9.5, color=C_GRAY)

def footer(slide, txt='Анализ DML | Парк 730E | QSK50 MCRS | май 2026'):
    add_rect(slide, 0, 18.55, 33.86, 0.5, fill=C_DARK2)
    add_text(slide, txt, 0.5, 18.6, 33, 0.4, fs=7, color=C_GRAY, align=PP_ALIGN.CENTER)

# rolling smooth helper
def smooth(v, w=5):
    s = pd.Series(v)
    return s.rolling(w, center=True, min_periods=1).mean().values

# ─── Slide 1: Title ───────────────────────────────────────────────────────────
def slide1(prs):
    sl = blank_slide(prs)
    set_bg(sl)
    add_rect(sl, 0, 0, 33.86, 19.05, fill=C_BG)
    add_rect(sl, 0, 0, 8, 19.05, fill=C_DARK2)
    add_rect(sl, 7.8, 0, 0.22, 19.05, fill=C_GREEN)
    add_rect(sl, 1, 2, 5, 0.08, fill=C_GREEN)
    add_text(sl, 'DML', 1, 2.3, 6, 3, fs=60, bold=True, color=C_GREEN)
    add_text(sl, 'LOG-FILES', 1, 5.5, 6, 1.5, fs=22, color=C_WHITE)
    add_text(sl, 'ANALYSIS', 1, 6.8, 6, 1.5, fs=22, color=C_WHITE)

    add_text(sl, 'АНАЛИЗ DML ЛОГ-ФАЙЛОВ', 9, 3, 23, 2, fs=32, bold=True, color=C_WHITE)
    add_text(sl, 'ПАРК 730E', 9, 5.3, 23, 2, fs=32, bold=True, color=C_GREEN)
    add_rect(sl, 9, 8.0, 23, 0.08, fill=C_GRAY)
    add_text(sl, '9 агрегатов  |  май 2026  |  Cummins QSK50 MCRS', 9, 8.3, 23, 1, fs=13, color=C_GRAY)

    boxes = [('9', 'агрегатов'), ('3', 'критических'), ('QSK50', 'MCRS двигатель'), ('DML', 'реал-тайм лог')]
    for i, (val, lbl) in enumerate(boxes):
        bx = 9 + i * 5.8
        add_rect(sl, bx, 10, 5.3, 2.8, fill=C_DARK2, line=C_GREEN, lw=1)
        add_text(sl, val, bx+0.2, 10.2, 5, 1.4, fs=26, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, bx+0.2, 11.5, 5, 0.8, fs=10, color=C_GRAY, align=PP_ALIGN.CENTER)

    add_text(sl, '730E  •  Cummins QSK50 MCRS  •  16 цилиндров  •  ECM0 / ECM1 / ECM144',
             9, 14, 23, 1, fs=10, color=C_GRAY, italic=True)
    add_text(sl, 'Данные: INSITE DML  •  16–21 мая 2026', 9, 15, 23, 1, fs=10, color=C_GRAY, italic=True)
    footer(sl)

# ─── Slide 2: Stat cards ──────────────────────────────────────────────────────
def slide2(prs):
    sl = blank_slide(prs)
    set_bg(sl)
    header_bar(sl, 'ОБЗОР ПАРКА 730E — СТАТУС АГРЕГАТОВ',
               'Ключевые параметры DML-сессий | 9 агрегатов | май 2026')

    order = sorted(DATA.keys(), key=lambda u: (
        0 if DATA[u] and DATA[u]['risk'] == 'КРИТИЧНО' else
        1 if DATA[u] and DATA[u]['risk'] == 'ВНИМАНИЕ' else 2))

    cols = [0.3, 11.5, 22.7]
    rows = [2.3, 8.0, 13.7]

    for i, uid in enumerate(order[:9]):
        d = DATA[uid]
        row, col_i = divmod(i, 3)
        x, y = cols[col_i], rows[row]
        risk = d['risk'] if d else 'НОРМА'
        rc = RISK_COLOR[risk]

        add_rect(sl, x, y, 10.8, 4.7, fill=C_DARK2, line=rc, lw=1.5)
        # risk badge
        add_rect(sl, x+0.2, y+0.2, 3.5, 0.65, fill=rc)
        add_text(sl, risk, x+0.2, y+0.15, 3.5, 0.75, fs=8, bold=True,
                 color=C_BG, align=PP_ALIGN.CENTER)
        # unit label
        add_text(sl, f'730E {UNIT_LABELS[uid]}', x+4.0, y+0.2, 6.5, 0.7,
                 fs=12, bold=True, color=rc)
        add_text(sl, d['date'], x+4.0, y+0.9, 6.5, 0.5, fs=9, color=C_GRAY)

        if d:
            kpis = [
                ('ЕГТ avg', f"{d['egt_avg']:.0f}°C" if not np.isnan(d['egt_avg']) else 'N/A', rc if (not np.isnan(d.get('egt_avg', np.nan)) and d['egt_avg'] > 450) else C_WHITE),
                ('ЕГТ max', f"{d['egt_max']:.0f}°C" if not np.isnan(d['egt_max']) else 'N/A', rc if (not np.isnan(d.get('egt_max', np.nan)) and d['egt_max'] >= 500) else C_WHITE),
                ('RPM avg', f"{d['rpm_avg']:.0f}", C_BLUE),
                ('Нагрузка', f"{d['load']:.0f}%" if not np.isnan(d['load']) else 'N/A', C_YELLOW),
                ('Т° масла', f"{d['oil_max']:.0f}°C" if not np.isnan(d['oil_max']) else 'N/A',
                 C_RED if (not np.isnan(d.get('oil_max', np.nan)) and d['oil_max'] >= 110) else
                 C_ORANGE if (not np.isnan(d.get('oil_max', np.nan)) and d['oil_max'] >= 105) else C_GREEN),
                ('Длит.', f"{d['dur']:.1f} мин", C_GRAY),
            ]
            for j, (lbl, val, vc) in enumerate(kpis):
                ky = y + 1.1 + j * 0.57
                add_text(sl, lbl + ':', x+0.3, ky, 3.5, 0.7, fs=9, color=C_GRAY)
                add_text(sl, val, x+4.2, ky, 6.0, 0.7, fs=10, bold=True, color=vc, align=PP_ALIGN.RIGHT)

    footer(sl)

# ─── Slide 3: EGT comparison bar chart ───────────────────────────────────────
def slide3(prs):
    uids = [u for u in DATA if DATA[u] and not np.isnan(DATA[u]['egt_avg']) and DATA[u]['egt_avg'] > 150]
    uids_sorted = sorted(uids, key=lambda u: DATA[u]['egt_max'], reverse=True)

    labels = [UNIT_LABELS[u] for u in uids_sorted]
    avgs   = [DATA[u]['egt_avg'] for u in uids_sorted]
    maxs   = [DATA[u]['egt_max'] for u in uids_sorted]
    colors = [RISK_COLOR[DATA[u]['risk']] for u in uids_sorted]

    x = np.arange(len(labels))
    fig, ax = plt.subplots(figsize=(14, 7))
    fig.patch.set_facecolor(C_BG)
    ax.set_facecolor(C_AXES)

    bar_w = 0.35
    b1 = ax.bar(x - bar_w/2, avgs, bar_w, label='Avg ЕГТ (°C)', color=C_BLUE, alpha=0.85)
    b2 = ax.bar(x + bar_w/2, maxs, bar_w, label='Max ЕГТ (°C)', color=[c for c in colors], alpha=0.9)

    ax.axhline(450, color=C_ORANGE, lw=1.5, ls='--', alpha=0.8, label='450°C (внимание)')
    ax.axhline(500, color=C_RED, lw=1.5, ls='--', alpha=0.8, label='500°C (критично)')
    ax.set_ylim(0, 650)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, rotation=0, fontsize=9, color=C_WHITE)
    ax.set_ylabel('ЕГТ (°C)', color=C_WHITE, fontsize=10)
    ax.set_title('Средняя и максимальная ЕГТ по агрегатам (из данных цилиндров)', color=C_WHITE, fontsize=11, fontweight='bold')
    ax.legend(fontsize=9, framealpha=0.4, facecolor=C_AXES, edgecolor=C_GRAY, labelcolor=C_WHITE, loc='upper right')
    ax.grid(True, alpha=0.3, axis='y')
    ax.tick_params(colors=C_WHITE, labelsize=9)
    for sp in ax.spines.values(): sp.set_edgecolor(C_GRAY)

    for b, v in zip(b1, avgs):
        ax.text(b.get_x()+b.get_width()/2, v+5, f'{v:.0f}', ha='center', va='bottom', fontsize=8, color=C_WHITE)
    for b, v, c in zip(b2, maxs, colors):
        ax.text(b.get_x()+b.get_width()/2, v+5, f'{v:.0f}', ha='center', va='bottom', fontsize=8, color=c, fontweight='bold')

    plt.tight_layout()
    img = save_fig(fig)

    sl = blank_slide(prs)
    set_bg(sl)
    header_bar(sl, 'ЕГТ — СРАВНЕНИЕ ПО АГРЕГАТАМ',
               'Средняя и максимальная температура выхлопных газов (из данных 16 цилиндров)')
    add_img(sl, img, 0.5, 2.2, 32.8, 15.7)
    footer(sl)

# ─── Slide 4: Multi-parameter grid ───────────────────────────────────────────
def slide4(prs):
    uids = [u for u in DATA if DATA[u]]
    labels = [UNIT_LABELS[u] for u in uids]
    x = np.arange(len(uids))

    def bar_vals(key, fallback=0):
        return [DATA[u][key] if DATA[u] and not np.isnan(DATA[u][key]) else fallback for u in uids]

    fig = plt.figure(figsize=(16, 10))
    fig.patch.set_facecolor(C_BG)
    gs = gridspec.GridSpec(2, 3, figure=fig, hspace=0.45, wspace=0.35)
    axes = [fig.add_subplot(gs[r, c]) for r in range(2) for c in range(3)]

    panels = [
        ('load', 'Нагрузка (%)', C_YELLOW, None, None),
        ('rpm_max', 'Макс. RPM', C_BLUE, None, None),
        ('egt_avg', 'Avg ЕГТ (°C)', C_ORANGE, 450, 500),
        ('oil_max', 'Макс. Т° масла (°C)', '#90EE90', 105, 110),
        ('boost_avg', 'Среднее давление наддува (кПа)', C_BLUE, None, None),
        ('oil_p_avg', 'Ср. давление масла (кПа)', C_GREEN, None, None),
    ]

    for ax, (key, title, color, warn, crit) in zip(axes, panels):
        ax.set_facecolor(C_AXES)
        vals = bar_vals(key)
        bar_colors = []
        for v in vals:
            if crit and v >= crit: bar_colors.append(C_RED)
            elif warn and v >= warn: bar_colors.append(C_ORANGE)
            else: bar_colors.append(color)
        bars = ax.barh(labels, vals, color=bar_colors, alpha=0.85)
        if warn: ax.axvline(warn, color=C_ORANGE, lw=1.2, ls='--', alpha=0.8)
        if crit: ax.axvline(crit, color=C_RED,    lw=1.2, ls='--', alpha=0.8)
        ax.set_title(title, color=C_WHITE, fontsize=9.5, fontweight='bold')
        ax.tick_params(colors=C_WHITE, labelsize=7.5)
        ax.grid(True, alpha=0.3, axis='x')
        for sp in ax.spines.values(): sp.set_edgecolor(C_GRAY)
        for bar, v in zip(bars, vals):
            ax.text(v + max(vals)*0.01, bar.get_y()+bar.get_height()/2,
                    f'{v:.0f}', va='center', fontsize=7.5, color=C_WHITE)

    fig.suptitle('Ключевые параметры — все агрегаты 730E', color=C_WHITE, fontsize=13, fontweight='bold', y=1.01)
    img = save_fig(fig)

    sl = blank_slide(prs)
    set_bg(sl)
    header_bar(sl, 'СВОДКА ПАРАМЕТРОВ — ВСЕ АГРЕГАТЫ',
               'Нагрузка, RPM, ЕГТ, температура масла, наддув, давление масла')
    add_img(sl, img, 0.3, 2.1, 33.2, 16.3)
    footer(sl)

# ─── Slide 5: Critical unit deep-dive ────────────────────────────────────────
def slide5(prs):
    # Unit #24 — highest EGT avg 515°C
    uid = '24'
    d = DATA[uid]
    if not d: return

    fig, axes = plt.subplots(3, 1, figsize=(15, 10))
    fig.patch.set_facecolor(C_BG)
    fig.suptitle(f'730E {UNIT_LABELS[uid]}  |  {d["date"]}  |  Детальный анализ (КРИТИЧНО — max ЕГТ: {d["egt_max"]:.0f}°C)',
                 color=C_WHITE, fontsize=11, fontweight='bold', y=1.0)

    sec = d['sec']

    for ax, (key, label, color, ylim, ref) in zip(axes, [
        ('rpm_ts',   'RPM двигателя',    C_BLUE,   (1000, 2200), 2100),
        ('load_ts',  'Нагрузка (%)',      C_YELLOW, (0, 115),     None),
        ('oil_ts',   'Т° масла (°C)',     '#90EE90',(50, 130),    110),
    ]):
        ax.set_facecolor(C_AXES)
        v = smooth(d[key])
        ax.fill_between(sec, v, alpha=0.15, color=color)
        ax.plot(sec, v, color=color, lw=1.8, alpha=0.9)
        if ref:
            ax.axhline(ref, color=C_RED, lw=1.2, ls='--', alpha=0.7)
        ax.set_ylim(ylim)
        ax.set_ylabel(label, color=C_WHITE, fontsize=9)
        ax.grid(True, alpha=0.3)
        ax.tick_params(colors=C_WHITE, labelsize=8)
        for sp in ax.spines.values(): sp.set_edgecolor(C_GRAY)
        if ax != axes[-1]: ax.set_xticklabels([])

    axes[-1].set_xlabel('Время (сек)', color=C_WHITE, fontsize=9)
    plt.tight_layout()
    img = save_fig(fig)

    sl = blank_slide(prs)
    set_bg(sl)
    header_bar(sl, f'КРИТИЧНО: 730E {UNIT_LABELS[uid]}  |  {d["date"]}  |  ЕГТ avg {d["egt_avg"]:.0f}°C  max {d["egt_max"]:.0f}°C',
               f'Длит. {d["dur"]:.1f} мин  |  RPM avg {d["rpm_avg"]:.0f}  |  Нагрузка {d["load"]:.0f}%  |  Цил.{d["hot_cyl"]} — {d["egt_max"]:.0f}°C')
    add_img(sl, img, 0.3, 2.1, 33.2, 15.5)

    add_rect(sl, 0, 17.7, 33.86, 1.35, fill=C_DARK2, line=C_RED, lw=1.2)
    add_text(sl, f'ВЫВОД: Агрегат 730E {UNIT_LABELS[uid]} — ЕГТ цил.{d["hot_cyl"]} достигает {d["egt_max"]:.0f}°C. Нагрузка {d["load"]:.0f}%  |  Масло {d["oil_max"]:.0f}°C.',
             0.5, 17.75, 32.5, 1.25, fs=10, bold=True, color=C_RED)
    footer(sl)

# ─── Slide 6: Per-cylinder EGT bars for critical units ───────────────────────
def slide6(prs):
    # Unit 24 per-cyl EGT bars
    uid = '24'
    d = DATA[uid]
    if not d: return

    cyls = sorted(d['cyl_avgs'].keys())
    avgs = [d['cyl_avgs'][n] for n in cyls]
    maxs = [d['cyl_maxes'][n] for n in cyls]
    bar_colors = [C_RED if m >= 500 else C_ORANGE if m >= 450 else C_BLUE for m in maxs]

    fig, ax = plt.subplots(figsize=(14, 6))
    fig.patch.set_facecolor(C_BG)
    ax.set_facecolor(C_AXES)

    x = np.arange(len(cyls))
    ax.bar(x - 0.2, avgs, 0.38, label='Avg ЕГТ', color=[c+'99' for c in bar_colors], alpha=0.9)
    ax.bar(x + 0.2, maxs, 0.38, label='Max ЕГТ', color=bar_colors, alpha=0.95)
    ax.axhline(450, color=C_ORANGE, lw=1.5, ls='--', alpha=0.8, label='450°C (внимание)')
    ax.axhline(500, color=C_RED, lw=1.5, ls='--', alpha=0.8, label='500°C (критично)')
    ax.set_xticks(x)
    ax.set_xticklabels([f'Ц{n}' for n in cyls], fontsize=9, color=C_WHITE)
    ax.set_ylabel('ЕГТ (°C)', color=C_WHITE, fontsize=10)
    ax.set_ylim(200, 620)
    ax.set_title(f'730E {UNIT_LABELS[uid]} — ЕГТ по цилиндрам  |  {d["date"]}', color=C_WHITE, fontsize=11, fontweight='bold')
    ax.legend(fontsize=9, framealpha=0.4, facecolor=C_AXES, edgecolor=C_GRAY, labelcolor=C_WHITE)
    ax.grid(True, alpha=0.3, axis='y')
    ax.tick_params(colors=C_WHITE, labelsize=9)
    for sp in ax.spines.values(): sp.set_edgecolor(C_GRAY)
    for xi, (a, m) in enumerate(zip(avgs, maxs)):
        ax.text(xi+0.2, m+4, f'{m:.0f}', ha='center', va='bottom', fontsize=8, color=C_WHITE, fontweight='bold')

    plt.tight_layout()
    img = save_fig(fig)

    sl = blank_slide(prs)
    set_bg(sl)
    header_bar(sl, f'730E {UNIT_LABELS[uid]} — ЕГТ ПО ЦИЛИНДРАМ',
               f'{d["date"]}  |  Горячий цил.: Ц{d["hot_cyl"]} ({d["egt_max"]:.0f}°C)')
    add_img(sl, img, 0.3, 2.2, 33.2, 15.5)

    hot_bank = 'В-банк (чётные)' if d['hot_cyl'] % 2 == 0 else 'А-банк (нечётные)'
    add_rect(sl, 0, 17.7, 33.86, 1.35, fill=C_DARK2, line=C_ORANGE, lw=1.2)
    add_text(sl, f'ВЫВОД: Цил.{d["hot_cyl"]} ({hot_bank}) — {d["egt_max"]:.0f}°C максимум. '
             f'Средняя ЕГТ по парку цилиндров: {d["egt_avg"]:.0f}°C.',
             0.5, 17.75, 32.5, 1.25, fs=10, bold=True, color=C_ORANGE)
    footer(sl)

# ─── Slide 7: Unit 28 long-session deep-dive ─────────────────────────────────
def slide7(prs):
    uid = '28'
    d = DATA[uid]
    if not d: return

    fig, axes = plt.subplots(2, 1, figsize=(15, 8))
    fig.patch.set_facecolor(C_BG)
    fig.suptitle(f'730E {UNIT_LABELS[uid]}  |  {d["date"]}  |  27.6 мин  |  max ЕГТ {d["egt_max"]:.0f}°C',
                 color=C_WHITE, fontsize=11, fontweight='bold', y=1.0)

    sec = d['sec']
    cyls_hot = sorted(d['cyl_maxes'], key=d['cyl_maxes'].get, reverse=True)[:6]

    ax0, ax1 = axes

    # RPM + Load
    ax0.set_facecolor(C_AXES)
    ax0t = ax0.twinx()
    ax0.fill_between(sec, smooth(d['rpm_ts']), alpha=0.1, color=C_BLUE)
    ax0.plot(sec, smooth(d['rpm_ts']), color=C_BLUE, lw=1.8, label='RPM')
    ax0t.plot(sec, smooth(d['load_ts']), color=C_YELLOW, lw=1.5, ls='--', alpha=0.8, label='Нагрузка %')
    ax0.set_ylabel('RPM', color=C_BLUE, fontsize=9)
    ax0t.set_ylabel('Нагрузка (%)', color=C_YELLOW, fontsize=9)
    ax0.set_ylim(800, 2200)
    ax0t.set_ylim(0, 120)
    ax0.grid(True, alpha=0.3)
    ax0.tick_params(colors=C_WHITE, labelsize=8)
    ax0t.tick_params(colors=C_YELLOW, labelsize=8)
    for sp in ax0.spines.values(): sp.set_edgecolor(C_GRAY)
    ax0.set_xticklabels([])

    # Hot cylinder EGT time series
    ax1.set_facecolor(C_AXES)
    colors = plt.cm.hot(np.linspace(0.3, 0.9, len(cyls_hot)))
    for cyl, col in zip(cyls_hot, colors):
        if cyl in d['cyl_egt']:
            ax1.plot(sec, smooth(d['cyl_egt'][cyl]), color=col, lw=1.5,
                     label=f'Ц{cyl} ({d["cyl_maxes"][cyl]:.0f}°C)', alpha=0.9)
    ax1.axhline(450, color=C_ORANGE, lw=1.3, ls='--', alpha=0.8)
    ax1.axhline(500, color=C_RED,    lw=1.3, ls='--', alpha=0.8, label='500°C')
    ax1.set_ylim(200, 620)
    ax1.set_ylabel('ЕГТ (°C)', color=C_WHITE, fontsize=9)
    ax1.set_xlabel('Время (сек)', color=C_WHITE, fontsize=9)
    ax1.legend(fontsize=8, ncol=3, loc='upper left', framealpha=0.35,
               facecolor=C_AXES, edgecolor=C_GRAY, labelcolor=C_WHITE)
    ax1.grid(True, alpha=0.3)
    ax1.tick_params(colors=C_WHITE, labelsize=8)
    for sp in ax1.spines.values(): sp.set_edgecolor(C_GRAY)

    plt.tight_layout()
    img = save_fig(fig)

    sl = blank_slide(prs)
    set_bg(sl)
    header_bar(sl, f'730E {UNIT_LABELS[uid]} — ДЛИТЕЛЬНАЯ СЕССИЯ  |  {d["date"]}',
               f'27.6 мин  |  RPM avg {d["rpm_avg"]:.0f}  |  Нагрузка {d["load"]:.0f}%  |  Цил.{d["hot_cyl"]} max {d["egt_max"]:.0f}°C')
    add_img(sl, img, 0.3, 2.1, 33.2, 15.5)

    add_rect(sl, 0, 17.7, 33.86, 1.35, fill=C_DARK2, line=C_ORANGE, lw=1.2)
    add_text(sl, f'ВЫВОД: Наиболее горячий цил.{d["hot_cyl"]} — {d["egt_max"]:.0f}°C. Мониторинг В-банка (чётные): Ц11 и Ц13 — лидеры нагрева.',
             0.5, 17.75, 32.5, 1.25, fs=10, bold=True, color=C_ORANGE)
    footer(sl)

# ─── Slide 8: Sensor fault analysis ─────────────────────────────────────────
def slide8(prs):
    # Compare units with good EGT vs sensor-fault units
    good_uids = [u for u in DATA if DATA[u] and not DATA[u]['egt_fault'] and DATA[u]['egt_avg'] > 200]
    fault_uids = [u for u in DATA if DATA[u] and DATA[u]['egt_fault']]

    fig, axes = plt.subplots(1, 2, figsize=(15, 7))
    fig.patch.set_facecolor(C_BG)
    fig.suptitle('Диагностика: рабочие агрегаты vs. отказ датчиков ЕГТ', color=C_WHITE, fontsize=12, fontweight='bold')

    # Left: EGT profile for good units
    ax0 = axes[0]
    ax0.set_facecolor(C_AXES)
    colors_map = {u: c for u, c in zip(good_uids, plt.cm.plasma(np.linspace(0.2, 0.9, len(good_uids))))}
    for u in good_uids:
        d = DATA[u]
        cyls = sorted(d['cyl_avgs'].keys())
        avgs = [d['cyl_avgs'][n] for n in cyls]
        ax0.plot([f'Ц{n}' for n in cyls], avgs, '-o', lw=1.5, ms=4,
                 color=colors_map[u], label=f'730E {UNIT_LABELS[u]}', alpha=0.85)
    ax0.axhline(450, color=C_ORANGE, lw=1.2, ls='--', alpha=0.8)
    ax0.axhline(500, color=C_RED, lw=1.2, ls='--', alpha=0.8)
    ax0.set_ylim(100, 650)
    ax0.set_title('Нормальный профиль ЕГТ\n(рабочие агрегаты)', color=C_WHITE, fontsize=10, fontweight='bold')
    ax0.set_ylabel('Avg ЕГТ (°C)', color=C_WHITE, fontsize=9)
    ax0.legend(fontsize=8, framealpha=0.3, facecolor=C_AXES, edgecolor=C_GRAY, labelcolor=C_WHITE)
    ax0.grid(True, alpha=0.3)
    ax0.tick_params(colors=C_WHITE, labelsize=8, rotation=45)
    for sp in ax0.spines.values(): sp.set_edgecolor(C_GRAY)

    # Right: fault units
    ax1 = axes[1]
    ax1.set_facecolor(C_AXES)
    colors_map2 = {u: c for u, c in zip(fault_uids, plt.cm.cool(np.linspace(0.2, 0.9, max(len(fault_uids), 1))))}
    for u in fault_uids:
        d = DATA[u]
        cyls = sorted(d['cyl_avgs'].keys())
        avgs = [d['cyl_avgs'][n] for n in cyls]
        ax1.plot([f'Ц{n}' for n in cyls], avgs, '-o', lw=1.5, ms=4,
                 color=colors_map2[u], label=f'730E {UNIT_LABELS[u]}', alpha=0.85)
    ax1.axhline(50, color=C_GRAY, lw=1.5, ls=':', alpha=0.8, label='50°C (порог отказа)')
    ax1.set_ylim(0, 200)
    ax1.set_title('Отказ датчиков ЕГТ\n(значения ≈58°C = потеря сигнала)', color=C_WHITE, fontsize=10, fontweight='bold')
    ax1.set_ylabel('Avg ЕГТ (°C)', color=C_WHITE, fontsize=9)
    if fault_uids:
        ax1.legend(fontsize=8, framealpha=0.3, facecolor=C_AXES, edgecolor=C_GRAY, labelcolor=C_WHITE)
    ax1.grid(True, alpha=0.3)
    ax1.tick_params(colors=C_WHITE, labelsize=8, rotation=45)
    for sp in ax1.spines.values(): sp.set_edgecolor(C_GRAY)

    plt.tight_layout()
    img = save_fig(fig)

    sl = blank_slide(prs)
    set_bg(sl)
    header_bar(sl, 'ДИАГНОСТИКА ДАТЧИКОВ ЕГТ',
               'Рабочий профиль vs. отказ датчиков (значение 58°C = разрыв сигнала)')
    add_img(sl, img, 0.3, 2.1, 33.2, 15.5)

    fault_list = ', '.join([f'730E {UNIT_LABELS[u]}' for u in fault_uids]) if fault_uids else 'нет'
    add_rect(sl, 0, 17.7, 33.86, 1.35, fill=C_DARK2, line=C_ORANGE, lw=1.2)
    add_text(sl, f'ВЫВОД: Агрегаты с отказом датчиков ЕГТ: {fault_list}. '
             f'Все цилиндры показывают ≈58°C — признак разрыва цепи термопары.',
             0.5, 17.75, 32.5, 1.25, fs=10, bold=True, color=C_ORANGE)
    footer(sl)

# ─── Slide 9: Two sessions of unit 39 ────────────────────────────────────────
def slide9(prs):
    d_a = DATA.get('39a')
    d_b = DATA.get('39b')
    if not d_a or not d_b: return

    fig, axes = plt.subplots(2, 2, figsize=(15, 9))
    fig.patch.set_facecolor(C_BG)
    fig.suptitle('730E #39 — Две сессии: 16.05 vs 21.05.2026', color=C_WHITE, fontsize=12, fontweight='bold', y=1.01)

    for row, (uid, d) in enumerate([('39a', d_a), ('39b', d_b)]):
        sec = d['sec']
        color_ts = C_BLUE if uid == '39a' else C_GREEN
        date = d['date']
        dur = d['dur']

        # RPM time series
        ax0 = axes[row][0]
        ax0.set_facecolor(C_AXES)
        ax0.fill_between(sec, smooth(d['rpm_ts']), alpha=0.12, color=color_ts)
        ax0.plot(sec, smooth(d['rpm_ts']), color=color_ts, lw=1.8)
        ax0.set_ylim(800, 2200)
        ax0.set_title(f'{date}  |  {dur:.1f} мин  |  RPM', color=C_WHITE, fontsize=9.5, fontweight='bold')
        ax0.set_ylabel('RPM', color=C_WHITE, fontsize=9)
        ax0.grid(True, alpha=0.3)
        ax0.tick_params(colors=C_WHITE, labelsize=8)
        for sp in ax0.spines.values(): sp.set_edgecolor(C_GRAY)

        # Per-cyl avg EGT bars
        ax1 = axes[row][1]
        ax1.set_facecolor(C_AXES)
        cyls = sorted(d['cyl_avgs'].keys())
        avgs = [d['cyl_avgs'][n] for n in cyls]
        bar_colors = [C_RED if v >= 500 else C_ORANGE if v >= 450 else color_ts for v in avgs]
        ax1.bar([f'Ц{n}' for n in cyls], avgs, color=bar_colors, alpha=0.85)
        ax1.axhline(450, color=C_ORANGE, lw=1.2, ls='--', alpha=0.8)
        ax1.axhline(500, color=C_RED, lw=1.2, ls='--', alpha=0.8)
        ax1.set_title(f'Avg ЕГТ по цилиндрам | max: Ц{d["hot_cyl"]} {d["egt_max"]:.0f}°C',
                      color=C_WHITE, fontsize=9.5, fontweight='bold')
        ax1.set_ylim(100, 620)
        ax1.tick_params(colors=C_WHITE, labelsize=7.5, rotation=45)
        ax1.grid(True, alpha=0.3, axis='y')
        for sp in ax1.spines.values(): sp.set_edgecolor(C_GRAY)

    plt.tight_layout()
    img = save_fig(fig)

    sl = blank_slide(prs)
    set_bg(sl)
    header_bar(sl, '730E #39 — ДВЕ СЕССИИ: 16.05 И 21.05.2026',
               f'Сравнение RPM и ЕГТ по цилиндрам между сессиями')
    add_img(sl, img, 0.3, 2.1, 33.2, 16.0)

    add_rect(sl, 0, 18.2, 33.86, 0.85, fill=C_DARK2, line=C_GRAY, lw=0.8)
    add_text(sl, f'16.05: max ЕГТ {d_a["egt_max"]:.0f}°C (Ц{d_a["hot_cyl"]})  |  21.05: max ЕГТ {d_b["egt_max"]:.0f}°C (Ц{d_b["hot_cyl"]})  |  Масло 16.05: {d_a["oil_max"]:.0f}°C',
             0.5, 18.25, 32.5, 0.75, fs=9.5, bold=True, color=C_GRAY)
    footer(sl)

# ─── Slide 10: EGT Heatmap ────────────────────────────────────────────────────
def slide10(prs):
    uids = list(DATA.keys())
    rows_u = uids
    cols_c = list(range(1, 17))

    matrix = np.full((len(rows_u), 16), np.nan)
    for i, uid in enumerate(rows_u):
        d = DATA[uid]
        if not d: continue
        for j, n in enumerate(cols_c):
            if n in d['cyl_maxes']:
                matrix[i, j] = d['cyl_maxes'][n]

    fig, ax = plt.subplots(figsize=(15, 6))
    fig.patch.set_facecolor(C_BG)
    ax.set_facecolor(C_AXES)

    masked = np.ma.masked_invalid(matrix)
    cmap = plt.cm.RdYlGn_r
    cmap.set_bad(color='#333333')
    im = ax.imshow(masked, aspect='auto', cmap=cmap, vmin=50, vmax=650)
    plt.colorbar(im, ax=ax, label='Max ЕГТ (°C)', shrink=0.8)

    ax.set_xticks(range(16))
    ax.set_xticklabels([f'Ц{n}' for n in cols_c], fontsize=9, color=C_WHITE)
    ax.set_yticks(range(len(rows_u)))
    ax.set_yticklabels([UNIT_LABELS[u] for u in rows_u], fontsize=9, color=C_WHITE)
    ax.set_title('Тепловая карта: максимальная ЕГТ (°C) — агрегаты × цилиндры', color=C_WHITE, fontsize=11, fontweight='bold')

    for i in range(len(rows_u)):
        for j in range(16):
            v = matrix[i, j]
            if not np.isnan(v):
                c = 'white' if v > 400 and v < 600 else ('black' if v < 200 else 'white')
                ax.text(j, i, f'{v:.0f}', ha='center', va='center', fontsize=7, color=c)

    ax.tick_params(colors=C_WHITE, labelsize=9)
    for sp in ax.spines.values(): sp.set_edgecolor(C_GRAY)
    plt.tight_layout()
    img = save_fig(fig)

    sl = blank_slide(prs)
    set_bg(sl)
    header_bar(sl, 'ТЕПЛОВАЯ КАРТА ЕГТ — ВСЕ АГРЕГАТЫ × 16 ЦИЛИНДРОВ',
               'Максимальная температура выхлопных газов по каждому цилиндру')
    add_img(sl, img, 0.3, 2.1, 33.2, 16.0)
    footer(sl)

# ─── Slide 11: Anomaly matrix ─────────────────────────────────────────────────
def slide11(prs):
    sl = blank_slide(prs)
    set_bg(sl)
    header_bar(sl, 'МАТРИЦА АНОМАЛИЙ — ПАРК 730E',
               'Обнаруженные отклонения по каждому агрегату | май 2026')

    col_headers = ['ЕГТ >500°C', 'Масло >105°C', 'Датч. ЕГТ', 'ЕГТ max', 'Статус']
    col_x = [3.5, 9.0, 14.5, 20.0, 26.5]
    header_y = 2.25

    add_rect(sl, 0.2, header_y, 33.4, 0.9, fill=C_DARK2)
    add_text(sl, 'Агрегат', 0.4, header_y+0.1, 3, 0.7, fs=9, bold=True, color=C_GRAY)
    for cx, ch in zip(col_x, col_headers):
        add_text(sl, ch, cx, header_y+0.1, 5.5, 0.7, fs=8.5, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    row_y = 3.25
    for uid in DATA:
        d = DATA[uid]
        risk = d['risk'] if d else 'НОРМА'
        rc = RISK_COLOR[risk]

        add_rect(sl, 0.2, row_y, 33.4, 1.5, fill='#1a1f23', line=C_GRAY, lw=0.3)
        add_text(sl, f'730E {UNIT_LABELS[uid]}', 0.4, row_y+0.4, 3.2, 0.8, fs=10, bold=True, color=C_WHITE)

        if d:
            checks = [
                (not np.isnan(d['egt_max']) and d['egt_max'] >= 500, f'{d["egt_max"]:.0f}°C' if not np.isnan(d['egt_max']) else '-'),
                (not np.isnan(d['oil_max']) and d['oil_max'] >= 105, f'{d["oil_max"]:.0f}°C' if not np.isnan(d['oil_max']) else '-'),
                (d['egt_fault'], 'ОТКАЗ' if d['egt_fault'] else 'НОРМА'),
                (False, f'{d["egt_max"]:.0f}°C' if not np.isnan(d['egt_max']) else 'N/A'),
                (False, risk),
            ]
            for cx, (is_bad, val) in zip(col_x, checks):
                color = (C_RED if is_bad else C_GREEN) if not isinstance(val, str) or val not in [risk] else rc
                if val == 'ОТКАЗ': color = C_RED
                elif val == 'НОРМА' and cx == col_x[2]: color = C_GREEN
                elif val == risk: color = rc
                add_text(sl, val, cx, row_y+0.4, 5.5, 0.75, fs=10, bold=True, color=color, align=PP_ALIGN.CENTER)

        row_y += 1.6

    add_rect(sl, 0, 17.7, 33.86, 1.35, fill=C_DARK2, line=C_ORANGE, lw=1.0)
    critical = [UNIT_LABELS[u] for u in DATA if DATA[u] and DATA[u]['risk'] == 'КРИТИЧНО']
    add_text(sl, f'КРИТИЧНО: {", ".join(critical)}. Требуется немедленная проверка цилиндров с ЕГТ >500°C.',
             0.5, 17.75, 32.5, 1.25, fs=10, bold=True, color=C_ORANGE)
    footer(sl)

# ─── Slide 12: Conclusions ───────────────────────────────────────────────────
def slide12(prs):
    sl = blank_slide(prs)
    set_bg(sl)
    add_rect(sl, 0, 0, 33.86, 2.2, fill=C_DARK2)
    add_rect(sl, 0, 2.0, 33.86, 0.12, fill=C_GREEN)
    add_text(sl, 'ЗАКЛЮЧЕНИЕ И РЕКОМЕНДАЦИИ — ПАРК 730E', 0.5, 0.2, 32, 1.2, fs=20, bold=True)
    add_text(sl, 'На основе DML-логов | 9 агрегатов | 16–21 мая 2026', 0.5, 1.45, 32, 0.6, fs=10, color=C_GRAY)

    # Critical risks
    add_rect(sl, 0.3, 2.3, 16, 0.75, fill=C_RED)
    add_text(sl, 'КРИТИЧЕСКИЕ РИСКИ', 0.3, 2.3, 16, 0.75, fs=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    critical_uids = [u for u in DATA if DATA[u] and DATA[u]['risk'] == 'КРИТИЧНО']
    risks = []
    for uid in critical_uids:
        d = DATA[uid]
        risks.append((f'730E {UNIT_LABELS[uid]} — ЕГТ max {d["egt_max"]:.0f}°C',
                      f'Цил.{d["hot_cyl"]} = {d["egt_max"]:.0f}°C при нагрузке {d["load"]:.0f}%. Превышение порога 500°C.'))

    for i, (title, desc) in enumerate(risks[:5]):
        y = 3.15 + i * 2.1
        add_rect(sl, 0.3, y, 16, 1.8, fill=C_DARK2, line=C_RED, lw=1)
        add_text(sl, title, 0.5, y+0.1, 15.5, 0.8, fs=10, bold=True, color=C_RED)
        add_text(sl, desc, 0.5, y+0.85, 15.5, 0.85, fs=8.5, color=C_GRAY)

    # Recommendations
    add_rect(sl, 17, 2.3, 16.5, 0.75, fill=C_GREEN)
    add_text(sl, 'РЕКОМЕНДАЦИИ', 17, 2.3, 16.5, 0.75, fs=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    recs = [
        ('СРОЧНО: Инспекция цилиндров >500°C',
         f'730E #24 Ц9, #27 Ц15, #28 Ц11: эндоскопия клапанов.\nТепловой зазор клапанов, состояние форсунок.'),
        ('Диагностика датчиков ЕГТ',
         'Агрегаты #23, #40: все датчики ЕГТ показывают 58°C.\nПроверить разъёмы термопар, заменить датчики.'),
        ('Мониторинг масла: #19 и #39a',
         'Масло 107°C и 109°C — близко к пределу 110°C.\nПроверить систему охлаждения масла.'),
        ('Регулярный DML-мониторинг',
         'Съём DML-логов каждые 2 недели для контроля.\nПороги: ЕГТ >450°C = предупреждение, >500°C = СТОП.'),
        ('Сравнительный анализ с NTE200',
         'EGT 730E ниже NTE200 (нет TIB=200%).\nБазовый уровень EGT: 400–450°C при нагрузке 70–90%.'),
    ]

    for i, (title, desc) in enumerate(recs[:5]):
        y = 3.15 + i * 2.1
        add_rect(sl, 17, y, 16.5, 1.8, fill=C_DARK2, line=C_GREEN, lw=1)
        add_text(sl, title, 17.2, y+0.1, 16, 0.8, fs=10, bold=True, color=C_GREEN)
        add_text(sl, desc, 17.2, y+0.85, 16, 0.85, fs=8.5, color=C_GRAY)

    footer(sl)

# ─── Bank asymmetry slides ────────────────────────────────────────────────────
def make_bank_chart(uid, d):
    """Bank asymmetry chart with 3-panel delta visualization."""
    bank_a = [1, 3, 5, 7, 9, 11, 13, 15]
    bank_b = [2, 4, 6, 8, 10, 12, 14, 16]
    colors_a = plt.cm.YlOrRd(np.linspace(0.35, 0.95, 8))
    colors_b = plt.cm.cool(np.linspace(0.25, 0.85, 8))

    sec = d['sec']
    cyl_egt = d['cyl_egt']
    valid_sec = sec[np.isfinite(sec)]
    t_max = valid_sec[-1] if len(valid_sec) > 0 else 1

    fig = plt.figure(figsize=(17, 11))
    fig.patch.set_facecolor(C_BG)
    fig.suptitle(
        f'730E {UNIT_LABELS[uid]}  |  {d["date"]}  |  Температурная асимметрия банков  |  Δ между цилиндрами',
        color=C_WHITE, fontsize=11, fontweight='bold', y=1.0)

    gs = gridspec.GridSpec(3, 2, figure=fig,
                           height_ratios=[3.2, 1.6, 1.6],
                           hspace=0.55, wspace=0.28)
    ax_a  = fig.add_subplot(gs[0, 0])
    ax_b  = fig.add_subplot(gs[0, 1])
    ax_db = fig.add_subplot(gs[1, :])
    ax_dt = fig.add_subplot(gs[2, :])

    def bank_ts_fn(bank):
        stacks = [cyl_egt[n] for n in bank if n in cyl_egt and np.any(np.isfinite(cyl_egt[n]))]
        if not stacks: return np.full_like(sec, np.nan, dtype=float)
        return np.nanmean(np.vstack(stacks), axis=0)

    def bank_max_ts(bank):
        stacks = [cyl_egt[n] for n in bank if n in cyl_egt and np.any(np.isfinite(cyl_egt[n]))]
        if not stacks: return np.full_like(sec, np.nan, dtype=float)
        return np.nanmax(np.vstack(stacks), axis=0)

    def bank_min_ts(bank):
        stacks = [cyl_egt[n] for n in bank if n in cyl_egt and np.any(np.isfinite(cyl_egt[n]))]
        if not stacks: return np.full_like(sec, np.nan, dtype=float)
        return np.nanmin(np.vstack(stacks), axis=0)

    def sm(v, w=5):
        return pd.Series(v).rolling(w, center=True, min_periods=1).mean().values

    a_avg_ts = sm(bank_ts_fn(bank_a))
    b_avg_ts = sm(bank_ts_fn(bank_b))

    bank_stats = {}

    for ax_p, bank, colors, bname, ecm in [
        (ax_a, bank_a, colors_a, 'А-банк (нечётные)', 'ECM144'),
        (ax_b, bank_b, colors_b, 'В-банк (чётные)',   'ECM1'),
    ]:
        ax_p.set_facecolor(C_AXES)

        # Intra-bank spread band
        mx_ts = sm(bank_max_ts(bank))
        mn_ts = sm(bank_min_ts(bank))
        ax_p.fill_between(sec, mn_ts, mx_ts, alpha=0.12, color=C_GRAY, label='Разброс')

        # Individual cylinders
        for cyl, color in zip(bank, colors):
            if cyl not in cyl_egt: continue
            v_s = sm(cyl_egt[cyl])
            lw = 2.0 if cyl in [1, 2, 6, 9, 11, 15] else 1.3
            ax_p.plot(sec, v_s, color=color, lw=lw, label=f'Ц{cyl}', alpha=0.9)

        # Bank average
        avg_ts = sm(bank_ts_fn(bank))
        ax_p.plot(sec, avg_ts, color=C_WHITE, lw=1.8, ls='--', alpha=0.6,
                  label='Avg банка', zorder=5)

        ax_p.axhline(450, color=C_ORANGE, lw=1.3, ls='--', alpha=0.85, label='450°C')
        ax_p.axhline(500, color=C_RED,    lw=1.4, ls='--', alpha=0.85, label='500°C')
        ax_p.set_ylim(100, 650)
        ax_p.set_xlim(0, t_max)
        ax_p.set_title(f'{bname}  |  {ecm}', color=C_WHITE, fontsize=10, fontweight='bold')
        ax_p.set_ylabel('ЕГТ (°C)', color=C_WHITE, fontsize=8.5)
        ax_p.set_xlabel('Время (сек)', color=C_WHITE, fontsize=8.5)
        ax_p.legend(fontsize=6.5, ncol=3, loc='upper left', framealpha=0.35,
                    facecolor=C_AXES, edgecolor='#445055', labelcolor=C_WHITE)
        ax_p.grid(True, alpha=0.3)
        ax_p.tick_params(colors=C_WHITE, labelsize=7.5)
        for sp in ax_p.spines.values(): sp.set_edgecolor('#445055')

        avgs, maxs, mins = [], [], []
        for c in bank:
            if c in d['cyl_avgs']:
                avgs.append(d['cyl_avgs'][c])
                maxs.append(d['cyl_maxes'][c])
            if c in cyl_egt:
                finite = cyl_egt[c][np.isfinite(cyl_egt[c])]
                if len(finite) > 5: mins.append(np.min(finite))
        if avgs:
            b_avg = np.mean(avgs); b_max = np.max(maxs)
            intra = b_max - (np.min(mins) if mins else b_avg)
        else:
            b_avg = b_max = intra = 0.0

        box_c = C_RED if b_max >= 500 else (C_ORANGE if b_max >= 450 else C_GREEN)
        ax_p.text(0.98, 0.04,
                  f'Avg: {b_avg:.0f}°C  Max: {b_max:.0f}°C\nΔ внутри банка: {intra:.0f}°C',
                  transform=ax_p.transAxes, ha='right', va='bottom',
                  color=C_WHITE, fontsize=8.5, fontweight='bold',
                  bbox=dict(boxstyle='round,pad=0.4', facecolor=C_AXES,
                            edgecolor=box_c, alpha=0.95, linewidth=1.5))
        bank_stats[bname] = (b_avg, b_max, intra)

    # ── Per-cylinder Δ from bank average ────────────────────────────────────
    ax_db.set_facecolor(C_AXES)
    ax_db.set_title('Δ каждого цилиндра от среднего по своему банку (avg − bank_avg)',
                    color=C_WHITE, fontsize=9.5, fontweight='bold')

    cyl_delta = {}
    for n in range(1, 17):
        if n not in cyl_egt: continue
        finite = cyl_egt[n][np.isfinite(cyl_egt[n])]
        if len(finite) < 5: continue
        bank = bank_a if n % 2 == 1 else bank_b
        b_vals = [cyl_egt[c][np.isfinite(cyl_egt[c])].mean()
                  for c in bank if c in cyl_egt and np.any(np.isfinite(cyl_egt[c]))]
        if b_vals:
            cyl_delta[n] = np.mean(finite) - np.mean(b_vals)

    if cyl_delta:
        cyls = sorted(cyl_delta.keys())
        deltas = [cyl_delta[n] for n in cyls]
        bar_clr = []
        for n, dv in zip(cyls, deltas):
            if dv > 30: bar_clr.append(C_RED)
            elif dv > 10: bar_clr.append(C_ORANGE)
            elif dv < -15: bar_clr.append(C_GREEN)
            else: bar_clr.append('#4EC9F0' if n % 2 == 1 else '#7090B0')

        x_pos = np.arange(len(cyls))
        bars = ax_db.bar(x_pos, deltas, color=bar_clr, alpha=0.88, width=0.7, zorder=3)
        ax_db.axhline(0,   color=C_WHITE,  lw=1.2, alpha=0.6, zorder=4)
        ax_db.axhline(20,  color=C_ORANGE, lw=0.8, ls=':', alpha=0.6)
        ax_db.axhline(-20, color=C_GREEN,  lw=0.8, ls=':', alpha=0.6)
        ax_db.set_xticks(x_pos)
        ax_db.set_xticklabels([f'Ц{n}' for n in cyls], fontsize=8.5, color=C_WHITE)
        ax_db.set_ylabel('Δ°C от avg', color=C_WHITE, fontsize=8.5)
        ax_db.grid(True, alpha=0.3, axis='y', zorder=1)
        ax_db.tick_params(colors=C_WHITE, labelsize=8)
        for sp in ax_db.spines.values(): sp.set_edgecolor('#445055')

        for bar, dv, n in zip(bars, deltas, cyls):
            offset = 1.5 if dv >= 0 else -4.0
            va = 'bottom' if dv >= 0 else 'top'
            ax_db.text(bar.get_x() + bar.get_width()/2, dv + offset,
                       f'{dv:+.0f}°', ha='center', va=va, fontsize=7.5,
                       color=C_RED if abs(dv) > 25 else C_WHITE, fontweight='bold')

        ax_db.axvline(7.5, color=C_GRAY, lw=1.2, ls='--', alpha=0.5)
        ylim = ax_db.get_ylim()
        ax_db.text(3.5,  ylim[1]*0.85, 'А-банк (нечётные)', ha='center',
                   color=colors_a[4], fontsize=8, fontweight='bold', alpha=0.9)
        ax_db.text(11.5, ylim[1]*0.85, 'В-банк (чётные)', ha='center',
                   color=colors_b[4], fontsize=8, fontweight='bold', alpha=0.9)

    # ── Inter-bank delta over time ───────────────────────────────────────────
    ax_dt.set_facecolor(C_AXES)
    ax_dt.set_title('Межбанковая дельта: В-банк avg − А-банк avg  (по времени)',
                    color=C_WHITE, fontsize=9.5, fontweight='bold')

    inter = sm(b_avg_ts - a_avg_ts, w=9)
    ax_dt.fill_between(sec, inter, 0, where=inter > 0, alpha=0.35, color=C_RED,   label='В горячее')
    ax_dt.fill_between(sec, inter, 0, where=inter < 0, alpha=0.35, color=C_GREEN, label='А горячее')
    ax_dt.plot(sec, inter, color=C_ORANGE, lw=2.0, alpha=0.9, label='Δ (В−А)')
    ax_dt.axhline(0,   color=C_WHITE,  lw=1.2, alpha=0.5, ls='--')
    ax_dt.axhline(20,  color=C_ORANGE, lw=0.8, ls=':', alpha=0.6, label='+20°C')
    ax_dt.axhline(-20, color=C_GREEN,  lw=0.8, ls=':', alpha=0.6, label='-20°C')
    ax_dt.set_xlim(0, t_max)
    ax_dt.set_ylabel('Δ (°C)', color=C_WHITE, fontsize=8.5)
    ax_dt.set_xlabel('Время (сек)', color=C_WHITE, fontsize=8.5)
    ax_dt.legend(fontsize=7.5, ncol=4, loc='upper right', framealpha=0.35,
                 facecolor=C_AXES, edgecolor='#445055', labelcolor=C_WHITE)
    ax_dt.grid(True, alpha=0.3)
    ax_dt.tick_params(colors=C_WHITE, labelsize=7.5)
    for sp in ax_dt.spines.values(): sp.set_edgecolor('#445055')

    fd = inter[np.isfinite(inter)]
    if len(fd) > 0:
        mean_d = np.mean(fd)
        ax_dt.text(0.99, 0.85, f'Ср. Δ(В−А) = {mean_d:+.0f}°C',
                   transform=ax_dt.transAxes, ha='right', va='top',
                   color=C_RED if mean_d > 15 else (C_GREEN if mean_d < -15 else C_WHITE),
                   fontsize=9, fontweight='bold',
                   bbox=dict(boxstyle='round,pad=0.35', facecolor=C_AXES,
                             edgecolor=C_ORANGE, alpha=0.9))

    fig.tight_layout(rect=[0, 0, 1, 0.97])
    img = save_fig(fig)
    return img, bank_stats, cyl_delta

def add_bank_slide(prs, uid, d):
    img, bank_stats, cyl_delta = make_bank_chart(uid, d)

    a_n, b_n = 'А-банк (нечётные)', 'В-банк (чётные)'
    a_avg, a_max, a_intra = bank_stats.get(a_n, (0, 0, 0))
    b_avg, b_max, b_intra = bank_stats.get(b_n, (0, 0, 0))
    delta_ab = b_avg - a_avg

    if abs(delta_ab) > 10:
        asym = f'В-банк горячее А-банка на {abs(delta_ab):.0f}°C' if delta_ab > 0 \
               else f'А-банк горячее В-банка на {abs(delta_ab):.0f}°C'
    else:
        asym = 'Банки симметричны'

    intra_txt = f'  |  Разброс: А={a_intra:.0f}°C  В={b_intra:.0f}°C'
    if cyl_delta:
        hot_cyl = max(cyl_delta, key=lambda n: cyl_delta[n])
        cyl_txt = f'  |  Цил.{hot_cyl} горячее avg на {cyl_delta[hot_cyl]:+.0f}°C'
    else:
        cyl_txt = ''

    conclusion = f'ВЫВОД: {asym}.{intra_txt}.{cyl_txt}'

    sl = blank_slide(prs)
    set_bg(sl)
    add_rect(sl, 0, 0, 33.86, 2.0, fill=C_DARK2)
    add_rect(sl, 0, 1.85, 33.86, 0.12, fill=C_GREEN)
    add_text(sl, f'ТЕМПЕРАТУРА ЦИЛИНДРОВ — 730E {UNIT_LABELS[uid]}  |  {d["date"]}',
             0.5, 0.1, 32, 1.0, fs=15, bold=True)
    add_text(sl, 'Δ внутри банков  |  Δ каждого цилиндра от среднего  |  Межбанковая асимметрия по времени',
             0.5, 1.15, 32, 0.7, fs=9, color=C_GRAY)

    sl.shapes.add_picture(img, Cm(0.2), Cm(2.05), Cm(33.46), Cm(15.55))

    add_rect(sl, 0, 17.68, 33.86, 1.37, fill=C_DARK2, line=C_ORANGE, lw=1.2)
    add_text(sl, conclusion, 0.4, 17.73, 33.0, 1.27, fs=9.5, bold=True, color=C_ORANGE)
    footer(sl)

# ─── Section divider ──────────────────────────────────────────────────────────
def section_divider(prs, num, title, subtitle=''):
    sl = blank_slide(prs)
    set_bg(sl)
    add_rect(sl, 0, 6.5, 33.86, 6.0, fill=C_DARK2)
    add_rect(sl, 0, 6.5, 33.86, 0.3, fill=C_GREEN)
    add_text(sl, str(num), 1.5, 7.5, 5, 3, fs=72, bold=True, color=C_GREEN)
    add_text(sl, title, 7, 7.5, 25, 2.5, fs=40, bold=True, color=C_WHITE)
    if subtitle:
        add_text(sl, subtitle, 7, 11.5, 25, 1.2, fs=14, color=C_GRAY)

# ─── Main ─────────────────────────────────────────────────────────────────────
def main():
    out = '/home/user/NTE200/DML_Анализ_LogFiles_730E.pptx'
    prs = new_prs()

    print("Building slides...")
    slide1(prs);  print("  Slide 1: Title")
    slide2(prs);  print("  Slide 2: Stat cards")
    slide3(prs);  print("  Slide 3: EGT comparison")
    slide4(prs);  print("  Slide 4: Multi-parameter grid")
    slide5(prs);  print("  Slide 5: Unit 24 deep-dive")
    slide6(prs);  print("  Slide 6: Unit 24 per-cyl EGT")
    slide7(prs);  print("  Slide 7: Unit 28 long session")
    slide8(prs);  print("  Slide 8: Sensor fault analysis")
    slide9(prs);  print("  Slide 9: Unit 39 two sessions")
    slide10(prs); print("  Slide 10: EGT heatmap")
    slide11(prs); print("  Slide 11: Anomaly matrix")
    slide12(prs); print("  Slide 12: Conclusions")

    # Bank asymmetry section
    section_divider(prs, '09', 'ТЕМПЕРАТУРА', 'А-банк (ECM144) vs В-банк (ECM1)  |  Все агрегаты 730E')
    print("  Section divider: Bank asymmetry")

    for uid in ['24', '19', '23a', '23b', '27', '28', '39a', '39b', '40']:
        d = DATA.get(uid)
        if d:
            add_bank_slide(prs, uid, d)
            print(f"  Bank slide: unit {uid}")

    print(f"\nTotal slides: {len(prs.slides)}")
    prs.save(out)
    print(f"Saved: {out}")
    cleanup()

if __name__ == '__main__':
    main()
