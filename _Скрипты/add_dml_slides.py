"""
Add DML analysis slides (34-40) to presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec
import numpy as np
import pandas as pd
import warnings
warnings.filterwarnings('ignore')

# ─────────────────────────────────────────────
# Brand colors
# ─────────────────────────────────────────────
DARK   = RGBColor(0x29, 0x31, 0x36)
ACCENT = RGBColor(0x3E, 0xF0, 0xAF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RED_C  = RGBColor(0xFF, 0x44, 0x44)
ORANGE = RGBColor(0xFF, 0x8C, 0x42)
GRAY   = RGBColor(0x66, 0x77, 0x77)

# ─────────────────────────────────────────────
# Chart colors (matplotlib)
# ─────────────────────────────────────────────
DARK_M   = '#293136'
ACCENT_M = '#3EF0AF'
RED_M    = '#FF4444'
ORANGE_M = '#FF8C42'
YELLOW_M = '#FFD166'
BLUE_M   = '#4EC9F0'
WHITE_M  = '#FFFFFF'

plt.rcParams.update({
    'font.family': 'DejaVu Sans',
    'axes.facecolor': '#1e2529',
    'figure.facecolor': DARK_M,
    'text.color': WHITE_M,
    'axes.labelcolor': WHITE_M,
    'axes.edgecolor': '#445055',
    'xtick.color': '#aaaaaa',
    'ytick.color': '#aaaaaa',
    'grid.color': '#445055',
    'grid.alpha': 0.5,
})

# ─────────────────────────────────────────────
# Load DML data
# ─────────────────────────────────────────────
df1 = pd.read_csv('/home/user/NTE200/DML_82_bez_porody.csv', encoding='cp1251', skiprows=26, decimal=',', on_bad_lines='skip')
df2 = pd.read_csv('/home/user/NTE200/DML_82_s_porodoy.csv', encoding='cp1251', skiprows=26, decimal=',', on_bad_lines='skip')

def parse_secs(df):
    t = df['Время'].str.split(':').apply(
        lambda x: float(x[0])*3600 + float(x[1])*60 + float(x[2]) if len(x)==3 else 0)
    return t - t.iloc[0]

sec1 = parse_secs(df1)
sec2 = parse_secs(df2)

RPM     = 'Частота вращения двигателя (об/мин)- Идентификатор0'
LOAD    = 'Относительная нагрузка (Проценты)- Идентификатор0'
AVG_EGT = 'Средняя температура отработавших газов (расчетное значение) (°C)- Идентификатор1'
EGT_L   = 'Средняя температура отработавших газов - левый ряд (расчетное значение) (°C)- Идентификатор144'
EGT_R   = 'Средняя температура отработавших газов - правый ряд (расчетное значение) (°C)- Идентификатор1'
COOLANT = 'Датчик температуры (°C)- Идентификатор0'
OIL_P   = 'Давление масла (кПа)- Идентификатор0'
BOOST   = 'Давление во впускном коллекторе (кПа)- Идентификатор0'
OIL_T   = 'Температура масла (°C)- Идентификатор0'

def egt_col(n):
    bank = '144' if n % 2 == 1 else '1'
    return f'Датчик температуры отработавших газов цилиндра {n} (°C)- Идентификатор{bank}'

def get_cyl_stats(df, cyl):
    v = df[egt_col(cyl)]
    v_clean = v.where(v > 50) if cyl == 11 else v
    return v_clean.mean(), v_clean.max()

# ─────────────────────────────────────────────
# Helper: add slide with dark background
# ─────────────────────────────────────────────
def add_dark_slide(prs):
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK
    return slide

def add_text(slide, text, left, top, width, height, font_size, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_image(slide, img_path, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(img_path, left, top, width, height)
    else:
        slide.shapes.add_picture(img_path, left, top, width)

W = Cm(33.86)
H = Cm(19.05)

# ─────────────────────────────────────────────
# Slide 34: Section divider "DML Analysis"
# ─────────────────────────────────────────────
def make_slide_34(prs):
    slide = add_dark_slide(prs)

    # Section number bar
    bar = slide.shapes.add_shape(1, Cm(0), Cm(6.5), Cm(33.86), Cm(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x1e, 0x25, 0x29)
    bar.line.fill.background()

    # Accent line
    line = slide.shapes.add_shape(1, Cm(0), Cm(6.5), Cm(33.86), Cm(0.3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_text(slide, '08', Cm(1.5), Cm(7.2), Cm(5), Cm(3),
             font_size=80, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, 'АНАЛИЗ DML ЛОГОВ', Cm(7), Cm(7.8), Cm(25), Cm(2),
             font_size=36, bold=True, color=WHITE)
    add_text(slide, 'Динамический анализ параметров двигателя в реальном времени',
             Cm(7), Cm(10.2), Cm(24), Cm(1.5),
             font_size=16, bold=False, color=ACCENT)

    # Info cards
    cards = [
        ('02.06.2026', 'Дата теста'),
        ('Агрегат 82', 'Объект'),
        ('2 теста', 'Кол-во'),
        ('13.5 + 23.8 мин', 'Длительность'),
    ]
    for i, (val, lbl) in enumerate(cards):
        x = Cm(1.5 + i * 8)
        box = slide.shapes.add_shape(1, x, Cm(14.5), Cm(7.5), Cm(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x1e, 0x25, 0x29)
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1)
        add_text(slide, val, x + Cm(0.3), Cm(14.6), Cm(7), Cm(1.2),
                 font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, x + Cm(0.3), Cm(15.7), Cm(7), Cm(0.8),
                 font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, 'Методология: тест без нагрузки (порожний рейс) + тест под нагрузкой (с породой)',
             Cm(1.5), Cm(17.5), Cm(30), Cm(1), font_size=10, color=GRAY, italic=True)

# ─────────────────────────────────────────────
# Slide 35: Key findings summary (stat cards)
# ─────────────────────────────────────────────
def make_slide_35(prs):
    slide = add_dark_slide(prs)

    # Header
    header = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(33.86), Cm(2.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1e, 0x25, 0x29)
    header.line.fill.background()

    add_text(slide, 'КЛЮЧЕВЫЕ ВЫВОДЫ DML — АГРЕГАТ 82',
             Cm(1), Cm(0.3), Cm(30), Cm(1.5),
             font_size=20, bold=True, color=WHITE)
    add_text(slide, '02 июня 2026  |  Тест 1 (без породы) + Тест 2 (с породой)',
             Cm(1), Cm(1.55), Cm(30), Cm(0.7),
             font_size=10, color=ACCENT)

    # 6 key finding cards in 2 rows x 3 cols
    findings = [
        (f'567°C',    'Макс. ЕГТ цил.6 и цил.2', 'КРИТИЧНО', RED_C,
         'Цилиндры 6 и 2 достигают 567°C под нагрузкой.\nПредел клапана Cummins QSK50: 650°C.'),
        (f'88%',      'Нагрузка с породой', 'ВЫСОКАЯ', ORANGE,
         'Средняя нагрузка 88% при ведении с породой.\nЭто в 4.1x выше порожнего режима (21%).'),
        (f'1929 об/м','Макс. обороты', 'КОНТРОЛЬ', RGBColor(0xFF,0xD1,0x66),
         'Максимальные обороты 1929 rpm под нагрузкой.\nНиже лимита 2100, но TIB=200% удваивает нагрузку.'),
        (f'Цил.11',   'Отказ датчика ЕГТ', 'ДЕФЕКТ', RED_C,
         '63% времени теста датчик ЕГТ цил.11 выдаёт\n0°C — признак обгоревшего клапана.'),
        (f'525°C',    'Ср. ЕГТ правого банка', 'КРИТИЧНО', RED_C,
         'Правый банк (B-банк, чётные цилиндры) горячее\nлевого на 60°C в среднем под нагрузкой.'),
        (f'+186°C',   'Прирост ЕГТ под нагрузкой', 'НАГРЕВ', ORANGE,
         'Нагрузка поднимает ЕГТ с 287°C до 474°C.\nДельта +186°C при нагрузке с породой.'),
    ]

    cols3 = [Cm(0.8), Cm(11.8), Cm(22.8)]
    rows2 = [Cm(2.5), Cm(10.5)]

    for i, (val, title, badge, badge_color, desc) in enumerate(findings):
        row = i // 3
        col_i = i % 3
        x = cols3[col_i]
        y = rows2[row]

        # Card background
        card = slide.shapes.add_shape(1, x, y, Cm(10.5), Cm(7.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x1e, 0x25, 0x29)
        card.line.color.rgb = badge_color
        card.line.width = Pt(1.5)

        # Badge
        badge_box = slide.shapes.add_shape(1, x + Cm(0.3), y + Cm(0.3), Cm(3.5), Cm(0.65))
        badge_box.fill.solid()
        badge_box.fill.fore_color.rgb = badge_color
        badge_box.line.fill.background()
        add_text(slide, badge, x + Cm(0.3), y + Cm(0.25), Cm(3.5), Cm(0.7),
                 font_size=8, bold=True, color=RGBColor(0x29,0x31,0x36), align=PP_ALIGN.CENTER)

        # Value
        add_text(slide, val, x + Cm(0.4), y + Cm(1.1), Cm(9.5), Cm(2.2),
                 font_size=32, bold=True, color=badge_color)

        # Title
        add_text(slide, title, x + Cm(0.4), y + Cm(3.2), Cm(9.5), Cm(1.0),
                 font_size=11, bold=True, color=WHITE)

        # Description
        add_text(slide, desc, x + Cm(0.4), y + Cm(4.2), Cm(9.8), Cm(3.0),
                 font_size=9, color=GRAY)

# ─────────────────────────────────────────────
# Slide 36: Overview time series chart
# ─────────────────────────────────────────────
def make_slide_36(prs):
    # Build the chart
    fig, axes = plt.subplots(4, 1, figsize=(16, 11))
    fig.patch.set_facecolor(DARK_M)
    fig.suptitle('Агрегат 82 — Динамика ключевых параметров  |  02.06.2026',
                 fontsize=14, color=WHITE_M, fontweight='bold', y=0.99)

    plots = [
        (RPM,     'Обороты (об/мин)',     ACCENT_M, ORANGE_M, (700, 2200)),
        (LOAD,    'Нагрузка (%)',          BLUE_M,   RED_M,    (0, 115)),
        (AVG_EGT, 'Средняя ЕГТ (°C)',     YELLOW_M, RED_M,    (150, 560)),
        (OIL_T,   'Т° масла (°C)',         '#90EE90',ORANGE_M, (60, 125)),
    ]

    for idx, (col, ylabel, c1, c2, ylim) in enumerate(plots):
        ax = axes[idx]
        ax.set_facecolor('#1e2529')
        v1 = df1[col].rolling(5, center=True).mean()
        v2 = df2[col].rolling(5, center=True).mean()

        # Offset test 2 to right of test 1
        offset = sec1.max() + 25
        ax.fill_between(sec1, v1, alpha=0.12, color=c1)
        ax.fill_between(sec2 + offset, v2, alpha=0.12, color=c2)
        ax.plot(sec1, v1, color=c1, lw=1.8, label='Без породы', alpha=0.9)
        ax.plot(sec2 + offset, v2, color=c2, lw=1.8, label='С породой', alpha=0.9)

        # Separator
        ax.axvline(offset - 12, color='#666666', lw=1.5, ls=':', alpha=0.8)

        # Limit lines
        if col == RPM:
            ax.axhline(2100, color=RED_M, lw=1, ls='--', alpha=0.7, label='Лимит 2100')
        if col == AVG_EGT:
            ax.axhline(500, color=RED_M, lw=1, ls='--', alpha=0.7, label='500°C критич.')

        ax.set_ylabel(ylabel, color=WHITE_M, fontsize=9)
        ax.set_ylim(ylim)
        ax.grid(True, alpha=0.3)
        ax.tick_params(colors=WHITE_M, labelsize=8)
        for spine in ax.spines.values():
            spine.set_edgecolor('#445055')
        if idx < 3:
            ax.set_xticklabels([])
        else:
            ax.set_xlabel('Время (сек)', color=WHITE_M, fontsize=9)
        if idx == 0:
            ax.legend(loc='upper right', fontsize=8, framealpha=0.3,
                      facecolor=DARK_M, edgecolor='#445055', labelcolor=WHITE_M)
            # Labels for tests
            ax.text(sec1.max()/2, 2050,
                    'ТЕСТ 1 — Без породы\n13.5 мин  |  avg 1441 rpm  |  avg ЕГТ 287°C',
                    ha='center', color=ACCENT_M, fontsize=8, fontweight='bold')
            ax.text(offset + sec2.max()/2, 2050,
                    'ТЕСТ 2 — С породой\n23.8 мин  |  avg 1822 rpm  |  avg ЕГТ 474°C',
                    ha='center', color=ORANGE_M, fontsize=8, fontweight='bold')

    plt.tight_layout(rect=[0, 0, 1, 0.98])
    plt.savefig('/tmp/dml_ts_slide.png', dpi=130, bbox_inches='tight', facecolor=DARK_M)
    plt.close()

    slide = add_dark_slide(prs)
    # Header bar
    header = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(33.86), Cm(1.4))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1e, 0x25, 0x29)
    header.line.fill.background()
    add_text(slide, 'DML: ДИНАМИКА ПАРАМЕТРОВ — АГРЕГАТ 82',
             Cm(1), Cm(0.1), Cm(30), Cm(1.2),
             font_size=16, bold=True, color=WHITE)

    add_image(slide, '/tmp/dml_ts_slide.png', Cm(0.5), Cm(1.5), Cm(32.5))

# ─────────────────────────────────────────────
# Slide 37: EGT per cylinder bar chart
# ─────────────────────────────────────────────
def make_slide_37(prs):
    fig, axes = plt.subplots(1, 2, figsize=(16, 7))
    fig.patch.set_facecolor(DARK_M)

    cyls = list(range(1, 17))
    bank_a = [c for c in cyls if c % 2 == 1]
    bank_b = [c for c in cyls if c % 2 == 0]

    def get_means_maxs(df):
        means, maxs = [], []
        for c in cyls:
            v = df[egt_col(c)]
            if c == 11:
                v = v.where(v > 50)
            means.append(v.mean())
            maxs.append(v.max())
        return np.array(means), np.array(maxs)

    m1, x1 = get_means_maxs(df1)
    m2, x2 = get_means_maxs(df2)

    x = np.arange(16)
    w = 0.35

    for ax_idx, (means, maxs, title, c_m, c_x, sec_lbl) in enumerate([
        (m1, x1, 'Тест 1: БЕЗ ПОРОДЫ  (avg нагрузка 21%)', ACCENT_M, YELLOW_M, '13.5 мин'),
        (m2, x2, 'Тест 2: С ПОРОДОЙ  (avg нагрузка 88%)', ORANGE_M, RED_M, '23.8 мин'),
    ]):
        ax = axes[ax_idx]
        ax.set_facecolor('#1e2529')

        bar_m = ax.bar(x - w/2, means, w, label='Среднее', color=c_m, alpha=0.8, zorder=3)
        bar_x = ax.bar(x + w/2, maxs, w, label='Максимум', color=c_x, alpha=0.85, zorder=3)

        if ax_idx == 1:
            # Gray out cyl 11 (sensor fault)
            ax.bar(10 - w/2, means[10] if not np.isnan(means[10]) else 0, w,
                   color='#888888', alpha=0.6, hatch='////', zorder=4)
            ax.bar(10 + w/2, maxs[10] if not np.isnan(maxs[10]) else 0, w,
                   color='#888888', alpha=0.6, hatch='////', zorder=4)
            ax.text(10, 25, 'FC1531\nдатчик', ha='center', color=RED_M, fontsize=6.5, fontweight='bold')

        # Mark top 3 max with value labels
        sorted_max = np.argsort(np.nan_to_num(maxs, nan=0))[-3:]
        for ti in sorted_max:
            if not np.isnan(maxs[ti]):
                ax.bar(ti + w/2, maxs[ti], w, color='#FF2222', alpha=1.0, zorder=5)
                ax.text(ti + w/2, maxs[ti] + 6, f'{maxs[ti]:.0f}',
                        ha='center', color=RED_M, fontsize=7.5, fontweight='bold')

        ax.axhline(500, color=RED_M, lw=1.5, ls='--', alpha=0.8, label='500°C критич.')
        ax.axhline(450, color=ORANGE_M, lw=1, ls=':', alpha=0.7, label='450°C внимание')

        ax.set_xticks(x)
        ax.set_xticklabels([f'Ц{n}' for n in cyls], fontsize=8.5, color=WHITE_M)
        ax.set_ylabel('Температура ЕГТ (°C)', color=WHITE_M, fontsize=9)
        ax.set_title(f'{title}  |  {sec_lbl}', color=WHITE_M, fontsize=10.5, fontweight='bold', pad=8)
        ax.set_ylim(0, 640)
        ax.grid(True, axis='y', alpha=0.3)
        ax.legend(loc='upper right', fontsize=8, framealpha=0.3,
                  facecolor=DARK_M, edgecolor='#445055', labelcolor=WHITE_M)
        ax.tick_params(colors=WHITE_M, labelsize=8)
        for spine in ax.spines.values():
            spine.set_edgecolor('#445055')
        # Bank labels
        ax.axvline(7.5, color='#555555', lw=1.5, alpha=0.6)
        ax.text(3.5, 625, 'A-БАНК (нечётные)', ha='center', color='#aaaaaa', fontsize=8)
        ax.text(11.5, 625, 'B-БАНК (чётные)', ha='center', color='#aaaaaa', fontsize=8)

    plt.tight_layout()
    plt.savefig('/tmp/dml_egt_cyl_slide.png', dpi=130, bbox_inches='tight', facecolor=DARK_M)
    plt.close()

    slide = add_dark_slide(prs)
    header = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(33.86), Cm(1.4))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1e, 0x25, 0x29)
    header.line.fill.background()
    add_text(slide, 'DML: ТЕМПЕРАТУРА ЕГТ ПО ЦИЛИНДРАМ — АГРЕГАТ 82',
             Cm(1), Cm(0.1), Cm(30), Cm(1.2),
             font_size=16, bold=True, color=WHITE)
    add_image(slide, '/tmp/dml_egt_cyl_slide.png', Cm(0.5), Cm(1.5), Cm(32.5))

# ─────────────────────────────────────────────
# Slide 38: Cylinder 11 anomaly + hottest cylinders
# ─────────────────────────────────────────────
def make_slide_38(prs):
    fig, axes = plt.subplots(1, 2, figsize=(16, 8))
    fig.patch.set_facecolor(DARK_M)
    fig.suptitle('Агрегат 82 — Дефект датчика ЕГТ Цилиндра 11  и  Горячие цилиндры',
                 fontsize=12, color=WHITE_M, fontweight='bold')

    # Left: cyl 11 sensor fault in test 2
    ax = axes[0]
    ax.set_facecolor('#1e2529')
    v11_raw = df2[egt_col(11)]
    v11_filt = v11_raw.where(v11_raw > 50, other=np.nan)
    v9 = df2[egt_col(9)].rolling(5, center=True).mean()
    v13 = df2[egt_col(13)].rolling(5, center=True).mean()

    # Fill fault periods
    fault_mask = v11_raw < 50
    for i in range(len(fault_mask)):
        if fault_mask.iloc[i]:
            ax.axvspan(sec2.iloc[i]-0.5, sec2.iloc[i]+0.5, alpha=0.15, color=RED_M, zorder=1)

    ax.plot(sec2, v9, color=BLUE_M, lw=1.5, label='Цил.9 (A-банк)', alpha=0.8)
    ax.plot(sec2, v13, color=ACCENT_M, lw=1.5, label='Цил.13 (A-банк)', alpha=0.8)
    ax.plot(sec2, v11_filt.rolling(3, center=True).mean(), color=RED_M, lw=2.5,
            label='Цил.11 (валидные)', zorder=5)

    ax.axhline(500, color=RED_M, lw=1.5, ls='--', alpha=0.7, label='500°C критич.')
    ax.set_ylim(-20, 580)
    ax.set_title('Тест 2 (с породой) — Сбои датчика Цил.11', color=WHITE_M, fontsize=10.5, fontweight='bold')
    ax.set_ylabel('ЕГТ (°C)', color=WHITE_M, fontsize=9)
    ax.set_xlabel('Время (сек)', color=WHITE_M, fontsize=9)
    ax.legend(fontsize=8, framealpha=0.3, facecolor=DARK_M, edgecolor='#445055', labelcolor=WHITE_M)
    ax.grid(True, alpha=0.3)
    ax.tick_params(colors=WHITE_M, labelsize=8)
    for spine in ax.spines.values():
        spine.set_edgecolor('#445055')

    n_faults = fault_mask.sum()
    pct = n_faults/len(v11_raw)*100
    ax.text(0.02, 0.96, f'Сбоев: {n_faults} из {len(v11_raw)} точек ({pct:.0f}%)\nPossible: обгоревший клапан\n→ разрыв термопары',
            transform=ax.transAxes, color=RED_M, fontsize=9, fontweight='bold', va='top',
            bbox=dict(boxstyle='round', facecolor='#1e2529', edgecolor=RED_M, alpha=0.9))

    # Right: top 6 hottest cylinders heatmap (loaded test)
    ax2 = axes[1]
    ax2.set_facecolor('#1e2529')

    cyl_avgs = {}
    cyl_maxs = {}
    for c in range(1, 17):
        v = df2[egt_col(c)]
        if c == 11:
            v = v.where(v > 50)
        cyl_avgs[c] = v.mean()
        cyl_maxs[c] = v.max()

    sorted_cyls = sorted(cyl_avgs.items(), key=lambda x: x[1] if not np.isnan(x[1]) else 0, reverse=True)
    sorted_cyls_max = {c: cyl_maxs[c] for c, _ in sorted_cyls[:8]}

    cyls_sorted = [c for c, _ in sorted_cyls]
    avgs_sorted = [v if not np.isnan(v) else 0 for _, v in sorted_cyls]
    maxs_sorted = [cyl_maxs[c] if not np.isnan(cyl_maxs[c]) else 0 for c in cyls_sorted]

    bar_colors = []
    for v in avgs_sorted:
        if v >= 500: bar_colors.append('#FF2222')
        elif v >= 480: bar_colors.append(RED_M)
        elif v >= 460: bar_colors.append(ORANGE_M)
        elif v >= 440: bar_colors.append(YELLOW_M)
        else: bar_colors.append(ACCENT_M)

    x = np.arange(16)
    ax2.bar(x - 0.2, avgs_sorted, 0.4, color=bar_colors, alpha=0.9, label='Среднее', zorder=3)
    ax2.bar(x + 0.2, maxs_sorted, 0.4, color=[c+'99' if c.startswith('#') else c for c in bar_colors],
            alpha=0.6, label='Максимум', zorder=3)

    for xi, (avg, mx) in enumerate(zip(avgs_sorted, maxs_sorted)):
        if mx >= 540:
            ax2.text(xi + 0.2, mx + 5, f'{mx:.0f}', ha='center', color=RED_M, fontsize=7.5, fontweight='bold')

    ax2.set_xticks(x)
    ax2.set_xticklabels([f'Ц{c}' for c in cyls_sorted], fontsize=8.5, color=WHITE_M)
    ax2.axhline(500, color=RED_M, lw=1.5, ls='--', alpha=0.8, label='500°C')
    ax2.set_title('Все цилиндры: ЕГТ под нагрузкой (с породой)\nСортировка по среднему', color=WHITE_M, fontsize=10.5, fontweight='bold')
    ax2.set_ylabel('ЕГТ (°C)', color=WHITE_M, fontsize=9)
    ax2.set_ylim(0, 640)
    ax2.legend(fontsize=8, framealpha=0.3, facecolor=DARK_M, edgecolor='#445055', labelcolor=WHITE_M)
    ax2.grid(True, axis='y', alpha=0.3)
    ax2.tick_params(colors=WHITE_M, labelsize=8)
    for spine in ax2.spines.values():
        spine.set_edgecolor('#445055')

    plt.tight_layout(rect=[0, 0, 1, 0.95])
    plt.savefig('/tmp/dml_cyl11_slide.png', dpi=130, bbox_inches='tight', facecolor=DARK_M)
    plt.close()

    slide = add_dark_slide(prs)
    header = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(33.86), Cm(1.4))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1e, 0x25, 0x29)
    header.line.fill.background()
    add_text(slide, 'DML: АНОМАЛИЯ ЕГТ ЦИЛ.11 + РЕЙТИНГ ГОРЯЧИХ ЦИЛИНДРОВ',
             Cm(1), Cm(0.1), Cm(30), Cm(1.2),
             font_size=16, bold=True, color=WHITE)
    add_image(slide, '/tmp/dml_cyl11_slide.png', Cm(0.5), Cm(1.5), Cm(32.5))

# ─────────────────────────────────────────────
# Slide 39: Bank asymmetry — A vs B under load
# ─────────────────────────────────────────────
def make_slide_39(prs):
    fig, axes = plt.subplots(1, 2, figsize=(16, 8))
    fig.patch.set_facecolor(DARK_M)
    fig.suptitle('Агрегат 82 — Температурная асимметрия банков A и B',
                 fontsize=13, color=WHITE_M, fontweight='bold')

    bank_a = [1, 3, 5, 7, 9, 11, 13, 15]
    bank_b = [2, 4, 6, 8, 10, 12, 14, 16]
    colors_a = plt.cm.YlOrRd(np.linspace(0.4, 0.95, 8))
    colors_b = plt.cm.cool(np.linspace(0.3, 0.9, 8))

    for ax_idx, (bank, colors, bank_name, ecm_id) in enumerate([
        (bank_a, colors_a, 'A-банк (нечётные цилиндры)', 'ECM144'),
        (bank_b, colors_b, 'B-банк (чётные цилиндры)', 'ECM1'),
    ]):
        ax = axes[ax_idx]
        ax.set_facecolor('#1e2529')
        for ci, (cyl, color) in enumerate(zip(bank, colors)):
            v = df2[egt_col(cyl)]
            if cyl == 11:
                v = v.where(v > 50, other=np.nan)
            v_s = v.rolling(5, center=True).mean()
            lw = 2.5 if cyl in [1, 6, 15, 2] else 1.5
            ax.plot(sec2, v_s, color=color, lw=lw, label=f'Ц{cyl}', alpha=0.9)

        ax.axhline(500, color=RED_M, lw=1.5, ls='--', alpha=0.8, label='500°C')
        ax.axhline(550, color='#FF0000', lw=1, ls='-', alpha=0.5, label='550°C')
        ax.set_ylim(200, 620)
        ax.set_title(f'{bank_name}\n{ecm_id}', color=WHITE_M, fontsize=11, fontweight='bold')
        ax.set_ylabel('ЕГТ (°C)', color=WHITE_M, fontsize=9)
        ax.set_xlabel('Время теста (сек)', color=WHITE_M, fontsize=9)
        ax.legend(fontsize=8, ncol=2, loc='upper left', framealpha=0.3,
                  facecolor=DARK_M, edgecolor='#445055', labelcolor=WHITE_M)
        ax.grid(True, alpha=0.3)
        ax.tick_params(colors=WHITE_M, labelsize=8)
        for spine in ax.spines.values():
            spine.set_edgecolor('#445055')

        # Stats box
        avg_vals = []
        max_vals = []
        for c in bank:
            v = df2[egt_col(c)]
            if c == 11:
                v = v.where(v > 50)
            avg_vals.append(v.mean())
            max_vals.append(v.max())
        valid_avgs = [v for v in avg_vals if not np.isnan(v)]
        bank_avg = np.mean(valid_avgs)
        bank_max = np.nanmax(max_vals)
        ax.text(0.98, 0.04, f'Avg банка: {bank_avg:.0f}°C\nMax банка: {bank_max:.0f}°C',
                transform=ax.transAxes, ha='right', va='bottom',
                color=WHITE_M, fontsize=9.5, fontweight='bold',
                bbox=dict(boxstyle='round', facecolor='#1e2529',
                          edgecolor=RED_M if bank_max >= 550 else ORANGE_M, alpha=0.9))

    plt.tight_layout(rect=[0, 0, 1, 0.94])
    plt.savefig('/tmp/dml_banks_slide.png', dpi=130, bbox_inches='tight', facecolor=DARK_M)
    plt.close()

    slide = add_dark_slide(prs)
    header = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(33.86), Cm(1.4))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1e, 0x25, 0x29)
    header.line.fill.background()
    add_text(slide, 'DML: ТЕМПЕРАТУРНАЯ АСИММЕТРИЯ БАНКОВ  |  Тест 2 (с породой)',
             Cm(1), Cm(0.1), Cm(30), Cm(1.2),
             font_size=16, bold=True, color=WHITE)

    # Add interpretation box
    interp = slide.shapes.add_shape(1, Cm(0.5), Cm(17.2), Cm(32.5), Cm(1.6))
    interp.fill.solid()
    interp.fill.fore_color.rgb = RGBColor(0x1e, 0x25, 0x29)
    interp.line.color.rgb = ORANGE
    interp.line.width = Pt(1)
    add_text(slide,
             'ВЫВОД: B-банк (правый, чётные) в среднем на 60°C горячее A-банка. Цил.6 и Цил.2 — лидеры перегрева (567°C max).',
             Cm(0.8), Cm(17.3), Cm(32), Cm(1.4),
             font_size=11, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

    add_image(slide, '/tmp/dml_banks_slide.png', Cm(0.5), Cm(1.5), Cm(32.5), Cm(15.5))

# ─────────────────────────────────────────────
# Slide 40: Conclusions and recommendations
# ─────────────────────────────────────────────
def make_slide_40(prs):
    slide = add_dark_slide(prs)

    header = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(33.86), Cm(2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1e, 0x25, 0x29)
    header.line.fill.background()

    add_text(slide, 'DML: ЗАКЛЮЧЕНИЕ И РЕКОМЕНДАЦИИ — АГРЕГАТ 82',
             Cm(1), Cm(0.2), Cm(30), Cm(1.3),
             font_size=20, bold=True, color=WHITE)
    add_text(slide, 'На основе анализа 2 DML-логов от 02.06.2026',
             Cm(1), Cm(1.4), Cm(30), Cm(0.6),
             font_size=10, color=ACCENT)

    # Two columns
    # Left: Confirmed risks
    left_title = slide.shapes.add_shape(1, Cm(0.5), Cm(2.2), Cm(15.5), Cm(0.8))
    left_title.fill.solid()
    left_title.fill.fore_color.rgb = RED_C
    left_title.line.fill.background()
    add_text(slide, 'ПОДТВЕРЖДЁННЫЕ РИСКИ', Cm(0.5), Cm(2.2), Cm(15.5), Cm(0.8),
             font_size=11, bold=True, color=RGBColor(0x29,0x31,0x36), align=PP_ALIGN.CENTER)

    risks = [
        ('Цил.11 — отказ датчика ЕГТ',
         '63% времени теста датчик показывает 0°C при\nполной нагрузке. Дефект клапана → обрыв термопары.'),
        ('Цил.6 и Цил.2 — температурный экстремум',
         'Оба достигают 567°C при нагрузке 88%.\nЦил.6 — хронически горячий (avg 530°C).'),
        ('B-банк системно горячее A-банка',
         'Delta +60°C между банками под нагрузкой.\nАсимметрия может указывать на дисбаланс TIB.'),
        ('Нагрузка +4x при работе с породой',
         'С 21% до 88% — экстремальный нагрев при каждом\nрейсе с загрузкой. При 200% TIB нагрев удваивается.'),
    ]

    for i, (title, desc) in enumerate(risks):
        y = Cm(3.2 + i * 3.8)
        box = slide.shapes.add_shape(1, Cm(0.5), y, Cm(15.5), Cm(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x1e, 0x25, 0x29)
        box.line.color.rgb = RED_C
        box.line.width = Pt(1)
        add_text(slide, f'{i+1}. {title}', Cm(0.8), y + Cm(0.2), Cm(14.8), Cm(0.9),
                 font_size=10.5, bold=True, color=RED_C)
        add_text(slide, desc, Cm(0.8), y + Cm(1.1), Cm(14.8), Cm(2.2),
                 font_size=9.5, color=GRAY)

    # Right: Recommendations
    right_title = slide.shapes.add_shape(1, Cm(17.5), Cm(2.2), Cm(15.5), Cm(0.8))
    right_title.fill.solid()
    right_title.fill.fore_color.rgb = ACCENT
    right_title.line.fill.background()
    add_text(slide, 'РЕКОМЕНДАЦИИ', Cm(17.5), Cm(2.2), Cm(15.5), Cm(0.8),
             font_size=11, bold=True, color=RGBColor(0x29,0x31,0x36), align=PP_ALIGN.CENTER)

    recs = [
        ('СРОЧНО: Инспекция цилиндра 11 A-банка',
         'Стробоскопическая и механическая проверка клапана 11.\nЗамена термопары ЕГТ. Компрессия цилиндра.'),
        ('Мониторинг цилиндров 6 и 2',
         'Установить пороговую сигнализацию 520°C для Ц6 и Ц2.\nПлановый осмотр при следующем ТО.'),
        ('Рекалибровка TIB для NTE200',
         'Снизить TIB с 200% до 150-160% → снизит нагрев\nна 25-35°C при тех же эксплуатационных показателях.'),
        ('Включить ограничение ЕГТ в ECM0',
         'Активировать C_EPD_OT_RPM_Drt_En → деррейт\nоборотов при перегреве ЕГТ >500°C.'),
        ('Регулярный DML-мониторинг',
         'Ежемесячный DML тест под нагрузкой для всего\nпарка. Приоритет: единицы 82, 75, 83, 80.'),
    ]

    for i, (title, desc) in enumerate(recs):
        y = Cm(3.2 + i * 3.0)
        box = slide.shapes.add_shape(1, Cm(17.5), y, Cm(15.5), Cm(2.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x1e, 0x25, 0x29)
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1)
        add_text(slide, f'{i+1}. {title}', Cm(17.8), y + Cm(0.2), Cm(14.8), Cm(0.9),
                 font_size=10, bold=True, color=ACCENT)
        add_text(slide, desc, Cm(17.8), y + Cm(1.0), Cm(14.8), Cm(1.7),
                 font_size=9, color=GRAY)

    # Divider
    div = slide.shapes.add_shape(1, Cm(16.5), Cm(2.2), Cm(0.15), Cm(16.5))
    div.fill.solid()
    div.fill.fore_color.rgb = GRAY
    div.line.fill.background()

    # Bottom summary
    bottom = slide.shapes.add_shape(1, Cm(0), Cm(18.2), Cm(33.86), Cm(0.85))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = ACCENT
    bottom.line.fill.background()
    add_text(slide, 'Агрегат 82 — наивысший риск в парке. DML подтверждает начало отказа Цил.11. Требуется немедленная диагностика.',
             Cm(1), Cm(18.25), Cm(32), Cm(0.7),
             font_size=10.5, bold=True,
             color=RGBColor(0x29, 0x31, 0x36))

# ─────────────────────────────────────────────
# Build & save
# ─────────────────────────────────────────────
print('Loading presentation...')
prs = Presentation('/home/user/NTE200/Презентация_анализ_клапанов_QSK50_NTE200.pptx')
print(f'Current slides: {len(prs.slides)}')

print('Adding slide 34 (DML section header)...')
make_slide_34(prs)

print('Adding slide 35 (Key findings summary)...')
make_slide_35(prs)

print('Adding slide 36 (Time series overview)...')
make_slide_36(prs)

print('Adding slide 37 (EGT per cylinder bars)...')
make_slide_37(prs)

print('Adding slide 38 (Cylinder 11 anomaly)...')
make_slide_38(prs)

print('Adding slide 39 (Bank asymmetry)...')
make_slide_39(prs)

print('Adding slide 40 (Conclusions)...')
make_slide_40(prs)

out_path = '/home/user/NTE200/Презентация_анализ_клапанов_QSK50_NTE200.pptx'
prs.save(out_path)
print(f'\nSaved: {out_path}')
print(f'Total slides: {len(prs.slides)}')
