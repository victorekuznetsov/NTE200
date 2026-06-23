#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Рендер pptx -> PDF через matplotlib (без LibreOffice). Рисует прямоугольники, текст, картинки по их позициям."""
import sys, io
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle
from matplotlib.backends.backend_pdf import PdfPages
from PIL import Image
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE

SRC=sys.argv[1] if len(sys.argv)>1 else "05 Отчёты/Презентации/Отчёт_Развитие_NTE200_ГБЦ_QSK50.pptx"
OUT=sys.argv[2] if len(sys.argv)>2 else "/tmp/claude-0/-home-user-NTE200/26b161fb-3e80-5997-bb3a-c0a0b40e0304/scratchpad/deck.pdf"
prs=Presentation(SRC)
W=prs.slide_width/914400.0; H=prs.slide_height/914400.0
def hexcol(rgb):
    try: return '#%02X%02X%02X'%(rgb[0],rgb[1],rgb[2])
    except: return None
def fill_color(sh):
    try:
        if sh.fill.type==1:  # solid
            return hexcol(sh.fill.fore_color.rgb)
    except Exception: pass
    return None

with PdfPages(OUT) as pdf:
    for slide in prs.slides:
        fig=plt.figure(figsize=(W,H)); ax=fig.add_axes([0,0,1,1]); ax.set_xlim(0,W); ax.set_ylim(0,H); ax.axis('off')
        # background: first full-size rect
        for sh in slide.shapes:
            try:
                l=sh.left/914400.0; t=sh.top/914400.0; w=sh.width/914400.0; h=sh.height/914400.0
            except Exception:
                continue
            y=H-t-h
            # picture
            if sh.shape_type==MSO_SHAPE_TYPE.PICTURE:
                try:
                    im=Image.open(io.BytesIO(sh.image.blob))
                    ax.imshow(im, extent=[l,l+w,y,y+h], aspect='auto', zorder=2)
                except Exception: pass
                continue
            # rectangle fill
            fc=fill_color(sh)
            if fc:
                ax.add_patch(Rectangle((l,y),w,h,facecolor=fc,edgecolor='none',zorder=1))
            # text
            if sh.has_text_frame:
                tf=sh.text_frame
                va=tf.vertical_anchor
                # gather paragraphs
                paras=tf.paragraphs
                # vertical placement
                n=len([p for p in paras])
                for i,p in enumerate(paras):
                    runs=p.runs
                    if not runs: continue
                    txt=''.join(r.text for r in runs)
                    if not txt.strip(): continue
                    for a,b in [('🔴','●'),('🟠','◆'),('🟢','○'),('✅','✓'),('❌','✗'),('⚠️','!'),('⚠','!'),('🌡','t°'),('⛔','✗'),('🧠',''),('🛢️',''),('▲','▲'),('•','•'),('→','→')]:
                        txt=txt.replace(a,b)
                    r0=runs[0]; fs=(r0.font.size.pt if r0.font.size else 14)
                    col=hexcol(r0.font.color.rgb) if (r0.font.color and r0.font.color.type is not None) else '#293136'
                    bold=bool(r0.font.bold)
                    al=p.alignment
                    if al==PP_ALIGN.CENTER: ha='center'; tx=l+w/2
                    elif al==PP_ALIGN.RIGHT: ha='right'; tx=l+w-0.05
                    else: ha='left'; tx=l+0.07
                    # y position: stack lines from top of box
                    lh=fs/72.0*1.25
                    if va==MSO_ANCHOR.MIDDLE: ty=y+h/2 + (n/2-i-0.5)*lh
                    else: ty=y+h-0.07-(i+0.7)*lh
                    ax.text(tx,ty,txt,ha=ha,va='center',fontsize=fs,color=col or '#293136',
                            fontweight='bold' if bold else 'normal',wrap=True,zorder=3)
        pdf.savefig(fig); plt.close(fig)
print("saved",OUT)
