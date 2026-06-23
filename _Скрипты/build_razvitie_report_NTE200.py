#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Светлая презентация «Развитие» — отчёт ТС и отдела надёжности по NTE200 (ГБЦ QSK50).
Палитра АО Развитие (light): фон #F5F5F5, текст #293136, акцент #3EF0AF."""
import os, tempfile
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── палитра (light) ──
BG=RGBColor(0xF5,0xF5,0xF5); INK=RGBColor(0x29,0x31,0x36); ACC=RGBColor(0x3E,0xF0,0xAF)
WHITE=RGBColor(0xFF,0xFF,0xFF); RED=RGBColor(0xE0,0x52,0x52); ORG=RGBColor(0xE0,0x8C,0x2E)
GRN=RGBColor(0x2E,0xA0,0x43); BLU=RGBColor(0x4E,0xC9,0xF0); GRY=RGBColor(0x70,0x78,0x80)
HX=lambda c:'#%02X%02X%02X'%(c[0],c[1],c[2])
EMU=914400
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
OUT="05 Отчёты/Презентации/Отчёт_Развитие_NTE200_ГБЦ_QSK50.pptx"
tmp=tempfile.mkdtemp()

def slide(bg=BG):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    r.shadow.inherit=False
    return s
def rect(s,x,y,w,h,fill,line=None,lw=0):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; return sh
def txt(s,x,y,w,h,text,size=14,bold=False,color=INK,align=PP_ALIGN.LEFT,italic=False,anchor=MSO_ANCHOR.TOP,font='Calibri'):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    lines=text.split('\n')
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; f=r.font; f.size=Pt(size); f.bold=bold; f.italic=italic
        f.color.rgb=color; f.name=font
    return tb
def footer(s):
    rect(s,0,7.12,13.333,0.38,INK)
    txt(s,0.4,7.16,12.5,0.3,"АО Развитие   •   QSK50 / NTE200   •   Полюс Магадан   •   2026   •   Конфиденциально",9,False,RGBColor(0xCF,0xD6,0xDA))
def header(s,title,kicker=None):
    rect(s,0,0,13.333,1.15,WHITE); rect(s,0,1.15,13.333,0.06,ACC)
    rect(s,0.0,0,0.22,1.15,ACC)
    if kicker: txt(s,0.5,0.18,12,0.3,kicker,11,True,GRY)
    txt(s,0.5,0.45,12.3,0.6,title,24,True,INK)
def chip(s,x,y,w,text,color):
    rect(s,x,y,w,0.34,color); txt(s,x+0.08,y+0.03,w-0.12,0.3,text,11,True,WHITE)

def table(s,x,y,w,col_w,rows,header_fill=INK,fsize=11,rowh=0.34,first_bold=True):
    # rows: list of list[str]; col_w fractions sum to 1
    cx=x
    th=rowh*len(rows)
    for ci,frac in enumerate(col_w):
        cw=w*frac; cy=y
        for ri,row in enumerate(rows):
            fill=header_fill if ri==0 else (WHITE if ri%2 else RGBColor(0xEC,0xEF,0xF1))
            rect(s,cx,cy,cw,rowh,fill)
            col=WHITE if ri==0 else INK
            txt(s,cx+0.06,cy+0.02,cw-0.1,rowh-0.02,str(row[ci]),fsize,(ri==0 and first_bold),col,anchor=MSO_ANCHOR.MIDDLE)
            cy+=rowh
        cx+=cw

def chart_egt(path):
    cyl=list(range(1,17));
    egt={2:567,6:567,4:555}; base={1:455,3:470,5:470,7:480,9:475,13:485,15:478,8:520,10:515,12:525,14:500,16:495}
    vals=[]; cols=[]
    for c in cyl:
        v=egt.get(c, base.get(c, 490 if c not in (11,) else 98)); vals.append(v)
        cols.append(HX(RED) if v>=540 else (HX(ORG) if v>=510 else (HX(GRY) if v<200 else HX(INK))))
    fig,ax=plt.subplots(figsize=(7.6,3.0),dpi=150); fig.patch.set_facecolor(HX(BG)); ax.set_facecolor(HX(BG))
    ax.bar([str(c) for c in cyl],vals,color=cols)
    ax.axhline(520,color=HX(RED),ls='--',lw=1.4,label='Лимит Cummins 520 °C')
    ax.set_ylabel('EGT, °C',color=HX(INK)); ax.set_title('EGT по цилиндрам — DML #82 (рейс с породой)',color=HX(INK),fontsize=11,weight='bold')
    ax.tick_params(colors=HX(INK)); [sp.set_color(HX(GRY)) for sp in ax.spines.values()]
    ax.legend(fontsize=8,facecolor=HX(WHITE),edgecolor=HX(GRY)); ax.text(10,120,'Ц11: датчик FC1531',color=HX(GRY),fontsize=8)
    plt.tight_layout(); plt.savefig(path,facecolor=HX(BG)); plt.close()

def chart_1010(path):
    cats=['Среднее','>110 %','>120 %']; nte=[99.7,6.0,0.01]; pol=[100,10,0.0]
    fig,ax=plt.subplots(figsize=(5.4,3.0),dpi=150); fig.patch.set_facecolor(HX(BG)); ax.set_facecolor(HX(BG))
    import numpy as np; x=np.arange(len(cats)); w=0.36
    ax.bar(x-w/2,nte,w,color=HX(ACC),label='NTE200 (от 180 т)')
    ax.bar(x+w/2,pol,w,color=HX(GRY),label='Лимит 10/10/20')
    ax.set_xticks(x); ax.set_xticklabels(cats,color=HX(INK)); ax.tick_params(colors=HX(INK))
    [sp.set_color(HX(GRY)) for sp in ax.spines.values()]
    ax.set_title('Соответствие политике 10/10/20',color=HX(INK),fontsize=11,weight='bold')
    ax.legend(fontsize=8,facecolor=HX(WHITE),edgecolor=HX(GRY))
    for i,v in enumerate(nte): ax.text(i-w/2,v+1,f'{v}',ha='center',color=HX(INK),fontsize=8)
    plt.tight_layout(); plt.savefig(path,facecolor=HX(BG)); plt.close()

def chart_oil(path):
    import numpy as np; labels=['P50','P90','P95']; nte=[0,5,6]; k730=[4,14,20]
    fig,ax=plt.subplots(figsize=(5.4,3.0),dpi=150); fig.patch.set_facecolor(HX(BG)); ax.set_facecolor(HX(BG))
    x=np.arange(3); w=0.36
    ax.bar(x-w/2,nte,w,color=HX(ACC),label='NTE200')
    ax.bar(x+w/2,k730,w,color=HX(BLU),label='730E')
    ax.set_xticks(x); ax.set_xticklabels(labels,color=HX(INK)); ax.tick_params(colors=HX(INK))
    [sp.set_color(HX(GRY)) for sp in ax.spines.values()]
    ax.set_ylabel('Fe, ppm',color=HX(INK)); ax.set_title('Износ Fe в масле двигателя (ниже = лучше)',color=HX(INK),fontsize=11,weight='bold')
    ax.legend(fontsize=8,facecolor=HX(WHITE),edgecolor=HX(GRY))
    plt.tight_layout(); plt.savefig(path,facecolor=HX(BG)); plt.close()

p_egt=os.path.join(tmp,'egt.png'); chart_egt(p_egt)
p_10=os.path.join(tmp,'l10.png'); chart_1010(p_10)
p_oil=os.path.join(tmp,'oil.png'); chart_oil(p_oil)

# 1 Title
s=slide(INK)
rect(s,0,3.05,13.333,0.06,ACC)
txt(s,0.8,1.5,11.7,0.6,"АО «РАЗВИТИЕ»",16,True,ACC)
txt(s,0.8,2.15,11.7,1.4,"Отказы ГБЦ и клапанов Cummins QSK50\nна самосвалах NTE200",30,True,WHITE)
txt(s,0.8,3.25,11.7,0.8,"Отчёт технической службы и отдела надёжности",16,False,RGBColor(0xCF,0xD6,0xDA))
txt(s,0.8,4.2,11.7,0.5,"Полюс Магадан · номинал 180 т · контрольная группа 730E-DC",13,False,GRY)
txt(s,0.8,6.4,11.7,0.4,"2026 · Конфиденциально · все числа из первичных данных",11,True,GRY)

# 2 Резюме
s=slide(); header(s,"Резюме: коренная причина и факторы","РЕЗЮМЕ")
rect(s,0.5,1.45,12.33,1.0,WHITE); rect(s,0.5,1.45,0.12,1.0,RED)
txt(s,0.75,1.55,12,0.85,"Коренная причина — системное снижение тепловой защиты в калибровке ЭБУ NTE200 (15 параметров):\nTIB 200% vs 100%, отключены триггеры перегрева (масло/ОЖ/наддув/топливо), задержки derate 43–51 с, обороты 2100 vs 1990.\nУсугубляющий конструктивный фактор — отсутствие наплавки Stellite на выпускном клапане (исполнение NHL).",13,False,INK)
rows=[["Фактор","Вердикт","Основание"],
["Калибровка ЭБУ (TIB + защиты)","Коренная причина","15 параметров ✅"],
["Отсутствие Stellite","Вторичный (конструкция)","письма NHL/Cummins ✅"],
["EGT 567 °C (лимит 520)","Прямое доказательство","DML #82"],
["Систематическая перегрузка","НЕ подтверждена","10/10/20 соблюдается ✅"],
["Марка масла","Не дифференцирует","= 730E (CI-4) ✅"],
["Топливо (сера)","Исключено","Е5, 2–10 мг/кг"]]
table(s,0.5,2.65,12.33,[0.36,0.30,0.34],rows,fsize=12,rowh=0.5); footer(s)

# 3 Section: ТС
s=slide(INK); rect(s,0,3.2,13.333,0.06,ACC)
txt(s,0.8,2.6,11.7,1,"ЧАСТЬ I",16,True,ACC)
txt(s,0.8,3.3,11.7,1,"Техническая служба — наблюдаемые факты",26,True,WHITE)

# 4 Парк и статус ECM
s=slide(); header(s,"Парк NTE200 и статус двигателей (ECM)","ОБЪЕКТ")
txt(s,0.5,1.4,12.3,0.6,"NTE200: карьерный самосвал 180 т, двигатель Cummins QSK50 MCRS (V16, Tier 2, без DPF/EGR), электропривод GE. Контроль — 730E-DC.",13)
rows=[["Статус по ECM (38 машин)","Машин"],["Отказ","10"],["EGT/связанные","10"],["Норма","17"],["Катастрофический (#83)","1"]]
table(s,0.5,2.3,5.6,[0.62,0.38],rows,fsize=13,rowh=0.5)
rect(s,6.6,2.3,6.2,2.5,WHITE); rect(s,6.6,2.3,0.12,2.5,ORG)
txt(s,6.85,2.45,5.8,2.3,"21 из 38 машин (55%) — в отказном или EGT-связанном статусе.\n\nПрошивки ЭБУ AQ60809.07 (23) и .08 (13) сосуществуют в парке.\n\nМакс. обороты по ECM ~2030–2060 (что согласуется с повышенным лимитом 2100 vs 1990 на 730E).",13,False,INK); footer(s)

# 5 Реестр отказов
s=slide(); header(s,"Реестр отказов ГБЦ / клапанов","РЕМОНТЫ")
rows=[["Маш","Дата","Наработка, ч","Операция"],
["#51","17.07.2025","~10 783","Замена ГБЦ (трещина → ОЖ в масло)"],
["#58","10.08.2025","~10 702","Клапан/седло + толкатель"],
["#47","11.08 / 07.12.2025","~12 890","Замена ГБЦ (даты расходятся)"],
["#48","23.11 / 17.12.2025","~14 997","Замена ГБЦ"],
["#43","21.01.2026","—","Клапаны + форсунки"],
["#83","~03.2026","2 776","Замена ГБЦ (ранний, «Катастроф.»)"],
["#48 / #58","03.2026","~15 126","Повторные ремонты"],
["#78 / #55 / #69","05.2026","—","Клапаны"]]
table(s,0.5,1.5,12.33,[0.14,0.22,0.18,0.46],rows,fsize=12,rowh=0.52)
txt(s,0.5,6.35,12.3,0.5,"Всего ~11 машин / ~13 событий. Прежняя оценка «30 событий / 24+ ед.» не подтверждается.",11,True,RED); footer(s)

# 6 DML EGT
s=slide(); header(s,"DML — температуры выхлопа (EGT)","ДИАГНОСТИКА · DML")
s.shapes.add_picture(p_egt,Inches(0.4),Inches(1.4),width=Inches(7.7))
rect(s,8.3,1.5,4.6,4.6,WHITE); rect(s,8.3,1.5,0.12,4.6,RED)
txt(s,8.5,1.65,4.3,4.4,"#82, рейс с породой:\n\n• Ц2 и Ц6 = 567 °C при лимите 520 °C (+47).\n\n• Банк B горячее банка A на ~60 °C — дисбаланс корректирующей дозы (TIB).\n\n• Датчик Ц11 ≈ 98 °C (FC1531) — обрыв термопары, индикатор обгоревшего клапана.\n\n• Без нагрузки 287 °C → с нагрузкой 474 °C (рост +187).",12.5,False,INK); footer(s)

# 7 ECM faults
s=slide(); header(s,"Ошибки ECM (Fault Codes) — флот","ДИАГНОСТИКА · ECM")
rows=[["FC","Машин","Описание"],
["611","46","Остановка горячего двигателя 🌡"],
["426","46","Канал J1939 — хаотичный характер"],
["234","14","Высокая частота вращения (Red)"],
["146","9","Температура ОЖ высокая 🌡"],
["197/235","18/13","Уровень охлаждающей жидкости"]]
table(s,0.5,1.5,7.0,[0.18,0.2,0.62],rows,fsize=12,rowh=0.5)
rect(s,7.8,1.5,5.0,4.4,WHITE); rect(s,7.8,1.5,0.12,4.4,ACC)
txt(s,8.0,1.65,4.7,4.2,"АКТИВНЫЕ неисправности:\n\n• #61 → FC 628 «Малое отклонение EGT» — маркер теплового дисбаланса.\n\n• #67 → FC 1619 «датчик EGT цилиндра» — отказ датчика.\n\n• #51 → FC 1549/323/1557 — цепи форсунок.\n\n361 файл INSITE, 106 кодов. «Горячий стоп» (FC611) — на 46 машинах.",12.5,False,INK); footer(s)

# 8 Топливо и масло
s=slide(); header(s,"Топливо и масло","ЭКСПЛУАТАЦИЯ")
s.shapes.add_picture(p_oil,Inches(0.4),Inches(1.5),width=Inches(5.5))
rect(s,6.2,1.5,6.6,4.5,WHITE); rect(s,6.2,1.5,0.12,4.5,GRN)
txt(s,6.42,1.62,6.3,4.3,"Топливо: 25+ паспортов, сера 2–10 мг/кг (Е5, в 2,4× ниже порога Cummins). Исключено.\n\nМасло ДВС (2024–2026, 1827 проб NTE200): Газпромнефть G-Profi MSI Plus + TEBOIL Super HPD Plus, оба 15W-40 CI-4. = как на 730E → не дифференцирует.\n\nСостояние масла отличное: TBN 9,7 (норма ≥2,5), САЖА 0,3 (низкая), окисление 0,1, Fe P95=6.\n\nНизкая сажа ОПРОВЕРГАЕТ версию NHL про «сажа/зола → износ клапанов». Отказ ТЕРМИЧЕСКИЙ.",12,False,INK); footer(s)

# 8b Чистота топлива (ISO 4406) — технический анализ
s=slide(); header(s,"Чистота топлива (ISO 4406) — технический анализ","ТОПЛИВО · ЧИСТОТА")
chip(s,0.5,1.35,3.0,"Cummins (бак): 18/16/13",INK)
chip(s,3.7,1.35,3.9,"Вход в цилиндры: 12/9/6",ORG)
txt(s,7.8,1.36,5.0,0.34,"MCRS: 2100 бар, зазор форсунок 2–5 мкм",10,True,GRY)
rows=[["Маш","Бак","После грубой","После тонкой","vs 18/16/13","vs 12/9/6"],
["№46","19/16/15","20/17/16","17/13/13","✅","❌"],
["№47","19/17/16","21/19/18","19/16/14","❌","❌"],
["№81","20/16/15","21/19/17","21/18/17","❌","❌"],
["№82","21/17/15","23/20/19","21/16/14","❌","❌"]]
table(s,0.5,1.85,8.2,[0.12,0.2,0.24,0.24,0.1,0.1],rows,fsize=12,rowh=0.5)
rect(s,8.9,1.85,3.95,4.3,WHITE); rect(s,8.9,1.85,0.12,4.3,ORG)
txt(s,9.1,1.98,3.6,4.15,"• Топливо в баке грязнее нормы Cummins на всех машинах.\n\n• После фильтрации в 18/16/13 укладывается только №46; цель 12/9/6 не достигается нигде (≈32× грязнее по >4 мкм).\n\n• Грубый фильтр местами УХУДШАЕТ чистоту (№82: 21→23) — насыщен/деградирован.\n\n• Топливо одно для NTE200 и 730E → не дифференцирует отказ, но угрожает форсункам MCRS (отказы форсунок #51/#47/#43).",11.5,False,INK)
txt(s,0.5,6.5,8.2,0.4,"Вывод: бортовая фильтрация недостаточна; требуется доп. фильтрация до бака (Cummins) и тонкая под 12/9/6 (Fleetguard).",10.5,True,RED); footer(s)

# 9 Section: надёжность
s=slide(INK); rect(s,0,3.2,13.333,0.06,ACC)
txt(s,0.8,2.6,11.7,1,"ЧАСТЬ II",16,True,ACC)
txt(s,0.8,3.3,11.7,1,"Отдел надёжности — анализ",26,True,WHITE)

# 10 10/10/20
s=slide(); header(s,"Грузоподъёмность и политика 10/10/20","НАГРУЗКА")
s.shapes.add_picture(p_10,Inches(0.4),Inches(1.5),width=Inches(5.6))
rows=[["Метрика","Ср.",">110%",">120%","Макс"],
["FinalPayload/180 т","99,7%","6,0%","0,01%","133%"],
["Родное LoadPercentage","102,8%","8,2%","0,10%","145%"]]
table(s,6.2,1.6,6.7,[0.40,0.15,0.15,0.15,0.15],rows,fsize=11,rowh=0.5)
rect(s,6.2,3.0,6.7,3.0,WHITE); rect(s,6.2,3.0,0.12,3.0,GRN)
txt(s,6.42,3.12,6.3,2.8,"Вывод: парк СООТВЕТСТВУЕТ 10/10/20 (среднее ≈100%, >120% практически нет).\n\nСистематическая перегрузка как значимый фактор НЕ подтверждается.\n\nПрежний тезис «макс 295 т / нарушены все критерии» — ошибка (номинал брался 220 т вместо 180 т). Корректно: макс 133% ≈ 239 т.",12.5,False,INK); footer(s)

# 11 Производительность
s=slide(); header(s,"Производительность парка","ПРОИЗВОДИТЕЛЬНОСТЬ")
kpis=[("65,5 млн т","перевезено всего"),("365 069","рейсов (131 файл)"),("30,6 мин","средний цикл"),
("4,32 км","гружёный пробег"),("31 км/ч","скорость гружёного"),("6 888","рейсов/машину")]
x=0.5; y=1.6
for i,(v,l) in enumerate(kpis):
    cx=0.5+(i%3)*4.25; cy=1.6+(i//3)*2.3
    rect(s,cx,cy,4.0,2.0,WHITE); rect(s,cx,cy,4.0,0.1,ACC)
    txt(s,cx+0.2,cy+0.4,3.6,0.8,v,28,True,INK)
    txt(s,cx+0.2,cy+1.35,3.6,0.5,l,13,False,GRY)
footer(s)

# 12 Анализ всех отказов
s=slide(); header(s,"Анализ всех отказов и возможных причин","АНАЛИЗ")
rows=[["Возможная причина","Вердикт"],
["Калибровка ЭБУ (TIB + отключённые защиты)","🔴 Коренная"],
["Отсутствие Stellite (NHL)","🟠 Вторичный"],
["Перегрузка","🟢 Не значима (10/10/20)"],
["Абразивный/зольный износ (версия NHL)","🟢 Ослаблена (низкий Fe)"],
["Марка/класс масла","🟢 Не дифференцирует"],
["Топливо (сера)","🟢 Исключено"],
["Естественный износ","🟢 Исключён (#83 при 2 776 ч)"]]
table(s,0.5,1.5,12.33,[0.6,0.4],rows,fsize=13,rowh=0.55); footer(s)

# 13 + 14 ECM 15 params
def egt_param_slide(part):
    s=slide(); header(s,f"15 параметров ЭБУ: NTE200 vs 730E ({part})","КАЛИБРОВКА ЭБУ")
    return s
rows1=[["Параметр","730E","NTE200","Влияние"],
["TIB верхний предел коррекции топлива","100%","200%","переобогащение / поздний впрыск"],
["Снижение оборотов по T масла","Вкл","ВЫКЛ","нет защиты от перегрева масла"],
["Триггер high T наддувочного воздуха","54 °C","ВЫКЛ","нет тревоги перегрева наддува"],
["Триггер high T ОЖ","85 °C","ВЫКЛ","нет сигнала перегрева ОЖ"],
["Триггер high T топлива","68 °C","ВЫКЛ","нет тревоги перегрева топлива"],
["Задержка снижения мощности (масло)","0 с","51 с","полная мощность в перегреве"],
["Задержка снижения мощности (ОЖ)","0 с","43 с","43 с без снижения нагрузки"]]
s=egt_param_slide("1/2"); table(s,0.5,1.5,12.33,[0.40,0.16,0.16,0.28],rows1,fsize=11.5,rowh=0.62); footer(s)
rows2=[["Параметр","730E","NTE200","Влияние"],
["Макс. рабочие обороты","1990","2100","EGT +25–35 °C, удары клапанов"],
["Макс. обороты (реф. 2/3)","1990","2225,8","раскрутка в отд. режимах"],
["Макс. наклон регулятора (droop)","0%","20%","±20% колебаний оборотов"],
["Горячий стоп — момент (верх.)","5509","3650 Н·м","срабатывает раньше"],
["Горячий стоп — момент (нижн.)","4722","1850 Н·м","нестабильный гистерезис"],
["Топливный код (карта впрыска)","FC0BU52","FC0WH95","иная модель сгорания"],
["Пуск вентилятора охл. наддува","50 °C","46 °C","компенсация, не защита"]]
s=egt_param_slide("2/2"); table(s,0.5,1.5,12.33,[0.40,0.16,0.16,0.28],rows2,fsize=11.5,rowh=0.62)
txt(s,0.5,6.5,12.3,0.5,"Вывод: NTE200 отключает все триггеры перегрева и допускает более высокие обороты/нагрузку.",12,True,RED); footer(s)

# 14b Аналитическое сравнение параметров, влияющих на отказ
s=slide(); header(s,"Аналитическое сравнение параметров, влияющих на отказ","СИНТЕЗ")
rows=[["Параметр","NTE200","730E","Дифференцирует отказ?"],
["TIB коррекция топлива","200%","100%","🔴 ДА — калибровка"],
["Защита по T масла (derate)","ВЫКЛ","ВКЛ","🔴 ДА — калибровка"],
["Триггеры перегрева (ОЖ/наддув/топл)","ВЫКЛ","ВКЛ","🔴 ДА — калибровка"],
["Макс. обороты, об/мин","2100","1990","🔴 ДА — калибровка"],
["EGT пиковая (цил.), °C","567","~540","🔴 следствие"],
["Наплавка Stellite клапана","нет","нет*","🟠 фон (конструкция)"],
["Нагрузка (10/10/20)","соблюд.","соблюд.","🟢 нет"],
["Масло: марка / класс","15W-40 CI-4","15W-40 CI-4","🟢 нет"],
["Сажа в масле (ST)","0,3","0,4","🟢 нет"],
["TBN / износ Fe P95","9,7 / 6","9,5 / 16","🟢 нет (Fe ниже)"],
["Топливо (сера), мг/кг","≤10","≤10","🟢 нет"]]
table(s,0.4,1.45,12.5,[0.34,0.18,0.18,0.30],rows,fsize=11,rowh=0.45)
txt(s,0.4,6.55,12.5,0.5,"Все дифференцирующие параметры (🔴) — в калибровке ЭБУ. Эксплуатация (масло/нагрузка/топливо, 🟢) — не различает парки.",11,True,RED); footer(s)

# 15 Письма NHL/Cummins
s=slide(); header(s,"Оценка писем NHL и Cummins","ПЕРЕПИСКА")
rows=[["Тезис письма","Вердикт"],
["NHL: износ от зольности масла (SAPS)","❌ смещает на эксплуатацию (низкий Fe)"],
["NHL: пыль из карьера","❌ переоценка (Si подчинён)"],
["NHL: Stellite не нужен «для rated»","❌ ложный (NTE200 ≠ rated: 567 °C)"],
["NHL: регулировка клапанов на 1500/5000 ч","⚠️ косвенное признание проблемы"],
["Cummins: масло CES20078 достаточно","❌ уход от корня"],
["Cummins: на NTE200 защита по T масла отключена","✅ ПОДТВЕРЖДАЕТ нашу находку"],
["Cummins: «не поставляет Stellite»","❌ уход от ответственности"]]
table(s,0.5,1.5,12.33,[0.55,0.45],rows,fsize=12,rowh=0.55); footer(s)

# 16 Достоверность
s=slide(); header(s,"Достоверность и пробелы данных","ОБЪЕКТИВНОСТЬ")
rows=[["Пункт","Статус"],
["Даты ремонтов #47, #48","❌ расходятся между источниками"],
["Счёт ремонтов («30/24+» vs ~13)","❌ прежняя оценка завышена"],
["«0 отказов 730E»","⚠️ нет полного реестра ремонтов 730E"],
["Цитаты Cummins 378-010/011, лимит 520 °C","⚠️ без дословного первоисточника"],
["EGT 567 (#82 DML) vs 329 (#83 INSITE)","⚠️ разные машины/условия"],
["Тоннаж 220 т / «295 т»","❌ исправлено на 180 т (239 т)"],
["Период проб масла","✅ 2024-02 — 2025-12"]]
table(s,0.5,1.5,12.33,[0.6,0.4],rows,fsize=12,rowh=0.55); footer(s)

# 17 Рекомендации
s=slide(); header(s,"Рекомендации","ДЕЙСТВИЯ")
chip(s,0.5,1.4,2.2,"НЕМЕДЛЕННО",RED)
txt(s,0.5,1.85,12.3,1.5,"• Перекалибровать ЭБУ NTE200 под профиль 730E: TIB 100%, включить защиту по T масла и триггеры перегрева, убрать задержки derate, снизить max обороты до 1990.\n• Контрольный DML после перекалибровки — подтвердить снижение EGT.\n• Заменить датчики EGT с FC1531/1619 (#82, #67).",13)
chip(s,0.5,3.5,2.6,"СРЕДНЕСРОЧНО",ORG)
txt(s,0.5,3.95,12.3,1.3,"• Запрос NHL/Cummins: обоснование исключения Stellite и отключения защит; поставка клапанов с наплавкой.\n• Унифицировать прошивку ЭБУ (.07/.08).\n• Сопоставить FC611/EGT-коды с историей ремонтов.",13)
chip(s,0.5,5.3,2.4,"ДОЛГОСРОЧНО",GRN)
txt(s,0.5,5.75,12.3,1.0,"• Мониторинг EGT и FC611 в реальном времени.\n• Устранить нарушение по паспорту топлива №246.",13); footer(s)

# 18 Closing
s=slide(INK); rect(s,0,3.5,13.333,0.06,ACC)
txt(s,0.8,2.7,11.7,1,"Коренная причина — калибровка ЭБУ (TIB + отключённые защиты),",20,True,WHITE)
txt(s,0.8,3.7,11.7,1,"усугублённая отсутствием наплавки Stellite.",20,True,WHITE)
txt(s,0.8,4.8,11.7,0.5,"Перегрузка и масло — не дифференцирующие факторы (подтверждено данными).",14,False,ACC)
txt(s,0.8,6.5,11.7,0.4,"АО «Развитие» · техническая служба + отдел надёжности · 2026",11,True,GRY)

prs.save(OUT)
print("saved:",OUT,"| слайдов:",len(prs.slides._sldIdLst))
