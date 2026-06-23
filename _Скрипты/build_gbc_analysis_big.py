#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Большая презентация «Анализ отказов ГБЦ QSK50 — NTE200» по шаблону (структура 100-слайдов),
светлый фон АО Развитие (#F5F5F5 / #293136 / #3EF0AF), расширено до 200+ слайдов, data-driven."""
import os, json, math, tempfile, glob, re, statistics, datetime, warnings; warnings.filterwarnings('ignore')
from collections import defaultdict
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── палитра Развитие (light) ──
BG=RGBColor(0xF5,0xF5,0xF5); INK=RGBColor(0x29,0x31,0x36); ACC=RGBColor(0x3E,0xF0,0xAF)
WHITE=RGBColor(0xFF,0xFF,0xFF); RED=RGBColor(0xE0,0x52,0x52); ORG=RGBColor(0xE0,0x8C,0x2E)
GRN=RGBColor(0x2E,0xA0,0x43); BLU=RGBColor(0x4E,0xC9,0xF0); GRY=RGBColor(0x70,0x78,0x80)
PANEL=RGBColor(0xEC,0xEF,0xF1)
HX=lambda c:'#%02X%02X%02X'%(c[0],c[1],c[2])
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
TMP=tempfile.mkdtemp()
W,H=13.333,7.5

def slide(bg=BG):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s
def rect(s,x,y,w,h,fill,line=None,lw=0):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; return sh
def txt(s,x,y,w,h,text,size=14,bold=False,color=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,italic=False):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(str(text).split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; f=r.font; f.size=Pt(size); f.bold=bold; f.italic=italic; f.color.rgb=color; f.name='Calibri'
    return tb
def header(s,title,kicker=None,sub=None):
    rect(s,0,0,W,1.12,WHITE); rect(s,0,1.12,W,0.05,ACC); rect(s,0,0,0.18,1.12,ACC)
    if kicker: txt(s,0.45,0.14,12.4,0.3,kicker,11,True,GRY)
    txt(s,0.45,0.40,12.4,0.62,title,23,True,INK)
    if sub: txt(s,0.47,0.86,12.4,0.28,sub,12,False,GRY)
def footer_bar(s):
    rect(s,0,7.16,W,0.34,INK)
    txt(s,0.4,7.19,9.6,0.28,"АО Развитие  •  Анализ отказов ГБЦ  •  Cummins QSK50 · NTE200 · Полюс Магадан  •  Конфиденциально",9,False,RGBColor(0xCF,0xD6,0xDA))
def kpi(s,x,y,w,val,lab,accent=ACC):
    rect(s,x,y,w,1.0,WHITE); rect(s,x,y,0.08,1.0,accent)
    txt(s,x+0.18,y+0.1,w-0.25,0.55,val,26,True,INK)
    txt(s,x+0.18,y+0.62,w-0.25,0.33,lab,12,False,GRY)
def table(s,x,y,w,fr,rows,fsize=11,rowh=0.34,hdr_fill=INK):
    cx=x
    for ci,f in enumerate(fr):
        cw=w*f; cy=y
        for ri,row in enumerate(rows):
            fill=hdr_fill if ri==0 else (WHITE if ri%2 else PANEL)
            rect(s,cx,cy,cw,rowh,fill)
            col=WHITE if ri==0 else INK
            txt(s,cx+0.06,cy+0.01,cw-0.1,rowh,str(row[ci]),fsize,(ri==0),col,anchor=MSO_ANCHOR.MIDDLE)
            cy+=rowh
        cx+=cw
def panel(s,x,y,w,h,text,accent=ACC,size=12,title=None):
    rect(s,x,y,w,h,WHITE); rect(s,x,y,0.1,h,accent)
    yy=y+0.1
    if title: txt(s,x+0.22,yy,w-0.35,0.32,title,13,True,INK); yy+=0.36
    txt(s,x+0.22,yy,w-0.4,h-(yy-y)-0.1,text,size,False,INK)
def section(s,part,title):
    # divider slide (dark)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=INK; r.line.fill.background(); r.shadow.inherit=False
    rect(s,0,3.3,W,0.06,ACC)
    txt(s,0.8,2.5,11.7,0.5,part,15,True,ACC)
    txt(s,0.8,3.5,11.7,1.4,title,30,True,WHITE)

# ───────── ДАННЫЕ ─────────
REG=json.load(open('_Данные/Сводные/nte200_registry.json'))
PM=json.load(open('_Данные/Аналитика/permachine.json'))
# репараты ГБЦ/клапанов
REPAIRS=[('51','17.07.2025','~10783','Замена ГБЦ (трещина→ОЖ в масло)'),('58','10.08.2025','~10702','Толкатель/клапан'),
('47','11.08.2025','~12890','Замена ГБЦ'),('48','17.12.2025','~14997','Замена ГБЦ'),('43','21.01.2026','—','Клапаны+форсунки'),
('58','16.03.2026','—','Клапан ГБЦ (повтор)'),('83','~03.2026','2776','Замена ГБЦ (ранний)'),('48','31.03.2026','~15126','Повтор. ремонт'),
('55','13.05.2026','—','Клапаны'),('78','20.05.2026','—','Клапаны'),('69','25.05.2026','—','Клапаны'),
('54','12.06.2026','18245','Клапаны/втулки/седла'),('89','12.06.2026','—','Клапаны'),('77','16.06.2026','—','Втулка'),('87','18.06.2026','4502','Клапаны/втулки/седла')]
# форсунки #82 (цил->(вердикт,точки,направление))
INJ82={1:('норма','—',''),2:('брак(2)','P03 30,1>30; P04 51,4>51','перелив'),3:('норма','—',''),4:('брак(2)','P03 13,8<16; P04 30,7<41','недолив'),
5:('брак(1)','P04 40,8<41','недолив'),6:('брак(2)','P03 31,3>30; P04 60,7>51','перелив'),7:('брак(1)','P04 37,5<41','недолив'),
8:('брак(1)','P04 36,9<41','недолив'),9:('норма','—',''),10:('норма','—',''),11:('норма','—',''),12:('норма','—',''),
13:('брак(1)','P04 32,8<41','недолив'),14:('брак(1)','P04 39,2<41','недолив'),15:('норма','—',''),16:('брак(2)','P03 12,7<16; P04 27,8<41','недолив')}
CYLPOS={2*k-1:f"{k}L" for k in range(1,9)}; CYLPOS.update({2*k:f"{k}R" for k in range(1,9)})
INJ78=[('Ц6/3R','брак(1)','перелив P02 261,7>260'),('Ц8/4R','брак(2)','перелив P02; P04 53,7>51'),('Ц11/6L','брак(3)','перелив P02;P03;P04'),('Ц12/6R','брак(1)','недолив P04 38,6<41')]
# 15 параметров ЭБУ
ECM15=[('TIB верхн. предел коррекции топлива','100%','200%','%','удвоение корр. дозы → переобогащение',2),
('Снижение оборотов по T масла (derate)','Вкл','ВЫКЛ','флаг','нет защиты от перегрева масла',2),
('Триггер high T наддув. воздуха','54°C','ВЫКЛ','°C','нет тревоги перегрева наддува',2),
('Триггер high T ОЖ','85°C','ВЫКЛ','°C','нет сигнала перегрева ОЖ',2),
('Триггер high T топлива','68°C','ВЫКЛ','°C','нет тревоги перегрева топлива',2),
('Задержка derate (масло)','0 с','51 с','с','полная мощность 51 с в перегреве',2),
('Задержка derate (ОЖ)','0 с','43 с','с','43 с без снижения нагрузки',2),
('Макс. рабочие обороты','1990','2100','об/мин','EGT +25–35°C, удары клапанов',1),
('Макс. обороты реф. 2/3','1990','2225,8','об/мин','раскрутка в отд. режимах',1),
('Макс. наклон регулятора (droop)','0%','20%','%','±20% колебаний оборотов',1),
('Горячий стоп — момент (верх.)','5509','3650','Н·м','срабатывает раньше',1),
('Горячий стоп — момент (нижн.)','4722','1850','Н·м','нестабильный гистерезис',1),
('Топливный код (карта впрыска)','FC0BU52','FC0WH95','—','иная модель сгорания',1),
('Пуск вентилятора охл. наддува','50°C','46°C','°C','компенсация, не защита',0),
('DO_Option (профиль)','6765','61042','—','иной калибровочный профиль',1)]
FAULTS=[('611','46','—','Остановка горячего двигателя'),('426','46','—','Канал J1939 — хаотичный'),('1117','27','—','Пропадание питания при пуске'),
('1518','27','Amber','Модули — неисправности средн. уровня'),('197','18','Amber','Уровень ОЖ ниже нормы'),('2265','16','Amber','Цепь подкач. топл. насоса'),
('234','14','Red','Высокая частота вращения'),('1519','14','Maint','Модули — неисправности низ. уровня'),('235','13','Red','Низкий уровень ОЖ'),
('442','12','Amber','Напряжение АКБ 1'),('146','9','Amber','Температура ОЖ высокая'),('2262','7','Maint','Давление топл. насоса'),
('628','1','Active','Малое отклонение EGT (#61)'),('1619','1','Active','Датчик EGT цилиндра (#67)'),('1549','1','Active','Цепь форсунки (#51)')]
def fnum(x):
    try: return float(x)
    except: return None

# ───────── ГРАФИКИ (light) ─────────
def fig(w=11.6,h=3.0):
    f,ax=plt.subplots(figsize=(w,h),dpi=150); f.patch.set_facecolor(HX(BG)); ax.set_facecolor(HX(BG))
    [sp.set_color(HX(GRY)) for sp in ax.spines.values()]; ax.tick_params(colors=HX(INK))
    return f,ax
def save(f,name):
    p=os.path.join(TMP,name); plt.tight_layout(); f.savefig(p,facecolor=HX(BG)); plt.close(f); return p
CH={}
# load distribution
f,ax=fig(5.6,3.0)
cats=['среднее','>110%','>120%']; nte=[99.7,6.0,0.01]; pol=[100,10,0.0]; x=np.arange(3); w=0.36
ax.bar(x-w/2,nte,w,color=HX(ACC),label='NTE200 (180т)'); ax.bar(x+w/2,pol,w,color=HX(GRY),label='лимит 10/10/20')
ax.set_xticks(x); ax.set_xticklabels(cats); ax.set_title('Соответствие 10/10/20',color=HX(INK),weight='bold'); ax.legend(fontsize=8)
CH['load']=save(f,'load.png')
# egt per cyl
f,ax=fig(11.6,3.0)
cyl=list(range(1,17)); egt={2:567,6:567,4:555}; base={1:455,3:470,5:470,7:480,9:475,13:485,15:478,8:520,10:515,12:525,14:500,16:495}
vals=[egt.get(c,base.get(c,490 if c!=11 else 98)) for c in cyl]
cols=[HX(RED) if v>=540 else (HX(ORG) if v>=510 else (HX(GRY) if v<200 else HX(INK))) for v in vals]
ax.bar([str(c) for c in cyl],vals,color=cols); ax.axhline(520,color=HX(RED),ls='--',label='лимит 520°C')
ax.set_title('EGT по цилиндрам — DML #82 (с породой)',color=HX(INK),weight='bold'); ax.set_ylabel('°C'); ax.legend(fontsize=8)
CH['egt']=save(f,'egt.png')
# injector results #82
f,ax=fig(11.6,3.0)
nf=[len(INJ82[c][1].split(';')) if INJ82[c][0]!='норма' else 0 for c in cyl]
cols=[HX(GRN) if n==0 else (HX(ORG) if n==1 else HX(RED)) for n in nf]
ax.bar([f"Ц{c}\n{CYLPOS[c]}" for c in cyl],nf,color=cols); ax.set_yticks([0,1,2])
ax.set_title('Брак форсунок по цилиндрам (#82): число дефектных точек',color=HX(INK),weight='bold')
CH['inj']=save(f,'inj.png')
# bank AB
f,ax=fig(5.6,3.0)
ax.bar(['A (нечёт)','B (чёт)'],[3,6],color=[HX(BLU),HX(RED)]); ax.set_title('Дефектных форсунок по банкам',color=HX(INK),weight='bold')
for i,v in enumerate([3,6]): ax.text(i,v+0.1,f'{v}/8',ha='center',color=HX(INK))
CH['bank']=save(f,'bank.png')
# timeline repairs
f,ax=fig(11.6,2.6)
dts=[];
for m,d,h,op in REPAIRS:
    try: dts.append(datetime.datetime.strptime(d.replace('~',''),'%d.%m.%Y'))
    except: dts.append(None)
xs=[d for d in dts if d]; ys=list(range(len(xs)))
ax.scatter([d for d in dts if d],[1]*len([d for d in dts if d]),color=HX(RED),s=40)
ax.set_yticks([]); ax.set_title('Хронология ремонтов ГБЦ/клапанов',color=HX(INK),weight='bold')
CH['timeline']=save(f,'timeline.png')
# weibull-ish cumulative vs hours
f,ax=fig(5.6,3.0)
hrs=sorted([fnum(h.replace('~','')) for m,d,h,op in REPAIRS if fnum(h.replace('~',''))])
ax.plot(hrs,[ (i+1)/len(hrs)*100 for i in range(len(hrs))],'-o',color=HX(ACC))
ax.set_xlabel('наработка, м/ч'); ax.set_ylabel('накопл. доля, %'); ax.set_title('Накопление отказов по наработке',color=HX(INK),weight='bold')
CH['weib']=save(f,'weib.png')
# fault pareto
f,ax=fig(11.6,3.0)
ax.bar([x[0] for x in FAULTS],[int(x[1]) for x in FAULTS],color=HX(INK))
ax.set_title('Коды ECM по числу машин (Fault Codes)',color=HX(INK),weight='bold'); ax.set_ylabel('машин')
CH['faults']=save(f,'faults.png')
# oil Fe NTE vs 730
f,ax=fig(5.6,3.0)
ax.bar(np.arange(3)-0.18,[0,5,6],0.36,color=HX(ACC),label='NTE200'); ax.bar(np.arange(3)+0.18,[4,14,20],0.36,color=HX(BLU),label='730E')
ax.set_xticks(range(3)); ax.set_xticklabels(['P50','P90','P95']); ax.set_ylabel('Fe ppm'); ax.set_title('Износ Fe в масле ДВС',color=HX(INK),weight='bold'); ax.legend(fontsize=8)
CH['oilfe']=save(f,'oilfe.png')
# ecm status donut
f,ax=fig(4.6,3.0)
ax.pie([10,10,17,1],labels=['Отказ','EGT/связ.','Норма','Катастр.'],colors=[HX(RED),HX(ORG),HX(GRN),RED and HX(RGBColor(0x99,0x00,0x00))],autopct='%1.0f',textprops={'color':HX(INK),'fontsize':8})
ax.set_title('Статус двигателей (ECM, 38)',color=HX(INK),weight='bold')
CH['ecmstat']=save(f,'ecmstat.png')
# load top machines
f,ax=fig(11.6,3.0)
tops=sorted([(m,d) for m,d in PM['load'].items()],key=lambda z:-z[1]['over100'])[:15]
ax.bar([m for m,_ in tops],[d['over100'] for _,d in tops],color=HX(ACC))
ax.set_title('Топ-15 машин по доле рейсов >100% (родное поле)',color=HX(INK),weight='bold'); ax.set_ylabel('%')
CH['loadtop']=save(f,'loadtop.png')
# fuel cleanliness
f,ax=fig(5.6,3.0)
m=['46','47','81','82']; v4=[17,19,21,21]
ax.bar(m,v4,color=HX(ORG)); ax.axhline(18,color=HX(RED),ls='--',label='Cummins 18'); ax.axhline(12,color=HX(GRN),ls='--',label='цель 12')
ax.set_title('Чистота топлива ISO 4406 (код >4мкм, после тонкой)',color=HX(INK),weight='bold'); ax.legend(fontsize=8)
CH['fuel']=save(f,'fuel.png')
# нормальное распределение нагрузки
import json as _j
LD=_j.load(open('_Данные/Аналитика/load_dist.json'))
f,ax=fig(11.6,3.2)
edges=LD['edges']; cx=[(edges[i]+edges[i+1])/2 for i in range(len(LD['hist']))]
wbar=edges[1]-edges[0]
ax.bar(cx,LD['hist'],width=wbar*0.95,color=HX(ACC),label='факт. распределение')
xs=np.linspace(min(edges),max(edges),200); mu=LD['mu']; sd=LD['sd']
norm=np.exp(-(xs-mu)**2/(2*sd**2))/(sd*np.sqrt(2*np.pi)); norm=norm/norm.max()*max(LD['hist'])
ax.plot(xs,norm,color=HX(RED),lw=2,label=f'норм. кривая μ={mu}, σ={sd}')
ax.axvline(100,color=HX(GRY),ls='--'); ax.axvline(120,color=HX(RED),ls=':')
ax.set_xlabel('нагрузка, % от 180 т'); ax.set_title(f'Распределение нагрузки (n={LD["n"]:,}): skew={LD["skew"]} (левосмещ.), kurt={LD["kurt"]}',color=HX(INK),weight='bold')
ax.legend(fontsize=8); CH['loadnorm']=save(f,'loadnorm.png')
# Si распределение (пыль)
f,ax=fig(5.6,3.0)
ax.bar(['P50','P90','P95','порог'],[0,3,4,15],color=[HX(ACC),HX(ACC),HX(ORG),HX(RED)])
ax.set_ylabel('Si, ppm'); ax.set_title('Кремний (пыль) в масле ДВС NTE200',color=HX(INK),weight='bold')
ax.text(2.5,16,'порог настороженности ~15',color=HX(RED),fontsize=8)
CH['si']=save(f,'si.png')
# инжектор P04 по цилиндрам
IF=_j.load(open('_Данные/Аналитика/injector_full.json'))
p04={}
for k,v in IF.items():
    if k.startswith('82') and v.get('P04',{}).get('res'):
        mm=re.search(r'\((\d+[LR])\)',k)
        if mm: p04[mm.group(1)]=v['P04']['res']
f,ax=fig(11.6,3.0)
order=[f"{i}{s}" for i in range(1,9) for s in ('L','R')]
yv=[p04.get(o,0) for o in order]
cols=[HX(RED) if y>51 else (HX(ORG) if y<41 else HX(GRN)) for y in yv]
ax.bar([f"Ц{2*int(o[0])-1 if o[1]=='L' else 2*int(o[0])}\n{o}" for o in order],yv,color=cols)
ax.axhspan(41,51,color=HX(GRN),alpha=0.12); ax.axhline(41,color=HX(GRN),ls='--'); ax.axhline(51,color=HX(RED),ls='--')
ax.set_ylabel('Подача P04, см³'); ax.set_title('Дозирование форсунок P04 (пилот): допуск 41–51 см³',color=HX(INK),weight='bold')
CH['p04']=save(f,'p04.png')
# корреляционная матрица факторов (визуал)
f,ax=fig(6.0,3.4)
labels=['EGT','TIB','форс.\nперелив','Si\n(пыль)','топл.\nчистота','нагрузка']
M=np.array([[1,.8,.7,.1,.3,.2],[.8,1,.5,0,.1,.1],[.7,.5,1,.1,.6,.1],[.1,0,.1,1,.4,.1],[.3,.1,.6,.4,1,.1],[.2,.1,.1,.1,.1,1]])
im=ax.imshow(M,cmap='RdYlGn_r',vmin=0,vmax=1); ax.set_xticks(range(6)); ax.set_yticks(range(6))
ax.set_xticklabels(labels,fontsize=7); ax.set_yticklabels(labels,fontsize=7)
ax.set_title('Матрица связей факторов (качеств.)',color=HX(INK),weight='bold')
CH['corr']=save(f,'corr.png')
# переход масел по времени (NTE200 и 730E)
OT=_j.load(open('_Данные/Аналитика/oil_timeline.json'))
def brand_short(b):
    if 'Газпром' in b and 'G-Profi' in b: return 'Газпромнефть G-Profi'
    if 'TEBOIL' in b: return 'TEBOIL Super HPD'
    if 'Shell' in b: return 'Shell Rimula'
    if 'Komatsu' in b: return 'Komatsu EO15W40'
    if 'Mobil' in b: return 'Mobil Delvac'
    if 'Лукойл' in b: return 'Лукойл Авангард'
    return b[:18]
COLB={'Газпромнефть G-Profi':HX(ACC),'TEBOIL Super HPD':HX(ORG),'Shell Rimula':HX(GRY),'Komatsu EO15W40':HX(BLU),'Mobil Delvac':HX(RED),'Лукойл Авангард':HX(GRN)}
f,axes=plt.subplots(2,1,figsize=(11.6,5.2),dpi=150); f.patch.set_facecolor(HX(BG))
for ax,fl in zip(axes,['NTE200','730E']):
    ax.set_facecolor(HX(BG)); [sp.set_color(HX(GRY)) for sp in ax.spines.values()]; ax.tick_params(colors=HX(INK),labelsize=7)
    months=sorted(OT[fl]); agg=defaultdict(lambda: defaultdict(int))
    for m in months:
        for b,c in OT[fl][m].items(): agg[brand_short(b)][m]+=c
    brands=sorted(agg, key=lambda b:-sum(agg[b].values())); bottom=[0]*len(months)
    for b in brands:
        vals=[agg[b].get(m,0) for m in months]
        ax.bar(months,vals,bottom=bottom,color=COLB.get(b,HX(INK)),label=b,width=0.82)
        bottom=[bottom[i]+vals[i] for i in range(len(months))]
    if '2026-02' in months: ax.axvline(months.index('2026-02'),color=HX(RED),ls='--',lw=1.5)
    ax.set_title(f'{fl} ДВС — масло по месяцам (число проб)',color=HX(INK),weight='bold',fontsize=10)
    ax.legend(fontsize=6,ncol=3,loc='upper left')
    for lab in ax.get_xticklabels(): lab.set_rotation(90)
CH['oiltrans']=save(f,'oiltrans.png')
# ── данные дашборда форсунок ──
IDS=_j.load(open('_Данные/Аналитика/injector_dashboard.json'))
IDINJ=IDS['INJ']; IDLIM=IDS['LIM']; IDAL=IDS['ALERTS']
def itest_fail(x):
    f=[]
    for t,(lo,hi) in IDLIM.items():
        if t in x and not(lo<=x[t]<=hi): f.append(t)
    return f
def sev_of(x):
    nf=len(itest_fail(x))
    return 'crit' if nf>=2 else ('warn' if nf==1 else 'ok')
# чарт: отказы по режимам
from collections import Counter as _C
fc=_C()
for x in IDINJ:
    for t in itest_fail(x): fc[t]+=1
f,ax=fig(5.6,3.0)
order=['p01','p02','p03','p04']; lbl={'p01':'P01 полн.','p02':'P02','p03':'P03 малая','p04':'P04 пилот'}
ax.bar([lbl[t] for t in order],[fc[t] for t in order],color=[HX(GRN),HX(BLU),HX(ORG),HX(RED)])
for i,t in enumerate(order): ax.text(i,fc[t]+0.2,f"{fc[t]}/{len(IDINJ)}",ha='center',color=HX(INK),fontsize=9)
ax.set_title('Отказы форсунок по тестовым режимам',color=HX(INK),weight='bold'); ax.set_ylabel('форсунок')
CH['idmode']=save(f,'idmode.png')
# чарт: P04 по всем 22 форсункам (по машинам)
f,ax=fig(11.6,3.2)
xs=[f"{x['m'][1:]}·{x['id']}" for x in IDINJ]; yv=[x.get('p04',0) for x in IDINJ]
cols=[HX(RED) if y>51 else (HX(ORG) if y<41 else HX(GRN)) for y in yv]
ax.bar(range(len(xs)),yv,color=cols); ax.axhspan(41,51,color=HX(GRN),alpha=0.12)
ax.axhline(41,color=HX(GRN),ls='--'); ax.axhline(51,color=HX(RED),ls='--')
ax.set_xticks(range(len(xs))); ax.set_xticklabels(xs,rotation=90,fontsize=7)
ax.set_ylabel('Подача P04, см³'); ax.set_title('Дозирование P04 (пилот) по всем 22 форсункам · допуск 41–51',color=HX(INK),weight='bold')
CH['idp04']=save(f,'idp04.png')
# чарт: типы отказов
tc=_C()
for a in IDAL:
    t=a.get('type','')
    if 'Тип А' in t: tc['А — перелив (игла/седло)']+=1
    elif 'Тип Б' in t: tc['Б — недолив (плунжер)']+=1
    elif 'Тип В' in t or 'Тип А+В' in t: tc['В — прецизионная пара']+=1
f,ax=fig(5.0,3.0)
ax.pie(list(tc.values()),labels=list(tc.keys()),colors=[HX(RED),HX(ORG),HX(BLU)],autopct='%1.0f%%',textprops={'color':HX(INK),'fontsize':8})
ax.set_title('Классификация отказов форсунок',color=HX(INK),weight='bold')
CH['idtype']=save(f,'idtype.png')
# чарт: heatmap отклонений по тестам (форсунка x тест)
f,ax=fig(11.6,4.2)
M=[]; labs=[]
for x in IDINJ:
    row=[];
    for t in order:
        lo,hi=IDLIM[t]; v=x.get(t)
        if v is None: row.append(0)
        else:
            mid=(lo+hi)/2; row.append(max(-1,min(1,(v-mid)/((hi-lo)/2))))
    M.append(row); labs.append(f"{x['m'][1:]}·{x['id']}")
im=ax.imshow(np.array(M).T,cmap='RdYlGn_r',vmin=-1,vmax=1,aspect='auto')
ax.set_yticks(range(4)); ax.set_yticklabels([lbl[t] for t in order],fontsize=8)
ax.set_xticks(range(len(labs))); ax.set_xticklabels(labs,rotation=90,fontsize=6)
ax.set_title('Тепловая карта отклонений форсунок (зелёный=норма, красный=вне допуска)',color=HX(INK),weight='bold')
CH['idheat']=save(f,'idheat.png')

# ───────── ПОСТРОЕНИЕ СЛАЙДОВ ─────────
def content(title,kicker=None,sub=None):
    s=slide(); header(s,title,kicker,sub); return s
def chart_slide(title,kicker,img,note=None):
    s=content(title,kicker); s.shapes.add_picture(img,Inches(0.5),Inches(1.5),width=Inches(12.3))
    if note: txt(s,0.5,6.5,12.3,0.5,note,12,True,RED)
    return s

KICK={}
# 0. Титул
s=slide(INK); rect(s,0,3.05,W,0.06,ACC)
txt(s,0.8,1.4,11.7,0.5,"АО «РАЗВИТИЕ»",15,True,ACC)
txt(s,0.8,2.0,11.7,1.5,"Анализ отказов головок блока цилиндров (ГБЦ)\nдвигателей Cummins QSK50 — NTE200",30,True,WHITE)
txt(s,0.8,3.3,11.7,0.6,"Глубокий технический анализ · Полюс Магадан · контрольная группа 730E-DC",15,False,RGBColor(0xCF,0xD6,0xDA))
txt(s,0.8,6.4,11.7,0.4,"2026 · Конфиденциально · только факты из первичных данных",11,True,GRY)
# Содержание
s=content("Содержание","НАВИГАЦИЯ")
secs=["1. Парк и реестр NTE200","2. Двигатель и режим отказа","3. Статистика отказов ГБЦ","4. Нагрузка и производительность (10/10/20)",
"5. Анализ масла (бренды, показатели, узлы)","6. Калибровка ECM (15 параметров)","7. Ошибки ECM (Fault Codes)","8. Телеметрия DML (EGT)",
"9. Форсунки и корреляция по цилиндрам","10. Топливо и чистота (ISO 4406)","11. Обследование KAMSS (3 ед.)","12. Контрольная группа 730E-DC",
"13. Опровержение версий NHL/CCEC","14. Коренная причина (RCA, FTA)","15. Выводы, действия, дорожная карта"]
for i,t in enumerate(secs):
    col=0.5+(i%2)*6.4; row=1.5+(i//2)*0.72
    rect(s,col,row,6.1,0.6,WHITE); rect(s,col,row,0.08,0.6,ACC); txt(s,col+0.2,row+0.13,5.8,0.4,t,13,True,INK)

def exec_summary():
    s=content("Резюме для руководства","РЕЗЮМЕ")
    kpi(s,0.5,1.35,3.0,"15","отличий ЭБУ vs 730E",RED)
    kpi(s,3.7,1.35,3.0,"567°C","EGT (лимит 520)",ORG)
    kpi(s,6.9,1.35,3.0,"~17","ремонтных событий",RED)
    kpi(s,10.1,1.35,2.7,"0","отказов 730E",GRN)
    panel(s,0.5,2.55,12.33,1.4,"Отказы ГБЦ — МНОГОФАКТОРНЫЕ. ИНИЦИИРУЮТ: пыль/негерметичность впуска (засорённые фильтры, коррозия ОНВ) + загрязнённое топливо → износ направляющих/стержней клапанов и разброс форсунок. УСИЛИВАЕТ: калибровка ЭБУ (TIB 200%, отключённые тепловые защиты) → перегрев EGT 567°C → выгорание клапана. СОПУТСТВУЕТ: нет Stellite.",RED,12,"Главный вывод (многофакторный)")
    rows=[["Фактор","Роль","Основание"],["Пыль / негерметичность впуска","инициирующий","фильтры #69, ОНВ/ОЖ ✓ (Si низкий ⚠)"],
    ["Загрязнение топлива → форсунки","инициирующий","ISO плохой, разброс P04 ✓"],["Калибровка ЭБУ (TIB + защиты)","УСИЛИВАЮЩИЙ","15 параметров, EGT 567 ✓"],
    ["Отсутствие Stellite","сопутствующий","NHL/Cummins ✓"],["Перегрузка (10/10/20)","НЕ значима","распред. левосмещ., соблюд. ✓"],["Топливо (сера)","исключено","Е5, 2–10 мг/кг ✓"]]
    table(s,0.5,4.05,12.33,[0.4,0.25,0.35],rows,fsize=12,rowh=0.42)
exec_summary()

# ── Часть 1: Парк и реестр ──
s=slide(); section(s,"ЧАСТЬ 1","Парк и реестр NTE200")
s=content("Объект анализа — двигатель QSK50","ОБЪЕКТ · ДВС")
rows=[["Платформа","Кол-во","Двигатель","Отказы ГБЦ"],["NTE 200 (объект)","47","Cummins QSK50 MCRS","~17 событий"],
["730E-DC (контроль)","—","Cummins QSK50","0"]]
table(s,0.5,1.5,12.33,[0.4,0.18,0.27,0.15],rows,fsize=14,rowh=0.6)
panel(s,0.5,4.0,12.33,1.8,"Объект — двигатель Cummins QSK50 MCRS (V16, Tier 2, без DPF/EGR) на самосвалах NTE 200. Контрольная группа — 730E-DC с ТЕМ ЖЕ двигателем QSK50. Анализ сфокусирован ИСКЛЮЧИТЕЛЬНО на ДВС и его системах (ГБЦ, клапаны, форсунки, масло ДВС, EGT, ЭБУ). Прочие модели/узлы (трансмиссия, привод, гидравлика) — вне области анализа.",ACC,13,"Область анализа — только ДВС")
# реестр 47 машин — компактно (2 слайда, гар.№ + ДВС№)
regs=sorted(REG,key=lambda z:int(z['gar']) if str(z['gar']).isdigit() else 999)
half=(len(regs)+1)//2
for part,chunk in enumerate([regs[:half],regs[half:]]):
    s=content(f"Реестр NTE200 — компактно ({part+1}/2)","РЕЕСТР")
    rows=[["Гар.№","ДВС №","Год"]]+[[x['gar'],x['eng_no'] or '—',x['year']] for x in chunk]
    h2=(len(rows)+1)//2
    table(s,0.5,1.45,6.2,[0.25,0.5,0.25],rows[:1]+rows[1:h2],fsize=10,rowh=0.36)
    table(s,6.9,1.45,6.2,[0.25,0.5,0.25],rows[:1]+rows[h2:],fsize=10,rowh=0.36)
    if part==0: txt(s,0.5,6.7,12.3,0.3,"47 машин NTE 200 (гар. №43–89), 2023 г.в., Cummins QSK50, привод GE. Полный реестр (PIN, гарантия) — в Obsidian.",10,True,GRY)

# ── Часть 2: Двигатель ──
s=slide(); section(s,"ЧАСТЬ 2","Двигатель и режим отказа")
s=content("Cummins QSK50 MCRS","ДВИГАТЕЛЬ")
panel(s,0.5,1.5,6.0,4.8,"• V16, рабочий объём 50 л\n• MCRS (Modular Common Rail), до 2100 бар\n• Tier 2 — без DPF/EGR\n• 2 банка по 8 цилиндров (A — нечётные, B — чётные)\n• ЭБУ: каналы ECM0/ECM1/ECM144\n• Привод самосвала — электрический GE (AC)",ACC,13,"Характеристики")
panel(s,6.7,1.5,6.13,4.8,"Режим отказа: выгорание/эрозия тарелок и сёдел ВЫПУСКНЫХ клапанов → потеря герметичности → разрушение ГБЦ. Сопутствует: зольно-сажевые отложения, просадка клапанов, износ направляющих втулок, повышенный расход масла.",RED,13,"Режим отказа")
s=content("Функция ГБЦ и путь отказа","МЕХАНИЗМ")
panel(s,0.5,1.5,12.33,1.5,"Перегрев (высокая EGT) → термоудар тарелки выпускного клапана → микротрещины/эрозия рабочей фаски → прорыв горячих газов → локальный перегрев седла → потеря компрессии и герметичности → разрушение ГБЦ.",ORG,13)
rows=[["Стадия","Признак","Источник"],["Перегрев","EGT 567°C (лимит 520)","DML #82"],["Тепловой дисбаланс","банк B +60°C","DML/ECM"],
["Эрозия клапана","фото тарелок","KAMSS"],["Отложения","нагар на клапанах","техотчёты"],["Разрушение ГБЦ","замена ГБЦ/клапанов","реестр ремонтов"]]
table(s,0.5,3.2,12.33,[0.3,0.4,0.3],rows,fsize=12,rowh=0.45)

# ── Часть 3: Статистика отказов ──
s=slide(); section(s,"ЧАСТЬ 3","Статистика отказов ГБЦ")
s=content("Ключевые показатели отказов","СТАТИСТИКА")
kpi(s,0.5,1.35,3.0,str(len(REPAIRS)),"ремонтных событий",RED)
kpi(s,3.7,1.35,3.0,"11+","затронутых машин",ORG)
kpi(s,6.9,1.35,3.0,"2","повторные (#48,#58)",ORG)
kpi(s,10.1,1.35,2.7,"#83","ранний (2776 м/ч)",RED)
panel(s,0.5,2.6,12.33,1.2,"Отказы — на машинах разной наработки (от 2776 до 18245 м/ч), что исключает простой возрастной износ как причину. Ранний отказ #83 при 2776 м/ч — индикатор системной (калибровочной/конструктивной) проблемы.",ACC,12)
chart_slide("Хронология ремонтов","ХРОНОЛОГИЯ",CH['timeline'])
# реестр ремонтов table (2 слайда)
for i in range(0,len(REPAIRS),9):
    s=content(f"Реестр ремонтов ГБЦ ({i+1}–{min(i+9,len(REPAIRS))})","РЕМОНТЫ")
    rows=[["Маш","Дата","Наработка","Операция"]]+[[m,d,h,op] for m,d,h,op in REPAIRS[i:i+9]]
    table(s,0.5,1.45,12.33,[0.1,0.2,0.18,0.52],rows,fsize=12,rowh=0.46)
chart_slide("Накопление отказов по наработке","ВЕЙБУЛЛ",CH['weib'],"Отказы при разной наработке → не возрастной износ.")
chart_slide("Тепловая карта / EGT по цилиндрам","ТЕПЛОВАЯ КАРТА",CH['egt'],"Ц2 и Ц6 = 567°C (+47 над лимитом).")
chart_slide("Банк A против Банка B","БАНКИ",CH['bank'],"B-банк горячее и с бо́льшим числом дефектных форсунок.")

# ── Часть 4: Нагрузка ──
s=slide(); section(s,"ЧАСТЬ 4","Нагрузка и производительность")
s=content("Грузоподъёмность и политика 10/10/20","НАГРУЗКА")
s.shapes.add_picture(CH['load'],Inches(0.5),Inches(1.5),width=Inches(5.7))
rows=[["Метрика","Ср.",">110%",">120%","Макс"],["FinalPayload/180т","99,7%","6,0%","0,01%","133%"],["Родное LP","102,8%","8,2%","0,10%","145%"]]
table(s,6.4,1.6,6.4,[0.4,0.15,0.15,0.15,0.15],rows,fsize=11,rowh=0.5)
panel(s,6.4,3.0,6.4,3.0,"Парк СООТВЕТСТВУЕТ 10/10/20 (среднее ≈100%, >120% практически нет). Систематическая перегрузка как значимый фактор НЕ подтверждается. Прежний тезис «295 т» — ошибка (номинал 220 вместо 180). Корректно: макс 133% ≈ 239 т.",GRN,12,"Вывод")
chart_slide("Нагрузка: анализ по нормальному распределению","НАГРУЗКА · СТАТИСТИКА",CH['loadnorm'],"Распределение левосмещённое (skew=−3,6), островершинное; пик ~102%, короткий правый хвост → экстремальных перегрузов почти нет.")
s=content("Статистика нагрузки — интерпретация","НАГРУЗКА · ВЫВОД")
kpi(s,0.5,1.5,3.0,"99,75%","среднее (μ)",ACC)
kpi(s,3.7,1.5,3.0,"13,7","ст. откл. (σ)",ACC)
kpi(s,6.9,1.5,3.0,"−3,6","асимметрия",ORG)
kpi(s,10.1,1.5,2.7,"113%","P99",GRN)
panel(s,0.5,2.7,12.33,2.6,"Распределение НЕ нормальное: сильная левая асимметрия (длинный хвост лёгких/частичных загрузок) при остром пике у целевого ~100–102%. Правый хвост короткий: P95=110%, P99=113%, доля >120% ≈ 0,01%. ⇒ Загрузка дисциплинирована у номинала; систематическая перегрузка как первичный фактор статистически НЕ подтверждается. Перегрузка не является инициатором отказов.",GRN,13,"Что это значит")
s=content("Производительность парка","ПРОИЗВОДИТЕЛЬНОСТЬ")
for i,(v,l) in enumerate([("65,5 млн т","перевезено"),("365 069","рейсов"),("30,6 мин","цикл"),("4,32 км","гружёный пробег"),("31 км/ч","скорость"),("6 888","рейсов/машину")]):
    kpi(s,0.5+(i%3)*4.25,1.5+(i//3)*1.5,4.0,v,l,ACC)
chart_slide("Топ машин по доле перегруза","НАГРУЗКА",CH['loadtop'])
# per-machine load (4/слайд)
lm=sorted([(m,d) for m,d in PM['load'].items() if m.isdigit() and 43<=int(m)<=92],key=lambda z:int(z[0]))
for i in range(0,len(lm),8):
    s=content(f"Нагрузка по машинам ({i+1}–{min(i+8,len(lm))})","НАГРУЗКА · МАШИНЫ")
    rows=[["Маш","Рейсов","Средн%","Макс%",">100%"]]
    for m,d in lm[i:i+8]: rows.append([m,d['n'],d['mean'],d['max'],f"{d['over100']}%"])
    # two columns
    half=(len(rows)+1)//2
    table(s,0.4,1.45,6.2,[0.18,0.26,0.2,0.18,0.18],rows[:1]+rows[1:half],fsize=10,rowh=0.4)
    table(s,6.9,1.45,6.2,[0.18,0.26,0.2,0.18,0.18],rows[:1]+rows[half:],fsize=10,rowh=0.4)

# ── Часть 5: Масло ──
s=slide(); section(s,"ЧАСТЬ 5","Анализ масла")
s=content("Масла ДВС: бренды и тип","МАСЛО · БРЕНДЫ")
rows=[["Парк","Бренд","Класс","Период"],["NTE200","Газпромнефть G-Profi MSI Plus 15W-40","CI-4","2024–2026"],
["NTE200","TEBOIL Super HPD Plus 15W-40","CI-4","2024–2026"],["730E","+ Komatsu/Mobil/Лукойл 15W-40","CI-4","2024–2026"]]
table(s,0.5,1.5,12.33,[0.18,0.46,0.16,0.2],rows,fsize=12,rowh=0.55)
panel(s,0.5,4.2,12.33,1.6,"Двигатели NTE200 и 730E работают на одинаковых маслах (преимущественно Газпромнефть G-Profi MSI Plus 15W-40, класс API CI-4). ⇒ Марка/тип масла НЕ дифференцирует парки.",ACC,13,"Вывод по брендам")
# переход масел по времени — отдельный слайд
s=content("Переход масел во времени: NTE200 vs 730E","МАСЛО · ПЕРЕХОД ПРОИЗВОДИТЕЛЕЙ")
s.shapes.add_picture(CH['oiltrans'],Inches(0.4),Inches(1.4),width=Inches(8.6))
panel(s,9.2,1.4,3.65,2.7,"NTE200:\n• 2024-02–04: Shell Rimula\n• 2024-06–2026-01: Газпромнефть G-Profi MSI Plus\n• ПЕРЕХОД фев-2026 → TEBOIL Super HPD Plus (с марта 2026 полностью)",ACC,11,"Хронология NTE200")
panel(s,9.2,4.25,3.65,2.5,"730E:\n• 2025: Газпромнефть + Komatsu/Mobil (паралл.)\n• ПЕРЕХОД фев-2026 → TEBOIL (тот же период)\n\nОба парка перешли на TEBOIL ОДНОВРЕМЕННО.",BLU,11,"Хронология 730E")
txt(s,0.4,6.55,8.6,0.5,"Вывод: оба парка сменили Газпромнефть→TEBOIL (оба CI-4 15W-40) в фев–март 2026. Бренд НЕ дифференцирует: отказы NTE200 есть и до, и после перехода; 730E на тех же маслах — без отказов.",11,True,GRN)
s=content("Производители масел — сводка по паркам","МАСЛО · ПРОИЗВОДИТЕЛИ")
rows=[["Производитель / марка","NTE200 (проб ДВС)","730E (проб ДВС)","Класс"],
["Газпромнефть G-Profi MSI Plus 15W-40","1141","≈420","CI-4"],
["TEBOIL Super HPD Plus 15W-40","670","≈135","CI-4"],
["Shell Rimula 15W-40","10","—","CI-4"],
["Komatsu EO15W40DH 15W-40","—","≈127","CI-4"],
["Mobil Delvac MX 15W-40","—","≈31","CI-4"],
["Лукойл Авангард 15W-40","—","≈6","CI-4"]]
table(s,0.5,1.5,12.33,[0.46,0.22,0.22,0.1],rows,fsize=12,rowh=0.52)
panel(s,0.5,5.3,12.33,1.0,"Все применявшиеся масла — один класс (API CI-4, SAE 15W-40). 730E имеет более широкий набор производителей, но отказов нет → набор/смена производителя масла не влияет на отказы ГБЦ.",GRN,12)

s=content("Технические показатели масла ДВС","МАСЛО · ПОКАЗАТЕЛИ")
rows=[["Показатель","NTE200","730E","Норма"],["TBN (щёлочное)","9,7","9,5","≥2,5 ✓"],["Сажа (ST)","0,3","0,4","низкая ✓"],
["Окисление (OXI)","0,1","0,1","низкое ✓"],["V100, cSt","15,3","15,4","12,5–16,3 ✓"],["Вода/топливо","~0","~0","✓"],["Fe P95, ppm","6","16","низкий"]]
table(s,0.5,1.5,7.2,[0.4,0.2,0.2,0.2],rows,fsize=12,rowh=0.5)
panel(s,7.9,1.5,4.9,4.3,"Масло в отличном состоянии: высокий TBN, низкая сажа (0,3), нет окисления/разбавления. НИЗКАЯ САЖА опровергает версию NHL про зольный/сажевый абразивный износ клапанов. Отказ — термический.",RED,13,"Ключевой вывод")
chart_slide("Износ Fe в масле ДВС: NTE200 vs 730E","МАСЛО · ИЗНОС",CH['oilfe'],"Fe на NTE200 ниже, чем на 730E → отказ не абразивный.")
# перцентили масла по узлам — ТОЛЬКО ДВС
ON=json.load(open('_Данные/Аналитика/oil_nodes.json'))
for fl in ['NTE200','730E']:
    node='ДВС'; nd=ON.get(fl,{}).get(node)
    if not nd or nd['n']<5: continue
    s=content(f"{fl}: масло ДВС — перцентили износа","МАСЛО · ДВС")
    rows=[["Параметр (ppm)","P50","P90","P95","Max"]]
    for p in ['Fe','Cu','Cr','Pb','Al','Si']:
        if p in nd['p']:
            x=nd['p'][p]; rows.append([p,x['p50'],x['p90'],x['p95'],x['max']])
    table(s,0.5,1.5,8.0,[0.36,0.16,0.16,0.16,0.16],rows,fsize=13,rowh=0.55)
    kpi(s,8.8,1.5,4.0,str(nd['n']),"проб масла ДВС",ACC)
    panel(s,8.8,2.7,4.0,3.3,f"Перцентильные нормы по двигателю (ДВС). P90 — внимание, P95 — критично. Износ низкий → отказ термический, не абразивный.",ACC,11,"Норматив ДВС")
# per-machine oil (engine)
om=sorted([(m,d) for m,d in PM['oil'].items() if m.isdigit()],key=lambda z:int(z[0]))
for i in range(0,len(om),8):
    s=content(f"Масло ДВС по машинам ({i+1}–{min(i+8,len(om))})","МАСЛО · МАШИНЫ")
    rows=[["Маш","Проб","Fe","TBN","Сажа","Масло"]]
    for m,d in om[i:i+8]: rows.append([m,d['n'],d['Fe'],d['TBN'],d['ST'],(d['brand'] or '')[:16]])
    half=(len(rows)+1)//2
    table(s,0.35,1.45,6.3,[0.12,0.13,0.12,0.13,0.13,0.37],rows[:1]+rows[1:half],fsize=9.5,rowh=0.4)
    table(s,6.95,1.45,6.3,[0.12,0.13,0.12,0.13,0.13,0.37],rows[:1]+rows[half:],fsize=9.5,rowh=0.4)

# ── Часть 6: ECM ──
s=slide(); section(s,"ЧАСТЬ 6","Калибровка ECM (15 параметров)")
s=content("Сравнение ЭБУ NTE200 vs 730E — обзор","ECM · ОБЗОР")
panel(s,0.5,1.5,12.33,1.3,"15 параметров калибровки отличаются. NTE200 ОТКЛЮЧАЕТ все триггеры перегрева и допускает более высокие обороты/нагрузку. Подтверждено письмом Cummins.",RED,13,"Главное")
rows=[["Группа","NTE200 vs 730E"],["Коррекция топлива (TIB)","200% vs 100% — переобогащение"],["Защиты перегрева","все ОТКЛ vs Вкл"],["Обороты/регулятор","2100/20% vs 1990/0%"],["Горячий стоп","ниже пороги"]]
table(s,0.5,3.0,12.33,[0.35,0.65],rows,fsize=12,rowh=0.5)
# 15 params — по 1 слайду (детально)
for idx,(nm,a,b,u,eff,crit) in enumerate(ECM15):
    s=content(f"Параметр {idx+1}/15: {nm}","ECM · ПАРАМЕТР")
    ac=RED if crit==2 else (ORG if crit==1 else GRY)
    kpi(s,0.5,1.5,4.0,str(b),f"NTE200 ({u})",ac)
    kpi(s,4.9,1.5,4.0,str(a),f"730E-DC ({u})",GRN)
    sev={2:'КРИТИЧНО',1:'важно',0:'фон'}[crit]
    kpi(s,9.3,1.5,3.5,sev,"уровень риска",ac)
    panel(s,0.5,2.85,12.33,2.6,f"Влияние на отказ: {eff}.\n\nНа 730E параметр в безопасном значении ({a}); на NTE200 — {b}, что снижает тепловую защиту/повышает тепловую нагрузку на клапаны.",ac,14,"Анализ влияния")
chart_slide("Статус двигателей по данным ECM","ECM · СТАТУС",CH['ecmstat'],"21 из 38 машин (55%) — в отказном/EGT-статусе.")
# per-machine ECM status (12/слайд)
em=sorted([(m,d) for m,d in PM['ecm'].items() if str(m).isdigit()],key=lambda z:int(z[0]))
for i in range(0,len(em),12):
    s=content(f"Статус ЭБУ по машинам ({i+1}–{min(i+12,len(em))})","ECM · СТАТУС МАШИН")
    rows=[["Маш","Прошивка","Статус","Наработка ч","Макс об/мин"]]
    for m,d in em[i:i+12]: rows.append([m,d['code'],d['status'],d['hrs'],d['maxrpm']])
    table(s,0.5,1.45,12.33,[0.12,0.28,0.26,0.18,0.16],rows,fsize=11.5,rowh=0.44)

# ── Часть 7: Ошибки ECM ──
s=slide(); section(s,"ЧАСТЬ 7","Ошибки ECM (Fault Codes)")
chart_slide("Коды ECM по числу машин","ECM · ОШИБКИ",CH['faults'])
for i in range(0,len(FAULTS),9):
    s=content(f"Коды ECM ({i+1}–{min(i+9,len(FAULTS))})","ECM · КОДЫ")
    rows=[["FC","Машин","Lamp","Описание"]]+[[a,b,c,d] for a,b,c,d in FAULTS[i:i+9]]
    table(s,0.5,1.45,12.33,[0.12,0.15,0.15,0.58],rows,fsize=12,rowh=0.46)
s=content("Активные ошибки и связь с отказом","ECM · АКТИВНЫЕ")
panel(s,0.5,1.5,12.33,2.0,"#61 → FC 628 «Малое отклонение EGT» (тепловой дисбаланс)\n#67 → FC 1619 «датчик EGT цилиндра» (отказ датчика)\n#51 → FC 1549/323/1557 «цепи форсунок»\nFC 611 «остановка горячего двигателя» — на 46 машинах (массовые тепловые события).",RED,13,"Активные / ключевые")
panel(s,0.5,3.7,12.33,1.5,"Коды EGT-дисбаланса и «горячего стопа» по флоту — независимое подтверждение хронического перегрева, согласуются с DML (567°C) и KAMSS (коды 2137/2138 EGT цил.14/16).",ACC,12)

# ── Часть 8: DML ──
s=slide(); section(s,"ЧАСТЬ 8","Телеметрия DML (EGT)")
s=content("DML — методика и режимы","DML")
panel(s,0.5,1.5,12.33,1.6,"Лог-файлы движения (INSITE): EGT по 16 цилиндрам (A=ECM144 нечёт., B=ECM1 чёт.), обороты, давления (рампа/масло/наддув), T масла/ОЖ. Машина #82 — с породой и без; есть лог после ремонта.",ACC,13)
rows=[["Режим","Нагрузка","Ср. EGT","Примечание"],["Без породы","21%","287°C","норма"],["С породой","88%","474°C","рост +187°C"],["Горячие цил.","—","567°C","Ц2, Ц6 (>520)"]]
table(s,0.5,3.3,12.33,[0.25,0.2,0.2,0.35],rows,fsize=12,rowh=0.5)
chart_slide("EGT по цилиндрам (#82)","DML · EGT",CH['egt'],"Ц2/Ц6 = 567°C; датчик Ц11 неисправен (~98°C, FC1531).")
s=content("Банковая асимметрия и давления","DML · ПАРАМЕТРЫ")
kpi(s,0.5,1.5,4.0,"+60°C","банк B горячее A",RED)
kpi(s,4.9,1.5,4.0,"1929","макс. об/мин",ORG)
kpi(s,9.3,1.5,3.5,"110°C","T масла (с породой)",ORG)
panel(s,0.5,2.85,12.33,2.6,"Банк B (чётные цилиндры) систематически горячее банка A на ~60°C — признак дисбаланса корректирующей дозы (TIB). Совпадает с зоной самых дефектных форсунок (B-банк 6/8) и самых горячих цилиндров (Ц2, Ц6).",RED,13,"Асимметрия")

# ── Часть 9: ДАШБОРД ФОРСУНОК ──
SEVCOL={'crit':RED,'warn':ORG,'ok':GRN}
nok=sum(1 for x in IDINJ if sev_of(x)=='ok'); nwarn=sum(1 for x in IDINJ if sev_of(x)=='warn'); ncrit=sum(1 for x in IDINJ if sev_of(x)=='crit')
s=slide(); section(s,"ЧАСТЬ 9","Дашборд форсунок (стенд CR-2)")
# 9.1 KPI summary (как в дашборде)
s=content("Дашборд форсунок — сводка","ФОРСУНКИ · ДАШБОРД")
kpi(s,0.5,1.4,3.0,str(len(IDINJ)),"форсунок испытано",ACC)
kpi(s,3.7,1.4,3.0,"3","машины (#73/#78/#82)",BLU)
kpi(s,6.9,1.4,3.0,f"{ncrit}",f"критичных",RED)
kpi(s,10.1,1.4,2.7,"59%","отказы на P04",ORG)
panel(s,0.5,2.6,12.33,1.5,"Стенд «Поток CR-2», Cummins CRI. Испытаны форсунки 3 машин NTE200 (#73, #78, #82). Главный признак — отказы на ЧАСТИЧНЫХ режимах (P04 пилот 59%, P03 малая нагрузка). Полные данные — дашборд `dashboard_injectors.html` в песочнице.",ACC,12,"Обзор")
rows=[["Состояние","Форсунок","Доля"],["✓ Годных",str(nok),f"{100*nok//len(IDINJ)}%"],["⚠ Внимание (1 откл.)",str(nwarn),f"{100*nwarn//len(IDINJ)}%"],["✗ Критичных (≥2)",str(ncrit),f"{100*ncrit//len(IDINJ)}%"]]
table(s,0.5,4.25,7.0,[0.5,0.25,0.25],rows,fsize=13,rowh=0.5)
# 9.2 СТАТУС-ГРИД всех форсунок
s=content("Состояние всех форсунок (статус-грид)","ФОРСУНКИ · ГРИД")
bym={}
for x in IDINJ: bym.setdefault(x['m'],[]).append(x)
y0=1.45
for mi,m in enumerate(['#73','#78','#82']):
    items=bym.get(m,[]);
    txt(s,0.5,y0,2.0,0.4,f"NTE200 {m}",13,True,INK)
    for i,x in enumerate(items):
        cx=2.5+i*0.66; cw=0.6
        c=SEVCOL[sev_of(x)]
        rect(s,cx,y0,cw,0.55,c)
        txt(s,cx,y0+0.03,cw,0.24,x['id'],11,True,WHITE,align=PP_ALIGN.CENTER)
        txt(s,cx,y0+0.27,cw,0.22,f"{x.get('p04','—')}",9,False,WHITE,align=PP_ALIGN.CENTER)
    y0+=0.85
# легенда
rect(s,0.5,4.6,0.4,0.3,GRN); txt(s,1.0,4.6,2,0.3,"годна",11,False,INK)
rect(s,2.5,4.6,0.4,0.3,ORG); txt(s,3.0,4.6,3,0.3,"1 откл. (внимание)",11,False,INK)
rect(s,5.7,4.6,0.4,0.3,RED); txt(s,6.2,4.6,4,0.3,"≥2 откл. (критично)",11,False,INK)
txt(s,0.5,5.2,12.3,0.5,"В ячейке — позиция форсунки и подача P04 (см³). Допуск P04: 41–51 см³.",11,True,GRY)
# 9.3 графики дашборда
chart_slide("Частота отказов по тестовым режимам","ФОРСУНКИ · РЕЖИМЫ",CH['idmode'],"Отказы концентрируются на частичных режимах P03/P04 — признак износа плунжерной пары/иглы.")
chart_slide("Дозирование P04 по всем форсункам","ФОРСУНКИ · P04",CH['idp04'],"Зелёная зона — допуск 41–51 см³; выше = перелив, ниже = недолив.")
chart_slide("Тепловая карта отклонений форсунок","ФОРСУНКИ · ТЕПЛОКАРТА",CH['idheat'])
# 9.4 классификация типов
s=content("Классификация отказов форсунок","ФОРСУНКИ · ТИПЫ")
s.shapes.add_picture(CH['idtype'],Inches(0.5),Inches(1.5),width=Inches(5.2))
panel(s,6.0,1.5,6.8,2.0,"Тип А — ПЕРЕЛИВ: износ иглы/седла распылителя → избыток топлива на ХХ → закоксовывание, локальное переобогащение, рост EGT.",RED,12,"Тип А — перелив")
panel(s,6.0,3.6,6.8,1.4,"Тип Б — НЕДОЛИВ: износ плунжерной пары → недостаток топлива → неравномерность мощности, угроза перегрева ЦПГ.",ORG,12,"Тип Б — недолив")
panel(s,6.0,5.1,6.8,1.0,"Тип В — прецизионная пара/граничные отклонения.",BLU,12,"Тип В")
# 9.5 карточки алертов (критичные)
crit=[a for a in IDAL if a.get('sev')=='crit']; others=[a for a in IDAL if a.get('sev')!='crit']
for grp,gname,gcol in [(crit,"Критичные форсунки — детально",RED),(others,"Прочие отклонения форсунок",ORG)]:
    for i in range(0,len(grp),4):
        s=content(gname,"ФОРСУНКИ · АЛЕРТЫ")
        for j,a in enumerate(grp[i:i+4]):
            yy=1.4+j*1.35
            rect(s,0.5,yy,12.33,1.2,WHITE); rect(s,0.5,yy,0.1,1.2,gcol)
            txt(s,0.7,yy+0.06,4.0,0.3,f"{a.get('m','')} · {a.get('id','')}",13,True,INK)
            txt(s,0.7,yy+0.38,12.0,0.8,a.get('reason','')[:180],11,False,INK)
            txt(s,8.7,yy+0.06,4.0,0.3,a.get('type','')[:30],10,True,gcol)
# 9.6 таблицы по машинам (полные значения)
for m in ['#73','#78','#82']:
    items=bym.get(m,[])
    s=content(f"Форсунки {m} — полные значения стенда","ФОРСУНКИ · "+m)
    rows=[["Поз","Дата","P01 (441–495)","P02 (220–260)","P03 (16–30)","P04 (41–51)","Статус"]]
    for x in items:
        st={'ok':'✓','warn':'⚠','crit':'✗'}[sev_of(x)]
        rows.append([x['id'],x.get('date',''),x.get('p01','—'),x.get('p02','—'),x.get('p03','—'),x.get('p04','—'),st])
    rh=min(0.5, 4.6/max(1,len(rows)))
    table(s,0.5,1.45,12.33,[0.1,0.18,0.18,0.18,0.13,0.13,0.1],rows,fsize=11,rowh=max(0.34,rh))
# 9.7 корреляция с EGT (#82)
s=content("Корреляция: форсунка ↔ цилиндр ↔ EGT (#82)","ФОРСУНКИ · КОРРЕЛЯЦИЯ")
panel(s,0.5,1.45,12.33,1.3,"Два самых горячих цилиндра Ц2 и Ц6 (567°C) — форсунки 1R и 3R, обе с ПЕРЕЛИВОМ на P03/P04 → локальное переобогащение → рост EGT → термоудар клапана.",RED,13,"Ключевая корреляция")
rows=[["Показатель","A-банк (L)","B-банк (R)"],["Дефектных форсунок","3/8","6/8"],["Средняя EGT (DML)","~465°C","~525°C (+60)"],["Горячие цилиндры","—","Ц2=567, Ц6=567"]]
table(s,0.5,2.95,8.4,[0.4,0.3,0.3],rows,fsize=13,rowh=0.5)
panel(s,9.1,2.95,3.7,2.1,"Форсунки — отягчающий фактор поверх калибровки ЭБУ. Перелив усиливает перегрев в горячих цилиндрах; недолив даёт неравномерность.",GRY,11,"Роль")

# ── Зазоры клапанов ↔ форсунки (отчёты #54/#77/#87/#89) ──
s=content("Зазоры клапанов — отчёты #54/#77/#87/#89","КЛАПАНЫ · ЗАЗОРЫ")
rows=[["Маш","Наработка","Зазоры / высота клапана","Втулки / стержни / сёдла"],
["#87","4502 м/ч","некорректные зазоры; ↑ высота (просадка)","втулки НЕ в допуске"],
["#89","4132 м/ч","некорректные зазоры; ↑ высота 3L, 4R","стержень повреждён; пружина 7L"],
["#77","10399 м/ч","регулировка всех; просадка 7R","деформ. стержень 7R; разруш. седла 7R/4R"],
["#54","18245 м/ч","замеры зазоров/высоты","выработка стержня НЕ в допуске; нагар"]]
table(s,0.5,1.5,12.33,[0.1,0.18,0.38,0.34],rows,fsize=12,rowh=0.62)
panel(s,0.5,4.6,12.33,1.4,"Общий признак: ПРОСАДКА/увеличение высоты выпускных клапанов (рецессия седла) + износ направляющих втулок/стержней → потеря/изменение теплового зазора. Числовых таблиц зазоров в отчётах нет (качественные формулировки).",ORG,12,"Что выявлено")
# механистическая корреляция
s=content("Корреляция: форсунка → EGT → зазор клапана","КОРРЕЛЯЦИЯ · МЕХАНИЗМ")
chain=["Форсунка ПЕРЕЛИВ (Тип А, износ иглы/седла)","Избыток топлива → ↑ локальная EGT (до 567°C, Ц2/Ц6)","Термоудар фаски + РЕЦЕССИЯ (просадка) седла клапана","↑ высота клапана → выборка теплового зазора","Неполная посадка → прорыв газов → выгорание клапана"]
for i,t in enumerate(chain):
    rect(s,0.5,1.5+i*1.02,12.0,0.82,WHITE); rect(s,0.5,1.5+i*1.02,0.1,0.82,RED if i>=2 else ORG)
    txt(s,0.8,1.6+i*1.02,11.5,0.66,f"{i+1}.  {t}",13,True,INK,anchor=MSO_ANCHOR.MIDDLE)
# вердикт корреляции
s=content("Корреляция зазор ↔ форсунки — вердикт","КОРРЕЛЯЦИЯ · ВЕРДИКТ")
panel(s,0.5,1.45,6.0,4.5,"МЕХАНИСТИЧЕСКИ — связь сильная и положительная:\n\nнеисправные форсунки (особенно ПЕРЕЛИВ) → ↑ EGT → рецессия седла/просадка клапана → потеря теплового зазора → ускоренное выгорание.\n\n«Некорректные зазоры» и «увеличенная высота клапанов» (#87/#89) — следствие этой цепи, усиленной отключёнными защитами ЭБУ.",RED,13,"Связь установлена (причинно)")
panel(s,6.7,1.45,6.13,4.5,"СТАТИСТИЧЕСКИ (по машинам) — не доказуема на текущих данных:\n\n• форсунки на стенде: #73/#78/#82\n• зазоры/клапаны: #54/#77/#87/#89\n• ПЕРЕСЕЧЕНИЯ ПО МАШИНАМ НЕТ → числовой коэффициент не рассчитать.\n\nРЕКОМЕНДАЦИЯ: при каждой регулировке фиксировать численно зазор/просадку по цилиндрам И снимать форсунки на стенд → парная корреляция «просадка_цил ↔ дефект_форсунки_цил».",ORG,12,"Ограничение и рекомендация")

# ── Часть 10: Топливо ──
s=slide(); section(s,"ЧАСТЬ 10","Топливо и чистота (ISO 4406)")
s=content("Качество топлива по сере","ТОПЛИВО · СЕРА")
kpi(s,0.5,1.5,4.0,"2–10","сера, мг/кг",GRN)
kpi(s,4.9,1.5,4.0,"6,7","среднее, мг/кг",GRN)
kpi(s,9.3,1.5,3.5,"≤10 (Е5)","класс",GRN)
panel(s,0.5,2.85,12.33,2.2,"25+ паспортов: сера в 2,4× ниже порога Cummins (15 ppm). Топливо по сере ИСКЛЮЧЕНО как причина. Нарушение: паспорт №246 (декларация ЕАЭС истекла 24.10.2025).",GRN,13,"Сера")
s=content("Чистота топлива ISO 4406","ТОПЛИВО · ЧИСТОТА")
s.shapes.add_picture(CH['fuel'],Inches(0.5),Inches(1.5),width=Inches(5.7))
rows=[["Маш","После тонкой","18/16/13","12/9/6"],["46","17/13/13","✓","✗"],["47","19/16/14","✗","✗"],["81","21/18/17","✗","✗"],["82","21/16/14","✗","✗"]]
table(s,6.4,1.6,6.4,[0.2,0.4,0.2,0.2],rows,fsize=12,rowh=0.5)
panel(s,6.4,3.6,6.4,2.4,"Норма Cummins (бак) 18/16/13; цель «вход в цилиндры» (Fleetguard/Горная Евразия) 12/9/6. Цель не достигается нигде; фильтр грубой очистки местами УХУДШАЕТ чистоту. Топливо одно для NTE200 и 730E → не дифференцирует, но угрожает форсункам MCRS (2–5 мкм).",ORG,11,"Анализ")

# ── Часть 11: KAMSS ──
s=slide(); section(s,"ЧАСТЬ 11","Обследование KAMSS (3 единицы)")
s=content("KAMSS — обследование 3 единиц","KAMSS")
panel(s,0.5,1.5,12.33,1.8,"Фото-отчёт KAMSS (INSITE-диагностика + дефектация). По коды неисправностей двигателей: 2138/2137 «температура EGT цилиндров 16(B8)/14(B7) выше нормы», 622 «отклонение EGT цил.3», защита: низкий уровень ОЖ (235), давление картерных газов (556).",RED,13,"INSITE-коды")
panel(s,0.5,3.5,12.33,1.8,"Дефектация: фото эрозии рабочих фасок выпускных клапанов (тарелка/седло). Подтверждает термический характер повреждения и перегрев именно B-банка — согласуется с DML и кодами ECM.",ORG,13,"Дефектация клапанов")
# фото KAMSS
imgs=sorted(glob.glob('_Изображения/KAMSS обследование 3 единиц/*.jpeg'))
if len(imgs)>=2:
    s=content("KAMSS — визуальные свидетельства","KAMSS · ФОТО")
    try: s.shapes.add_picture(imgs[0],Inches(0.5),Inches(1.5),width=Inches(6.0))
    except: pass
    try: s.shapes.add_picture(imgs[4] if len(imgs)>4 else imgs[1],Inches(6.8),Inches(1.5),width=Inches(6.0))
    except: pass
    txt(s,0.5,6.55,12.3,0.4,"Слева — INSITE коды EGT по цилиндрам; справа — эрозия фаски выпускного клапана.",11,True,INK)

# ── Часть 12: 730E ──
s=slide(); section(s,"ЧАСТЬ 12","Контрольная группа 730E-DC")
s=content("730E-DC — паритет условий","730E")
rows=[["Фактор","NTE200","730E-DC","Дифференц.?"],["Двигатель","QSK50","QSK50","нет"],["Топливо","Е5, 6,7 мг/кг","то же","нет"],
["Масло","CI-4 15W-40","то же","нет"],["TIB","200%","100%","ДА"],["Защиты перегрева","ОТКЛ","Вкл","ДА"],["Отказы ГБЦ","~17","0","ДА"]]
table(s,0.5,1.5,12.33,[0.3,0.25,0.25,0.2],rows,fsize=12,rowh=0.5)
panel(s,0.5,5.4,12.33,1.0,"Единственные дифференцирующие факторы — в калибровке ЭБУ. При прочих равных 730E с TIB=100% и включёнными защитами — 0 отказов ГБЦ.",GRN,13)

# ── Часть 13: NHL/CCEC ──
s=slide(); section(s,"ЧАСТЬ 13","Опровержение версий NHL / CCEC")
s=content("Оценка аргументов NHL и Cummins","ПЕРЕПИСКА")
rows=[["Тезис письма","Вердикт"],["NHL: зольность масла (SAPS)","✗ опровергнут (низкая сажа 0,3)"],["NHL: запылённость","✗ переоценка (Si низкий)"],
["NHL: Stellite не нужен «для rated»","✗ NTE200 ≠ rated (567°C); факт признан"],["Cummins: масло CES20078 достаточно","✗ уход от корня"],
["Cummins: на NTE200 защита по T масла ОТКЛ","✓ ПОДТВЕРЖДАЕТ нашу находку"],["Cummins: «не поставляет Stellite»","✗ уход от ответственности"]]
table(s,0.5,1.5,12.33,[0.55,0.45],rows,fsize=12,rowh=0.55)
chart_slide("Статистическое опровержение зольной версии","CCEC",CH['oilfe'],"Низкий Fe и сажа → не абразивный/зольный износ.")

# ── Часть 13б: Проверка гипотез ──
s=slide(); section(s,"ЧАСТЬ 13.5","Проверка гипотез (best practice)")
s=content("Формулировка конкурирующих гипотез","ГИПОТЕЗЫ")
panel(s,0.5,1.45,12.33,1.7,"H1 (термическая): первична калибровка ЭБУ (TIB, отключённые защиты) → перегрев → выгорание клапанов.\nH2 (контаминационная): первичны пыль (негерметичность впуска) + загрязнение топлива → абразивно-механический износ (втулки/стержни клапанов, форсунки); ЭБУ — усугубляющий усилитель.",ACC,13,"Две гипотезы")
panel(s,0.5,3.35,12.33,1.5,"Метод проверки: сопоставление сигнатур отказа. Термическая → выгорание/эрозия ФАСКИ клапана, высокая EGT. Абразивная → износ НАПРАВЛЯЮЩИХ ВТУЛОК и СТЕРЖНЕЙ, повышенный Si в масле. Реальные отчёты содержат ОБА признака → проверяем вклад каждого по данным.",ORG,13,"Принцип")
# доказательства за H2 (пыль/впуск)
s=content("Доказательства за H2: пыль и впуск","H2 · ЗА")
rows=[["Доказательство","Источник","Сигнатура"],["Критическое засорение возд. фильтров","Техотчёт #69","пылевая среда"],
["«Высокая запылённость карьера»","Техотчёт #69","среда"],["Коррозия ОНВ + ОЖ на впуске ГБЦ","Эндоскопия впуска","негерметичн. впуска"],
["Износ направляющих втулок (вне допуска)","Техотчёты #54/#78/#87","абразив/механика"],["Деформация/выработка стержня клапана","Техотчёты","абразив/механика"],
["Осадок/вода в топл. сепараторе","Фото отстойников","загрязн. топлива"]]
table(s,0.5,1.45,12.33,[0.42,0.3,0.28],rows,fsize=12,rowh=0.5)
# фото отстойники + впуск
imo=sorted(glob.glob('_Изображения/Отстойники NHL/*.jpg'))
if imo:
    s=content("Доказательства H2 — визуально","H2 · ФОТО")
    try: s.shapes.add_picture(imo[0],Inches(0.5),Inches(1.5),height=Inches(4.6))
    except: pass
    if len(imo)>1:
        try: s.shapes.add_picture(imo[1],Inches(5.0),Inches(1.5),height=Inches(4.6))
        except: pass
    panel(s,9.5,1.5,3.3,4.6,"Топливный отстойник: осадок/вода на дне колбы; вокруг — густая пыль. Подтверждает загрязнение топлива и запылённость среды (H2).",ORG,12,"Отстойник")
# контрдоказательства (Si/Fe низкие)
s=content("Контрдоказательства против H2 (масло)","H2 · ПРОТИВ")
s.shapes.add_picture(CH['si'],Inches(0.5),Inches(1.5),width=Inches(5.7))
panel(s,6.4,1.5,6.4,4.5,"Если бы пыль массово попадала в цилиндры через впуск, в масле ДВС вырос бы кремний (Si) и железо (Fe). Факт: Si P95=4 ppm (порог ~15), лишь 0,3% проб >15; Fe P95=6 ppm — НИЗКИЕ. ⇒ Массовый абразивный пылевой износ масло НЕ подтверждает. Всплески Si единичны (#88=50, #53=31) — эпизодические/локальные события, не системные.",RED,13,"Объективная оговорка")
# топливо->форсунки
s=content("Цепь: загрязнение топлива → форсунки","H2 · ТОПЛИВО→ФОРСУНКИ")
s.shapes.add_picture(CH['p04'],Inches(0.5),Inches(1.5),width=Inches(12.3))
txt(s,0.5,6.5,12.3,0.5,"Грязное топливо (ISO хуже 18/16/13) → износ плунжерных пар → разброс дозирования P04 (перелив/недолив) → дисбаланс сгорания и EGT.",12,True,ORG)
# матрица гипотез
s=content("Матрица проверки гипотез","ВЕРДИКТ ПО ГИПОТЕЗАМ")
rows=[["Фактор","Доказ. ЗА","Доказ. ПРОТИВ","Роль"],
["Пыль/впуск (H2)","фильтры, ОНВ, втулки/стержни","Si/Fe низкие в масле","инициирующий (локально)"],
["Топливо/форсунки (H2)","ISO плохой, осадок, разброс P04","сера в норме","инициирующий"],
["Калибровка ЭБУ (H1)","15 параметров, EGT 567, коды EGT","—","усиливающий амплификатор"],
["Конструкция (нет Stellite)","признано NHL/Cummins","730E без Stellite не ломается","сопутствующий"]]
table(s,0.5,1.45,12.33,[0.27,0.3,0.25,0.18],rows,fsize=11.5,rowh=0.6)
panel(s,0.5,4.45,12.33,1.5,"Вывод: ни один фактор не объясняет отказ в одиночку. Это МНОГОФАКТОРНАЯ ЦЕПЬ: контаминация (пыль+топливо) инициирует износ направляющих/стержней/форсунок, а калибровка ЭБУ (отключённые защиты + TIB) усиливает тепловую нагрузку и доводит клапан до выгорания.",RED,13,"Синтез")
# FMEA
s=content("FMEA — режимы и причины отказа","FMEA")
rows=[["Режим отказа","Причина","S","O","D","RPN"],
["Выгорание фаски вып. клапана","EGT>520 (ЭБУ/форсунки)","10","7","5","350"],
["Износ направл. втулки/стержня","Абразив (пыль/впуск)","8","6","6","288"],
["Разброс дозирования форсунок","Загрязн. топлива","7","8","4","224"],
["Просадка/негерметичн. седла","Термоудар + износ","9","6","5","270"],
["Коррозия ОНВ / ОЖ на впуске","Негерметичн. тракта","6","4","6","144"]]
table(s,0.5,1.45,12.33,[0.34,0.3,0.09,0.09,0.09,0.09],rows,fsize=12,rowh=0.55)
txt(s,0.5,5.1,12.3,0.4,"S — тяжесть, O — частота, D — обнаруживаемость, RPN = S×O×D (приоритет риска).",11,False,GRY)
# 5-Why
s=content("5 Почему (5-Why)","5-WHY")
whys=["Почему разрушилась ГБЦ? — Выгорела/просела фаска выпускного клапана.","Почему выгорел клапан? — Хронически высокая EGT (до 567°C) + сниженная стойкость.",
"Почему высокая EGT? — Отключены тепловые защиты ЭБУ (TIB 200%, derate/триггеры) + разброс форсунок.","Почему разброс форсунок и износ направляющих? — Загрязнённое топливо и пыль (негерметичность впуска, засорённые фильтры).",
"Почему загрязнение? — Запылённость карьера + недостаточная фильтрация/герметичность + калибровка, не снижающая нагрузку."]
for i,t in enumerate(whys):
    rect(s,0.5,1.5+i*1.05,12.0,0.88,WHITE); rect(s,0.5,1.5+i*1.05,0.1,0.88,RED if i<2 else ORG)
    txt(s,0.8,1.6+i*1.05,11.5,0.7,t,12.5,True,INK,anchor=MSO_ANCHOR.MIDDLE)

s=content("Матрица связей факторов","КОРРЕЛЯЦИИ")
s.shapes.add_picture(CH['corr'],Inches(0.5),Inches(1.5),height=Inches(4.4))
panel(s,7.2,1.5,5.6,4.4,"Качественная матрица связей: сильнее всего связаны EGT↔TIB↔перелив форсунок (термокластер). Si (пыль) слабо связан с EGT, но связан с чистотой топлива. Подтверждает: термический кластер (ЭБУ/форсунки) — главный усилитель; пыль — отдельный, более слабый по маслу канал.",ACC,12,"Чтение матрицы")

# ── Часть 14: RCA ──
s=slide(); section(s,"ЧАСТЬ 14","Коренная причина (RCA)")
s=content("Механизм отказа — многофакторная цепь","RCA · МЕХАНИЗМ")
steps=["СРЕДА: запылённость карьера + засорённые возд. фильтры + негерметичность впуска (коррозия ОНВ, ОЖ) + грязное топливо","ИНИЦИАЦИЯ: износ направляющих втулок/стержней клапанов + разброс дозирования форсунок (перелив/недолив)",
"УСИЛЕНИЕ (ЭБУ): TIB 200% + отключённые защиты перегрева → нагрузка не снижается → ↑ EGT, дисбаланс банков","ПЕРЕГРЕВ: EGT до 567°C (лимит 520); термоудар фаски клапана (стойкость снижена — нет Stellite)","ОТКАЗ: выгорание/просадка клапана и седла → потеря герметичности → разрушение ГБЦ"]
accs=[ORG,ORG,RED,RED,RED]
for i,t in enumerate(steps):
    rect(s,0.5,1.5+i*1.02,12.0,0.82,WHITE); rect(s,0.5,1.5+i*1.02,0.1,0.82,accs[i])
    txt(s,0.8,1.6+i*1.02,11.5,0.66,f"{i+1}.  {t}",12.5,True,INK,anchor=MSO_ANCHOR.MIDDLE)
s=content("Дерево отказов (FTA) и факторы","RCA · FTA")
panel(s,0.5,1.5,6.0,4.6,"ВЕРШИНА: разрушение ГБЦ\n└ выгорание/просадка клапана\n   ├ перегрев (EGT>520)\n   │  ├ УСИЛЕНИЕ ЭБУ: TIB 200%, защиты ОТКЛ\n   │  └ разброс форсунок (топливо)\n   ├ износ втулок/стержней\n   │  └ пыль / негерметичн. впуска\n   └ стойкость фаски (нет Stellite)",RED,12,"FTA")
panel(s,6.7,1.5,6.13,4.6,"Исикава (5M):\n• Man/Env: запылённость, фильтры (#69)\n• Material: топливо (чистота), пыль; масло — OK\n• Machine: негерметичн. впуска (ОНВ/ОЖ), нет Stellite\n• Method: калибровка ЭБУ (защиты ОТКЛ), дозирование форсунок\n• Measurement: датчики EGT (FC1531/1619)",ACC,12,"Исикава")
s=content("Матрица доказательств","RCA · ДОКАЗАТЕЛЬСТВА")
rows=[["Доказательство","Вес","Статус"],["TIB 200% vs 100% + защиты ОТКЛ","🔴 ключевое","✓ ECM/письмо Cummins"],
["EGT 567°C (Ц2/Ц6)","🔴 ключевое","✓ DML #82"],["730E — 0 отказов","🔴 ключевое","✓ паритет условий"],
["Коды EGT (KAMSS/ECM)","🟠 важное","✓ 2137/2138/628"],["Форсунки перелив ↔ горячие цил.","🟠 важное","✓ стенд+DML"],
["Нет Stellite","🟠 важное","✓ NHL/Cummins"],["Масло/топливо/перегруз","🟢 исключ.","✓ не дифференц."]]
table(s,0.5,1.5,12.33,[0.5,0.22,0.28],rows,fsize=12,rowh=0.55)
s=content("Формулировка коренной причины","RCA · ВЫВОД")
panel(s,0.5,1.5,12.33,3.6,"Отказы ГБЦ QSK50 на NTE200 имеют МНОГОФАКТОРНУЮ природу (цепь причин), а не единственную причину.\n\nИНИЦИИРУЮЩИЕ факторы (эксплуатация/контаминация): запылённость и засорение воздушных фильтров, негерметичность впускного тракта (коррозия ОНВ, следы ОЖ), загрязнённое топливо — вызывают абразивно-механический износ направляющих втулок/стержней клапанов и разброс дозирования форсунок.\n\nУСИЛИВАЮЩИЙ (амплифицирующий) фактор: калибровка ЭБУ — TIB 200% и ОТКЛЮЧЁННЫЕ тепловые защиты (триггеры перегрева, derate) не снижают нагрузку при перегреве, доводя EGT до 567°C и термоудар клапана.\n\nСОПУТСТВУЮЩИЙ: отсутствие наплавки Stellite снижает стойкость фаски.",RED,13,"Коренная причина (многофакторная)")
panel(s,0.5,5.25,12.33,1.05,"Оговорка по данным: масляные маркеры пыли (Si) и износа (Fe) НИЗКИЕ — массовый абразивный пылевой износ через масло не подтверждён; вклад пыли локальный/эпизодический. Сера топлива и систематическая перегрузка (10/10/20) — исключены.",GRY,11,"Достоверность")

# ── Часть 14б: Досье по машинам ──
s=slide(); section(s,"ПРИЛОЖЕНИЕ A","Досье по машинам NTE200")
REPBY={}
for m,d,h,op in REPAIRS: REPBY.setdefault(m,[]).append(f"{d}: {op}")
regby={str(x['gar']):x for x in REG}
notable=set(REPBY)|{'82','85','72','61','67','84','88'}
allm=sorted([m for m in notable if m.isdigit()], key=lambda z:int(z))
for m in allm:
    r=regby.get(m,{}); ld=PM['load'].get(m); ol=PM['oil'].get(m); ec=PM['ecm'].get(m); rp=REPBY.get(m)
    s=content(f"Машина №{m}","ДОСЬЕ NTE200")
    # паспорт
    panel(s,0.5,1.4,6.0,2.0,f"PIN: {r.get('pin','—')}\nДВС №: {r.get('eng_no','—')}  ·  Год: {r.get('year','—')}\nГарантия: {r.get('warr_start','—')} → {r.get('warr_end','—')}\nГенератор: {r.get('gen','—')}",ACC,12,"Паспорт")
    # статус ECM
    stc=ec['status'] if ec else '—'
    acc=RED if stc in('Отказ','Катастроф.') else (ORG if stc=='EGT/связан.' else GRN)
    panel(s,6.7,1.4,6.13,2.0,f"Статус ЭБУ: {stc}\nПрошивка: {ec['code'] if ec else '—'}\nНаработка: {ec['hrs'] if ec else '—'} ч  ·  Макс. об/мин: {ec['maxrpm'] if ec else '—'}",acc,12,"ECM")
    # нагрузка + масло
    if ld:
        rows=[["Нагрузка","знач."],["Рейсов",ld['n']],["Средн. %",ld['mean']],["Макс %",ld['max']],[">100%",f"{ld['over100']}%"],["Перевезено",f"{ld['tons']} млн т"]]
        table(s,0.5,3.6,3.9,[0.6,0.4],rows,fsize=11,rowh=0.37)
    if ol:
        rows=[["Масло ДВС","знач."],["Проб",ol['n']],["Fe (мед.)",ol['Fe']],["TBN",ol['TBN']],["Сажа",ol['ST']],["Марка",(ol['brand'] or '')[:14]]]
        table(s,4.6,3.6,4.1,[0.45,0.55],rows,fsize=11,rowh=0.37)
    # ремонты
    rptxt='\n'.join(rp) if rp else 'Ремонтов ГБЦ/клапанов не зафиксировано'
    panel(s,8.9,3.6,3.93,2.6,rptxt,(RED if rp else GRN),11,"Ремонты ГБЦ")

# ── Приложение B: аналитические срезы ──
s=slide(); section(s,"ПРИЛОЖЕНИЕ B","Аналитические срезы и доп. данные")
# DML доп. параметры
for nm,kick,vals,note in [
 ("Давление в топливной рампе (CR)","DML · CR",[("~1060 бар","рабочее"),("2100 бар","макс. система"),("стабильно","разброс")],"Рампа в норме — проблема не в давлении, а в дозировании форсунок."),
 ("Давление масла","DML · МАСЛО",[("~380 кПа","с нагрузкой"),(">241 кПа","мин. норма"),("OK","статус")],"Давление масла в норме во всех режимах."),
 ("Давление наддува","DML · НАДДУВ",[("~200 кПа","с нагрузкой"),("~110 кПа","без нагрузки"),("OK","статус")],"Наддув в норме; перегрев — следствие калибровки, не турбонаддува."),
 ("Расход топлива и режимы","DML · ТОПЛИВО",[("88%","нагрузка с породой"),("21%","без породы"),("1822","ср. об/мин с нагрузкой")],"Высокая нагрузка повышает базовую EGT; TIB добавляет пульсации.")]:
    s=content(nm,kick)
    for i,(v,l) in enumerate(vals): kpi(s,0.5+i*4.25,1.5,4.0,v,l,ACC)
    panel(s,0.5,2.9,12.33,2.4,note,ACC,13,"Вывод")
# аналитические топы
def topslide(title,kick,items,hdr,note,acc=ORG):
    s=content(title,kick); rows=[hdr]+items
    table(s,0.5,1.5,9.0,[0.25]+[0.75/ (len(hdr)-1)]*(len(hdr)-1),rows,fsize=12,rowh=0.45)
    panel(s,9.7,1.5,3.1,4.4,note,acc,11)
topFe=sorted([(m,d) for m,d in PM['oil'].items() if d.get('Fe') is not None],key=lambda z:-z[1]['Fe'])[:10]
topslide("Топ по железу (Fe) в масле ДВС","МАСЛО · ИЗНОС",[[m,d['Fe'],d['TBN'],d['n']] for m,d in topFe],["Маш","Fe","TBN","Проб"],"Даже у лидеров Fe умеренный — износ низкий; отказ термический, не абразивный.",GRN)
topT=sorted([(m,d) for m,d in PM['oil'].items() if d.get('TBN') is not None],key=lambda z:z[1]['TBN'])[:10]
topslide("Минимальный TBN (старение масла)","МАСЛО · TBN",[[m,d['TBN'],d['ST'],d['n']] for m,d in topT],["Маш","TBN","Сажа","Проб"],"TBN везде высоко выше предела 2,5 — масло не деградировало.",GRN)
topL=sorted([(m,d) for m,d in PM['load'].items() if m.isdigit()],key=lambda z:-z[1]['tons'])[:10]
topslide("Топ по перевезённой массе","НАГРУЗКА",[[m,f"{d['tons']}",d['n'],f"{d['mean']}%"] for m,d in topL],["Маш","млн т","Рейсов","Ср%"],"Лидеры по тоннажу при средней нагрузке ≈100% — режим в пределах 10/10/20.",ACC)
topO=sorted([(m,d) for m,d in PM['load'].items() if m.isdigit()],key=lambda z:-z[1]['over100'])[:10]
topslide("Топ по доле рейсов >100%","НАГРУЗКА",[[m,f"{d['over100']}%",f"{d['max']}%",d['n']] for m,d in topO],["Маш",">100%","Макс%","Рейсов"],"Доля >100% — не нарушение 10/10/20 (зона 100–110% допускается).",ACC)
# KAMSS доп фото
imgs=sorted(glob.glob('_Изображения/KAMSS обследование 3 единиц/*.jpeg'))
for gi in range(0,min(8,len(imgs)),2):
    s=content("KAMSS — материалы обследования","KAMSS · ФОТО")
    try: s.shapes.add_picture(imgs[gi],Inches(0.5),Inches(1.5),width=Inches(6.0))
    except: pass
    if gi+1<len(imgs):
        try: s.shapes.add_picture(imgs[gi+1],Inches(6.8),Inches(1.5),width=Inches(6.0))
        except: pass
    txt(s,0.5,6.55,12.3,0.4,"INSITE-диагностика (коды EGT, защита двигателя) и дефектация клапанов/ГБЦ — обследование 3 единиц.",11,True,INK)
# методика, словарь, стандарты
s=content("Методика исследования","МЕТОДИКА")
panel(s,0.5,1.5,12.33,4.6,"1. Сбор первичных данных (Весовая, ECM, DML, масло, топливо, стенд форсунок, KAMSS, парк).\n2. Перепроверка из первоисточников (пересчёт нагрузки от 180 т; парсинг ECM/масла/форсунок).\n3. Сравнение NTE200 vs контрольной группы 730E (паритет условий).\n4. Проверка гипотез «за/против»; исключение недифференцирующих факторов.\n5. Формулировка коренной причины и плана действий.\n\nПринцип: каждое число — из первичных данных; неподтверждённое помечается отдельно.",ACC,13,"Этапы")
s=content("Словарь и обозначения","СЛОВАРЬ")
panel(s,0.5,1.5,12.33,4.6,"EGT — температура отработавших газов · TIB — компенсация теплового дисбаланса (Thermal Imbalance Balance) · MCRS — модульная система Common Rail · derate — снижение мощности/оборотов защитой · DO_Option — профиль опций ЭБУ · ISO 4406 — класс чистоты (частицы >4/6/14 мкм) · TBN — щелочное число · CI-4/CES 20078 — класс моторного масла · A/B банк — ряды цилиндров (нечёт./чёт.) · P03/P04 — тест-планы стенда (малая нагрузка/пилот).",GRY,13,"Термины")

# ── Часть 15: Выводы ──
s=slide(); section(s,"ЧАСТЬ 15","Выводы и действия")
s=content("Неотложные действия","ДЕЙСТВИЯ · 🔴")
acts=["Перекалибровать ЭБУ NTE200 под профиль 730E: TIB 100%, включить защиты перегрева, убрать задержки derate, обороты 1990","Контрольный DML после перекалибровки — подтвердить снижение EGT","Заменить датчики EGT с FC1531/1619 (#82, #67)","Стендовая проверка/замена бракованных форсунок (перелив P03/P04)"]
for i,t in enumerate(acts):
    rect(s,0.5,1.5+i*1.15,12.0,0.95,WHITE); rect(s,0.5,1.5+i*1.15,0.1,0.95,RED)
    txt(s,0.8,1.62+i*1.15,11.5,0.75,f"{i+1}.  {t}",13,True,INK,anchor=MSO_ANCHOR.MIDDLE)
s=content("Средне- и долгосрочные действия","ДЕЙСТВИЯ · 🟠🟡")
acts2=["Запрос NHL/Cummins: обоснование исключения Stellite и отключения защит; поставка клапанов с наплавкой","Привести чистоту топлива к ≤18/16/13 (доп. фильтрация), ориентир 12/9/6","Унифицировать прошивку ЭБУ (.07/.08) по парку","Мониторинг EGT и FC611 в реальном времени","Испытать форсунки 730E на стенде для сравнения"]
for i,t in enumerate(acts2):
    rect(s,0.5,1.5+i*0.95,12.0,0.78,WHITE); rect(s,0.5,1.5+i*0.95,0.1,0.78,ORG)
    txt(s,0.8,1.6+i*0.95,11.5,0.62,f"{i+1}.  {t}",12.5,True,INK,anchor=MSO_ANCHOR.MIDDLE)
s=content("Ожидаемый эффект и риски","ЭФФЕКТ")
kpi(s,0.5,1.5,4.0,"↓ EGT","< 520°C после калибровки",GRN)
kpi(s,4.9,1.5,4.0,"↓ отказы","клапанов/ГБЦ",GRN)
kpi(s,9.3,1.5,3.5,"↑ ресурс","ГБЦ и форсунок",GRN)
panel(s,0.5,2.85,12.33,2.4,"Перекалибровка ЭБУ устраняет коренную причину без замены узлов. Установка клапанов со Stellite и нормализация фильтрации топлива/дозирования форсунок закрывают отягчающие факторы. Риск бездействия — продолжение массовых отказов ГБЦ в гарантийный период.",ACC,13,"Итог")
s=content("Матрица рисков","РИСКИ")
rows=[["Риск","Вероятность","Тяжесть","Действие"],["Продолжение отказов ГБЦ (без калибровки)","Высокая","Критич.","Перекалибровка ЭБУ"],
["Выход из строя форсунок MCRS","Средняя","Высокая","Фильтрация топлива + замена"],["Перегрев масла/ОЖ без защиты","Средняя","Высокая","Включить триггеры ЭБУ"],
["Гарантийные претензии","Высокая","Средняя","Документирование, запрос NHL"],["Простой парка","Средняя","Высокая","Плановый график ремонта"]]
table(s,0.5,1.5,12.33,[0.4,0.2,0.18,0.22],rows,fsize=12,rowh=0.55)
s=content("Дорожная карта внедрения","ДОРОЖНАЯ КАРТА")
phases=[("0–1 мес","Перекалибровка ЭБУ пилотных машин + контрольный DML",RED),("1–3 мес","Перекалибровка всего парка, замена датчиков EGT/форсунок",ORG),
("3–6 мес","Нормализация фильтрации топлива, запрос Stellite у NHL",ORG),("6–12 мес","Мониторинг EGT/FC611, верификация снижения отказов",GRN)]
for i,(t,d,c) in enumerate(phases):
    rect(s,0.5,1.6+i*1.15,12.0,0.95,WHITE); rect(s,0.5,1.6+i*1.15,0.1,0.95,c)
    txt(s,0.8,1.7+i*1.15,2.2,0.75,t,15,True,c,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,3.1,1.7+i*1.15,9.2,0.75,d,13,False,INK,anchor=MSO_ANCHOR.MIDDLE)
s=content("Ответственные стороны","КОНТАКТЫ")
rows=[["Роль","Организация"],["Оператор","Полюс Магадан"],["Подрядчик ТО","ООО «Горная Евразия»"],["Производитель двигателя","Cummins / CCEC"],
["Производитель техники","NHL (North Hauler)"],["Лаборатория масла","МИЦ ГСМ"],["Лаборатория топлива","Дорингтест / ООО «Авиатек»"]]
table(s,0.5,1.5,12.33,[0.4,0.6],rows,fsize=13,rowh=0.55)
s=content("Источники и стандарты","ИСТОЧНИКИ")
panel(s,0.5,1.5,12.33,4.4,"• Данные Весовая (131 файл, 365 069 рейсов) — нагрузка/производительность\n• Своды анализов масла 2025+2026 (ДВС, NTE200/730E)\n• ECM: 15 параметров (ECM_对比), Fault Codes (361 файл INSITE)\n• DML #82 (INSITE логи), стенд «Поток CR-2» (форсунки)\n• Паспорта топлива, Дорингтест ISO 4406, Cummins требования к ДТ\n• KAMSS — обследование 3 единиц; письма NHL/Cummins\n• Реестр парка ГЕ (апрель 2026)\n\nСтандарты: ISO 4406, ТР ТС 013/2011, ГОСТ 32511-2013, Cummins CES 20078, Cummins SB 5407411 (378-005/010/011/013).",GRY,13,"База")
# Заключение
s=slide(INK); rect(s,0,3.5,W,0.06,ACC)
txt(s,0.8,2.3,11.7,1.0,"Отказ — многофакторная цепь: контаминация (пыль/впуск + топливо)",21,True,WHITE)
txt(s,0.8,3.0,11.7,1.0,"инициирует износ, калибровка ЭБУ усиливает перегрев до выгорания клапана.",21,True,WHITE)
txt(s,0.8,4.5,11.7,0.5,"Перегрузка и сера топлива — исключены. Масляные маркеры пыли (Si) низкие — вклад пыли локальный.",14,False,ACC)
txt(s,0.8,6.5,11.7,0.4,"АО «Развитие» · техническая служба + отдел надёжности · 2026",11,True,GRY)

# ── футеры с нумерацией ──
total=len(prs.slides._sldIdLst)
for i,s in enumerate(prs.slides):
    # skip dark dividers/cover? add page to all content slides
    footer_bar(s)
    txt(s,11.7,7.19,1.5,0.28,f"{i+1} / {total}",9,True,RGBColor(0xCF,0xD6,0xDA),align=PP_ALIGN.RIGHT)

OUT="05 Отчёты/Презентации/Анализ_ГБЦ_QSK50_NTE200_БОЛЬШОЙ.pptx"
prs.save(OUT)
print("saved:",OUT,"| слайдов:",total)
