#!/usr/bin/env python3
"""
Presentation: Analysis of exhaust valve failures on QSK50 engines, NTE200 dump trucks
АО Развитие brand style — #293136 dark / #F5F5F5 light / #3EF0AF accent
"""
import os, io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import fitz

WORK_DIR = "/home/user/NTE200"

# ── АО Развитие Brand Palette ─────────────────────────────────────────────────
C_DARK   = RGBColor(0x29, 0x31, 0x36)   # #293136 — primary dark bg
C_DARK2  = RGBColor(0x1F, 0x26, 0x2B)   # deeper dark for headers
C_LIGHT  = RGBColor(0xF5, 0xF5, 0xF5)   # #F5F5F5 — content slide bg
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_ACC    = RGBColor(0x3E, 0xF0, 0xAF)   # #3EF0AF — Развитие teal accent
C_GRN2   = RGBColor(0x10, 0xB9, 0x81)   # darker teal
C_AMB    = RGBColor(0xF5, 0x9E, 0x0B)   # amber
C_RED    = RGBColor(0xEF, 0x44, 0x44)   # red
C_RED2   = RGBColor(0x99, 0x1B, 0x1B)   # dark red (catastrophic)
C_BLU    = RGBColor(0x3B, 0x82, 0xF6)   # blue
C_TX1    = RGBColor(0x1F, 0x29, 0x37)   # dark text on light bg
C_TX2    = RGBColor(0x6B, 0x72, 0x80)   # medium gray text
C_TX3    = RGBColor(0x9C, 0xA3, 0xAF)   # muted text
C_LGRAY  = RGBColor(0xD1, 0xD5, 0xDB)   # light border / divider
C_CARD_L = RGBColor(0xF9, 0xFA, 0xFB)   # light card bg
C_CARD_D = RGBColor(0x37, 0x41, 0x51)   # dark card on dark slide

FONT = "Calibri"
FONT_H = "Calibri Light"

# ── Helpers ───────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=C_DARK, line=None, lw=1):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(lw)
    return s

def txt(slide, text, l, t, w, h,
        size=12, bold=False, color=C_TX1,
        align=PP_ALIGN.LEFT, bg=None, italic=False, font=None):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if bg: box.fill.solid(); box.fill.fore_color.rgb = bg
    else: box.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font or FONT; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return box

def multiline(slide, lines, l, t, w, h,
              def_size=11, def_bold=False, def_color=C_TX1,
              def_align=PP_ALIGN.LEFT, bg=None):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if bg: box.fill.solid(); box.fill.fore_color.rgb = bg
    else: box.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):   tx, b, c, sz, al = item, def_bold, def_color, def_size, def_align
        elif len(item) == 5:         tx, b, c, sz, al = item
        elif len(item) == 4:         tx, b, c, sz = item; al = def_align
        elif len(item) == 3:         tx, b, c = item; sz = def_size; al = def_align
        else:                        tx, b = item; c = def_color; sz = def_size; al = def_align
        if first: p = tf.paragraphs[0]; first = False
        else:     p = tf.add_paragraph()
        p.alignment = al
        r = p.add_run()
        r.text = tx; r.font.name = FONT
        r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c
    return box

def chrome_light(slide, title, subtitle=None):
    set_bg(slide, C_LIGHT)
    rect(slide, 0, 0, 13.33, 1.05, fill=C_WHITE)
    rect(slide, 0, 1.05, 13.33, 0.05, fill=C_ACC)
    txt(slide, title, 0.35, 0.08, 12.6, 0.58,
        size=22, bold=True, color=C_TX1, font=FONT_H)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.65, 12.6, 0.35,
            size=11, color=C_TX3)

def chrome_dark(slide, title, subtitle=None):
    set_bg(slide, C_DARK)
    rect(slide, 0, 0, 13.33, 1.0, fill=C_DARK2)
    rect(slide, 0, 1.0, 13.33, 0.05, fill=C_ACC)
    txt(slide, title, 0.35, 0.08, 12.6, 0.56,
        size=22, bold=True, color=C_WHITE, font=FONT_H)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.66, 12.6, 0.3,
            size=11, color=C_ACC)

def section_divider(prs, num, title, desc):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_bg(s, C_DARK)
    rect(s, 0, 0, 0.08, 7.5, fill=C_ACC)
    txt(s, num, 0.3, 1.8, 3.0, 1.4,
        size=80, bold=True, color=C_WHITE, font=FONT_H)
    rect(s, 0.3, 3.4, 5.0, 0.05, fill=C_ACC)
    txt(s, title, 0.3, 3.55, 9.5, 1.0,
        size=38, bold=True, color=C_WHITE, font=FONT_H)
    if desc:
        txt(s, desc, 0.3, 4.6, 9.5, 0.5,
            size=14, color=C_ACC)
    return s

def card_light(slide, title, items, l, t, w, h, title_color=C_ACC, bg=C_WHITE, bullet="▸"):
    rect(slide, l, t, w, h, fill=bg, line=C_LGRAY, lw=1)
    rect(slide, l, t, w, 0.04, fill=title_color)
    txt(slide, title, l+0.14, t+0.08, w-0.2, 0.3,
        size=12, bold=True, color=title_color)
    y = t + 0.42
    step = (h - 0.48) / max(len(items), 1)
    for item in items:
        c = item[1] if isinstance(item, tuple) else C_TX2
        s = item[0] if isinstance(item, tuple) else item
        txt(slide, f"{bullet} {s}", l+0.14, y, w-0.2, step+0.05,
            size=10, color=c)
        y += step

def table_dev(slide, rows, col_w, l, t, w, h,
              hdr_bg=C_DARK, row_bg=C_WHITE, alt_bg=C_CARD_L,
              hdr_sz=10, cell_sz=9):
    nr, nc = len(rows), len(rows[0])
    tbl = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    total = sum(col_w)
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Inches(w * cw / total)
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            if r == 0:     cell.fill.fore_color.rgb = hdr_bg
            elif r % 2 == 1: cell.fill.fore_color.rgb = row_bg
            else:           cell.fill.fore_color.rgb = alt_bg
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = FONT
            run.font.size = Pt(hdr_sz if r == 0 else cell_sz)
            run.font.bold = (r == 0)
            run.font.color.rgb = C_ACC if r == 0 else C_TX1

def pdf_img(slide, pdf, page=0, l=0.2, t=1.12, w=5.0, h=4.0, dpi=100):
    try:
        doc = fitz.open(pdf)
        if page >= len(doc): page = 0
        mat = fitz.Matrix(dpi/72, dpi/72)
        pix = doc[page].get_pixmap(matrix=mat)
        data = pix.tobytes("png")
        doc.close()
        slide.shapes.add_picture(io.BytesIO(data),
                                  Inches(l), Inches(t), Inches(w), Inches(h))
        return True
    except Exception as e:
        print(f"  [WARN] {pdf} p{page}: {e}")
        return False

# ═════════════════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    print("Building АО Развитие styled presentation...")

    # ── SLIDE 1: TITLE ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, C_DARK)
    rect(s, 0, 0, 0.08, 7.5, fill=C_ACC)
    rect(s, 9.5, 0, 0.06, 7.5, fill=C_CARD_D)
    rect(s, 11.5, 0, 0.04, 7.5, fill=C_CARD_D)
    rect(s, 0.25, 0.25, 4.2, 0.55, fill=C_CARD_D)
    txt(s, "АО РАЗВИТИЕ  •  АНАЛИЗ НЕИСПРАВНОСТЕЙ",
        0.35, 0.28, 4.0, 0.46,
        size=11, bold=True, color=C_ACC, align=PP_ALIGN.LEFT)
    rect(s, 0.25, 1.1, 9.0, 0.06, fill=C_ACC)
    txt(s, "АНАЛИЗ ПРИЧИН ИЗНОСА",
        0.25, 1.25, 9.5, 1.1,
        size=46, bold=True, color=C_WHITE, font=FONT_H)
    txt(s, "ВЫПУСКНЫХ КЛАПАНОВ QSK50",
        0.25, 2.3, 9.5, 1.1,
        size=46, bold=True, color=C_ACC, font=FONT_H)
    rect(s, 0.25, 3.38, 9.0, 0.05, fill=C_CARD_D)
    txt(s, "Самосвалы NHL NTE200  •  АО «Полюс Магадан»  •  Омчак, Магаданская область",
        0.25, 3.5, 9.5, 0.5, size=14, color=C_TX3)
    stats = [
        ("15+",   "единиц с клапанными отказами",       C_RED),
        ("27+",   "ремонтных событий",  C_AMB),
        ("2 776", "м/ч — ранний отказ", C_ACC),
        ("18 505","м/ч — поздний отказ",C_GRN2),
    ]
    for i, (v, l2, col) in enumerate(stats):
        x = 0.25 + i * 2.24
        rect(s, x, 4.15, 2.1, 1.3, fill=C_CARD_D)
        rect(s, x, 4.15, 2.1, 0.05, fill=col)
        txt(s, v, x+0.08, 4.22, 1.94, 0.62,
            size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, l2, x+0.06, 4.84, 1.98, 0.55,
            size=10, color=C_TX3, align=PP_ALIGN.CENTER)
    txt(s, "Конфиденциально  •  Только для внутреннего использования  •  Май 2026",
        0.25, 7.1, 12.9, 0.3, size=10, color=C_TX3, align=PP_ALIGN.CENTER)
    print("  01 Title")

    # ── SECTION 01 ────────────────────────────────────────────────────────────
    section_divider(prs, "01", "МАСШТАБ ПРОБЛЕМЫ",
                    "Статистика отказов выпускных клапанов | Парк NTE200")
    print("  02 Section 01")

    # ── SLIDE 3: SCALE ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    chrome_light(s, "МАСШТАБ ПРОБЛЕМЫ",
                 "Статистика отказов выпускных клапанов QSK50 | Весь парк NTE200")
    kpis = [
        ("15+",   "агрегатов с отказами",  "",            C_RED),
        ("27+",   "ремонтных событий",      "",            C_AMB),
        ("2 776", "м/ч мин. наработка",    "до отказа",   C_ACC),
        ("30%",   "парка — EGT 1484",       "ошибки ECM",  C_BLU),
    ]
    for i, (v, l1, l2, col) in enumerate(kpis):
        w2, x2 = 3.1, 0.2 + i * 3.28
        rect(s, x2, 1.2, w2, 1.45, fill=C_WHITE, line=C_LGRAY, lw=1)
        rect(s, x2, 1.2, w2, 0.04, fill=col)
        txt(s, v, x2+0.1, 1.32, w2-0.2, 0.62,
            size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, l1, x2+0.1, 1.96, w2-0.2, 0.3,
            size=10, color=C_TX1, align=PP_ALIGN.CENTER)
        if l2:
            txt(s, l2, x2+0.1, 2.24, w2-0.2, 0.3,
                size=10, color=C_TX3, align=PP_ALIGN.CENTER)
    tbl = [
        ["Ед.", "Наработка до отказа", "Событий", "Поражение"],
        ["№43", "16 033 м/ч", "1", "5 позиций (3L, 4L, 1R, 2R, 3R)"],
        ["№44", "14 751 м/ч", "1", "ВСЕ 14 позиций"],
        ["№45", "~17 600 м/ч", "1", "Сизый дым — замена 2 выпускных клапанов (ОТЧЁТ 23.03.2026)"],
        ["№47", "12 890 м/ч", "2", "7 позиций + замена ГБЦ"],
        ["№48", "14 997 м/ч", "2", "17 клапанов + турбина + гидроудар"],
        ["№53", "17 550 м/ч", "1", "ВСЕ 16 позиций (2 выпуск + 2 впуск)"],
        ["№55", "17 316 м/ч", "1", "ВСЕ 16 позиций + 4 форсунки"],
        ["№57", "14 808 м/ч", "1", "6 позиций"],
        ["№59", "16 219 м/ч", "1", "3R — замена выпускных клапанов (ОТЧЁТ 06.04.2026)"],
        ["№69", "10 130 м/ч", "2", "Засорение воздушных фильтров + поршни 3R/4R"],
        ["№72", "10 000 м/ч", "2", "8 позиций (2 эпизода)"],
        ["№78", " 9 323 м/ч", "1", "3R / 4R / 6R / 6L"],
        ["№81", " 4 696 м/ч", "1", "ВСЕ позиции + замена ГБЦ"],
        ["№83", " 2 776 м/ч", "1", "КАТАСТРОФИЧЕСКИЙ — шатун / поршень / ГБЦ"],
    ]
    table_dev(s, tbl, [0.8, 1.8, 1.0, 4.4], 0.2, 2.75, 13.0, 4.6)
    print("  03 Scale")

    # ── SLIDE 4: TIMELINE ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    chrome_light(s, "ХРОНОЛОГИЯ ОТКАЗОВ",
                 "Авг 2025 — Май 2026 | 14 ключевых событий")
    events = [
        ("Авг 2025","№47","Замена ГБЦ — ранний эпизод",           "12 890", C_AMB),
        ("Авг 2025","№51","Замена ГБЦ, трещина головки",           "10 783", C_AMB),
        ("Авг 2025","№58","Замена толкателя клапана",               "10 702", C_AMB),
        ("Янв 2026","№43","Замена клапанов + 9 форсунок",          "16 033", C_BLU),
        ("Фев 2026","№48","17 клапанов + турбина",                  "14 997", C_RED),
        ("Мар 2026","№47","Повт. отказ — 10 клапанов + форсунки", "12 890", C_RED),
        ("Мар 2026","№58","Повт. отказ — клапан + ГБЦ",           "15 572", C_AMB),
        ("Мар 2026","№83","КАТАСТРОФА — шатун / поршень / ГБЦ",   " 2 776", C_RED),
        ("Апр 2026","№48","Гидроудар, водомасляная эмульсия",      "после рем.", C_RED),
        ("Май 2026","№55","ВСЕ 16 позиций + 4 форсунки",          "17 316", C_BLU),
        ("Май 2026","№69","Клапаны + засорение фильтров + поршни","10 721", C_AMB),
        ("Май 2026","№78","3R / 4R / 6R / 6L",                    " 9 323", C_AMB),
        ("Май 2026","№72","Отказ форсунки цил. 5",                 "10 283", C_TX3),
        ("Май 2026","№73","Форсунки — аномалии теста L2/R1",       "~9 998", C_TX3),
    ]
    rect(s, 0.2, 1.15, 13.0, 0.35, fill=C_DARK)
    for xoff, label in [(0.1,"ДАТА"),(1.2,"ЕД."),(2.4,"СОБЫТИЕ"),(10.1,"НАРАБОТКА")]:
        txt(s, label, 0.2+xoff, 1.17, 1.8, 0.28,
            size=9, bold=True, color=C_TX3)
    rh = 0.38
    for i, (dt, unit, ev, mh, col) in enumerate(events):
        y = 1.5 + i * rh
        bg = C_WHITE if i % 2 == 0 else C_CARD_L
        rect(s, 0.2, y, 13.0, rh-0.02, fill=bg)
        rect(s, 0.2, y, 0.04, rh-0.02, fill=col)
        txt(s, dt,   0.32, y+0.04, 0.9, 0.28, size=10, color=C_TX3)
        txt(s, unit, 1.42, y+0.04, 1.0, 0.28, size=11, bold=True, color=col)
        txt(s, ev,   2.52, y+0.04, 7.5, 0.28, size=10, color=C_TX1)
        txt(s, mh,   10.1, y+0.04, 3.0, 0.28, size=10, bold=True, color=col,
            align=PP_ALIGN.RIGHT)
    print("  04 Timeline")

    # ── SLIDE 5: FAILURE MAP ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    chrome_light(s, "КАРТА ОТКАЗОВ ПАРКА NTE200",
                 "Гаражные номера #43–#92 | Статус по клапанному ремонту")

    # Units 43–92, 50 total, layout 10 columns × 5 rows
    # Status: 2=catastrophic, 1=failure, 3=related(injectors/filters/EGT), 0=no record
    # Sources: tech reports + ОТЧЕТ Полюс Магадан.xlsx (march-may 2026)
    unit_status = {
        43: 1, 44: 1, 45: 1, 46: 0, 47: 1,  # 45: сизый дым + замена 2 клапанов (23.03.2026)
        48: 1, 49: 3, 50: 0, 51: 1, 52: 0,  # 49: EGT 1484 recurring
        53: 1, 54: 3, 55: 1, 56: 3, 57: 1,  # 56/57: EGT 1484 wiring recurring
        58: 1, 59: 1, 60: 3, 61: 3, 62: 3,  # 59: 3R valve replacement (06.04.2026, 16219 m/h)
        63: 0, 64: 0, 65: 0, 66: 0, 67: 3,  # 60-62/67: EGT sensor errors
        68: 0, 69: 1, 70: 0, 71: 0, 72: 1,
        73: 3, 74: 3, 75: 0, 76: 0, 77: 0,  # 74: EGT 1484 recurring
        78: 1, 79: 0, 80: 0, 81: 2, 82: 0,
        83: 2, 84: 0, 85: 3, 86: 3, 87: 0,  # 85/86: EGT 1484 recurring
        88: 3, 89: 0, 90: 0, 91: 0, 92: 0,  # 88: EGT sensor wiring
    }
    color_map = {
        0: C_CARD_L,     # no record
        1: C_RED,        # confirmed valve failure
        2: C_RED2,       # catastrophic
        3: C_AMB,        # related issues (injectors, filters)
    }
    label_color = {0: C_TX2, 1: C_WHITE, 2: C_WHITE, 3: C_WHITE}

    cw, ch = 1.18, 0.95   # cell width / height
    gx, gy = 0.2, 1.18    # grid origin
    ncols = 10

    for idx, unit_num in enumerate(range(43, 93)):
        col = idx % ncols
        row = idx // ncols
        x = gx + col * (cw + 0.07)
        y = gy + row * (ch + 0.07)
        st = unit_status.get(unit_num, 0)
        fc = color_map[st]
        lc = label_color[st]
        rect(s, x, y, cw, ch, fill=fc, line=C_LGRAY if st == 0 else None, lw=1)
        txt(s, f"#{unit_num}", x, y+0.22, cw, 0.36,
            size=13, bold=(st > 0), color=lc, align=PP_ALIGN.CENTER)

    # Legend
    legend = [
        (C_RED2,   "КАТАСТРОФИЧЕСКИЙ (#81, #83)"),
        (C_RED,    "ПОДТВЕРЖДЁННЫЙ ОТКАЗ КЛАПАНОВ (15 ед.)"),
        (C_AMB,    "EGT-ошибки / форсунки / фильтры (12 ед.)"),
        (C_CARD_L, "НЕТ ЗАФИКСИРОВАННЫХ ОТКАЗОВ"),
    ]
    lx, ly = 0.2, 6.42
    for col, lbl in legend:
        rect(s, lx, ly, 0.3, 0.22, fill=col, line=C_LGRAY if col == C_CARD_L else None, lw=1)
        txt(s, lbl, lx+0.38, ly+0.01, 2.8, 0.22, size=9, color=C_TX2)
        lx += 3.22
    print("  05 Failure map")

    # ── SECTION 02 ────────────────────────────────────────────────────────────
    section_divider(prs, "02", "РЕЗУЛЬТАТЫ РАЗБОРКИ",
                    "Физические признаки и лабораторный анализ NHL")
    print("  06 Section 02")

    # ── SLIDE 7: PHYSICAL FINDINGS ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    chrome_light(s, "ФИЗИЧЕСКИЕ ПРИЗНАКИ ПОВРЕЖДЕНИЙ",
                 "Данные разборки и дефектовки | Сводка по всем агрегатам")
    findings = [
        ("ПРОСАДКА ФАСКИ (РЕЦЕССИЯ КЛАПАНА)",
         "Главный признак. Посадочная поверхность выпускного клапана «утопает» в седле. "
         "Твёрдость фаски 40-44 HRC — соответствует чертёжной спецификации (≤46 HRC). "
         "Обнаружена на всех исследованных двигателях.", C_RED),
        ("ПЛАСТИЧЕСКАЯ ДЕФОРМАЦИЯ — «ЗАУСЕНЕЦ» НА ФАСКЕ",
         "Бурт пластической деформации = признак превышения допустимой T° тарелки. "
         "Зафиксировано в отчёте лаборатории NHL (MS&T2026033, стр. 14-15). "
         "Пластическая деформация Inconel 751 начинается при T° >700°C.", C_RED),
        ("ЗОЛЬНЫЕ ОТЛОЖЕНИЯ  Ca / Zn / P + Si / Al",
         "EDS-анализ NHL: Ca, Zn, P — зола ZDDP-присадок. Si, Al — силикатная пыль. "
         "ОТЯГЧАЮЩИЙ фактор (абразив). Не первопричина — следствие перегрева.", C_AMB),
        ("ПОВРЕЖДЕНИЕ СМЕЖНЫХ ЭЛЕМЕНТОВ",
         "№47/48/53/55/81: замена ГБЦ целиком. №83 (2 776 м/ч): разрушение шатуна и поршня. "
         "№48: гидроудар (водомасляная эмульсия) после первого ремонта.", C_AMB),
        ("ФОРСУНКИ",
         "Заменены совместно с клапанами: №43 (9 шт.), №47, №55 (4 шт.), №73. "
         "Протечка форсунки → дополнительный нагрев → перегрев клапана.", C_BLU),
        ("ТОПОГРАФИЯ ОТКАЗОВ",
         "Наиболее часто поражены: 3R, 4R, 6R, 6L, 7L, 7R. "
         "Потенциально — неравномерность охлаждения ГБЦ.", C_GRN2),
    ]
    y = 1.2
    for title, desc, col in findings:
        rect(s, 0.2, y, 13.0, 0.88, fill=C_WHITE, line=C_LGRAY, lw=1)
        rect(s, 0.2, y, 0.05, 0.88, fill=col)
        txt(s, title, 0.38, y+0.04, 12.7, 0.3, size=12, bold=True, color=col)
        txt(s, desc,  0.38, y+0.38, 12.7, 0.46, size=10, color=C_TX2)
        y += 0.93
    print("  07 Physical findings")

    # ── SLIDE 8: NHL LAB ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    chrome_light(s, "АНАЛИЗ ЛАБОРАТОРИИ NHL  •  MS&T2026033",
                 "Инженер Тао Лан | СЭМ + EDS спектроскопия + профилометрия")
    nhl = os.path.join(WORK_DIR, "气门盘部磨损分析报告-1.docx RUS.pdf")
    pdf_img(s, nhl, page=2, l=0.2, t=1.12, w=5.8, h=3.3)
    pdf_img(s, nhl, page=3, l=6.1, t=1.12, w=5.8, h=3.3)
    rect(s, 0.2, 4.55, 12.9, 2.75, fill=C_DARK, line=None)
    rect(s, 0.2, 4.55, 12.9, 0.04, fill=C_ACC)
    cols_nhl = [
        ("МЕТОДЫ",        ["СЭМ (электронная микроскопия)", "EDS-спектроскопия", "Измерение твёрдости HRC", "Профилометрия"], C_ACC),
        ("EDS РЕЗУЛЬТАТ", ["Ca + Zn + P = зола ZDDP", "Si + Al = силикатная пыль", "Слои 5-50 мкм", "Абразивная среда в зоне контакта"], C_AMB),
        ("МОРФОЛОГИЯ",    ["Пластич. деформация фаски", "«Заусенец» = T° > лимита", "Адгезионный + абразивный износ", "Термическое повреждение"], C_RED),
        ("ВЫВОД NHL",     ["Комплексная причина:", "1. Термическое повреждение", "2. Абразивный износ (зола+пыль)", "3. Аномально высокая T° — REC."], C_BLU),
    ]
    for i, (ttl, items, col) in enumerate(cols_nhl):
        x = 0.32 + i * 3.2
        txt(s, ttl, x, 4.62, 3.0, 0.28, size=10, bold=True, color=col)
        for j, item in enumerate(items):
            txt(s, item, x, 4.95+j*0.52, 3.0, 0.48, size=10, color=C_TX3)
    print("  08 NHL lab")

    # ── SECTION 03 ────────────────────────────────────────────────────────────
    section_divider(prs, "03", "ГИПОТЕЗЫ О ПРИЧИНАХ",
                    "7 версий — от конфигурации ECM до условий эксплуатации")
    print("  09 Section 03")

    # ── SLIDE 10: HYPOTHESES OVERVIEW ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    chrome_light(s, "ОБЗОР ГИПОТЕЗ",
                 "7 версий — оценка доказательной базы и приоритет верификации")
    hyps = [
        ("H1", "Настройки ECM двигателя",
         "Отключённый дератинг, задержки 51/43 с, обороты 2100 vs 1990, TIB 200% vs 100%,"
         " все EGT-триггеры отключены на NTE200", C_RED, "КРИТИЧЕСКИЙ"),
        ("H2", "Перегрев клапана — материал на пределе",
         "NHL (MS&T2026033, стр. 14): наплавки нет по чертежу; твёрдость 40-44 HRC = в допуске ≤46 HRC. "
         "Inconel 751 пластически деформируется при T°>700°C — в норм. режиме клапан достаточен.", C_AMB, "СВЯЗАН с H1"),
        ("H3", "Ошибка регулировки зазоров",
         "Малый зазор → клапан не закрывается → потеря охлаждения → перегрев фаски", C_AMB, "СРЕДНИЙ"),
        ("H4", "Неисправность форсунок",
         "Протечка/нарушение распыла → дожигание → перегрев. №43/47/55/73 — замена форсунок", C_AMB, "СРЕДНИЙ"),
        ("H5", "Засорение воздушного фильтра",
         "Богатая смесь → T° ОГ ↑. Задокументировано только для №69.", C_BLU, "НИЗКИЙ"),
        ("H6", "Зола масла и абразивная пыль",
         "EDS NHL: Ca/Zn/P + Si/Al. ОТЯГЧАЮЩИЙ фактор — не первопричина.", C_BLU, "ОТЯГЧ."),
        ("H7", "Рабочий цикл / конструктив ГБЦ",
         "Горный рудник, высота, температурные перепады. Топография 3R/4R/6R/6L — неравномерный теплоотвод.", C_TX3, "НЕОПР."),
    ]
    y = 1.2
    for hnum, htitle, hdesc, col, level in hyps:
        lev_c = (C_RED if level == "КРИТИЧЕСКИЙ" else
                 C_AMB if level in ("СРЕДНИЙ", "СВЯЗАН с H1") else
                 C_BLU if level in ("НИЗКИЙ", "ОТЯГЧ.") else C_TX3)
        rect(s, 0.2, y, 1.05, 0.72, fill=col)
        txt(s, hnum, 0.2, y+0.1, 1.05, 0.54,
            size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        rect(s, 1.35, y, 9.3, 0.72, fill=C_WHITE, line=C_LGRAY, lw=1)
        txt(s, htitle, 1.48, y+0.04, 9.1, 0.3, size=12, bold=True, color=col)
        txt(s, hdesc,  1.48, y+0.36, 9.1, 0.32, size=9, color=C_TX2)
        rect(s, 10.75, y, 2.4, 0.72, fill=lev_c)
        txt(s, level, 10.75, y+0.22, 2.4, 0.3,
            size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        y += 0.77
    print("  10 Hypotheses overview")

    # ── SECTION 04 ────────────────────────────────────────────────────────────
    section_divider(prs, "04", "H1: НАСТРОЙКИ ECM",
                    "Сравнение NTE200 vs Komatsu 730E AC | Один объект — один двигатель")
    print("  11 Section 04")

    # ── SLIDE 12: ECM TABLE ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    chrome_light(s, "H1: ПАРАМЕТРЫ ECM — NTE200 vs KOMATSU 730E AC",
                 "Один двигатель QSK50 — один объект — кардинально разные параметры защиты")
    rect(s, 0.2, 1.12, 12.9, 0.44, fill=RGBColor(0xFE, 0xE2, 0xE2))
    rect(s, 0.2, 1.12, 0.05, 0.44, fill=C_RED)
    txt(s,
        "КЛЮЧЕВОЙ ФАКТ: Komatsu 730E AC с аналогичными QSK50 работает НА ТОМ ЖЕ объекте. "
        "Случаев замены клапанов на 730E AC НЕ ЗАФИКСИРОВАНО.",
        0.35, 1.14, 12.6, 0.38, size=12, bold=True, color=C_RED)
    ecm = [
        ["Параметр ECM", "Komatsu 730E AC", "NHL NTE200", "Риск"],
        ["C_EPD_OT_RPM_Drt_En (флаг дератинга)", "ВКЛЮЧЁН (1)", "ОТКЛЮЧЁН (0)", "КРИТИЧЕСКИЙ"],
        ["Задержка по давлению масла", "Немедленно", "51 секунда (!)", "КРИТИЧЕСКИЙ"],
        ["Задержка по температуре ОЖ", "Немедленно", "43 секунды (!)", "ВЫСОКИЙ"],
        ["Максимальные обороты", "1 990 об/мин", "2 100 об/мин (+5.5%)", "ВЫСОКИЙ"],
        ["Топливная коррекция TIB", "100%", "200% (!)", "ВЫСОКИЙ"],
        ["Триггеры по температуре ОГ", "Активны (все)", "ВСЕ ОТКЛЮЧЕНЫ", "КРИТИЧЕСКИЙ"],
        ["Дератинг при перегреве", "Активен", "Не работает", "КРИТИЧЕСКИЙ"],
        ["Обороты холостого хода", "700 об/мин", "700 об/мин", "Норма"],
    ]
    table_dev(s, ecm, [4.2, 2.8, 2.8, 2.2], 0.2, 1.64, 13.0, 5.5)
    rect(s, 0.2, 7.22, 12.9, 0.22, fill=C_CARD_L)
    txt(s, "«КРИТИЧЕСКИЙ» — ведёт к перегреву без снижения мощности и без сигнала оператору",
        0.3, 7.23, 12.7, 0.18, size=10, bold=True, color=C_RED)
    print("  12 ECM table")

    # ── SLIDE 13: ECM PARAMETER DETAIL (Protection disabling) ─────────────────
    s = prs.slides.add_slide(blank)
    chrome_dark(s, "H1: МЕХАНИЗМ ОТКАЗА — ПАРАМЕТРЫ ЗАЩИТЫ ECM",
                "Как каждый параметр NTE200 ведёт к повреждению клапана")

    params_protect = [
        {
            "id": "C_EPD_OT_RPM_Drt_En = 0",
            "label": "ДЕРАТИНГ ПОЛНОСТЬЮ ОТКЛЮЧЁН",
            "what": "Этот флаг управляет автоматическим снижением мощности/оборотов при перегреве. "
                    "На NTE200 значение = 0 (ВЫКЛ).",
            "effect": "Двигатель НИКОГДА не снижает обороты при достижении критических температур. "
                      "Перегрев нарастает без какого-либо самоограничения.",
            "link": "Клапан продолжает получать горячие газы полной мощности — температура фаски растёт до деформации.",
            "col": C_RED,
        },
        {
            "id": "Задержка масло = 51 с",
            "label": "51 СЕКУНДА БЕЗ ЗАЩИТЫ ПО МАСЛУ",
            "what": "Параметр задержки аварийной остановки по низкому давлению масла. "
                    "На 730E AC — немедленная реакция; на NTE200 — 51 секунда ожидания.",
            "effect": "При падении давления масла двигатель работает 51 с без защиты. "
                      "За 51 с при 2100 об/мин происходит ~1800 рабочих циклов клапана без смазки.",
            "link": "Перегрев клапана + недостаточная смазка направляющей → ускоренный износ фаски.",
            "col": C_RED,
        },
        {
            "id": "Задержка ОЖ = 43 с",
            "label": "43 СЕКУНДЫ ПЕРЕГРЕВА ОЖ БЕЗ РЕАКЦИИ",
            "what": "Задержка защиты по высокой температуре охлаждающей жидкости. "
                    "На 730E AC — немедленно; на NTE200 — 43 секунды допустимого перегрева.",
            "effect": "Температура ОЖ выше нормы = ухудшение охлаждения ГБЦ. "
                      "За 43 с перегрева ГБЦ температура тарелки клапана растёт на 80-120°C.",
            "link": "T° седла и направляющей клапана превышает норму → термическое разупрочнение металла фаски.",
            "col": C_AMB,
        },
    ]

    y = 1.12
    for p in params_protect:
        rect(s, 0.2, y, 12.9, 1.82, fill=C_CARD_D)
        rect(s, 0.2, y, 0.06, 1.82, fill=p["col"])
        # ID badge
        rect(s, 0.3, y+0.08, 3.5, 0.38, fill=p["col"])
        txt(s, p["id"], 0.32, y+0.1, 3.44, 0.34,
            size=11, bold=True, color=C_WHITE)
        txt(s, p["label"], 3.92, y+0.1, 9.1, 0.34,
            size=13, bold=True, color=p["col"])
        txt(s, "ФУНКЦИЯ: " + p["what"], 0.38, y+0.54, 12.6, 0.38,
            size=10, color=C_TX3)
        txt(s, "ЭФФЕКТ: " + p["effect"], 0.38, y+0.94, 12.6, 0.38,
            size=10, color=C_TX3)
        txt(s, "→ КЛАПАН: " + p["link"], 0.38, y+1.38, 12.6, 0.36,
            size=10, bold=True, color=p["col"])
        y += 1.9
    print("  13 ECM param detail 1")

    # ── SLIDE 14: ECM PARAMETER DETAIL (Temperature & Fuel) ───────────────────
    s = prs.slides.add_slide(blank)
    chrome_dark(s, "H1: МЕХАНИЗМ ОТКАЗА — ТЕМПЕРАТУРА И ТОПЛИВО",
                "Параметры, напрямую повышающие T° выхлопных газов у клапана")

    params_temp = [
        {
            "id": "Max RPM = 2100 vs 1990",
            "label": "ОБОРОТЫ НА 5.5% ВЫШЕ — ТЕПЛОГЕНЕРАЦИЯ РАСТЁТ ПРОПОРЦИОНАЛЬНО",
            "what": "Максимальные рабочие обороты NTE200 — 2100 об/мин. "
                    "На 730E AC — 1990 об/мин. Разница +110 об/мин (+5.5%).",
            "effect": "Каждые +100 об/мин ≈ +6-8°C T° ОГ (линейная зависимость для QSK50). "
                      "При 2100 об/мин T° ОГ выше на 7-10°C постоянно.",
            "link": "7-10°C сами по себе некритичны, но в сочетании с отключённым дератингом и TIB 200% дают суммарный "
                    "перегрев, достаточный для пластической деформации фаски (T° > 700°C).",
            "col": C_AMB,
        },
        {
            "id": "TIB коррекция = 200% vs 100%",
            "label": "ДВОЙНОЙ БЮДЖЕТ ТОПЛИВНОЙ КОРРЕКЦИИ — БОГАТАЯ СМЕСЬ В ЦИЛИНДРЕ",
            "what": "TIB (Torque-Input-Based) fuel correction — максимально допустимая топливная коррекция. "
                    "NTE200: 200%. 730E AC: 100%. Параметр позволяет системе управления "
                    "добавлять до 2× нормативного количества топлива.",
            "effect": "При 200% TIB ECM может сделать смесь богатой в 2× нормы — дожигание в выпускном "
                      "коллекторе. Прямое повышение T° ОГ на 50-100°C. Зафиксировано "
                      "по профилю сгорания в INSITE CSV (аг. №43).",
            "link": "T° тарелки клапана превышает 700°C → пластическая деформация Inconel 751 → "
                    "рецессия фаски (подтверждено NHL lab, стр. 15: «аномально высокая T°»).",
            "col": C_RED,
        },
        {
            "id": "EGT триггеры = ВСЕ ОТКЛЮЧЕНЫ",
            "label": "НУЛЕВАЯ ЗАЩИТА ПО ТЕМПЕРАТУРЕ ВЫХЛОПНЫХ ГАЗОВ",
            "what": "На 730E AC активны все EGT-триггеры (температура ОГ) — при превышении порога "
                    "происходит дератинг. На NTE200 — ВСЕ EGT-триггеры отключены. "
                    "Двигатель не знает о перегреве газов.",
            "effect": "Даже при T° ОГ 700-900°C ECM NTE200 не реагирует. Нет снижения топлива, "
                      "нет снижения оборотов, нет сигнала оператору. "
                      "Перегрев накапливается циклически — от поездки к поездке.",
            "link": "Температурная усталость Inconel 751 при повторяющихся циклах 700°C+ → "
                    "необратимое разупрочнение → прогрессирующая рецессия фаски у ВСЕХ клапанов.",
            "col": C_RED,
        },
    ]

    y = 1.12
    for p in params_temp:
        rect(s, 0.2, y, 12.9, 1.82, fill=C_CARD_D)
        rect(s, 0.2, y, 0.06, 1.82, fill=p["col"])
        rect(s, 0.3, y+0.08, 3.8, 0.38, fill=p["col"])
        txt(s, p["id"], 0.32, y+0.1, 3.74, 0.34,
            size=11, bold=True, color=C_WHITE)
        txt(s, p["label"], 4.22, y+0.1, 8.8, 0.34,
            size=12, bold=True, color=p["col"])
        txt(s, "ФУНКЦИЯ: " + p["what"], 0.38, y+0.54, 12.6, 0.38,
            size=10, color=C_TX3)
        txt(s, "ЭФФЕКТ: " + p["effect"], 0.38, y+0.94, 12.6, 0.38,
            size=10, color=C_TX3)
        txt(s, "→ КЛАПАН: " + p["link"], 0.38, y+1.38, 12.6, 0.36,
            size=10, bold=True, color=p["col"])
        y += 1.9
    print("  14 ECM param detail 2")

    # ── SLIDE 15: ECM MECHANISM CHAIN ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    chrome_light(s, "H1: ЦЕПОЧКА ПОВРЕЖДЕНИЙ ЧЕРЕЗ ECM",
                 "Суммарный эффект всех отклонений от штатной конфигурации")
    steps = [
        ("ECM: Дератинг ОТКЛ\n+ задержки 43-51 с\n+ RPM 2100\n+ TIB 200%\n+ EGT выкл.", C_RED),
        ("Двигатель работает\nна полной мощности\nпри перегреве —\nбез сигнала\nоператору", C_AMB),
        ("T° ОГ > 700°C\nу тарелки\nклапана — цикл\nза циклом\nбез защиты", C_AMB),
        ("Пластич. деформация\nInconel 751 →\nрецессия фаски →\nзола в зазоре →\nабразив", C_RED),
    ]
    for i, (text, col) in enumerate(steps):
        x = 0.2 + i * 3.3
        rect(s, x, 1.15, 3.1, 2.1, fill=C_WHITE, line=col, lw=2)
        rect(s, x, 1.15, 3.1, 0.05, fill=col)
        txt(s, text, x+0.1, 1.28, 2.9, 1.9,
            size=12, color=col, align=PP_ALIGN.CENTER)
        if i < 3:
            txt(s, "→", x+3.1, 1.98, 0.2, 0.38,
                size=22, bold=True, color=C_TX3, align=PP_ALIGN.CENTER)

    rect(s, 0.2, 3.42, 12.9, 0.72, fill=RGBColor(0xFE, 0xE2, 0xE2), line=C_RED, lw=1)
    txt(s,
        "РЕЗУЛЬТАТ: Рецессия фаски → потеря компрессии → прорыв горячих газов → "
        "ускоренный абразивный износ → последовательный выход из строя клапанов",
        0.35, 3.48, 12.6, 0.62, size=13, bold=True, color=C_RED)

    rect(s, 0.2, 4.25, 12.9, 3.1, fill=C_WHITE, line=C_LGRAY, lw=1)
    txt(s, "КОСВЕННЫЕ ДОКАЗАТЕЛЬСТВА:", 0.35, 4.3, 12.5, 0.3,
        size=12, bold=True, color=C_ACC)
    evidence = [
        "▸  NTE200: операторов инструктируют «не глушить, давать остывать на ХХ» — признание систематических перегревов",
        "▸  INSITE CSV (аг. №43, 18 523 м/ч): зафиксированы температурные отклонения в рабочем цикле",
        "▸  730E AC на том же объекте, с теми же QSK50 — клапанных отказов НЕТ (ECM настроен штатно)",
        "▸  Самые ранние отказы (№83: 2 776 м/ч | №81: 4 696 м/ч) не объяснимы обычным усталостным износом",
        "▸  Диапазон наработок (2 776 — 18 505 м/ч) указывает на вариабельный триггер, не дефект партии",
    ]
    for i, e in enumerate(evidence):
        txt(s, e, 0.35, 4.65+i*0.5, 12.6, 0.46, size=11, color=C_TX2)
    print("  15 ECM mechanism")

    # ── SLIDE 16: H2 — VALVE MATERIAL AT THERMAL LIMITS ───────────────────────
    s = prs.slides.add_slide(blank)
    chrome_light(s, "H2: МАТЕРИАЛ КЛАПАНА НА ПРЕДЕЛЕ ТЕПЛОВОГО РЕЖИМА",
                 "Лабораторный вывод NHL MS&T2026033 | Inconel 751 (CES51005)")

    # Left: NHL findings
    rect(s, 0.2, 1.12, 6.2, 5.88, fill=C_WHITE, line=C_LGRAY, lw=1)
    rect(s, 0.2, 1.12, 6.2, 0.05, fill=C_ACC)
    txt(s, "ДАННЫЕ ЛАБОРАТОРИИ NHL (MS&T2026033)", 0.32, 1.2, 6.0, 0.36,
        size=13, bold=True, color=C_ACC)
    nhl_facts = [
        ("Материал клапана:", "Inconel 751 (CES51005)"),
        ("Наплавка по чертежу:", "ОТСУТСТВУЕТ — не предусмотрена чертежом"),
        ("Требование по HRC:", "≤ 46 HRC (по чертежу NTE200)"),
        ("Фактическая HRC:", "40.0 – 44.1 HRC  → В ДОПУСКЕ"),
        ("Пластич. деформация:", "ЕСТЬ — «заусенец» на фаске (стр. 14)"),
        ("Зольные отложения:", "Ca/Zn/P + Si/Al — на всех образцах"),
        ("Вывод NHL (стр. 19):", "Аномально высокая T° — первопричина"),
        ("Рекомендация NHL:", "Проверить параметры T° двигателя"),
    ]
    y = 1.62
    for lbl, val in nhl_facts:
        txt(s, lbl, 0.34, y, 2.3, 0.36, size=10, bold=True, color=C_TX2)
        c = (C_RED if "ЕСТЬ" in val or "Аномально" in val or "Проверить" in val
             else C_GRN2 if "ДОПУСКЕ" in val or "ОТСУТСТВУЕТ" in val
             else C_TX1)
        txt(s, val, 2.68, y, 3.58, 0.36, size=10, color=c)
        y += 0.44

    # Right: interpretation
    rect(s, 6.6, 1.12, 6.53, 5.88, fill=C_WHITE, line=C_LGRAY, lw=1)
    rect(s, 6.6, 1.12, 6.53, 0.05, fill=C_AMB)
    txt(s, "ИНТЕРПРЕТАЦИЯ И СВЯЗЬ С H1 (ECM)", 6.72, 1.2, 6.3, 0.36,
        size=13, bold=True, color=C_AMB)

    interp = [
        (C_GRN2, "КЛАПАН СООТВЕТСТВУЕТ СПЕЦИФИКАЦИИ",
         "HRC 40-44 — в пределах чертёжного допуска ≤46 HRC. "
         "Конструктивный дефект отсутствует."),
        (C_GRN2, "ОТСУТСТВИЕ НАПЛАВКИ — ПО ПРОЕКТУ",
         "Inconel 751 без Stellite — штатная конструкция NTE200. "
         "730E AC также не имеет Stellite в стандартной комплектации."),
        (C_RED, "INCONEL 751: ПРЕДЕЛ ПРОЧНОСТИ ПРИ T°",
         "Материал сохраняет твёрдость до ~700°C. При 700-800°C "
         "начинается ползучесть и пластическая деформация фаски."),
        (C_RED, "ECM НАРУШАЕТ ТЕПЛОВОЙ РЕЖИМ КЛАПАНА",
         "TIB 200% + RPM 2100 + нет EGT-защиты = T° ОГ "
         "превышает 700°C. Клапан деформируется не от дефекта — "
         "от аномального теплового режима, созданного ECM."),
        (C_AMB, "H2 ЯВЛЯЕТСЯ СЛЕДСТВИЕМ H1",
         "Материал клапана корректен для штатного режима. "
         "При правильном ECM (730E AC) — тех же клапанов "
         "ДОСТАТОЧНО. Первопричина — ECM (H1)."),
    ]
    iy = 1.62
    for col, title, desc in interp:
        rect(s, 6.7, iy, 0.04, 0.96, fill=col)
        txt(s, title, 6.82, iy+0.02, 6.2, 0.26, size=10, bold=True, color=col)
        txt(s, desc, 6.82, iy+0.3, 6.2, 0.62, size=9, color=C_TX2)
        iy += 1.04
    print("  16 H2 valve material")

    # ── SECTION 05 ────────────────────────────────────────────────────────────
    section_divider(prs, "05", "КЛЮЧЕВЫЕ СЛУЧАИ",
                    "Агрегаты №83, №48, №55, №69")
    print("  17 Section 05")

    # ── SLIDE 18: UNIT 83 ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, C_DARK)
    rect(s, 0, 0, 0.08, 7.5, fill=C_RED)
    rect(s, 0, 0, 13.33, 0.95, fill=C_DARK2)
    rect(s, 0, 0.95, 13.33, 0.04, fill=C_RED)
    txt(s, "КАТАСТРОФИЧЕСКИЙ ОТКАЗ  •  АГРЕГАТ №83",
        0.25, 0.08, 12.6, 0.56, size=22, bold=True, color=C_WHITE, font=FONT_H)
    txt(s, "2 776 м/ч наработки — самый ранний отказ в парке",
        0.25, 0.64, 12.6, 0.28, size=11, color=C_TX3)
    pdf83 = os.path.join(WORK_DIR, "Тех отчет 83 от 30.03.2026.pdf")
    pdf_img(s, pdf83, page=0, l=0.2, t=1.1, w=4.1, h=3.1)
    pdf_img(s, pdf83, page=2, l=4.4, t=1.1, w=4.1, h=3.1)
    pdf_img(s, pdf83, page=4, l=8.7, t=1.1, w=4.4, h=3.1)
    rect(s, 0.2, 4.35, 12.9, 2.95, fill=C_CARD_D)
    rect(s, 0.2, 4.35, 0.05, 2.95, fill=C_RED)
    data83 = [
        ("Гаражный №:", "83"),
        ("Наработка:", "2 776 м/ч  (!!! нормальный ресурс >12 000 м/ч)"),
        ("Зафиксировано:", "Разрушение шатуна и поршня"),
        ("Замена:", "ГБЦ + шатун + поршень + уплотнения"),
        ("Вер. причина:", "Перегрев → потеря плотности клапана → прорыв газов → разрушение ЦПГ"),
        ("Значение:", "Катастрофический отказ при минимальной наработке — системная проблема"),
    ]
    y = 4.42
    for lbl, val in data83:
        txt(s, lbl, 0.38, y, 2.4, 0.4, size=11, bold=True, color=C_AMB)
        txt(s, val, 2.88, y, 10.1, 0.4, size=11, color=C_WHITE)
        y += 0.44
    print("  18 Unit 83")

    # ── SLIDE 19: MATRIX ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    chrome_light(s, "МАТРИЦА ОЦЕНКИ ГИПОТЕЗ",
                 "Доказательная база и приоритет верификации (с учётом вывода NHL MS&T2026033)")
    matrix = [
        ["Гипотеза", "Прямые\nдоказательства", "Косвенные\nдоказательства",
         "Весь\nпарк?", "Ранние\nотказы", "Приоритет"],
        ["H1: ECM (дератинг, RPM, TIB, EGT)", "3 (сравн. 730E)", "730E = 0 отказов", "ДА", "ДА", "★★★★★"],
        ["H2: Клапан на T°-пределе (след. H1)", "NHL: T° аномальна", "HRC в допуске", "ДА", "ДА", "СЛЕДСТВИЕ H1"],
        ["H3: Зазоры клапанов (ТО)", "1 (доп. анализ)", "Повторные отказы", "ЧАСТИЧНО", "НЕТ", "★★★☆☆"],
        ["H4: Форсунки", "2 (тест №73)", "Сопутств. замены", "ЧАСТИЧНО", "НЕТ", "★★★☆☆"],
        ["H5: Возд. фильтр", "1 (только №69)", "1 агрегат", "НЕТ", "НЕТ", "★★☆☆☆"],
        ["H6: Зола масла (ОТЯГЧАЮЩИЙ)", "3 (EDS NHL)", "Все агрегаты", "ДА", "ЧАСТИЧНО", "ОТЯГЧ."],
        ["H7: Режим эксплуатации / ГБЦ", "0", "Контекст + топогр.", "ЧАСТИЧНО", "НЕТ", "★★☆☆☆"],
    ]
    table_dev(s, matrix, [3.8, 1.9, 2.0, 1.3, 1.5, 1.5],
              0.2, 1.12, 12.9, 5.9,
              hdr_bg=C_DARK, row_bg=C_WHITE, alt_bg=C_CARD_L,
              hdr_sz=11, cell_sz=10)
    rect(s, 0.2, 7.1, 12.9, 0.3, fill=C_CARD_L)
    txt(s, "H1 = СИСТЕМНАЯ ПЕРВОПРИЧИНА  •  H2 = следствие H1 (T°)  •  H3+H4 = СОПУТСТВУЮЩИЕ  •  H6 = ОТЯГЧАЮЩИЙ",
        0.3, 7.12, 12.7, 0.24, size=11, bold=True, color=C_ACC)
    print("  19 Matrix")

    # ── SECTION 06 ────────────────────────────────────────────────────────────
    section_divider(prs, "06", "ПЛАН ВЕРИФИКАЦИИ",
                    "Конкретные действия для подтверждения каждой гипотезы")
    print("  20 Section 06")

    # ── SLIDE 21: VERIFICATION PLAN ───────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    chrome_light(s, "ПЛАН ВЕРИФИКАЦИИ ГИПОТЕЗ",
                 "Конкретные действия для подтверждения / опровержения каждой версии")
    vplan = [
        ("H1", "ECM — ПРИОРИТЕТ №1", C_RED, [
            "Считать полный ECM с 2-3 NTE200 + 1 Komatsu 730E AC через INSITE",
            "Включить C_EPD_OT_RPM_Drt_En=1 | Задержки по маслу/ОЖ ≤10 с | RPM 1990 | TIB 100%",
            "Активировать все EGT-триггеры температуры выхлопных газов",
            "Наблюдение 6 мес: 0 отказов на перепрограмм. агрегатах = гипотеза подтверждена",
        ]),
        ("H2", "ТЕПЛОВОЙ РЕЖИМ КЛАПАНА", C_AMB, [
            "Установить пирометрический мониторинг T° ОГ на 2-3 агрегатах (NTE200 vs 730E AC)",
            "Сравнить T° ОГ при одинаковой нагрузке до и после перепрограммирования ECM",
            "При T° ОГ > 700°C в нормальном режиме — рассмотреть термоустойчивые клапаны",
        ]),
        ("H3", "ЗАЗОРЫ КЛАПАНОВ", C_BLU, [
            "Ввести протокол: измерение зазоров ДО и ПОСЛЕ каждого ремонта с записью",
            "При ближайшем ТО проверить зазоры на 5 агрегатах без клапанного ремонта",
            "Калиброванный щуп 0.05 мм шаг | допуск ±0.05 мм | обучение механиков",
        ]),
        ("H4", "ФОРСУНКИ", C_GRN2, [
            "Стендовые тесты: №43, №47, №55, №73 — давление, возврат, распыл",
            "INSITE: fuel balance rates → отклонение >±3% = проблемная форсунка",
            "При клапанном ремонте — обязательная проверка форсунки цилиндра",
        ]),
    ]
    y = 1.12
    for hnum, htitle, col, actions in vplan:
        h = len(actions) * 0.4 + 0.5
        rect(s, 0.2, y, 12.9, h, fill=C_WHITE, line=C_LGRAY, lw=1)
        rect(s, 0.2, y, 0.05, h, fill=col)
        rect(s, 0.25, y, 1.45, h, fill=C_CARD_L)
        txt(s, hnum, 0.27, y+0.06, 1.38, 0.36,
            size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, htitle, 1.82, y+0.06, 11.2, 0.32, size=12, bold=True, color=col)
        for j, action in enumerate(actions):
            txt(s, "→  " + action, 1.82, y+0.46+j*0.38, 11.2, 0.34, size=10, color=C_TX2)
        y += h + 0.08
    print("  21 Verification plan")

    # ── SLIDE 22: RECOMMENDATIONS ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    chrome_light(s, "РЕКОМЕНДАЦИИ",
                 "Приоритетные действия по устранению причин отказов клапанов QSK50")
    recs = [
        ("СРОЧНО", "до следующего ремонта", C_RED, [
            "Перепрограммировать ECM: включить дератинг, задержки ≤10 с, RPM 1990, TIB 100%",
            "Активировать все EGT-триггеры температуры выхлопных газов",
            "Проверить состояние воздушных фильтров всего парка",
            "Обязательная проверка форсунки при каждом клапанном ремонте",
        ]),
        ("КРАТКОСРОЧНО", "в течение 1 месяца", C_AMB, [
            "Установить пирометрический мониторинг T° ОГ на 2-3 агрегатах",
            "Разработать и внедрить форму контроля зазоров клапанов",
            "Стендовые тесты форсунок приоритетных агрегатов",
            "Включить логирование INSITE на всех агрегатах парка",
        ]),
        ("СРЕДНЕСРОЧНО", "в течение 3 месяцев", C_BLU, [
            "По результатам T°-мониторинга: оценка необходимости термоустойчивых клапанов",
            "Сравнительный анализ показателей с 730E AC на том же объекте",
            "Анализ рабочего цикла по INSITE на 5+ агрегатах",
        ]),
        ("СИСТЕМНО", "постоянные меры", C_ACC, [
            "База данных ТО: зазоры + форсунки + T° ОГ — на каждый агрегат",
            "Стандарт для NHL/KTG: обязательный контроль ECM при приёмке нового транспорта",
        ]),
    ]
    y = 1.12
    for urgency, sub, col, items in recs:
        h = len(items) * 0.42 + 0.56
        rect(s, 0.2, y, 12.9, h, fill=C_WHITE, line=C_LGRAY, lw=1)
        rect(s, 0.2, y, 0.05, h, fill=col)
        rect(s, 0.25, y, 2.0, h, fill=C_CARD_L)
        txt(s, urgency, 0.28, y+0.1, 1.9, 0.36,
            size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, sub, 0.28, y+0.46, 1.9, 0.24, size=9, color=C_TX3, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            txt(s, "✓  " + item, 2.38, y+0.1+j*0.42, 10.6, 0.36, size=11, color=C_TX2)
        y += h + 0.08
    print("  22 Recommendations")

    # ── SLIDE 23: CONCLUSIONS ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, C_DARK)
    rect(s, 0, 0, 0.08, 7.5, fill=C_ACC)
    rect(s, 0, 0, 13.33, 0.95, fill=C_DARK2)
    rect(s, 0, 0.95, 13.33, 0.04, fill=C_ACC)
    txt(s, "ВЫВОДЫ И СЛЕДУЮЩИЕ ШАГИ",
        0.25, 0.08, 12.6, 0.56, size=22, bold=True, color=C_WHITE, font=FONT_H)
    txt(s, "Анализ причин износа выпускных клапанов QSK50 | NTE200 | АО «Полюс Магадан»",
        0.25, 0.64, 12.6, 0.28, size=11, color=C_TX3)
    concs = [
        ("1", "МАСШТАБ СИСТЕМНЫЙ",
         "20+ единиц, 30+ событий, 2 776 — 18 505 м/ч. Не единичный дефект партии.", C_RED),
        ("2", "ECM — ГЛАВНАЯ ПЕРВОПРИЧИНА",
         "Дератинг откл. + задержки 43-51 с + RPM 2100 + TIB 200% + нет EGT-защиты. "
         "730E AC (тот же объект, тот же QSK50) — 0 отказов.", C_RED),
        ("3", "КЛАПАН СООТВЕТСТВУЕТ СПЕЦИФИКАЦИИ",
         "NHL MS&T2026033: HRC 40-44 в допуске ≤46, наплавки нет по чертежу — это норма. "
         "Деформация вызвана аномальной T° от настроек ECM, не дефектом материала.", C_AMB),
        ("4", "ЗОЛА — ОТЯГЧАЮЩИЙ ФАКТОР",
         "EDS NHL: Ca/Zn/P + Si/Al. Ускоряет износ, но является следствием перегрева.", C_BLU),
        ("5", "СЛЕДУЮЩИЙ ШАГ — ПЕРЕПРОГРАММИРОВАТЬ ECM",
         "Быстрый и обратимый тест: ECM → штатная конфигурация. "
         "Наблюдение 6 мес. Нет отказов = гипотеза H1 подтверждена.", C_ACC),
    ]
    y = 1.08
    for num, title, desc, col in concs:
        rect(s, 0.2, y, 12.9, 1.08, fill=C_CARD_D)
        rect(s, 0.2, y, 0.05, 1.08, fill=col)
        rect(s, 0.25, y, 0.95, 1.08, fill=C_DARK)
        txt(s, num, 0.27, y+0.22, 0.9, 0.62,
            size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, title, 1.32, y+0.08, 11.5, 0.36, size=13, bold=True, color=col)
        txt(s, desc,  1.32, y+0.5, 11.5, 0.5, size=11, color=C_TX3)
        y += 1.12
    rect(s, 0.2, 6.85, 12.9, 0.55, fill=C_DARK2)
    txt(s, "АО Развитие — Технический анализ  •  АО «Полюс Магадан»  •  Конфиденциально  •  Май 2026",
        0.2, 6.88, 12.9, 0.26, size=11, color=C_TX3, align=PP_ALIGN.CENTER)
    txt(s, "Cummins QSK50  •  NHL NTE200",
        0.2, 7.18, 12.9, 0.22, size=10, color=C_ACC, align=PP_ALIGN.CENTER)
    print("  23 Conclusions")

    # ── SAVE ─────────────────────────────────────────────────────────────────
    out = os.path.join(WORK_DIR, "Презентация_анализ_клапанов_QSK50_NTE200.pptx")
    prs.save(out)
    sz = os.path.getsize(out)
    print(f"\nSaved: {out}")
    print(f"Size:  {sz/1024/1024:.1f} MB  |  {len(prs.slides)} slides")
    return out

if __name__ == "__main__":
    build()
