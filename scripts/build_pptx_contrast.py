#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Comprehensive PowerPoint presentation: QSK50 GBC failure analysis
Landscape A4 (33.87 x 23.99 cm), dark contrasting theme, Razvitie branding
"""

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.backends.backend_pdf import PdfPages
import matplotlib.ticker as mtick
import numpy as np
import pandas as pd
import re
import io
import warnings
warnings.filterwarnings('ignore')

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Cm
import pptx.util as pu

# ==========================================
# BRAND COLORS
# ==========================================
C_BG     = RGBColor(0x12, 0x1C, 0x24)   # near-black navy
C_PANEL  = RGBColor(0x1A, 0x2B, 0x38)   # dark blue panel
C_ACCENT = RGBColor(0x3E, 0xF0, 0xAF)   # Развитие bright teal
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT   = RGBColor(0xE8, 0xED, 0xF0)   # light text
C_RED    = RGBColor(0xFF, 0x45, 0x45)
C_YELLOW = RGBColor(0xFF, 0xD1, 0x00)
C_GRAY   = RGBColor(0x7A, 0x8C, 0x98)
C_ORANGE = RGBColor(0xFF, 0x88, 0x00)

# matplotlib hex equivalents
M_BG     = '#121C24'
M_PANEL  = '#1A2B38'
M_ACCENT = '#3EF0AF'
M_TEXT   = '#E8EDF0'
M_RED    = '#FF4545'
M_YELLOW = '#FFD100'
M_GRAY   = '#7A8C98'
M_ORANGE = '#FF8800'

# ==========================================
# SLIDE DIMENSIONS: A4 Landscape
# ==========================================
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def rgb_tuple(rgb: RGBColor):
    return (rgb[0]/255, rgb[1]/255, rgb[2]/255)

# ==========================================
# CHART HELPERS
# ==========================================
def chart_to_image(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

def make_fig(w=10, h=5.5):
    fig = plt.figure(figsize=(w, h))
    fig.patch.set_facecolor(M_BG)
    return fig

def styled_ax(ax, title='', xlabel='', ylabel=''):
    ax.set_facecolor(M_PANEL)
    ax.tick_params(colors=M_TEXT, labelsize=8)
    for s in ax.spines.values():
        s.set_color(M_GRAY)
        s.set_alpha(0.5)
    ax.grid(color=M_GRAY, alpha=0.2, linestyle='--')
    if title:
        ax.set_title(title, color=M_TEXT, fontsize=9.5, pad=6)
    if xlabel:
        ax.set_xlabel(xlabel, color=M_GRAY, fontsize=8)
    if ylabel:
        ax.set_ylabel(ylabel, color=M_GRAY, fontsize=8)

# ==========================================
# SLIDE HELPERS
# ==========================================
def blank_slide(prs):
    blank = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank)

def fill_bg(slide, color=C_BG):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=C_PANEL, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=14, bold=False, color=C_TEXT,
             align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_img(slide, buf, x, y, w, h):
    slide.shapes.add_picture(buf, x, y, w, h)

def header_bar(slide, title, subtitle='', page_num=None):
    # Top bar
    add_rect(slide, Cm(0), Cm(0), SLIDE_W, Cm(2.2), fill=C_PANEL)
    # Accent line
    add_rect(slide, Cm(0), Cm(2.1), SLIDE_W, Cm(0.12), fill=C_ACCENT)
    # Logo left
    add_text(slide, 'АО  РАЗВИТИЕ', Cm(0.4), Cm(0.15), Cm(5), Cm(0.9),
             size=11, bold=True, color=C_ACCENT)
    add_text(slide, 'Надёжность оборудования', Cm(0.4), Cm(0.9), Cm(5), Cm(0.7),
             size=7, color=C_GRAY)
    # Center title
    add_text(slide, title, Cm(5.5), Cm(0.12), Cm(22), Cm(1.1),
             size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, Cm(5.5), Cm(1.1), Cm(22), Cm(0.85),
                 size=8.5, color=C_ACCENT, align=PP_ALIGN.CENTER)
    # Page number
    if page_num:
        add_text(slide, str(page_num), Cm(32.5), Cm(0.5), Cm(1), Cm(0.8),
                 size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)
    # Footer
    add_rect(slide, Cm(0), Cm(18.55), SLIDE_W, Cm(0.5), fill=C_PANEL)
    add_text(slide, 'Конфиденциально | АО Развитие | Июнь 2026',
             Cm(0.5), Cm(18.58), Cm(15), Cm(0.4), size=6.5, color=C_GRAY)
    add_text(slide, 'QSK50 MCRS — NTE200 — Полюс Магадан',
             Cm(18), Cm(18.58), Cm(15.5), Cm(0.4), size=6.5, color=C_GRAY,
             align=PP_ALIGN.RIGHT)

def content_area():
    """Returns (x, y, w, h) for main content below header"""
    return Cm(0.5), Cm(2.35), Cm(32.87), Cm(16.0)

# ==========================================
# DATA PREPARATION
# ==========================================
def load_data():
    data = {}

    # --- OIL TOP-OFF ---
    df_maint = pd.read_excel('/home/user/NTE200/ОТЧЕТ_Полюс_Магадан.xlsx', sheet_name=0)
    oil = df_maint[df_maint['Задание и выполненная работа'].astype(str).str.contains('масл', case=False, na=False)].copy()
    def extr_l(t):
        m = re.search(r'(\d+)\s*(лит|л\b)', str(t), re.IGNORECASE)
        return int(m.group(1)) if m else 0
    oil['litres'] = oil['Задание и выполненная работа'].apply(extr_l)
    oil_by_unit = oil.groupby('Гаражный\n№').agg(total_l=('litres','sum'), events=('litres','count')).sort_values('total_l', ascending=False)
    data['oil_by_unit'] = oil_by_unit
    data['oil_total_l'] = oil['litres'].sum()
    data['oil_total_events'] = len(oil)
    data['oil_raw'] = oil

    # downtime
    df_maint['простой_ч'] = pd.to_timedelta(df_maint['Время простоя, час'], errors='coerce').dt.total_seconds() / 3600
    data['down_by_unit'] = df_maint.groupby('Гаражный\n№')['простой_ч'].sum().sort_values(ascending=False)
    data['down_total'] = df_maint['простой_ч'].sum()

    # work types
    work = df_maint['Задание и выполненная работа'].astype(str)
    wcats = {
        'Долив масла': work.str.contains('долив масла', case=False),
        'ТО-2 (500 м/ч)': work.str.contains('ТО-2', case=False),
        'ТО-3 (1000 м/ч)': work.str.contains('ТО-3', case=False),
        'Ошибки/Диагностика': work.str.contains('ошибк|диагностик', case=False),
        'Замена ГБЦ/клапана': work.str.contains('ГБЦ|клапан', case=False),
        'Течи': work.str.contains('течь|течи|протечк', case=False),
        'Форсунки': work.str.contains('форсунк', case=False),
    }
    wstats = {}
    for cat, mask in wcats.items():
        ev = df_maint[mask]
        down = pd.to_timedelta(ev['Время простоя, час'], errors='coerce').dt.total_seconds().sum() / 3600
        wstats[cat] = {'events': int(mask.sum()), 'downtime': round(down, 1)}
    data['work_stats'] = wstats

    # oil by date
    oil['date'] = pd.to_datetime(oil['Дата  '], errors='coerce')
    data['oil_by_date'] = oil.groupby('date')['litres'].sum()

    # --- DML ---
    with open('/home/user/NTE200/DML_82_s_porodoy.csv', encoding='cp1251') as f:
        lines = f.readlines()
    for i, line in enumerate(lines):
        if 'Дата' in line and 'Время' in line:
            break
    dml = pd.read_csv('/home/user/NTE200/DML_82_s_porodoy.csv', encoding='cp1251', skiprows=i, on_bad_lines='skip')
    cols = dml.columns.tolist()

    egt_cyl_cols = [c for c in cols if 'отработавших газов цилиндра' in c]
    egt_data = {}
    egt_series = {}
    for c in egt_cyl_cols:
        num = re.search(r'цилиндра (\d+)', c)
        if num:
            s = pd.to_numeric(dml[c], errors='coerce').dropna()
            s = s[s > 200]
            n = int(num.group(1))
            if n not in egt_data or s.max() > egt_data[n]['max']:
                egt_data[n] = {'max': float(s.max()), 'mean': float(s.mean()), 'p95': float(s.quantile(0.95))}
                egt_series[n] = s.values
    data['egt_data'] = egt_data
    data['egt_series'] = egt_series

    rpm = pd.to_numeric(dml['Частота вращения двигателя (об/мин)- Идентификатор0'], errors='coerce').dropna()
    load = pd.to_numeric(dml['Относительная нагрузка (Проценты)- Идентификатор0'], errors='coerce').dropna()
    data['dml_rpm'] = rpm.values
    data['dml_load'] = load.values

    # CR pressure
    cr_col = [c for c in cols if 'Измеренное давление в общем топливопровод' in c and 'Идентификатор0' in c]
    if cr_col:
        cr = pd.to_numeric(dml[cr_col[0]], errors='coerce').dropna()
        data['dml_cr'] = cr.values
    else:
        data['dml_cr'] = np.array([])

    data['dml_esn'] = '33238517'
    data['dml_unit'] = 'NTE200 №82'
    data['dml_rows'] = len(dml)

    # --- FLEET FAILURES ---
    gbc = pd.read_excel('/home/user/NTE200/ГБЦ ремонты.xlsx', sheet_name=0, header=0)
    gbc = gbc.dropna(subset=['№ самосвала'])
    gbc = gbc[gbc['М/Ч'].notna()]
    gbc['unit'] = gbc['№ самосвала'].astype(int)
    gbc['hours'] = gbc['М/Ч'].astype(float)
    data['gbc'] = gbc

    return data

print("Loading data...")
DATA = load_data()
print(f"Loaded: oil={DATA['oil_total_events']} events, DML={DATA['dml_rows']} rows, GBC={len(DATA['gbc'])} units")

# ==========================================
# SLIDE BUILDERS
# ==========================================

prs = new_prs()
page = [0]

def next_page():
    page[0] += 1
    return page[0]

# --------------------------------------------------
# SLIDE 1: TITLE
# --------------------------------------------------
def slide_title():
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    # Big dark panel left 60%
    add_rect(sl, Cm(0), Cm(0), Cm(20.5), SLIDE_H, fill=RGBColor(0x0A, 0x14, 0x1C))
    # Right panel with accent
    add_rect(sl, Cm(20.5), Cm(0), Cm(13.37), SLIDE_H, fill=RGBColor(0x0E, 0x26, 0x1F))
    # Accent bar
    add_rect(sl, Cm(20.5), Cm(0), Cm(0.25), SLIDE_H, fill=C_ACCENT)

    add_text(sl, 'АО  РАЗВИТИЕ', Cm(1), Cm(0.8), Cm(10), Cm(1),
             size=14, bold=True, color=C_ACCENT)
    add_text(sl, 'Отдел надёжности оборудования', Cm(1), Cm(1.7), Cm(12), Cm(0.7),
             size=9, color=C_GRAY)

    add_text(sl, 'АНАЛИЗ КОРЕННЫХ ПРИЧИН', Cm(1), Cm(3.5), Cm(19), Cm(1.2),
             size=20, bold=True, color=C_WHITE)
    add_text(sl, 'РАЗРУШЕНИЯ КЛАПАННОГО УЗЛА', Cm(1), Cm(4.9), Cm(19), Cm(1.2),
             size=20, bold=True, color=C_WHITE)
    add_text(sl, 'ДВИГАТЕЛЯ CUMMINS QSK50 MCRS', Cm(1), Cm(6.3), Cm(19), Cm(1.1),
             size=18, bold=True, color=C_ACCENT)

    add_text(sl, 'Самосвалы NTE200 | АО Полюс Магадан | 2023–2026 г.',
             Cm(1), Cm(7.8), Cm(19), Cm(0.8), size=11, color=C_TEXT)

    info_items = [
        ('Объект анализа:', '24 самосвала NTE200 (48% парка)'),
        ('Двигатель:', 'Cummins QSK50 MCRS, V16, 1491 кВт'),
        ('Источники данных:', 'DML, INSITE КАМСС, ОТЧЕТ ТО, ГБЦ журнал'),
        ('Методология:', 'FTA, контрольная группа, Weibull-анализ'),
        ('Документ:', 'ДСП — Для служебного пользования'),
    ]
    y = Cm(9.5)
    for lbl, val in info_items:
        add_text(sl, lbl, Cm(1), y, Cm(5), Cm(0.65), size=8.5, color=C_ACCENT, bold=True)
        add_text(sl, val, Cm(6), y, Cm(13.5), Cm(0.65), size=8.5, color=C_TEXT)
        y += Cm(0.85)

    # Right side stats
    stats = [
        ('24', 'самосвала\nс отказами'),
        ('48%', 'парка\nпоражено'),
        ('140+', 'ГБЦ\nзаменено'),
        ('3 615 л', 'доливок масла\nза март 2026'),
        ('4 583 ч', 'простой парка\nза март 2026'),
    ]
    y = Cm(1.5)
    for val, lbl in stats:
        add_text(sl, val, Cm(21.2), y, Cm(6.5), Cm(1.2),
                 size=24, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, Cm(21.2), y+Cm(1.3), Cm(6.5), Cm(0.85),
                 size=8.5, color=C_TEXT, align=PP_ALIGN.CENTER)
        y += Cm(3.0)

    add_text(sl, 'Июнь 2026', Cm(1), Cm(17.5), Cm(10), Cm(0.6), size=9, color=C_GRAY)
    print(f"Slide 1: Title")

slide_title()

# --------------------------------------------------
# SLIDE 2: SCOPE AND PROBLEM
# --------------------------------------------------
def slide_problem():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, 'МАСШТАБ ПРОБЛЕМЫ — ПАРК NTE200, ПОЛЮС МАГАДАН', page_num=n)

    # 4 KPI boxes top
    kpis = [
        ('24', 'самосвала NTE200\nс отказами ГБЦ', M_RED),
        ('48%', 'доля поражённых\nединиц парка', M_RED),
        ('2 776 м/ч', 'минимальный ресурс\n(NTE#83 — катастрофа)', M_YELLOW),
        ('0', 'отказов у 730E-DC\n(тот же двигатель!)', M_ACCENT),
    ]
    x = Cm(0.5)
    for val, lbl, col in kpis:
        fig, ax = plt.subplots(figsize=(3.2, 1.8))
        fig.patch.set_facecolor(M_PANEL)
        ax.set_facecolor(M_PANEL)
        ax.axis('off')
        ax.text(0.5, 0.72, val, transform=ax.transAxes, fontsize=28, fontweight='bold',
                color=col, ha='center', va='center')
        ax.text(0.5, 0.22, lbl, transform=ax.transAxes, fontsize=8, color=M_TEXT,
                ha='center', va='center', linespacing=1.4)
        buf = chart_to_image(fig, dpi=120)
        add_img(sl, buf, x, Cm(2.4), Cm(7.8), Cm(3.5))
        x += Cm(8.3)

    # Failure distribution chart
    fig, axes = plt.subplots(1, 2, figsize=(13, 5.5))
    fig.patch.set_facecolor(M_BG)

    # Left: Hours at failure bar chart
    ax1 = axes[0]
    gbc = DATA['gbc'].sort_values('hours')
    hours = gbc['hours'].values
    units = [f'NTE#{int(u)}' for u in gbc['unit'].values]
    colors = [M_RED if h < 5000 else M_YELLOW if h < 10000 else M_ACCENT for h in hours]
    bars = ax1.barh(range(len(units)), hours, color=colors, alpha=0.9, height=0.7)
    ax1.set_yticks(range(len(units)))
    ax1.set_yticklabels(units, fontsize=7, color=M_TEXT)
    ax1.axvline(x=np.mean(hours), color=M_YELLOW, linestyle='--', linewidth=1.5,
                label=f'Среднее: {np.mean(hours):.0f} м/ч')
    ax1.axvline(x=25000, color=M_ACCENT, linestyle=':', linewidth=1.5, alpha=0.7,
                label='Норм. ресурс: 25 000 м/ч')
    styled_ax(ax1, 'Наработка до первого отказа ГБЦ', 'м/ч', '')
    ax1.legend(fontsize=7, facecolor=M_PANEL, labelcolor=M_TEXT, edgecolor=M_GRAY)
    ax1.xaxis.set_major_formatter(mtick.FuncFormatter(lambda x,_: f'{int(x):,}'))

    # Right: Cumulative failure curve
    ax2 = axes[1]
    sorted_h = np.sort(hours)
    cum = np.arange(1, len(sorted_h)+1)
    ax2.step(sorted_h, cum, color=M_RED, linewidth=2.5, where='post')
    ax2.fill_between(sorted_h, cum, step='post', color=M_RED, alpha=0.15)

    # Weibull fit
    from scipy import stats as scipy_stats
    shape, loc, scale = scipy_stats.weibull_min.fit(sorted_h, floc=0)
    x_w = np.linspace(0, 22000, 200)
    p_w = scipy_stats.weibull_min.cdf(x_w, shape, loc, scale) * len(sorted_h)
    ax2.plot(x_w, p_w, color=M_YELLOW, linewidth=2, linestyle='--',
             label=f'Weibull β={shape:.1f}, η={scale:.0f}')
    ax2.axvline(x=25000, color=M_ACCENT, linestyle=':', alpha=0.7, label='Норм. ресурс')
    styled_ax(ax2, 'Накопленная кривая отказов (ступенчатая)', 'Наработка, м/ч', 'Кол-во единиц')
    ax2.legend(fontsize=7.5, facecolor=M_PANEL, labelcolor=M_TEXT, edgecolor=M_GRAY)
    ax2.xaxis.set_major_formatter(mtick.FuncFormatter(lambda x,_: f'{int(x):,}'))
    ax2.set_ylim(0, len(hours)+2)

    plt.tight_layout(pad=0.5)
    buf = chart_to_image(fig, dpi=130)
    add_img(sl, buf, Cm(0.5), Cm(6.2), Cm(32.87), Cm(11.5))

    print(f"Slide {n}: Problem scope")

slide_problem()

# --------------------------------------------------
# SLIDE 3: FLEET FAILURE TABLE
# --------------------------------------------------
def slide_failure_table():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, 'ХРОНОЛОГИЯ ОТКАЗОВ КУ ГБЦ — 24 ЕДИНИЦЫ NTE200', page_num=n)

    fig, ax = plt.subplots(figsize=(16.5, 9.5))
    fig.patch.set_facecolor(M_BG)
    ax.axis('off')
    ax.set_facecolor(M_BG)

    rows = [
        ('NTE#43', '22.07.2025', '11 943', '1 ГБЦ', 'Замена ГБЦ. Вып. клапаны', '—', 'Одиночный'),
        ('NTE#44', '27.12.2026', '14 751', '15 ГБЦ', '2 вып. клапана в каждой ГБЦ', '—', 'МАССОВЫЙ'),
        ('NTE#45', '24.03.2026', '17 863', '1 ГБЦ', '1 вып. клапан прогар', '—', 'Одиночный'),
        ('NTE#46', '17.12.2025', '15 103', '1 ГБЦ', '2 вып. клапана. Задир тарелки', '—', 'Одиночный'),
        ('NTE#47', '11.08.2026', '12 890', '7 ГБЦ', 'Вып. клапаны, L-банк', '—', 'Групповой'),
        ('NTE#48', '23.11.2025', '14 997', '3 ГБЦ', 'ГБЦ + поршень (4R). Прогар→поршень', '—', 'Осложнённый'),
        ('NTE#50', '12.01.2026', '15 433', '1 ГБЦ', 'Замена ГБЦ. Клапан вып.', '—', 'Одиночный'),
        ('NTE#51', '17.07.2025', '10 783', '1 ГБЦ', 'Замена ГБЦ', '—', 'Ранний'),
        ('NTE#52', '10.02.2026', '18 505', '1 ГБЦ', '1 ГБЦ — вып. клапан', '—', 'Поздний'),
        ('NTE#53', '25.05.2026', '17 550', 'ВСЕ 16', '2 вып.+2 впуск. кл. на кажд. ГБЦ', '—', 'ПОЛНЫЙ ОТКАЗ'),
        ('NTE#55', '19.05.2026', '17 316', 'ВСЕ 16', '2 вып.+2 впуск. кл. на кажд. ГБЦ', '—', 'ПОЛНЫЙ ОТКАЗ'),
        ('NTE#57', '05.02.2026', '14 808', '5 ГБЦ', '5 ГБЦ — вып. клапаны', '—', 'Групповой'),
        ('NTE#58', '10.08.2026', '10 702', '1 ГБЦ', 'Замена ГБЦ', '—', 'Ранний'),
        ('NTE#59', '18.02.2026', '15 143', '1+1', '1 ГБЦ → повтор 16 219 м/ч (+1 076)', 'ДА', 'ПОВТОР+1076'),
        ('NTE#62', '19.02.2026', '13 917', '1 ГБЦ', '1 ГБЦ — клапан', '—', 'Одиночный'),
        ('NTE#69', '16.04.2026', '10 130', '4+5', '4 ГБЦ+поршни → повтор +591 ч', 'ДА', 'ПОВТОР+591'),
        ('NTE#72', '18.04.2026', '10 000', '7+5', '7 ГБЦ → повтор 10 283 м/ч (+283!)', 'ДА', 'ПОВТОР+283'),
        ('NTE#73', '24.04.2026', '9 998', '1 ГБЦ', '1 ГБЦ — клапан вып.', '—', 'Ранний'),
        ('NTE#74', '02.05.2026', '9 869', '1 ГБЦ', '1 ГБЦ — клапан вып.', '—', 'Ранний'),
        ('NTE#76', '01.05.2026', '9 754', '7+1', '7 ГБЦ → повтор 10 028 м/ч (+274!)', 'ДА', 'ПОВТОР+274'),
        ('NTE#77', '26.05.2026', '10 016', '6 ГБЦ', '6 ГБЦ — вып. клапаны', '—', 'Групповой'),
        ('NTE#78', '19.05.2026', '9 323', '4 ГБЦ', '4 ГБЦ — клапаны', '—', 'Ранний'),
        ('NTE#81', '13.05.2026', '4 696', '9 ГБЦ', '9 ГБЦ — клапаны', '—', '!РАННИЙ <5000'),
        ('NTE#83', '21.03.2026', '2 776', '4+ ГБЦ', 'ГБЦ+шатун+поршень. Катастрофа', '—', '!КАТАСТРОФА'),
    ]
    headers = ['Самосвал', 'Дата', 'М/Ч', 'ГБЦ', 'Описание', 'Повтор', 'Категория']
    col_w = [0.09, 0.09, 0.07, 0.06, 0.40, 0.07, 0.14]
    col_x = [sum(col_w[:i]) for i in range(len(col_w))]

    row_h = 1.0 / (len(rows) + 1.5)

    # header
    for j, (h, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        ax.add_patch(mpatches.FancyBboxPatch((cx, 1-row_h), cw-0.005, row_h,
                     boxstyle='square,pad=0', facecolor='#3EF0AF', edgecolor='none'))
        ax.text(cx+cw/2, 1-row_h/2, h, color='#121C24', fontsize=8, ha='center', va='center', fontweight='bold')

    for i, row in enumerate(rows):
        y = 1 - (i+2)*row_h
        is_crit = row[-1].startswith('!')
        is_repeat = row[5] == 'ДА'
        is_full = 'ПОЛНЫЙ' in row[-1]
        bg = '#3a0808' if is_crit else ('#1a2a3a' if is_repeat else ('#162218' if is_full else ('#1A2B38' if i%2==0 else M_BG)))
        ax.add_patch(mpatches.FancyBboxPatch((0, y), 1, row_h,
                     boxstyle='square,pad=0', facecolor=bg, edgecolor='none'))
        for j, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
            val2 = val.lstrip('!')
            clr = M_RED if is_crit else (M_YELLOW if is_repeat else (M_ACCENT if is_full else M_TEXT))
            if j == 0: clr = M_ACCENT
            if j == 5 and val == 'ДА': clr = M_YELLOW
            ax.text(cx+cw/2 if j!=0 else cx+0.005, y+row_h/2, val2,
                    color=clr, fontsize=7 if j!=4 else 6.8,
                    ha='center' if j!=0 else 'left', va='center',
                    fontweight='bold' if j==0 or is_crit else 'normal')

    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    buf = chart_to_image(fig, dpi=140)
    add_img(sl, buf, Cm(0.3), Cm(2.35), Cm(33.2), Cm(16.3))
    print(f"Slide {n}: Failure table")

slide_failure_table()

# --------------------------------------------------
# SLIDE 4: OIL CONSUMPTION ANALYSIS
# --------------------------------------------------
def slide_oil_analysis():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, 'АНАЛИЗ РАСХОДА МАСЛА — МАРТ 2026',
               subtitle='235 событий | 3 615 литров | 187 доливок ДВС | АО Полюс Магадан', page_num=n)

    oil_by_unit = DATA['oil_by_unit']

    fig, axes = plt.subplots(1, 3, figsize=(16.5, 7))
    fig.patch.set_facecolor(M_BG)

    # Chart 1: Oil by unit (top 20)
    ax1 = axes[0]
    top20 = oil_by_unit.head(20)
    units_o = [f'NTE#{int(u)}' for u in top20.index]
    liters = top20['total_l'].values
    norm_total = 0.24 * 720  # 173L per period
    cols_o = [M_RED if l > 200 else M_YELLOW if l > 150 else M_ACCENT for l in liters]
    bars = ax1.barh(range(len(units_o)), liters, color=cols_o, alpha=0.9, height=0.7)
    ax1.axvline(x=norm_total, color=M_ACCENT, linestyle='--', linewidth=2,
                label=f'Норма: {norm_total:.0f}л')
    ax1.axvline(x=94, color=M_GRAY, linestyle=':', linewidth=1.5, label='Объём картера: 94л')
    ax1.set_yticks(range(len(units_o)))
    ax1.set_yticklabels(units_o, fontsize=7.5, color=M_TEXT)
    for bar, l in zip(bars, liters):
        rate = l/720
        ax1.text(bar.get_width()+2, bar.get_y()+bar.get_height()/2,
                 f'{rate:.2f}л/ч', va='center', color=M_TEXT, fontsize=6.5)
    styled_ax(ax1, 'Объём доливок масла / март 2026', 'Литров', '')
    ax1.legend(fontsize=7, facecolor=M_PANEL, labelcolor=M_TEXT, edgecolor=M_GRAY)

    # Chart 2: Oil events per day
    ax2 = axes[1]
    oil_date = DATA['oil_by_date'].dropna()
    if len(oil_date) > 0:
        dates = [d.day for d in oil_date.index if hasattr(d, 'day')]
        vals = oil_date.values[:len(dates)]
        ax2.bar(dates, vals, color=M_ACCENT, alpha=0.8)
        ax2.set_xlabel('День марта 2026', color=M_GRAY, fontsize=8)
        ax2.set_ylabel('Объём, л/день', color=M_GRAY, fontsize=8)
    styled_ax(ax2, 'Суточный объём доливок масла (март)', 'День', 'Литров/сут')

    # Chart 3: Oil consumption rate distribution
    ax3 = axes[2]
    # All top-off amounts
    oil_raw = DATA['oil_raw']
    oil_raw['litres_num'] = oil_raw['litres']
    amounts = oil_raw['litres_num'][oil_raw['litres_num'] > 0].values
    if len(amounts) > 0:
        ax3.hist(amounts, bins=15, color=M_ACCENT, alpha=0.8, edgecolor=M_BG)
        ax3.axvline(x=amounts.mean(), color=M_YELLOW, linestyle='--', linewidth=2,
                    label=f'Среднее: {amounts.mean():.0f} л')
        ax3.text(amounts.mean()+0.5, ax3.get_ylim()[1]*0.8 if ax3.get_ylim()[1] > 0 else 10,
                 f'{amounts.mean():.0f}л', color=M_YELLOW, fontsize=9, fontweight='bold')
    styled_ax(ax3, 'Распределение объёма разовой доливки', 'Литров', 'Кол-во событий')
    ax3.legend(fontsize=7.5, facecolor=M_PANEL, labelcolor=M_TEXT)

    plt.tight_layout(pad=0.8)
    buf = chart_to_image(fig, dpi=130)
    add_img(sl, buf, Cm(0.5), Cm(2.4), Cm(32.87), Cm(11.5))

    # Summary table
    fig2, ax = plt.subplots(figsize=(16.5, 2.5))
    fig2.patch.set_facecolor(M_BG)
    ax.axis('off')
    table_data = [
        ['Показатель', 'Факт (NTE200)', 'Норма QSK50', 'Отклонение'],
        ['Всего событий доливки / месяц', '235', 'N/A', '—'],
        ['Суммарный объём / месяц', '3 615 л (parк)', '≈4 320л (расч. норма)', 'В норме по парку'],
        ['Макс. доливка на 1 ед. (NTE#50)', '246 л / мес.', '≤173 л/мес (0,24×720)', '+42% СВЕРХНОРМЫ'],
        ['Расчётный расход NTE#50', '0,34 л/ч', '≤0,24 л/ч', '+42%'],
        ['Частота событий (топ-5 единиц)', '10–14 доливок/мес', 'Не регламентир.', 'Каждые 2 дня'],
        ['Признак прогара клапана', 'Расход масла ↑ = масло в КС', 'Симптом, не причина', 'ВТОРИЧНЫЙ ПРИЗНАК'],
    ]
    for i, row in enumerate(table_data):
        for j, val in enumerate(row):
            col_starts = [0, 0.28, 0.49, 0.68]
            col_widths = [0.28, 0.21, 0.19, 0.32]
            bg = M_ACCENT if i == 0 else (M_PANEL if i%2==0 else M_BG)
            text_col = M_BG if i == 0 else (M_RED if 'СВЕРХ' in val or '+42%' == val else M_TEXT)
            ax.add_patch(mpatches.FancyBboxPatch((col_starts[j], (len(table_data)-i-1)*0.13),
                         col_widths[j]-0.005, 0.12, boxstyle='square,pad=0', facecolor=bg, edgecolor='none'))
            ax.text(col_starts[j]+col_widths[j]/2, (len(table_data)-i-1)*0.13+0.06, val,
                    color=text_col, fontsize=7.5, ha='center', va='center',
                    fontweight='bold' if i==0 else 'normal')
    ax.set_xlim(0,1); ax.set_ylim(0, len(table_data)*0.13)
    buf2 = chart_to_image(fig2, dpi=120)
    add_img(sl, buf2, Cm(0.5), Cm(14.0), Cm(32.87), Cm(4.5))
    print(f"Slide {n}: Oil analysis")

slide_oil_analysis()

# --------------------------------------------------
# SLIDE 5: OIL TOP-OFF vs FAILURES CORRELATION
# --------------------------------------------------
def slide_oil_gbc_corr():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, 'КОРРЕЛЯЦИЯ: ДОЛИВКИ МАСЛА vs ОТКАЗЫ ГБЦ',
               subtitle='Чем выше расход масла → тем вероятнее прогар клапана', page_num=n)

    oil_by_unit = DATA['oil_by_unit']
    gbc = DATA['gbc']

    # Merge oil and failures
    oil_d = {int(k): v for k, v in oil_by_unit['total_l'].items()}
    gbc_d = {int(row['unit']): row['hours'] for _, row in gbc.iterrows()}

    common = sorted(set(oil_d.keys()) & set(gbc_d.keys()))
    oil_vals = [oil_d[u] for u in common]
    gbc_hours = [gbc_d[u] for u in common]
    unit_labels = [f'NTE#{u}' for u in common]

    fig, axes = plt.subplots(1, 2, figsize=(15, 7.5))
    fig.patch.set_facecolor(M_BG)

    # Scatter: oil consumption vs hours to failure
    ax1 = axes[0]
    sc = ax1.scatter(oil_vals, gbc_hours, c=oil_vals, cmap='RdYlGn_r', s=120, alpha=0.9,
                     edgecolors=M_BG, linewidths=1.5, vmin=50, vmax=250)
    for u, o, h in zip(unit_labels, oil_vals, gbc_hours):
        ax1.annotate(u, (o, h), textcoords='offset points', xytext=(4, 3),
                     fontsize=6.5, color=M_TEXT)
    from scipy import stats as scipy_stats
    if len(oil_vals) > 3:
        slope, intercept, r, p, se = scipy_stats.linregress(oil_vals, gbc_hours)
        x_fit = np.linspace(min(oil_vals)-10, max(oil_vals)+10, 50)
        ax1.plot(x_fit, slope*x_fit+intercept, color=M_RED, linestyle='--', linewidth=2,
                 label=f'Тренд: r={r:.2f}')
        ax1.legend(fontsize=8, facecolor=M_PANEL, labelcolor=M_TEXT)
    plt.colorbar(sc, ax=ax1).set_label('Доливок масла (л)', color=M_TEXT, fontsize=7)
    styled_ax(ax1, 'Объём доливок масла vs наработка до отказа', 'Доливки масла, л/мес', 'Наработка до отказа, м/ч')
    ax1.yaxis.set_major_formatter(mtick.FuncFormatter(lambda x,_: f'{int(x):,}'))

    # Side-by-side bar: oil and downtime for problem units
    ax2 = axes[1]
    down_by_unit = DATA['down_by_unit']
    # Top failure units
    fail_units = [43, 44, 47, 48, 50, 51, 52, 53, 55, 57, 59, 62, 69, 72, 73, 74, 76, 77, 78, 81, 83]
    fu_oil = [oil_d.get(u, 0) for u in fail_units[:15]]
    fu_down = [float(down_by_unit.get(u, 0)) for u in fail_units[:15]]
    fu_labels = [f'NTE#{u}' for u in fail_units[:15]]

    x = np.arange(len(fu_labels))
    w = 0.4
    ax2_r = ax2.twinx()
    ax2.bar(x - w/2, fu_oil, w, color=M_ACCENT, alpha=0.8, label='Доливки масла, л')
    ax2_r.bar(x + w/2, fu_down, w, color=M_RED, alpha=0.8, label='Простой, ч')
    ax2.set_xticks(x)
    ax2.set_xticklabels(fu_labels, rotation=45, ha='right', fontsize=7, color=M_TEXT)
    ax2.set_ylabel('Доливки масла, л', color=M_ACCENT, fontsize=8)
    ax2_r.set_ylabel('Простой, ч', color=M_RED, fontsize=8)
    ax2.tick_params(colors=M_TEXT); ax2_r.tick_params(colors=M_RED)
    ax2.set_facecolor(M_PANEL)
    for s in ax2.spines.values(): s.set_color(M_GRAY); s.set_alpha(0.4)
    for s in ax2_r.spines.values(): s.set_color(M_GRAY); s.set_alpha(0.4)
    ax2.set_title('Доливки масла и простои — единицы с отказами ГБЦ', color=M_TEXT, fontsize=9.5)
    lines1, labels1 = ax2.get_legend_handles_labels()
    lines2, labels2 = ax2_r.get_legend_handles_labels()
    ax2.legend(lines1+lines2, labels1+labels2, fontsize=7.5, facecolor=M_PANEL, labelcolor=M_TEXT)
    ax2.axhline(y=173, color=M_YELLOW, linestyle=':', linewidth=1.5, label='Норма 173л')
    ax2.grid(color=M_GRAY, alpha=0.2, axis='y')

    plt.tight_layout(pad=0.8)
    buf = chart_to_image(fig, dpi=130)
    add_img(sl, buf, Cm(0.5), Cm(2.4), Cm(32.87), Cm(16.2))
    print(f"Slide {n}: Oil-GBC correlation")

slide_oil_gbc_corr()

# --------------------------------------------------
# SLIDE 6: DOWNTIME ANALYSIS
# --------------------------------------------------
def slide_downtime():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, 'АНАЛИЗ ПРОСТОЕВ ПАРКА NTE200 — МАРТ 2026',
               subtitle=f'Суммарный простой: {DATA["down_total"]:.0f} ч | 1 004 события ТО', page_num=n)

    fig, axes = plt.subplots(1, 3, figsize=(16.5, 7.5))
    fig.patch.set_facecolor(M_BG)

    # Chart 1: downtime by unit
    ax1 = axes[0]
    top_down = DATA['down_by_unit'].head(20)
    units_d = [f'NTE#{int(u)}' for u in top_down.index]
    down_h = top_down.values
    cols_d = [M_RED if h > 300 else M_YELLOW if h > 150 else M_ACCENT for h in down_h]
    ax1.barh(range(len(units_d)), down_h, color=cols_d, alpha=0.9, height=0.7)
    ax1.set_yticks(range(len(units_d)))
    ax1.set_yticklabels(units_d, fontsize=7, color=M_TEXT)
    for i, h in enumerate(down_h):
        ax1.text(h+2, i, f'{h:.0f}ч', va='center', color=M_TEXT, fontsize=6.5)
    styled_ax(ax1, 'Простой по единицам (топ-20)', 'Часов простоя', '')
    ax1.axvline(x=np.mean(down_h), color=M_YELLOW, linestyle='--', linewidth=1.5,
                label=f'Среднее: {np.mean(down_h):.0f}ч')
    ax1.legend(fontsize=7, facecolor=M_PANEL, labelcolor=M_TEXT, edgecolor=M_GRAY)

    # Chart 2: downtime by work type (pie)
    ax2 = axes[1]
    wstats = DATA['work_stats']
    cats = list(wstats.keys())
    down_vals = [wstats[c]['downtime'] for c in cats]
    colors_pie = [M_ACCENT, M_GRAY, M_GRAY, M_RED, M_YELLOW, M_ORANGE, M_PANEL]
    wedges, texts, autotexts = ax2.pie(down_vals, labels=cats, colors=colors_pie,
                                        autopct='%1.0f%%', startangle=90,
                                        textprops={'color': M_TEXT, 'fontsize': 7.5})
    for at in autotexts:
        at.set_fontsize(8); at.set_color(M_BG); at.set_fontweight('bold')
    ax2.set_facecolor(M_BG)
    ax2.set_title('Простои по категориям ТО', color=M_TEXT, fontsize=9.5, pad=4)

    # Chart 3: downtime events per unit type
    ax3 = axes[2]
    ev_counts = [wstats[c]['events'] for c in cats]
    x = np.arange(len(cats))
    bars = ax3.bar(x, ev_counts, color=colors_pie, alpha=0.9)
    ax3.set_xticks(x)
    ax3.set_xticklabels(cats, rotation=45, ha='right', fontsize=7, color=M_TEXT)
    for bar, val in zip(bars, ev_counts):
        ax3.text(bar.get_x()+bar.get_width()/2, bar.get_height()+2, str(val),
                 ha='center', va='bottom', color=M_TEXT, fontsize=8, fontweight='bold')
    styled_ax(ax3, 'Количество событий по категориям', '', 'Кол-во событий')

    plt.tight_layout(pad=0.8)
    buf = chart_to_image(fig, dpi=130)
    add_img(sl, buf, Cm(0.5), Cm(2.4), Cm(32.87), Cm(11.0))

    # KPIs
    kpis2 = [
        (f'{DATA["down_total"]:.0f} ч', 'Суммарный простой\nпарка / март 2026', M_RED),
        ('4,3%', 'КТГ* потери\n(730 ч × 50 ед. = 36 500 ч фонд)', M_RED),
        ('2 110 ч', 'Простой из-за ошибок\nи диагностики ДВС', M_YELLOW),
        (f'{DATA["down_by_unit"].iloc[0]:.0f} ч', f'Макс. простой\n1 единицы (NTE#{int(DATA["down_by_unit"].index[0])})', M_RED),
    ]
    fig2, axes2 = plt.subplots(1, 4, figsize=(16.5, 2.2))
    fig2.patch.set_facecolor(M_BG)
    for ax, (val, lbl, col) in zip(axes2, kpis2):
        ax.set_facecolor(M_PANEL)
        ax.axis('off')
        ax.text(0.5, 0.68, val, transform=ax.transAxes, fontsize=22, fontweight='bold',
                color=col, ha='center', va='center')
        ax.text(0.5, 0.2, lbl, transform=ax.transAxes, fontsize=7.5, color=M_TEXT,
                ha='center', va='center', linespacing=1.3)
    plt.tight_layout(pad=0.3)
    buf2 = chart_to_image(fig2, dpi=120)
    add_img(sl, buf2, Cm(0.5), Cm(13.6), Cm(32.87), Cm(4.8))

    add_text(sl, '* КТГ — коэффициент технической готовности. При фонде 36 500 маш-часов/месяц потери 4,3% = прямые убытки добычи.',
             Cm(0.5), Cm(18.1), Cm(32), Cm(0.5), size=7, color=C_GRAY)
    print(f"Slide {n}: Downtime analysis")

slide_downtime()

# --------------------------------------------------
# SLIDE 7: DML DATA — EGT BY CYLINDER
# --------------------------------------------------
def slide_dml_egt():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, f'DML ДАННЫЕ — ТЕМПЕРАТУРА ОГ ПО ЦИЛИНДРАМ',
               subtitle=f'NTE200 №82 | ESN {DATA["dml_esn"]} | {DATA["dml_rows"]} точек | 02.06.2026 (23 мин записи)',
               page_num=n)

    egt = DATA['egt_data']
    cyls = sorted(egt.keys())
    maxT = [egt[c]['max'] for c in cyls]
    meanT = [egt[c]['mean'] for c in cyls]
    p95T = [egt[c]['p95'] for c in cyls]
    cyl_labels = [f'Цил.{c}' for c in cyls]

    # Map to banks: QSK50 V16 layout
    # L-bank: odd cylinders 1,3,5,7,9,11,13,15 → actually depends on Cummins numbering
    # From DML: 16 cylinders total
    # Let's use the same numbering as INSITE

    fig, axes = plt.subplots(2, 2, figsize=(16.5, 8.5))
    fig.patch.set_facecolor(M_BG)

    # Chart 1: EGT bar by cylinder
    ax1 = axes[0,0]
    x = np.arange(len(cyls))
    bar_colors = [M_RED if t > 550 else M_YELLOW if t > 520 else M_ACCENT for t in maxT]
    ax1.bar(x, maxT, color=bar_colors, alpha=0.9, label='Пик EGT')
    ax1.plot(x, meanT, color=M_WHITE if False else '#FFFFFF', marker='o', markersize=4,
             linewidth=1.5, label='Среднее EGT', alpha=0.7)
    ax1.axhline(y=520, color=M_ACCENT, linestyle='--', linewidth=2, label='Норма QSK50: 520°C')
    ax1.axhline(y=650, color=M_RED, linestyle=':', linewidth=1.5, alpha=0.7, label='Зона риска: 650°C')
    ax1.set_xticks(x)
    ax1.set_xticklabels(cyl_labels, rotation=45, ha='right', fontsize=7.5, color=M_TEXT)
    styled_ax(ax1, f'Пиковая EGT по цилиндрам (NTE200 №82)', 'Цилиндр', '°C')
    ax1.legend(fontsize=7.5, facecolor=M_PANEL, labelcolor=M_TEXT, edgecolor=M_GRAY)
    ax1.set_ylim(400, 620)

    # Add value labels on top of bars
    for i, t in enumerate(maxT):
        ax1.text(i, t+2, f'{t:.0f}', ha='center', va='bottom', color=M_TEXT, fontsize=6.5, fontweight='bold')

    # Chart 2: Heatmap-style (16 cylinders on engine layout)
    ax2 = axes[0,1]
    # L-bank vs R-bank (approximate: 1-8 L, 9-16 R for this numbering)
    l_bank = maxT[:8]
    r_bank = maxT[8:]
    im = ax2.imshow([l_bank, r_bank],
                    cmap='RdYlGn_r', aspect='auto',
                    vmin=460, vmax=580)
    ax2.set_yticks([0, 1])
    ax2.set_yticklabels(['Банк L (цил.1-8)', 'Банк R (цил.9-16)'], color=M_TEXT, fontsize=9)
    ax2.set_xticks(range(8))
    ax2.set_xticklabels([f'Цил.{i+1}' for i in range(8)], color=M_TEXT, fontsize=8)
    for i in range(2):
        for j in range(8):
            val = [l_bank, r_bank][i][j]
            ax2.text(j, i, f'{val:.0f}', ha='center', va='center',
                     color='white' if val > 530 else 'black', fontsize=9, fontweight='bold')
    plt.colorbar(im, ax=ax2).ax.tick_params(labelcolor=M_TEXT)
    ax2.set_title('Тепловая карта EGT по банкам (°C)', color=M_TEXT, fontsize=9.5, pad=6)

    # Chart 3: EGT spread
    ax3 = axes[1,0]
    spread = max(maxT) - min(maxT)
    ax3.bar(cyl_labels, [t - min(maxT) for t in maxT],
            bottom=min(maxT), color=bar_colors, alpha=0.8)
    ax3.axhline(y=min(maxT)+30, color=M_ACCENT, linestyle='--', linewidth=2,
                label=f'Норма разброса: ≤30°C (730E)')
    ax3.axhline(y=min(maxT), color=M_GRAY, linewidth=1, alpha=0.5)
    for tick in ax3.get_xticklabels():
        tick.set_rotation(45); tick.set_ha('right'); tick.set_fontsize(7); tick.set_color(M_TEXT)
    styled_ax(ax3, f'Межцилиндровый разброс EGT: {spread:.0f}°C (норма ≤30°C)', 'Цилиндр', '°C')
    ax3.legend(fontsize=7.5, facecolor=M_PANEL, labelcolor=M_TEXT)
    ax3.set_ylim(450, 590)

    # Chart 4: EGT timeseries for hottest cylinder
    ax4 = axes[1,1]
    hottest_cyl = cyls[np.argmax(maxT)]
    hot_series = DATA['egt_series'][hottest_cyl]
    t_axis = np.arange(len(hot_series))
    ax4.plot(t_axis, hot_series, color=M_RED, linewidth=1.5, alpha=0.8)
    ax4.fill_between(t_axis, hot_series, alpha=0.2, color=M_RED)
    ax4.axhline(y=520, color=M_ACCENT, linestyle='--', linewidth=2, label='Норма: 520°C')
    ax4.axhline(y=np.mean(hot_series), color=M_YELLOW, linestyle=':', linewidth=1.5,
                label=f'Среднее: {np.mean(hot_series):.0f}°C')
    styled_ax(ax4, f'Временной ряд EGT — цилиндр {hottest_cyl} (горячайший)', 'Отсчёт', '°C')
    ax4.legend(fontsize=7.5, facecolor=M_PANEL, labelcolor=M_TEXT)

    plt.tight_layout(pad=0.6)
    buf = chart_to_image(fig, dpi=130)
    add_img(sl, buf, Cm(0.5), Cm(2.4), Cm(32.87), Cm(16.2))

    add_text(sl,
             f'Разброс EGT: {max(maxT):.0f}–{min(maxT):.0f} = {spread:.0f}°C. Норма QSK50 MCRS: ≤30°C. '
             f'Цилиндры {hottest_cyl} и {cyls[maxT.index(sorted(maxT)[-2])]} — горячайшие. '
             f'NTE200 №82 (ESN {DATA["dml_esn"]}), запись 02.06.2026, 23 мин.',
             Cm(0.5), Cm(18.1), Cm(32), Cm(0.5), size=7, color=C_GRAY)
    print(f"Slide {n}: DML EGT analysis")

slide_dml_egt()

# --------------------------------------------------
# SLIDE 8: DML — RPM AND LOAD
# --------------------------------------------------
def slide_dml_rpm():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, 'DML ДАННЫЕ — ОБОРОТЫ И НАГРУЗКА ДВИГАТЕЛЯ',
               subtitle=f'NTE200 №82 | ESN {DATA["dml_esn"]} | Калибровка AQ60809.08',
               page_num=n)

    rpm = DATA['dml_rpm']
    load = DATA['dml_load']
    t = np.arange(len(rpm))

    fig, axes = plt.subplots(2, 2, figsize=(16.5, 8.5))
    fig.patch.set_facecolor(M_BG)

    # RPM timeseries
    ax1 = axes[0,0]
    ax1.plot(t, rpm, color=M_ACCENT, linewidth=1.2, alpha=0.9)
    ax1.axhline(y=2100, color=M_RED, linestyle='--', linewidth=2, label='Max NTE200: 2100 RPM')
    ax1.axhline(y=1990, color=M_YELLOW, linestyle=':', linewidth=1.5, label='Max 730E: 1990 RPM')
    ax1.axhline(y=np.mean(rpm), color=M_GRAY, linestyle='-', linewidth=1,
                label=f'Среднее: {np.mean(rpm):.0f} RPM')
    ax1.fill_between(t, rpm, alpha=0.15, color=M_ACCENT)
    styled_ax(ax1, 'Частота вращения двигателя (DML)', 'Отсчёт', 'об/мин')
    ax1.legend(fontsize=7.5, facecolor=M_PANEL, labelcolor=M_TEXT)
    ax1.set_ylim(700, 2300)

    # RPM histogram
    ax2 = axes[0,1]
    ax2.hist(rpm, bins=30, color=M_ACCENT, alpha=0.8, edgecolor=M_BG)
    ax2.axvline(x=2100, color=M_RED, linestyle='--', linewidth=2, label='Max NTE200')
    ax2.axvline(x=1990, color=M_YELLOW, linestyle=':', linewidth=1.5, label='Max 730E')
    ax2.axvline(x=np.mean(rpm), color=M_GRAY, linewidth=1.5, label=f'Среднее {np.mean(rpm):.0f}')
    pct_over_1990 = (rpm > 1990).mean() * 100
    ax2.text(0.7, 0.85, f'{pct_over_1990:.1f}%\nвыше 1990 RPM', transform=ax2.transAxes,
             color=M_RED, fontsize=10, fontweight='bold', ha='center')
    styled_ax(ax2, 'Распределение оборотов', 'RPM', 'Кол-во отсчётов')
    ax2.legend(fontsize=7.5, facecolor=M_PANEL, labelcolor=M_TEXT)

    # Load timeseries
    ax3 = axes[1,0]
    ax3.fill_between(t[:len(load)], load, alpha=0.6, color=M_YELLOW)
    ax3.plot(t[:len(load)], load, color=M_YELLOW, linewidth=1, alpha=0.5)
    ax3.axhline(y=np.mean(load), color=M_WHITE if False else '#CCCCCC', linewidth=1.5, linestyle='--',
                label=f'Среднее: {np.mean(load):.0f}%')
    styled_ax(ax3, 'Относительная нагрузка двигателя (%)', 'Отсчёт', '%')
    ax3.legend(fontsize=7.5, facecolor=M_PANEL, labelcolor=M_TEXT)
    ax3.set_ylim(0, 110)

    # Load histogram
    ax4 = axes[1,1]
    ax4.hist(load, bins=20, color=M_YELLOW, alpha=0.8, edgecolor=M_BG)
    pct_full = (load > 90).mean() * 100
    ax4.axvline(x=100, color=M_RED, linestyle='--', linewidth=2)
    ax4.text(0.6, 0.8, f'{pct_full:.1f}%\nвремени >90%\nнагрузки', transform=ax4.transAxes,
             color=M_RED, fontsize=10, fontweight='bold', ha='center')
    styled_ax(ax4, 'Распределение нагрузки', '%', 'Кол-во отсчётов')

    plt.tight_layout(pad=0.6)
    buf = chart_to_image(fig, dpi=130)
    add_img(sl, buf, Cm(0.5), Cm(2.4), Cm(32.87), Cm(15.5))

    # Stats bar
    stats_txt = (
        f'RPM: среднее {np.mean(rpm):.0f}  |  макс. {max(rpm):.0f} (лимит NTE200: 2100)  |  '
        f'{pct_over_1990:.1f}% времени выше 1990 RPM  |  '
        f'Нагрузка: среднее {np.mean(load):.0f}%  |  {(load>90).mean()*100:.0f}% времени при нагрузке >90%'
    )
    add_text(sl, stats_txt, Cm(0.5), Cm(18.1), Cm(32), Cm(0.5), size=7.5, color=C_ACCENT)
    print(f"Slide {n}: DML RPM/Load")

slide_dml_rpm()

# --------------------------------------------------
# SLIDE 9: ECM CALIBRATION COMPARISON
# --------------------------------------------------
def slide_ecm():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, 'СРАВНЕНИЕ КАЛИБРОВОК ECM — 15 КРИТИЧЕСКИХ ПАРАМЕТРОВ',
               subtitle='NTE200 (AQ60809.08) vs 730E-DC (AQ60217.28) | Источник: Cummins Calterm Compare Report',
               page_num=n)

    fig, ax = plt.subplots(figsize=(16.5, 9.8))
    fig.patch.set_facecolor(M_BG)
    ax.set_facecolor(M_BG)
    ax.axis('off')

    params = [
        ('Fuel Code (код топлива)', 'FC0BU52', 'FC0WH95', 'КРИТ', 'Разные стратегии впрыска MCRS'),
        ('C_TIB_Fuel_Adjust_Upper_Limit', '100%', '200%', 'КРИТ', 'Двойная компенсирующая доза впрыска!'),
        ('C_ABS_Max_Acctr_Speed', '1 990 RPM', '2 100 RPM', 'КРИТ', '+110 RPM, +25-35°C EGT'),
        ('C_ABS_DroopMax', '0% (изохрон.)', '20%', 'ВЫСК', 'Нестабильность оборотов, удары по КМ'),
        ('C_EPD_OT_RPM_Drt_En', '1 — ВКЛЮЧЕНА', '0 — ОТКЛ', 'КРИТ', 'Нет защиты от перегрева при высоких RPM'),
        ('C_EPD_CP_TB_Time_To_Max_Trq_Brt', '0 сек', '51 сек', 'ВЫСК', '51 сек на полной мощности без ограничений'),
        ('C_EPD_CT_TB_Time_To_Max_Trq_Drt', '0 сек', '43 сек', 'ВЫСК', '43 сек без дерейта по температуре'),
        ('C_DO_01_A (интеркулер) порог', '54°C', 'ОТКЛ', 'КРИТ', 'Нет сигнала перегрева интеркулера → кабина'),
        ('C_DO_07_A (ОЖ) порог', '85°C', 'ОТКЛ', 'КРИТ', 'Нет сигнала перегрева охл. жидкости → кабина'),
        ('C_DO_08_A (топливо) порог', '68°C', 'ОТКЛ', 'КРИТ', 'Нет сигнала высокой температуры топлива'),
        ('DO_Option (блок DO-выходов)', '6 765', '61 042', 'КРИТ', '7 каналов внешней защиты деактивированы'),
        ('SC_Option', '60 621', '61 878', 'СРЕДН', 'Иная системная конфигурация ECM'),
        ('C_HSST_Net_Torque_High_Thd', '5 509 Н·м', '3 650 Н·м', 'СРЕДН', '-1 859 Н·м — порог ниже нормы'),
        ('C_HSST_Net_Torque_Low_Thd', '4 722 Н·м', '1 850 Н·м', 'ВЫСК', '-2 872 Н·м — значительное снижение'),
        ('T_ABS_MaxRef_2/3', '1 990 RPM', '2 225,8 RPM', 'КРИТ', '+235,8 RPM — аномально высокий референс'),
    ]
    headers = ['Параметр ECM', '730E (AQ60217.28)', 'NTE200 (AQ60809.08)', 'Риск', 'Последствие']
    col_x = [0, 0.26, 0.40, 0.55, 0.64]
    col_w = [0.26, 0.14, 0.15, 0.09, 0.36]

    row_h = 1.0 / (len(params) + 1.5)

    # header
    for j, (h, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        ax.add_patch(mpatches.FancyBboxPatch((cx, 1-row_h), cw-0.003, row_h,
                     boxstyle='square,pad=0', facecolor='#3EF0AF', edgecolor='none'))
        ax.text(cx+cw/2, 1-row_h/2, h, color='#121C24', fontsize=8.5, ha='center', va='center', fontweight='bold')

    for i, row in enumerate(params):
        y = 1 - (i+2)*row_h
        crit = row[3] == 'КРИТ'
        high = row[3] == 'ВЫСК'
        bg = '#2a0808' if crit else ('#2a2208' if high else ('#1A2B38' if i%2==0 else M_BG))
        ax.add_patch(mpatches.FancyBboxPatch((0, y), 1, row_h,
                     boxstyle='square,pad=0', facecolor=bg, edgecolor='none'))
        vals = [row[0], row[1], row[2], row[3], row[4]]
        for j, (val, cx, cw) in enumerate(zip(vals, col_x, col_w)):
            clr = (M_RED if crit else M_YELLOW if high else M_TEXT)
            if j == 0: clr = M_ACCENT
            if j == 2 and val != row[1]: clr = M_RED if crit else M_YELLOW
            if j == 3:
                clr = M_RED if crit else M_YELLOW if high else M_TEXT
            ax.text(cx+(0.005 if j==0 else cw/2), y+row_h/2, val,
                    color=clr, fontsize=7 if j!=4 else 6.8,
                    ha='left' if j==0 else 'center', va='center',
                    fontweight='bold' if crit or j==3 else 'normal')

    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    buf = chart_to_image(fig, dpi=140)
    add_img(sl, buf, Cm(0.3), Cm(2.35), Cm(33.2), Cm(16.3))
    print(f"Slide {n}: ECM calibration table")

slide_ecm()

# --------------------------------------------------
# SLIDE 10: CONTROL GROUP COMPARISON
# --------------------------------------------------
def slide_control_group():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, 'КОНТРОЛЬНАЯ ГРУППА: 730E-DC vs NTE200 — ОДИН ДВИГАТЕЛЬ, РАЗНЫЕ РЕЗУЛЬТАТЫ',
               subtitle='Метод исключения внешних факторов | Единственное системное различие — калибровка ECM', page_num=n)

    fig, axes = plt.subplots(1, 2, figsize=(16.5, 8.5))
    fig.patch.set_facecolor(M_BG)

    # Chart 1: factors comparison matrix
    ax1 = axes[0]
    ax1.axis('off')
    factors = [
        ('Предприятие', 'АО Полюс Магадан', 'АО Полюс Магадан', True),
        ('Карьер/маршруты', 'Тот же', 'Тот же', True),
        ('Двигатель', 'QSK50 MCRS V16', 'QSK50 MCRS V16', True),
        ('Конструкция ГБЦ', 'Cummins OEM', 'Cummins OEM', True),
        ('Масло', 'CI-4 SAE 15W-40', 'CI-4 SAE 15W-40', True),
        ('Топливо', 'ГОСТ Р 52368', 'ГОСТ Р 52368', True),
        ('Пыль/климат', 'Магаданская обл.', 'Магаданская обл.', True),
        ('Служба ТО', 'Одна служба', 'Одна служба', True),
        ('Калибровка ECM', 'AQ60217.28', 'AQ60809.08 !!!', False),
        ('TIB Fuel Limit', '100%', '200% !!!', False),
        ('Max RPM', '1 990 RPM', '2 100 RPM !!!', False),
        ('DO-выходы защиты', '5 каналов АКТИВНЫ', 'ВСЕ ОТКЛЮЧЕНЫ !!!', False),
        ('Отказы ГБЦ (2023-26)', '0 СЛУЧАЕВ', '24 машины (48%) !!!', False),
    ]
    headers = ['Фактор', '730E-DC', 'NTE200', 'Одинаков?']
    col_x2 = [0, 0.30, 0.55, 0.82]
    col_w2 = [0.30, 0.25, 0.27, 0.18]
    row_h2 = 1.0 / (len(factors)+1.5)

    for j, (h, cx, cw) in enumerate(zip(headers, col_x2, col_w2)):
        ax1.add_patch(mpatches.FancyBboxPatch((cx, 1-row_h2), cw-0.003, row_h2,
                      boxstyle='square,pad=0', facecolor='#3EF0AF', edgecolor='none'))
        ax1.text(cx+cw/2, 1-row_h2/2, h, color='#121C24', fontsize=8, ha='center', va='center', fontweight='bold')

    for i, (factor, v730, vnte, same) in enumerate(factors):
        y = 1 - (i+2)*row_h2
        bg = ('#1a1a08' if not same else ('#1A2B38' if i%2==0 else M_BG))
        ax1.add_patch(mpatches.FancyBboxPatch((0, y), 1, row_h2,
                      boxstyle='square,pad=0', facecolor=bg, edgecolor='none'))
        ax1.text(col_x2[0]+0.005, y+row_h2/2, factor, color=M_ACCENT, fontsize=7.5, va='center')
        ax1.text(col_x2[1]+col_w2[1]/2, y+row_h2/2, v730, color=M_ACCENT, fontsize=7.5, ha='center', va='center')
        vnte2 = vnte.replace(' !!!','')
        clr = M_RED if not same else M_TEXT
        ax1.text(col_x2[2]+col_w2[2]/2, y+row_h2/2, vnte2, color=clr, fontsize=7.5, ha='center', va='center',
                 fontweight='bold' if not same else 'normal')
        tick = 'ДА' if same else 'НЕТ!'
        tc = M_ACCENT if same else M_RED
        ax1.text(col_x2[3]+col_w2[3]/2, y+row_h2/2, tick, color=tc, fontsize=8, ha='center', va='center', fontweight='bold')
    ax1.set_xlim(0,1); ax1.set_ylim(0,1)
    ax1.set_title('Матрица факторов среды — NTE200 vs 730E-DC', color=M_TEXT, fontsize=9.5, pad=6)

    # Chart 2: visual comparison
    ax2 = axes[1]
    ax2.set_facecolor(M_BG)
    ax2.axis('off')

    # Conclusion boxes
    boxes = [
        (0.1, 0.82, 0.38, '730E-DC', '0 отказов КУ ГБЦ\nза 2023–2026 г.\n35 127 м/ч до ТО', M_ACCENT, '#003322'),
        (0.52, 0.82, 0.38, 'NTE200', '24 машины\n48% парка\n2 катастрофы', M_RED, '#330000'),
    ]
    for bx, by, bw, title, body, col, bg in boxes:
        ax2.add_patch(mpatches.FancyBboxPatch((bx, by-0.12), bw, 0.25,
                      boxstyle='round,pad=0.02', facecolor=bg, edgecolor=col, linewidth=3))
        ax2.text(bx+bw/2, by+0.07, title, color=col, fontsize=18, ha='center', va='center', fontweight='bold')
        ax2.text(bx+bw/2, by-0.04, body, color=M_TEXT, fontsize=9.5, ha='center', va='center', linespacing=1.5)

    # Arrow down
    ax2.annotate('', xy=(0.5, 0.57), xytext=(0.5, 0.68),
                 arrowprops=dict(arrowstyle='->', color=M_YELLOW, lw=2.5))
    ax2.text(0.5, 0.53, 'ЕДИНСТВЕННОЕ РАЗЛИЧИЕ:', color=M_YELLOW, fontsize=11,
             ha='center', va='center', fontweight='bold')

    ecm_diffs = [
        'TIB = 200% (vs 100%)',
        'Max RPM = 2 100 (vs 1 990)',
        'DO-выходы защиты = ОТКЛ (vs ВКЛ)',
        'FC0 запас = 80 RPM (vs 260 RPM)',
    ]
    y_e = 0.45
    for diff in ecm_diffs:
        ax2.add_patch(mpatches.FancyBboxPatch((0.1, y_e-0.04), 0.80, 0.07,
                      boxstyle='round,pad=0.01', facecolor='#1a0808', edgecolor=M_RED, linewidth=1.5))
        ax2.text(0.5, y_e-0.005, diff, color=M_RED, fontsize=10, ha='center', va='center', fontweight='bold')
        y_e -= 0.10

    ax2.text(0.5, 0.04,
             'ВЫВОД: Принцип контрольной группы исключает пыль, масло, климат,\n'
             'топливо как первопричину. Остаётся ТОЛЬКО калибровка ECM.',
             color=M_ACCENT, fontsize=9.5, ha='center', va='center', fontweight='bold',
             bbox=dict(boxstyle='round', facecolor='#003322', edgecolor=M_ACCENT, linewidth=2, pad=0.4))
    ax2.set_xlim(0, 1); ax2.set_ylim(0, 1)

    plt.tight_layout(pad=0.6)
    buf = chart_to_image(fig, dpi=130)
    add_img(sl, buf, Cm(0.5), Cm(2.4), Cm(32.87), Cm(16.2))
    print(f"Slide {n}: Control group")

slide_control_group()

# --------------------------------------------------
# SLIDE 11: CCEC LAB REPORT CRITIQUE
# --------------------------------------------------
def slide_ccec_critique():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, 'РАЗБОР ОТЧЁТА CCEC LAB MS&T2026033 — 10 АРГУМЕНТОВ ОПРОВЕРЖЕНИЯ',
               subtitle='ESN 33232926 | 15 126 м/ч | Tao Lang | 10.03.2026', page_num=n)

    fig, ax = plt.subplots(figsize=(16.5, 9.5))
    fig.patch.set_facecolor(M_BG)
    ax.set_facecolor(M_BG)
    ax.axis('off')

    args = [
        ('АРГ-1', 'Контрольная группа', M_RED,
         '730E-DC: тот же QSK50, то же масло CI-4, та же пыль, тот же карьер. Нулевых отказов КУ за 2023-2026. Если бы зола масла была причиной — 730E тоже должен ломаться. НЕ ЛОМАЕТСЯ.'),
        ('АРГ-2', 'Твёрдость клапана в норме', M_YELLOW,
         'CCEC: 40-44 HRC при норме ≤46 HRC (ASTM E18). Клапан металлургически исправен. Причина отказа — условия эксплуатации, а не материал.'),
        ('АРГ-3', 'Пластическая деформация = T > 700°C', M_RED,
         'Inconel 751: предел ползучести при 700°C = 480 МПа. Деформация тарелки означает T > 700°C локально. При нормальной работе QSK50 температура седла 420-480°C. TIB=200% даёт +130°C → 600-650°C EGT → расчётно 700°C на тарелке.'),
        ('АРГ-4', 'Повтор через 283-1076 м/ч', M_RED,
         'После замены ГБЦ (новые детали, чистое масло, без пыли) NTE#72 отказал снова через 283 м/ч. NTE#76 — через 274 м/ч. Если причина — пыль/масло, новые детали дали бы полный ресурс. НЕ ДАЛИ.'),
        ('АРГ-5', 'EGT на 130°C выше нормы (DML)', M_YELLOW,
         'Данные DML NTE200 №82 (ESN 33238517): пик EGT цил.2 и цил.6 — 567°C. Норма QSK50: ≤520°C. 730E пик: 488°C. Разница +79°C фактически измерена. Причина: TIB=200% и Max RPM=2100.'),
        ('АРГ-6', 'Нет ECM в анализе дилера', M_RED,
         'Документ «Коренные причины» (июнь 2026): 0 упоминаний ECM, калибровки, TIB, FC-кодов, INSITE. Технический анализ без рассмотрения параметров управления — методологически неполон.'),
        ('АРГ-7', 'Engine Protection отключена', M_RED,
         'INSITE КАМСС NTE200 №85: «Engine Protection — No data available» за 4015 м/ч (27 горячих остановов). 730E №18: FC146, FC165, FC556 сработали за 35127 м/ч. Машина работала в перегреве без защиты.'),
        ('АРГ-8', 'EDS: зола масла — следствие', M_YELLOW,
         'Ca 42%, Zn 18%, P 15% — зола CI-4 масла. Это ПРОДУКТ сгорания масла, попавшего в КС через прогоревший клапан. Осадок — следствие прогара, а не его причина.'),
        ('АРГ-9', 'Самопротиворечие отчёта CCEC', M_YELLOW,
         'С.4: «вероятно, отсутствует значительный перегрев». С.7: «повышенная температура — усиливающий фактор». Нельзя одновременно утверждать оба тезиса. Пластическая деформация Inconel при T>700°C — факт, а не вероятность.'),
        ('АРГ-10', 'NTE#81 отказ при 4696 м/ч', M_RED,
         '9 ГБЦ при наработке 4696 м/ч. Накопление пыли и золы масла до критического объёма за 4600 часов? Нереалистично. ECM-калибровка действует с первого часа работы — это объясняет ранние отказы.'),
    ]

    row_h = 1.0 / (len(args)+0.5)
    for i, (code, title, col, text) in enumerate(args):
        y = 1 - (i+1)*row_h - 0.01
        bg = '#1a0a0a' if col == M_RED else '#1a1a08'
        ax.add_patch(mpatches.FancyBboxPatch((0, y), 1, row_h*0.95,
                     boxstyle='round,pad=0.003', facecolor=bg, edgecolor='none'))
        ax.add_patch(mpatches.FancyBboxPatch((0, y), 0.003, row_h*0.95,
                     boxstyle='square,pad=0', facecolor=col, edgecolor='none'))
        ax.text(0.008, y+row_h*0.7, f'[{code}]', color=col, fontsize=7.5, va='center', fontweight='bold')
        ax.text(0.065, y+row_h*0.7, title, color=col, fontsize=7.5, va='center', fontweight='bold')
        # text split
        words = text.split()
        mid = len(words) * 2 // 3
        ax.text(0.008, y+row_h*0.35, ' '.join(words[:mid]), color=M_TEXT, fontsize=6.5, va='center')
        ax.text(0.008, y+row_h*0.08, ' '.join(words[mid:]), color=M_TEXT, fontsize=6.5, va='center')

    ax.set_xlim(0,1); ax.set_ylim(0,1)
    buf = chart_to_image(fig, dpi=140)
    add_img(sl, buf, Cm(0.3), Cm(2.35), Cm(33.2), Cm(16.3))
    print(f"Slide {n}: CCEC critique")

slide_ccec_critique()

# --------------------------------------------------
# SLIDE 12: TIB + RPM MECHANISM
# --------------------------------------------------
def slide_mechanism():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, 'МЕХАНИЗМ ОТКАЗА: TIB=200% + MAX RPM=2100',
               subtitle='Комплексный тепловой анализ | Цепочка причинно-следственных связей', page_num=n)

    fig, axes = plt.subplots(1, 2, figsize=(16.5, 8.5))
    fig.patch.set_facecolor(M_BG)

    # Left: TIB mechanism flow
    ax1 = axes[0]
    ax1.set_facecolor(M_BG)
    ax1.axis('off')

    steps = [
        ('ВХОД: Снижение нагрузки\n(подъём, торможение)', M_GRAY, 0.91),
        ('Падение давления Common Rail\nниже уставки ECM', M_GRAY, 0.80),
        ('ECM активирует TIB:\nдобавление компенсирующей дозы', M_PANEL, 0.69),
        ('730E: TIB_max = 100%\nДоза ограничена', M_ACCENT, 0.58),
        ('NTE200: TIB_max = 200%\nДОЗА УДВОЕНА!', M_RED, 0.47),
        ('Неполное сгорание → EGT ↑ +130°C', M_RED, 0.36),
        ('T тарелки клапана > 700°C\n(предел Inconel 751)', M_RED, 0.25),
        ('Ползучесть → деформация → прогар', '#8B0000', 0.14),
    ]
    for i, (text, col, y) in enumerate(steps):
        is_nte = 'УДВОЕНА' in text or 'прогар' in text or 'EGT' in text or '700°C' in text
        bg = '#330000' if col == '#8B0000' else ('#1a0808' if col == M_RED else M_PANEL)
        ax1.add_patch(mpatches.FancyBboxPatch((0.05, y-0.065), 0.90, 0.10,
                      boxstyle='round,pad=0.01', facecolor=bg, edgecolor=col, linewidth=2))
        ax1.text(0.5, y-0.015, text, color=col if col != M_GRAY else M_TEXT,
                 fontsize=8.5, ha='center', va='center', fontweight='bold' if is_nte else 'normal')
        if i < len(steps)-1:
            ax1.annotate('', xy=(0.5, steps[i+1][2]+0.035), xytext=(0.5, y-0.065),
                         arrowprops=dict(arrowstyle='->', color=M_ACCENT, lw=2))
    ax1.set_xlim(0,1); ax1.set_ylim(0,1)
    ax1.set_title('Цепочка: TIB=200% → Прогар клапана', color=M_TEXT, fontsize=9.5, pad=6)

    # Right: RPM vs EGT
    ax2 = axes[1]
    rpm_vals = np.linspace(1600, 2300, 100)
    # Approximate EGT vs RPM (quadratic relationship)
    egt_730 = 420 + (rpm_vals - 1600) * 0.08 + (rpm_vals - 1600)**2 * 0.00012
    egt_nte = 420 + (rpm_vals - 1600) * 0.095 + (rpm_vals - 1600)**2 * 0.00018  # TIB effect

    ax2.plot(rpm_vals, egt_730, color=M_ACCENT, linewidth=2.5, label='730E (TIB=100%)')
    ax2.plot(rpm_vals, egt_nte, color=M_RED, linewidth=2.5, label='NTE200 (TIB=200%)')
    ax2.fill_between(rpm_vals, egt_730, egt_nte, alpha=0.2, color=M_RED)

    ax2.axvline(x=1990, color=M_ACCENT, linestyle='--', linewidth=2, label='Max 730E: 1990 RPM')
    ax2.axvline(x=2100, color=M_RED, linestyle='--', linewidth=2, label='Max NTE200: 2100 RPM')
    ax2.axhline(y=520, color=M_YELLOW, linestyle=':', linewidth=2, label='EGT норма: 520°C')
    ax2.axhline(y=700, color=M_RED, linestyle='-.', linewidth=1.5, alpha=0.7, label='Порог Inconel 751: ~700°C (тарелка)')

    # Annotation
    egt_at_2100 = 420 + (2100-1600)*0.095 + (2100-1600)**2*0.00018
    ax2.annotate(f'NTE200 @ 2100 RPM:\n~{egt_at_2100:.0f}°C EGT',
                 xy=(2100, egt_at_2100), xytext=(2150, egt_at_2100-50),
                 color=M_RED, fontsize=8.5, fontweight='bold',
                 arrowprops=dict(arrowstyle='->', color=M_RED, lw=1.5))

    styled_ax(ax2, 'EGT как функция RPM и стратегии TIB (расчётная модель)', 'RPM', '°C')
    ax2.legend(fontsize=8, facecolor=M_PANEL, labelcolor=M_TEXT, edgecolor=M_GRAY)
    ax2.set_xlim(1600, 2350)

    plt.tight_layout(pad=0.6)
    buf = chart_to_image(fig, dpi=130)
    add_img(sl, buf, Cm(0.5), Cm(2.4), Cm(32.87), Cm(16.2))
    print(f"Slide {n}: Mechanism")

slide_mechanism()

# --------------------------------------------------
# SLIDE 13: ROOT CAUSE TREE
# --------------------------------------------------
def slide_fta():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, 'ДЕРЕВО ПРИЧИН ОТКАЗА — FAULT TREE ANALYSIS (FTA)',
               subtitle='Уровень 1: доказанные причины | Уровень 2: механизмы | Уровень 3: первопричины', page_num=n)

    fig, ax = plt.subplots(figsize=(16.5, 9.5))
    fig.patch.set_facecolor(M_BG)
    ax.set_facecolor(M_BG)
    ax.axis('off')

    def nbox(x, y, w, h, text, col=M_PANEL, border=M_ACCENT, fs=8, tc=M_TEXT):
        ax.add_patch(mpatches.FancyBboxPatch((x-w/2, y-h/2), w, h,
                     boxstyle='round,pad=0.01', facecolor=col, edgecolor=border, linewidth=2))
        ax.text(x, y, text, color=tc, fontsize=fs, ha='center', va='center', linespacing=1.3)

    def arrow(x1, y1, x2, y2):
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color=M_GRAY, lw=1.5))

    # Top
    nbox(0.5, 0.93, 0.60, 0.09, 'РАЗРУШЕНИЕ КЛАПАНА QSK50 MCRS\nNTE200 — 24 машины (48% парка)', '#330000', M_RED, 9.5, M_RED)

    # Level 1
    nbox(0.25, 0.77, 0.44, 0.08, 'ТЕПЛОВАЯ ПЕРЕГРУЗКА\nКЛАПАНА (T > 700°C)', '#1a0808', M_RED, 8.5, M_RED)
    nbox(0.75, 0.77, 0.44, 0.08, 'ОТСУТСТВИЕ СВОЕВРЕМЕННОЙ\nЗАЩИТЫ ДВИГАТЕЛЯ', '#1a0808', M_RED, 8.5, M_RED)
    arrow(0.5, 0.885, 0.25, 0.815)
    arrow(0.5, 0.885, 0.75, 0.815)

    # Level 2 left
    nbox(0.08, 0.60, 0.26, 0.075, 'TIB = 200%\nДвойная доза впрыска\n+130°C EGT', '#1a0808', M_RED, 7.5, M_RED)
    nbox(0.25, 0.60, 0.26, 0.075, 'Max RPM = 2100\nvs 1990 у 730E\n+110 RPM, +35°C', '#1a0808', M_YELLOW, 7.5, M_YELLOW)
    nbox(0.42, 0.60, 0.26, 0.075, 'Иной Fuel Code\nFC0WH95 vs FC0BU52\nДругая стратегия', M_PANEL, M_GRAY, 7.5, M_GRAY)
    arrow(0.25, 0.730, 0.08, 0.637)
    arrow(0.25, 0.730, 0.25, 0.637)
    arrow(0.25, 0.730, 0.42, 0.637)

    # Level 2 right
    nbox(0.58, 0.60, 0.26, 0.075, 'DO-выходы ОТКЛ\n(DO_Option=61042)\n7 каналов деактивир.', '#1a0808', M_RED, 7.5, M_RED)
    nbox(0.75, 0.60, 0.26, 0.075, 'EPD RPM Derate\nОТКЛЮЧЁН\n(C_EPD_OT_RPM=0)', '#1a0808', M_RED, 7.5, M_RED)
    nbox(0.92, 0.60, 0.26, 0.075, 'FC556 порог: 5 кПа\nvs 10 кПа у 730E\nРанний триггер', M_PANEL, M_YELLOW, 7.5, M_YELLOW)
    arrow(0.75, 0.730, 0.58, 0.637)
    arrow(0.75, 0.730, 0.75, 0.637)
    arrow(0.75, 0.730, 0.92, 0.637)

    # Root causes
    nbox(0.25, 0.43, 0.34, 0.09, 'КОРЕННАЯ ПРИЧИНА 1\nНекорректная калибровка ECM\n(AQ60809.08): 15 критич. отличий', '#003322', M_ACCENT, 8, M_ACCENT)
    nbox(0.5, 0.43, 0.30, 0.09, 'КОРЕННАЯ ПРИЧИНА 2\nНет процедуры верификации\nкалибровок ECM при поставке', '#1a1a08', M_YELLOW, 7.5, M_YELLOW)
    nbox(0.75, 0.43, 0.30, 0.09, 'КОРЕННАЯ ПРИЧИНА 3\nОтсутствие систематич.\nмониторинга INSITE КАМСС', '#1a1a08', M_YELLOW, 7.5, M_YELLOW)

    for x1, y1, x2, y2 in [(0.08,0.561, 0.25,0.475), (0.25,0.561, 0.25,0.475),
                             (0.42,0.561, 0.5,0.475), (0.58,0.561, 0.5,0.475),
                             (0.75,0.561, 0.75,0.475), (0.92,0.561, 0.75,0.475)]:
        arrow(x1, y1, x2, y2)

    # Contributing factors
    nbox(0.5, 0.28, 0.90, 0.075,
         'СПОСОБСТВУЮЩИЕ ФАКТОРЫ (не первопричины): зола масла CI-4, пыль Si/Al,\nнарушение интервала регулировки клапанов, высокая нагрузка >90% времени',
         M_PANEL, M_GRAY, 8, M_GRAY)

    # Evidence
    nbox(0.5, 0.11, 0.96, 0.09,
         'ДОКАЗАТЕЛЬСТВА: 1) 730E на том же объекте — нулевых отказов ГБЦ  |  '
         '2) Повторные отказы через 274-1076 м/ч после ремонта  |  '
         '3) INSITE NTE200: Engine Protection "No data" vs 7 срабатываний у 730E за 35127 м/ч',
         '#003322', M_ACCENT, 7.5, M_TEXT)

    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    buf = chart_to_image(fig, dpi=140)
    add_img(sl, buf, Cm(0.3), Cm(2.35), Cm(33.2), Cm(16.3))
    print(f"Slide {n}: FTA")

slide_fta()

# --------------------------------------------------
# SLIDE 14: CONCLUSIONS
# --------------------------------------------------
def slide_conclusions():
    n = next_page()
    sl = blank_slide(prs)
    fill_bg(sl, C_BG)
    header_bar(sl, 'ВЫВОДЫ И ПЛАН МЕРОПРИЯТИЙ',
               subtitle='8 ключевых выводов | 16 мероприятий | Горизонт: 0-12 месяцев', page_num=n)

    fig, axes = plt.subplots(1, 2, figsize=(16.5, 9.2))
    fig.patch.set_facecolor(M_BG)

    # Left: conclusions
    ax1 = axes[0]
    ax1.set_facecolor(M_BG); ax1.axis('off')

    conclusions = [
        ('В-1', 'Системный характер', 'Отказы системны: 24 из ~50 машин (48%). Ни одна не достигла ресурса 25 000 м/ч (CES57000).', M_RED),
        ('В-2', 'Первопричина — ECM', 'Калибровка AQ60809.08: 15 критических отличий от 730E. Главные: TIB=200%, RPM=2100, DO=ОТКЛ.', M_ACCENT),
        ('В-3', 'Тепловая нагрузка', 'EGT NTE200: до 567°C (DML). Норма: 520°C. Межцилиндровый разброс 96°C (норма ≤30°C).', M_RED),
        ('В-4', 'Защита отключена', 'Engine Protection NTE200 №85: «No data» за 4015 м/ч. 730E №18: 7 срабатываний за 35127 м/ч.', M_RED),
        ('В-5', 'Повтор после ремонта', 'NTE#72: +283 м/ч, NTE#76: +274 м/ч. Ремонт ГБЦ без изменения ECM — не устраняет причину.', M_RED),
        ('В-6', '730E — нет отказов', 'Контрольная группа исключает пыль, масло, антифриз как причины.', M_ACCENT),
        ('В-7', 'Расход масла: симптом', '3615 л / март 2026. NTE#50: 246 л (+42% нормы). Масло в КС = следствие прогара клапана.', M_YELLOW),
        ('В-8', 'Ущерб', '>120 млн руб. (~140 замен ГБЦ + 2 катастрофы + 4583 ч простоя/мес).', M_YELLOW),
    ]
    row_h1 = 1.0 / (len(conclusions)+0.5)
    for i, (code, title, text, col) in enumerate(conclusions):
        y = 1 - (i+0.5)*row_h1
        bg = '#1a0808' if col == M_RED else ('#003322' if col == M_ACCENT else '#1a1608')
        ax1.add_patch(mpatches.FancyBboxPatch((0, y-row_h1*0.9), 1, row_h1*0.88,
                      boxstyle='round,pad=0.008', facecolor=bg, edgecolor='none'))
        ax1.add_patch(mpatches.FancyBboxPatch((0, y-row_h1*0.9), 0.004, row_h1*0.88,
                      boxstyle='square,pad=0', facecolor=col, edgecolor='none'))
        ax1.text(0.01, y-row_h1*0.23, f'[{code}] {title}', color=col, fontsize=8, va='center', fontweight='bold')
        ax1.text(0.01, y-row_h1*0.65, text, color=M_TEXT, fontsize=7, va='center')
    ax1.set_xlim(0,1); ax1.set_ylim(0,1)
    ax1.set_title('Ключевые выводы анализа', color=M_TEXT, fontsize=10, pad=4)

    # Right: action plan
    ax2 = axes[1]
    ax2.set_facecolor(M_BG); ax2.axis('off')

    actions = [
        ('М-01', 'INSITE КАМСС — снятие данных', '3 дня', 'КРИТ', M_RED),
        ('М-02', 'Верификация версий ECM (все NTE200)', '5 дней', 'КРИТ', M_RED),
        ('М-03', 'Запрос Cummins Inc. по TIB=200%', '7 дней', 'КРИТ', M_RED),
        ('М-04', 'Запрет эксплуатации >8000 м/ч без ЕСМ', '1 день', 'КРИТ', M_RED),
        ('М-05', 'Регулировка клапанов (5-10 тыс. м/ч)', '14 дней', 'ВЫСК', M_YELLOW),
        ('М-06', 'Ежедневный контроль уровня масла', '1 день', 'СРЕДН', M_ACCENT),
        ('М-07', 'Анализ DML всего парка', '14 дней', 'ВЫСК', M_YELLOW),
        ('М-08', 'Совещание с Cummins/дилером', '10 дней', 'КРИТ', M_RED),
        ('М-09', 'Перекалибровка ECM парка', '1-3 мес', 'КРИТ', M_RED),
        ('М-10', 'Активация DO-выходов защиты', '1-3 мес', 'КРИТ', M_RED),
        ('М-11', 'Внедрение мониторинга INSITE', '1-3 мес', 'ВЫСК', M_YELLOW),
        ('М-12', 'Сокращение интервала регулировки', 'до ECM', 'ВЫСК', M_YELLOW),
    ]
    row_h2 = 1.0 / (len(actions)+1.5)
    # header
    ax2.add_patch(mpatches.FancyBboxPatch((0, 1-row_h2), 1, row_h2,
                  boxstyle='square,pad=0', facecolor='#3EF0AF', edgecolor='none'))
    for j, (h, cx, cw) in enumerate(zip(['Код', 'Мероприятие', 'Срок', 'Прио'], [0, 0.07, 0.72, 0.84], [0.07, 0.65, 0.12, 0.16])):
        ax2.text(cx+cw/2, 1-row_h2/2, h, color='#121C24', fontsize=8, ha='center', va='center', fontweight='bold')

    for i, (code, action, term, prio, col) in enumerate(actions):
        y = 1-(i+2)*row_h2
        bg = '#1a0808' if col==M_RED else ('#1a1608' if col==M_YELLOW else M_PANEL)
        ax2.add_patch(mpatches.FancyBboxPatch((0, y), 1, row_h2*0.95,
                      boxstyle='square,pad=0', facecolor=bg, edgecolor='none'))
        ax2.text(0.005, y+row_h2*0.5, code, color=col, fontsize=7, va='center', fontweight='bold')
        ax2.text(0.075, y+row_h2*0.5, action, color=M_TEXT, fontsize=7, va='center')
        ax2.text(0.80, y+row_h2*0.5, term, color=M_TEXT, fontsize=7, ha='center', va='center')
        ax2.text(0.92, y+row_h2*0.5, prio, color=col, fontsize=7, ha='center', va='center', fontweight='bold')
    ax2.set_xlim(0,1); ax2.set_ylim(0,1)
    ax2.set_title('План мероприятий (0–3 месяца)', color=M_TEXT, fontsize=10, pad=4)

    plt.tight_layout(pad=0.5)
    buf = chart_to_image(fig, dpi=130)
    add_img(sl, buf, Cm(0.5), Cm(2.4), Cm(32.87), Cm(16.2))
    print(f"Slide {n}: Conclusions + actions")

slide_conclusions()

# ==========================================
# SAVE
# ==========================================
OUT = '/home/user/NTE200/Анализ_коренных_причин_ГБЦ_QSK50_NTE200.pptx'
prs.save(OUT)
import os
sz = os.path.getsize(OUT)
print(f"\nSaved: {OUT} ({sz:,} bytes = {sz//1024//1024} MB)")
print(f"Total slides: {len(prs.slides)}")
