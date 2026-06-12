#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DML Log Analysis Presentation Builder
NHL NTE200 Dump Trucks — Cummins QSK50 MCRS
Generates: DML_Анализ_LogFiles_NTE200.pptx
"""

import os
import sys
import warnings
import tempfile
import numpy as np
import pandas as pd

warnings.filterwarnings('ignore')

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.gridspec as gridspec
from matplotlib.colors import LinearSegmentedColormap

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

# ─── Brand colors ────────────────────────────────────────────────────────────
C_BG      = '#293136'
C_AXES    = '#1e2529'
C_GREEN   = '#3EF0AF'
C_WHITE   = '#FFFFFF'
C_RED     = '#FF4444'
C_ORANGE  = '#FF8C42'
C_YELLOW  = '#FFD166'
C_BLUE    = '#4EC9F0'
C_GRAY    = '#667788'
C_DARK2   = '#20272b'

def hex2rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def rgb_color(h):
    r, g, b = hex2rgb(h)
    return RGBColor(r, g, b)

# ─── Matplotlib dark theme ────────────────────────────────────────────────────
plt.rcParams.update({
    'figure.facecolor':  C_BG,
    'axes.facecolor':    C_AXES,
    'axes.edgecolor':    C_GRAY,
    'axes.labelcolor':   C_WHITE,
    'xtick.color':       C_GRAY,
    'ytick.color':       C_GRAY,
    'text.color':        C_WHITE,
    'grid.color':        C_GRAY,
    'grid.alpha':        0.3,
    'legend.facecolor':  C_AXES,
    'legend.edgecolor':  C_GRAY,
    'font.family':       'DejaVu Sans',
    'font.size':         9,
    'axes.titlesize':    11,
    'axes.labelsize':    9,
})

# ─── Slide dimensions: 33.86cm × 19.05cm (widescreen 16:9) ───────────────────
SLIDE_W = Cm(33.86)
SLIDE_H = Cm(19.05)

# ─── File configuration ───────────────────────────────────────────────────────
LOG_DIR = '/home/user/NTE200/log_files_nte200'

FILE_CFG = {
    '43':  (f'{LOG_DIR}/NTE200 №43 16.05.2026.csv',  'cp1251', 'en'),
    '45':  (f'{LOG_DIR}/NTE200 №45 18.05.2026.csv',  'cp1251', 'ru'),
    '50':  (f'{LOG_DIR}/NTE200 №50 19.05.2026.csv',  'cp1251', 'ru'),
    '52':  (f'{LOG_DIR}/NTE200 №52 18.05.2026.csv',  'cp1251', 'ru'),
    '56':  (f'{LOG_DIR}/NTE200 №56 19.05.2026.csv',  'cp1251', 'ru'),
    '84a': (f'{LOG_DIR}/NTE200 №84 16.05.2026.csv',  'cp1251', 'en'),
    '84b': (f'{LOG_DIR}/NTE200 №84 21.05.2026.csv',  'utf-8',  'ru'),
    '85':  (f'{LOG_DIR}/NTE200 №85 19.05.2026.csv',  'cp1251', 'ru'),
    '87':  (f'{LOG_DIR}/NTE200 №87 21.05.2026.csv',  'utf-8',  'ru'),
}

# ─── Known stats (pre-computed) ───────────────────────────────────────────────
KNOWN_STATS = {
    '43':  {'dur': 19.8, 'rpm_avg': 1813, 'rpm_max': 1929, 'load': 85,
            'egt_avg': 495, 'egt_max': 520, 'oil_max': 80,  'risk': 'ВНИМАНИЕ'},
    '45':  {'dur':  2.5, 'rpm_avg': 1825, 'rpm_max': 1909, 'load': 90,
            'egt_avg': 454, 'egt_max': 477, 'oil_max': 97,  'risk': 'НОРМА'},
    '50':  {'dur': 21.8, 'rpm_avg': 1865, 'rpm_max': 1927, 'load': 92,
            'egt_avg': 507, 'egt_max': 528, 'oil_max': 110, 'risk': 'КРИТИЧНО'},
    '52':  {'dur':  2.0, 'rpm_avg': 1847, 'rpm_max': 1909, 'load': 95,
            'egt_avg': 474, 'egt_max': 490, 'oil_max': 84,  'risk': 'НОРМА'},
    '56':  {'dur': 20.3, 'rpm_avg': 1769, 'rpm_max': 1922, 'load': 74,
            'egt_avg': 474, 'egt_max': 526, 'oil_max': 106, 'risk': 'ВНИМАНИЕ'},
    '84a': {'dur': 17.8, 'rpm_avg': 1839, 'rpm_max': 1924, 'load': 88,
            'egt_avg': 473, 'egt_max': 498, 'oil_max': 106, 'risk': 'ВНИМАНИЕ'},
    '84b': {'dur':  9.5, 'rpm_avg': 1809, 'rpm_max': 1936, 'load': 83,
            'egt_avg': 458, 'egt_max': 506, 'oil_max': 104, 'risk': 'ВНИМАНИЕ'},
    '85':  {'dur': 31.8, 'rpm_avg': 1731, 'rpm_max': 1930, 'load': 72,
            'egt_avg': 474, 'egt_max': 566, 'oil_max': 107, 'risk': 'КРИТИЧНО'},
    '87':  {'dur':  9.5, 'rpm_avg': 1790, 'rpm_max': 1941, 'load': 82,
            'egt_avg': 450, 'egt_max': 514, 'oil_max': 102, 'risk': 'НОРМА'},
}

RISK_COLOR = {
    'КРИТИЧНО': C_RED,
    'ВНИМАНИЕ': C_ORANGE,
    'НОРМА':    C_GREEN,
}

UNIT_LABELS = {
    '43':  '#43 (клапан)',
    '45':  '#45',
    '50':  '#50',
    '52':  '#52',
    '56':  '#56',
    '84a': '#84 (16.05)',
    '84b': '#84 (21.05)',
    '85':  '#85',
    '87':  '#87',
}

# ─── Data loading ─────────────────────────────────────────────────────────────

def parse_time_seconds(series):
    """Convert HH:MM:SS.mmm time strings to seconds from start."""
    def to_sec(s):
        try:
            parts = str(s).split(':')
            return float(parts[0]) * 3600 + float(parts[1]) * 60 + float(parts[2])
        except Exception:
            return np.nan
    t = series.apply(to_sec)
    return t - t.iloc[0]

def get_col(df, *keywords, exclude=None):
    """Find column containing all keywords (case-insensitive)."""
    for col in df.columns:
        cl = col.lower()
        if all(k.lower() in cl for k in keywords):
            if exclude and any(e.lower() in cl for e in exclude):
                continue
            return col
    return None

def load_unit(uid):
    """Load a unit CSV and extract key time series."""
    path, enc, lang = FILE_CFG[uid]
    df = pd.read_csv(path, skiprows=26, encoding=enc, sep=',',
                     decimal=',', low_memory=False, on_bad_lines='skip')

    # Determine time column
    time_col = 'Time' if lang == 'en' else 'Время'
    if time_col not in df.columns:
        time_col = df.columns[1]

    sec = parse_time_seconds(df[time_col])

    def to_float(col):
        if col is None:
            return None
        s = df[col]
        if s.dtype == object:
            s = s.astype(str).str.replace(',', '.').str.strip()
            s = pd.to_numeric(s, errors='coerce')
        return s.values

    if lang == 'en':
        rpm_col       = get_col(df, 'Engine Speed (RPM)', 'Identifier0')
        load_col      = get_col(df, 'Percent Load', 'Identifier0')
        egt_avg_col   = get_col(df, 'Average Exhaust Temperature (Calculated)', 'Identifier1')
        egt_left_col  = get_col(df, 'Average Exhaust Temperature - Left Bank', 'Identifier144')
        egt_right_col = get_col(df, 'Average Exhaust Temperature - Right Bank', 'Identifier1')
        oil_temp_col  = get_col(df, 'Engine Oil Temperature', 'Identifier0', exclude=['Voltage'])
        oil_press_col = get_col(df, 'Engine Oil Pressure (kPa)', 'Identifier0')
        coolant_col   = get_col(df, 'Engine Coolant Temperature', 'Identifier0',
                                exclude=['Voltage', 'Level'])
        boost_col     = get_col(df, 'Intake Manifold Pressure (kPa)', 'Identifier0')
        crankcase_col = get_col(df, 'Crankcase Pressure (kPa)', 'Identifier0')
        cyl_cols = {}
        for n in range(1, 17):
            suffix = 'Identifier144' if n % 2 == 1 else 'Identifier1'
            c = get_col(df, f'Exhaust Temperature Sensor Cylinder {n} (', suffix)
            if c:
                cyl_cols[n] = c
    else:
        rpm_col       = get_col(df, 'Частота вращения двигателя', 'Идентификатор0')
        load_col      = get_col(df, 'Относительная нагрузка', 'Идентификатор0')
        egt_avg_col   = get_col(df, 'Средняя температура отработавших газов (расчетное',
                                'Идентификатор1')
        egt_left_col  = get_col(df, 'левый ряд', 'Идентификатор144')
        egt_right_col = get_col(df, 'правый ряд', 'Идентификатор1')
        oil_temp_col  = get_col(df, 'Температура масла', 'Идентификатор0')
        oil_press_col = get_col(df, 'Давление масла (кПа)', 'Идентификатор0')
        coolant_col   = get_col(df, 'Датчик температуры (°C)', 'Идентификатор0')
        boost_col     = get_col(df, 'Давление во впускном коллекторе', 'Идентификатор0')
        crankcase_col = get_col(df, 'Давление картерных газов', 'Идентификатор0')
        cyl_cols = {}
        for n in range(1, 17):
            suffix = 'Идентификатор144' if n % 2 == 1 else 'Идентификатор1'
            c = get_col(df, f'цилиндра {n} (°C)', suffix)
            if c:
                cyl_cols[n] = c

    result = {
        'sec':       sec.values,
        'rpm':       to_float(rpm_col),
        'load':      to_float(load_col),
        'egt_avg':   to_float(egt_avg_col),
        'egt_left':  to_float(egt_left_col),
        'egt_right': to_float(egt_right_col),
        'oil_temp':  to_float(oil_temp_col),
        'oil_press': to_float(oil_press_col),
        'coolant':   to_float(coolant_col),
        'boost':     to_float(boost_col),
        'crankcase': to_float(crankcase_col),
        'cyl_egt':   {},
    }
    for n, col in cyl_cols.items():
        vals = to_float(col)
        if vals is not None:
            vals_f = np.array(vals, dtype=float)
            # filter fault readings < 50°C
            vals_f = np.where(vals_f < 50, np.nan, vals_f)
            result['cyl_egt'][n] = vals_f

    return result

# ─── Pre-load all units ───────────────────────────────────────────────────────

print("Loading DML log files...")
DATA = {}
for uid in FILE_CFG:
    print(f"  Loading unit {uid}...")
    try:
        DATA[uid] = load_unit(uid)
    except Exception as e:
        print(f"  WARNING: Could not load unit {uid}: {e}")
        DATA[uid] = None
print("All files loaded.\n")

# ─── Helper: save fig to temp PNG ────────────────────────────────────────────

_tmp_files = []

def fig_to_tmp(fig):
    """Save matplotlib figure to temp file, return path."""
    fd, path = tempfile.mkstemp(suffix='.png')
    os.close(fd)
    fig.savefig(path, dpi=150, bbox_inches='tight',
                facecolor=C_BG, edgecolor='none')
    plt.close(fig)
    _tmp_files.append(path)
    return path

def cleanup_tmp():
    for f in _tmp_files:
        try:
            os.remove(f)
        except Exception:
            pass

# ─── PPTX helpers ────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def set_bg(slide, color_hex=C_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color(color_hex)

def add_rect(slide, x, y, w, h, fill_hex=None, line_hex=None, line_width=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Cm(x), Cm(y), Cm(w), Cm(h)
    )
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb_color(fill_hex)
    else:
        shape.fill.background()
    if line_hex and line_width > 0:
        shape.line.color.rgb = rgb_color(line_hex)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=12, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, font_name='DejaVu Sans'):
    txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb_color(color)
    run.font.name = font_name
    return txBox

def add_image(slide, img_path, x, y, w, h):
    slide.shapes.add_picture(img_path, Cm(x), Cm(y), Cm(w), Cm(h))

def add_header_bar(slide, title, subtitle=None):
    """Add top header bar."""
    add_rect(slide, 0, 0, 33.86, 2.2, fill_hex=C_DARK2)
    # Green accent line
    add_rect(slide, 0, 2.0, 33.86, 0.12, fill_hex=C_GREEN)
    add_text(slide, title, 0.5, 0.15, 28, 1.2,
             font_size=18, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.5, 1.2, 32, 0.8,
                 font_size=10, color=C_GRAY)

def add_footer(slide, text='Анализ DML | NTE200 | QSK50 MCRS | 02.06.2026'):
    add_rect(slide, 0, 18.55, 33.86, 0.5, fill_hex=C_DARK2)
    add_text(slide, text, 0.5, 18.6, 33, 0.4,
             font_size=7, color=C_GRAY, align=PP_ALIGN.CENTER)

# ─── Slide 1: Title ──────────────────────────────────────────────────────────

def make_slide1(prs):
    slide = blank_slide(prs)
    set_bg(slide, C_BG)

    # Background accent rectangles
    add_rect(slide, 0, 0, 33.86, 19.05, fill_hex=C_BG)
    add_rect(slide, 0, 0, 8, 19.05, fill_hex=C_DARK2)
    add_rect(slide, 7.8, 0, 0.25, 19.05, fill_hex=C_GREEN)

    # Logo area / decorative
    add_rect(slide, 1, 2, 5, 0.08, fill_hex=C_GREEN)
    add_text(slide, 'DML', 1, 2.3, 6, 3,
             font_size=60, bold=True, color=C_GREEN)
    add_text(slide, 'LOG-FILES', 1, 5.5, 6, 1.5,
             font_size=22, bold=False, color=C_WHITE)
    add_text(slide, 'ANALYSIS', 1, 6.8, 6, 1.5,
             font_size=22, bold=False, color=C_WHITE)

    # Main title
    add_text(slide, 'АНАЛИЗ DML ЛОГ-ФАЙЛОВ', 9, 3, 23, 3,
             font_size=34, bold=True, color=C_WHITE)
    add_text(slide, 'ПАРК NTE200', 9, 5.5, 23, 2,
             font_size=34, bold=True, color=C_GREEN)

    # Subtitle
    add_rect(slide, 9, 8.2, 23, 0.08, fill_hex=C_GRAY)
    add_text(slide, '9 агрегатов  |  май 2026  |  Cummins QSK50 MCRS', 9, 8.5, 23, 1,
             font_size=14, color=C_GRAY)

    # Info boxes
    boxes = [
        ('9', 'агрегатов'),
        ('3', 'критических'),
        ('QSK50', 'MCRS двигатель'),
        ('DML', 'реал-тайм лог'),
    ]
    for i, (val, lbl) in enumerate(boxes):
        bx = 9 + i * 5.8
        add_rect(slide, bx, 10, 5.2, 3, fill_hex=C_DARK2, line_hex=C_GREEN, line_width=0.5)
        add_text(slide, val, bx + 0.3, 10.3, 4.5, 1.5,
                 font_size=22, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, bx + 0.3, 11.7, 4.5, 1,
                 font_size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, '02 июня 2026', 9, 14, 10, 1,
             font_size=11, color=C_GRAY)
    add_text(slide, 'Горная Евразия — Техническая служба', 9, 15, 23, 1,
             font_size=11, color=C_GRAY)

# ─── Slide 2: Overview table ─────────────────────────────────────────────────

def make_slide2(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header_bar(slide, 'ОБЗОР ПАРКА — ВСЕ АГРЕГАТЫ',
                   'Сводная таблица параметров DML | 9 агрегатов | сортировка по уровню риска ЕГТ')
    add_footer(slide)

    # Sort units by EGT risk (КРИТИЧНО → ВНИМАНИЕ → НОРМА), then by avg EGT desc
    risk_order = {'КРИТИЧНО': 0, 'ВНИМАНИЕ': 1, 'НОРМА': 2}
    sorted_units = sorted(KNOWN_STATS.keys(),
                          key=lambda u: (risk_order[KNOWN_STATS[u]['risk']],
                                         -KNOWN_STATS[u]['egt_avg']))

    # 3 columns × 3 rows
    card_w = 10.5
    card_h = 4.5
    col_positions = [0.4, 11.4, 22.4]
    row_positions = [2.4, 7.1, 11.8]

    for idx, uid in enumerate(sorted_units):
        st = KNOWN_STATS[uid]
        col = idx % 3
        row = idx // 3
        cx = col_positions[col]
        cy = row_positions[row]
        risk = st['risk']
        rcol = RISK_COLOR[risk]

        # Card background
        add_rect(slide, cx, cy, card_w, card_h, fill_hex=C_DARK2,
                 line_hex=rcol, line_width=1.0)

        # Unit number header
        add_rect(slide, cx, cy, card_w, 0.9, fill_hex=C_AXES)
        label = UNIT_LABELS.get(uid, f'#{uid}')
        add_text(slide, f'АГРЕГАТ {label}', cx + 0.3, cy + 0.1, 7.5, 0.7,
                 font_size=11, bold=True, color=C_WHITE)

        # Risk badge
        add_rect(slide, cx + card_w - 2.8, cy + 0.1, 2.6, 0.7, fill_hex=rcol)
        add_text(slide, risk, cx + card_w - 2.8, cy + 0.12, 2.6, 0.65,
                 font_size=8, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

        # Stats
        stats_data = [
            ('Время теста',   f"{st['dur']:.1f} мин"),
            ('Ср. об/мин',    f"{st['rpm_avg']}"),
            ('Нагрузка',      f"{st['load']}%"),
            ('ЕГТ средн.',    f"{st['egt_avg']}°C"),
            ('ЕГТ макс.',     f"{st['egt_max']}°C"),
            ('Т масло макс.', f"{st['oil_max']}°C"),
        ]
        for si, (label_s, val_s) in enumerate(stats_data):
            row_y = cy + 0.95 + si * 0.57
            col_val_color = C_WHITE
            if label_s == 'ЕГТ средн.' and st['egt_avg'] >= 500:
                col_val_color = C_RED
            elif label_s == 'ЕГТ средн.' and st['egt_avg'] >= 470:
                col_val_color = C_ORANGE
            elif label_s == 'ЕГТ макс.' and st['egt_max'] >= 530:
                col_val_color = C_RED
            elif label_s == 'ЕГТ макс.' and st['egt_max'] >= 500:
                col_val_color = C_ORANGE
            elif label_s == 'Т масло макс.' and st['oil_max'] >= 105:
                col_val_color = C_RED
            elif label_s == 'Т масло макс.' and st['oil_max'] >= 100:
                col_val_color = C_ORANGE

            add_text(slide, label_s, cx + 0.3, row_y, 5.5, 0.55,
                     font_size=8.5, color=C_GRAY)
            add_text(slide, val_s, cx + 6.5, row_y, 3.7, 0.55,
                     font_size=8.5, bold=True, color=col_val_color,
                     align=PP_ALIGN.RIGHT)

# ─── Slide 3: EGT Comparison Bar Chart ───────────────────────────────────────

def make_slide3(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header_bar(slide, 'СРАВНЕНИЕ ЕГТ — ВСЕ АГРЕГАТЫ',
                   'Средняя и максимальная температура выхлопных газов | лимиты 450°C / 500°C')
    add_footer(slide)

    units_order = ['50', '43', '85', '56', '84a', '84b', '52', '87', '45']
    labels      = [UNIT_LABELS[u] for u in units_order]
    egt_avg     = [KNOWN_STATS[u]['egt_avg'] for u in units_order]
    egt_max     = [KNOWN_STATS[u]['egt_max'] for u in units_order]

    bar_colors_avg = []
    bar_colors_max = []
    for u in units_order:
        r = KNOWN_STATS[u]['risk']
        if r == 'КРИТИЧНО':
            bar_colors_avg.append(C_RED)
            bar_colors_max.append('#cc0000')
        elif r == 'ВНИМАНИЕ':
            bar_colors_avg.append(C_ORANGE)
            bar_colors_max.append('#cc6000')
        else:
            bar_colors_avg.append(C_GREEN)
            bar_colors_max.append('#00cc77')

    x = np.arange(len(units_order))
    w = 0.38

    fig, ax = plt.subplots(figsize=(14, 5.5))
    fig.patch.set_facecolor(C_BG)
    ax.set_facecolor(C_AXES)

    bars1 = ax.bar(x - w/2, egt_avg, w, color=bar_colors_avg,
                   label='Средн. ЕГТ', alpha=0.9, zorder=3)
    bars2 = ax.bar(x + w/2, egt_max, w, color=bar_colors_max,
                   label='Макс. ЕГТ', alpha=0.9, zorder=3, hatch='//')

    # Limit lines
    ax.axhline(500, color=C_RED, linewidth=1.5, linestyle='--', zorder=4,
               label='Лимит 500°C')
    ax.axhline(450, color=C_YELLOW, linewidth=1.2, linestyle=':', zorder=4,
               label='Предупр. 450°C')
    # Unit 82 reference
    ax.axhline(567, color=C_BLUE, linewidth=1.0, linestyle='-.', zorder=4,
               label='#82 (02.06) макс. 567°C')

    # Value labels on bars
    for bar in bars1:
        h = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2, h + 2, f'{int(h)}',
                ha='center', va='bottom', fontsize=7.5, color=C_WHITE)
    for bar in bars2:
        h = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2, h + 2, f'{int(h)}',
                ha='center', va='bottom', fontsize=7.5, color=C_WHITE)

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9)
    ax.set_ylabel('Температура ЕГТ (°C)', color=C_WHITE)
    ax.set_ylim(400, 600)
    ax.set_title('Температура выхлопных газов — сравнение агрегатов', color=C_WHITE, fontsize=12)
    ax.legend(loc='upper right', fontsize=8, framealpha=0.7)
    ax.grid(axis='y', alpha=0.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    # Color x-tick labels by risk
    for tick, u in zip(ax.get_xticklabels(), units_order):
        tick.set_color(RISK_COLOR[KNOWN_STATS[u]['risk']])

    img_path = fig_to_tmp(fig)
    add_image(slide, img_path, 0.5, 2.4, 32.8, 15.8)

# ─── Slide 4: Key Parameters Comparison ──────────────────────────────────────

def make_slide4(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header_bar(slide, 'КЛЮЧЕВЫЕ ПАРАМЕТРЫ — СРАВНЕНИЕ',
                   '6-панельный обзор: нагрузка, обороты, ЕГТ, температура масла, наддув, давление масла')
    add_footer(slide)

    units   = ['43', '45', '50', '52', '56', '84a', '84b', '85', '87']
    labels  = [UNIT_LABELS[u] for u in units]
    colors  = [RISK_COLOR[KNOWN_STATS[u]['risk']] for u in units]

    # Gather data
    load_vals     = [KNOWN_STATS[u]['load']    for u in units]
    rpm_max_vals  = [KNOWN_STATS[u]['rpm_max'] for u in units]
    egt_avg_vals  = [KNOWN_STATS[u]['egt_avg'] for u in units]
    oil_max_vals  = [KNOWN_STATS[u]['oil_max'] for u in units]

    # Boost avg from live data (or fallback)
    boost_vals = []
    oil_press_vals = []
    for u in units:
        d = DATA.get(u)
        if d and d.get('boost') is not None:
            b = np.array(d['boost'], dtype=float)
            b_valid = b[~np.isnan(b)]
            boost_vals.append(float(np.nanmean(b_valid)) if len(b_valid) > 0 else 0)
        else:
            boost_vals.append(0)
        if d and d.get('oil_press') is not None:
            op = np.array(d['oil_press'], dtype=float)
            op_valid = op[~np.isnan(op)]
            oil_press_vals.append(float(np.nanmin(op_valid)) if len(op_valid) > 0 else 0)
        else:
            oil_press_vals.append(0)

    panels = [
        ('Нагрузка ср. (%)',       load_vals,      C_BLUE,   None, None),
        ('Макс. об/мин (RPM)',     rpm_max_vals,   C_GREEN,  None, None),
        ('Ср. ЕГТ (°C)',           egt_avg_vals,   C_ORANGE, 500,  450),
        ('Макс. т. масла (°C)',    oil_max_vals,   C_YELLOW, 110,  105),
        ('Ср. наддув (кПа)',       boost_vals,     C_BLUE,   None, None),
        ('Мин. давл. масла (кПа)', oil_press_vals, C_GREEN,  None, None),
    ]

    fig = plt.figure(figsize=(14, 8))
    fig.patch.set_facecolor(C_BG)
    gs = gridspec.GridSpec(2, 3, figure=fig, hspace=0.5, wspace=0.4)

    for idx, (title, vals, color, limit1, limit2) in enumerate(panels):
        ax = fig.add_subplot(gs[idx // 3, idx % 3])
        ax.set_facecolor(C_AXES)

        y_pos = np.arange(len(labels))
        bar_colors = []
        for v, u in zip(vals, units):
            if limit1 and v >= limit1:
                bar_colors.append(C_RED)
            else:
                bar_colors.append(RISK_COLOR[KNOWN_STATS[u]['risk']])

        ax.barh(y_pos, vals, color=bar_colors, alpha=0.85, height=0.65)

        if limit1:
            ax.axvline(limit1, color=C_RED, linewidth=1.2, linestyle='--', alpha=0.8)
        if limit2:
            ax.axvline(limit2, color=C_YELLOW, linewidth=1.0, linestyle=':', alpha=0.8)

        ax.set_yticks(y_pos)
        ax.set_yticklabels(labels, fontsize=7)
        ax.set_title(title, color=C_WHITE, fontsize=9, fontweight='bold')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.tick_params(axis='both', which='both', length=2, colors=C_GRAY)

        max_v = max(vals) if vals else 1
        for bar_val, pos in zip(vals, y_pos):
            if bar_val > 0:
                ax.text(bar_val + max_v * 0.01, pos, f'{bar_val:.0f}',
                        va='center', fontsize=7, color=C_WHITE)

    fig.suptitle('Сравнительный анализ ключевых параметров', color=C_WHITE,
                 fontsize=12, y=1.01)

    img_path = fig_to_tmp(fig)
    add_image(slide, img_path, 0.5, 2.4, 32.8, 15.8)

# ─── Slide 5: Unit 50 Deep Dive ──────────────────────────────────────────────

def make_slide5(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header_bar(slide, 'АГРЕГАТ №50 — КРИТИЧЕСКИЙ АНАЛИЗ',
                   'ЕГТ ср. 507°C | Т масла 110°C (лимит) | Нагрузка 92% | 21.8 мин')
    add_footer(slide)

    d = DATA.get('50')
    if not d:
        add_text(slide, 'Данные не загружены', 5, 8, 20, 2, font_size=14, color=C_RED)
        return

    sec  = np.array(d['sec'],      dtype=float)
    rpm  = np.array(d['rpm'],      dtype=float) if d.get('rpm')      is not None else None
    load = np.array(d['load'],     dtype=float) if d.get('load')     is not None else None
    egt  = np.array(d['egt_avg'],  dtype=float) if d.get('egt_avg')  is not None else None
    oil  = np.array(d['oil_temp'], dtype=float) if d.get('oil_temp') is not None else None

    t_min = sec / 60.0

    fig = plt.figure(figsize=(14, 8))
    fig.patch.set_facecolor(C_BG)
    gs = gridspec.GridSpec(2, 2, figure=fig, hspace=0.42, wspace=0.35)

    panels_data = [
        ('Обороты двигателя (об/мин)', rpm,  C_BLUE,   1800, None),
        ('Нагрузка (%)',               load, C_GREEN,  90,   None),
        ('Ср. ЕГТ (°C)',               egt,  C_ORANGE, 500,  450),
        ('Температура масла (°C)',      oil,  C_YELLOW, 110,  105),
    ]

    for idx, (title, vals, color, limit1, limit2) in enumerate(panels_data):
        ax = fig.add_subplot(gs[idx // 2, idx % 2])
        ax.set_facecolor(C_AXES)

        if vals is not None:
            ax.plot(t_min, vals, color=color, linewidth=1.5, alpha=0.9)
            ax.fill_between(t_min, vals, alpha=0.15, color=color)

            if limit1:
                ax.axhline(limit1, color=C_RED, linewidth=1.2, linestyle='--',
                           label=f'Лимит {limit1}', alpha=0.9)
                valid_mask = ~np.isnan(vals)
                if np.any(valid_mask):
                    max_val = np.nanmax(vals)
                    max_idx = np.nanargmax(vals)
                    ax.annotate(f'Макс: {max_val:.0f}',
                                xy=(t_min[max_idx], max_val),
                                xytext=(t_min[max_idx] + 0.5, max_val + max_val * 0.02),
                                fontsize=7.5, color=C_RED,
                                arrowprops=dict(arrowstyle='->', color=C_RED, lw=1))
            if limit2:
                ax.axhline(limit2, color=C_YELLOW, linewidth=1.0, linestyle=':',
                           label=f'Предупр. {limit2}', alpha=0.8)

        ax.set_title(title, color=C_WHITE, fontsize=9, fontweight='bold')
        ax.set_xlabel('Время (мин)', color=C_GRAY, fontsize=7.5)
        ax.legend(fontsize=7, framealpha=0.6)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(alpha=0.25)

    fig.suptitle('Агрегат №50 — Временные ряды (КРИТИЧЕСКИЙ СТАТУС)', color=C_RED,
                 fontsize=12, fontweight='bold')

    img_path = fig_to_tmp(fig)
    add_image(slide, img_path, 0.5, 2.4, 32.8, 15.8)

# ─── Slide 6: Unit 50 Cylinder EGT ───────────────────────────────────────────

def make_slide6(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header_bar(slide, 'АГРЕГАТ №50 — ПРОФИЛЬ ЕГТ ПО ЦИЛИНДРАМ',
                   'Средняя и максимальная ЕГТ каждого цилиндра | Выявление горячих цилиндров')
    add_footer(slide)

    d = DATA.get('50')
    if not d or not d.get('cyl_egt'):
        add_text(slide, 'Данные цилиндров недоступны', 5, 8, 20, 2,
                 font_size=14, color=C_RED)
        return

    cyls = sorted(d['cyl_egt'].keys())
    cyl_avg = []
    cyl_max = []
    for n in cyls:
        vals = np.array(d['cyl_egt'][n], dtype=float)
        valid = vals[~np.isnan(vals)]
        cyl_avg.append(float(np.nanmean(valid)) if len(valid) > 0 else 0)
        cyl_max.append(float(np.nanmax(valid)) if len(valid) > 0 else 0)

    cyl_avg = np.array(cyl_avg)
    cyl_max = np.array(cyl_max)

    # Hottest cylinders
    hottest_avg_idx = set(np.argsort(cyl_avg)[::-1][:3].tolist())
    hottest_max_idx = set(np.argsort(cyl_max)[::-1][:3].tolist())

    bar_colors_avg = []
    bar_colors_max = []
    for i, n in enumerate(cyls):
        if i in hottest_avg_idx:
            bar_colors_avg.append(C_RED)
        elif cyl_avg[i] >= 500:
            bar_colors_avg.append(C_ORANGE)
        else:
            bar_colors_avg.append(C_BLUE if n % 2 == 0 else C_GREEN)

        if i in hottest_max_idx:
            bar_colors_max.append('#cc0000')
        elif cyl_max[i] >= 520:
            bar_colors_max.append(C_ORANGE)
        else:
            bar_colors_max.append('#0077aa' if n % 2 == 0 else '#00aa55')

    x = np.arange(len(cyls))

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 5.5))
    fig.patch.set_facecolor(C_BG)

    for ax in (ax1, ax2):
        ax.set_facecolor(C_AXES)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(axis='y', alpha=0.25)

    # Avg EGT
    bars = ax1.bar(x, cyl_avg, color=bar_colors_avg, alpha=0.9, zorder=3)
    ax1.axhline(500, color=C_RED, linewidth=1.5, linestyle='--', alpha=0.9, label='500°C лимит')
    ax1.axhline(450, color=C_YELLOW, linewidth=1.0, linestyle=':', alpha=0.8, label='450°C предупр.')
    for bar, val in zip(bars, cyl_avg):
        if val > 0:
            ax1.text(bar.get_x() + bar.get_width()/2, val + 2, f'{val:.0f}',
                     ha='center', va='bottom', fontsize=7, color=C_WHITE)
    ax1.set_xticks(x)
    ax1.set_xticklabels([f'Ц{n}' for n in cyls], fontsize=8)
    ax1.set_ylabel('Температура (°C)', color=C_WHITE)
    ax1.set_title('Средняя ЕГТ по цилиндрам', color=C_WHITE, fontsize=11, fontweight='bold')
    ax1.legend(fontsize=8)
    ax1.set_ylim(380, 560)

    # Max EGT
    bars2 = ax2.bar(x, cyl_max, color=bar_colors_max, alpha=0.9, zorder=3)
    ax2.axhline(520, color=C_RED, linewidth=1.5, linestyle='--', alpha=0.9, label='520°C крит.')
    ax2.axhline(500, color=C_ORANGE, linewidth=1.0, linestyle=':', alpha=0.8, label='500°C лимит')
    for bar, val in zip(bars2, cyl_max):
        if val > 0:
            ax2.text(bar.get_x() + bar.get_width()/2, val + 2, f'{val:.0f}',
                     ha='center', va='bottom', fontsize=7, color=C_WHITE)
    ax2.set_xticks(x)
    ax2.set_xticklabels([f'Ц{n}' for n in cyls], fontsize=8)
    ax2.set_ylabel('Температура (°C)', color=C_WHITE)
    ax2.set_title('Максимальная ЕГТ по цилиндрам', color=C_WHITE, fontsize=11, fontweight='bold')
    ax2.legend(fontsize=8)
    ax2.set_ylim(380, 570)

    # Legend
    bank_a = mpatches.Patch(color=C_GREEN, label='Ряд A (нечётн.)')
    bank_b = mpatches.Patch(color=C_BLUE,  label='Ряд B (чётн.)')
    hot_p  = mpatches.Patch(color=C_RED,   label='Горячий цилиндр')
    fig.legend(handles=[bank_a, bank_b, hot_p], loc='lower center',
               ncol=3, fontsize=9, framealpha=0.7, bbox_to_anchor=(0.5, -0.01))

    fig.suptitle('Агрегат №50 — Распределение ЕГТ по цилиндрам',
                 color=C_WHITE, fontsize=12, fontweight='bold')

    img_path = fig_to_tmp(fig)
    add_image(slide, img_path, 0.5, 2.4, 32.8, 15.8)

# ─── Slide 7: Unit 85 Deep Dive ──────────────────────────────────────────────

def make_slide7(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header_bar(slide, 'АГРЕГАТ №85 — КРИТИЧЕСКИЙ АНАЛИЗ (ПИКОВЫЙ ЕГТ 566°C)',
                   'Наивысший пиковый ЕГТ в парке | 31.8 мин | Сравнение левый/правый ряд')
    add_footer(slide)

    d = DATA.get('85')
    if not d:
        add_text(slide, 'Данные не загружены', 5, 8, 20, 2, font_size=14, color=C_RED)
        return

    sec   = np.array(d['sec'],       dtype=float)
    egt   = np.array(d['egt_avg'],   dtype=float) if d.get('egt_avg')   is not None else None
    left  = np.array(d['egt_left'],  dtype=float) if d.get('egt_left')  is not None else None
    right = np.array(d['egt_right'], dtype=float) if d.get('egt_right') is not None else None
    rpm   = np.array(d['rpm'],       dtype=float) if d.get('rpm')       is not None else None
    load  = np.array(d['load'],      dtype=float) if d.get('load')      is not None else None

    t_min = sec / 60.0

    fig = plt.figure(figsize=(14, 8))
    fig.patch.set_facecolor(C_BG)
    gs = gridspec.GridSpec(2, 2, figure=fig, hspace=0.42, wspace=0.35)

    # Panel 1 (top, spans both cols): EGT time series with peak annotation
    ax1 = fig.add_subplot(gs[0, :])
    ax1.set_facecolor(C_AXES)
    if egt is not None:
        ax1.plot(t_min, egt, color=C_ORANGE, linewidth=1.8, label='Ср. ЕГТ', zorder=3)
        ax1.fill_between(t_min, egt, alpha=0.15, color=C_ORANGE)
        peak_val = np.nanmax(egt)
        peak_idx = int(np.nanargmax(egt))
        ax1.axhline(peak_val, color=C_RED, linewidth=1.5, linestyle='--', alpha=0.9,
                    label=f'Пик {peak_val:.0f}°C')
        ax1.axhline(500, color=C_ORANGE, linewidth=1.2, linestyle=':', alpha=0.8,
                    label='Лимит 500°C')
        ax1.axhline(567, color=C_BLUE, linewidth=1.0, linestyle='-.', alpha=0.7,
                    label='#82 (02.06) 567°C')
        ax1.annotate(f'ПИКОВОЕ ЗНАЧЕНИЕ\n{peak_val:.0f}°C',
                     xy=(t_min[peak_idx], peak_val),
                     xytext=(max(t_min[peak_idx] - 6, 1), peak_val - 50),
                     fontsize=9, color=C_RED, fontweight='bold',
                     arrowprops=dict(arrowstyle='->', color=C_RED, lw=1.5))
    if left is not None:
        ax1.plot(t_min, left, color=C_BLUE, linewidth=1.2, linestyle='--',
                 label='Лев. ряд', alpha=0.8)
    if right is not None:
        ax1.plot(t_min, right, color=C_GREEN, linewidth=1.2, linestyle='--',
                 label='Прав. ряд', alpha=0.8)
    ax1.set_title('Температура ЕГТ: общая, левый и правый ряд', color=C_WHITE,
                  fontsize=10, fontweight='bold')
    ax1.set_xlabel('Время (мин)', color=C_GRAY, fontsize=8)
    ax1.set_ylabel('Температура (°C)', color=C_WHITE, fontsize=8)
    ax1.legend(fontsize=8, framealpha=0.7, ncol=5)
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.grid(alpha=0.25)

    # Panel 2 (bottom left): RPM
    ax2 = fig.add_subplot(gs[1, 0])
    ax2.set_facecolor(C_AXES)
    if rpm is not None:
        ax2.plot(t_min, rpm, color=C_BLUE, linewidth=1.5)
        ax2.fill_between(t_min, rpm, alpha=0.12, color=C_BLUE)
        ax2.axhline(1800, color=C_GRAY, linewidth=1.0, linestyle='--', alpha=0.7,
                    label='1800 RPM')
    ax2.set_title('Обороты двигателя (об/мин)', color=C_WHITE, fontsize=9)
    ax2.set_xlabel('Время (мин)', color=C_GRAY, fontsize=7.5)
    ax2.legend(fontsize=8)
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    ax2.grid(alpha=0.25)

    # Panel 3 (bottom right): Load
    ax3 = fig.add_subplot(gs[1, 1])
    ax3.set_facecolor(C_AXES)
    if load is not None:
        ax3.plot(t_min, load, color=C_GREEN, linewidth=1.5)
        ax3.fill_between(t_min, load, alpha=0.12, color=C_GREEN)
    ax3.set_title('Нагрузка (%)', color=C_WHITE, fontsize=9)
    ax3.set_xlabel('Время (мин)', color=C_GRAY, fontsize=7.5)
    ax3.spines['top'].set_visible(False)
    ax3.spines['right'].set_visible(False)
    ax3.grid(alpha=0.25)

    fig.suptitle('Агрегат №85 — Пиковый ЕГТ 566°C (наивысший в парке)',
                 color=C_RED, fontsize=12, fontweight='bold')

    img_path = fig_to_tmp(fig)
    add_image(slide, img_path, 0.5, 2.4, 32.8, 15.8)

# ─── Slide 8: Unit 43 Deep Dive ──────────────────────────────────────────────

def make_slide8(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header_bar(slide, 'АГРЕГАТ №43 — КЛАПАННАЯ НЕИСПРАВНОСТЬ (ЗАМЕНА 21.01.2026)',
                   'Анализ временных рядов | Профиль ЕГТ по цилиндрам | Сравнение с #50')
    add_footer(slide)

    d43 = DATA.get('43')
    d50 = DATA.get('50')

    fig = plt.figure(figsize=(14, 8))
    fig.patch.set_facecolor(C_BG)
    gs = gridspec.GridSpec(2, 2, figure=fig, hspace=0.45, wspace=0.35)

    # Panel 1: EGT time series for #43
    ax1 = fig.add_subplot(gs[0, 0])
    ax1.set_facecolor(C_AXES)
    if d43 and d43.get('egt_avg') is not None:
        sec43 = np.array(d43['sec'], dtype=float)
        egt43 = np.array(d43['egt_avg'], dtype=float)
        ax1.plot(sec43/60, egt43, color=C_ORANGE, linewidth=1.8, label='#43 Ср. ЕГТ')
        ax1.fill_between(sec43/60, egt43, alpha=0.15, color=C_ORANGE)
        ax1.axhline(500, color=C_RED, linewidth=1.2, linestyle='--', alpha=0.8, label='500°C')
        ax1.axhline(450, color=C_YELLOW, linewidth=1.0, linestyle=':', alpha=0.7, label='450°C')
        max_v = float(np.nanmax(egt43))
        max_i = int(np.nanargmax(egt43))
        ax1.annotate(f'{max_v:.0f}°C', xy=(sec43[max_i]/60, max_v),
                     xytext=(max(sec43[max_i]/60 - 3, 0.5), max_v + 8),
                     fontsize=8, color=C_RED,
                     arrowprops=dict(arrowstyle='->', color=C_RED, lw=1))
    ax1.set_title('№43 — Ср. ЕГТ (временной ряд)', color=C_WHITE, fontsize=9)
    ax1.set_xlabel('Время (мин)', color=C_GRAY, fontsize=7.5)
    ax1.legend(fontsize=7.5)
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.grid(alpha=0.25)

    # Panel 2: EGT comparison #43 vs #50
    ax2 = fig.add_subplot(gs[0, 1])
    ax2.set_facecolor(C_AXES)
    if d43 and d43.get('egt_avg') is not None:
        s43 = np.array(d43['sec'], dtype=float)
        e43 = np.array(d43['egt_avg'], dtype=float)
        ax2.plot(s43/60, e43, color=C_ORANGE, linewidth=1.5, label='#43 (клапан)', alpha=0.9)
    if d50 and d50.get('egt_avg') is not None:
        s50 = np.array(d50['sec'], dtype=float)
        e50 = np.array(d50['egt_avg'], dtype=float)
        ax2.plot(s50/60, e50, color=C_RED, linewidth=1.5, label='#50 (крит.)', alpha=0.9)
    ax2.axhline(500, color=C_GRAY, linewidth=1.0, linestyle='--', alpha=0.7, label='500°C')
    ax2.set_title('Сравнение ЕГТ: №43 vs №50', color=C_WHITE, fontsize=9)
    ax2.set_xlabel('Время (мин)', color=C_GRAY, fontsize=7.5)
    ax2.legend(fontsize=8)
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    ax2.grid(alpha=0.25)

    # Panel 3: Cylinder EGT profile #43
    ax3 = fig.add_subplot(gs[1, 0])
    ax3.set_facecolor(C_AXES)
    if d43 and d43.get('cyl_egt'):
        cyls = sorted(d43['cyl_egt'].keys())
        avg_vals = []
        for n in cyls:
            v = np.array(d43['cyl_egt'][n], dtype=float)
            valid = v[~np.isnan(v)]
            avg_vals.append(float(np.nanmean(valid)) if len(valid) > 0 else 0)
        bar_colors = []
        for n, v in zip(cyls, avg_vals):
            if v >= 500:
                bar_colors.append(C_RED)
            elif v >= 470:
                bar_colors.append(C_ORANGE)
            elif n % 2 == 0:
                bar_colors.append(C_BLUE)
            else:
                bar_colors.append(C_GREEN)
        ax3.bar(range(len(cyls)), avg_vals, color=bar_colors, alpha=0.88)
        ax3.axhline(500, color=C_RED, linewidth=1.2, linestyle='--', alpha=0.8, label='500°C')
        ax3.set_xticks(range(len(cyls)))
        ax3.set_xticklabels([f'Ц{n}' for n in cyls], fontsize=7)
        ax3.set_title('№43 — Ср. ЕГТ по цилиндрам', color=C_WHITE, fontsize=9)
        ax3.legend(fontsize=7.5)
        ax3.set_ylim(380, 560)
        for i, v in enumerate(avg_vals):
            if v > 0:
                ax3.text(i, v + 2, f'{v:.0f}', ha='center', fontsize=6.5, color=C_WHITE)
    ax3.spines['top'].set_visible(False)
    ax3.spines['right'].set_visible(False)
    ax3.grid(axis='y', alpha=0.25)

    # Panel 4: Cylinder EGT comparison #43 vs #50
    ax4 = fig.add_subplot(gs[1, 1])
    ax4.set_facecolor(C_AXES)

    cyls_all = list(range(1, 17))
    avg43 = []
    avg50 = []
    for n in cyls_all:
        v43 = d43['cyl_egt'].get(n) if (d43 and d43.get('cyl_egt')) else None
        v50 = d50['cyl_egt'].get(n) if (d50 and d50.get('cyl_egt')) else None
        for lst, v in [(avg43, v43), (avg50, v50)]:
            if v is not None:
                vf = np.array(v, dtype=float)
                valid = vf[~np.isnan(vf)]
                lst.append(float(np.nanmean(valid)) if len(valid) > 0 else 0)
            else:
                lst.append(0)

    x = np.arange(len(cyls_all))
    w = 0.38
    ax4.bar(x - w/2, avg43, w, color=C_ORANGE, alpha=0.85, label='#43')
    ax4.bar(x + w/2, avg50, w, color=C_RED,    alpha=0.85, label='#50 (крит.)')
    ax4.axhline(500, color=C_GRAY, linewidth=1.0, linestyle='--', alpha=0.7)
    ax4.set_xticks(x)
    ax4.set_xticklabels([f'Ц{n}' for n in cyls_all], fontsize=6.5)
    ax4.set_title('Профиль цилиндров: №43 vs №50', color=C_WHITE, fontsize=9)
    ax4.legend(fontsize=8)
    ax4.set_ylim(350, 560)
    ax4.spines['top'].set_visible(False)
    ax4.spines['right'].set_visible(False)
    ax4.grid(axis='y', alpha=0.25)

    fig.suptitle('Агрегат №43 — Анализ клапанной неисправности',
                 color=C_ORANGE, fontsize=12, fontweight='bold')

    img_path = fig_to_tmp(fig)
    add_image(slide, img_path, 0.5, 2.4, 32.8, 15.8)

# ─── Slide 9: Unit 84 Two Sessions ───────────────────────────────────────────

def make_slide9(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header_bar(slide, 'АГРЕГАТ №84 — ДВЕ СЕССИИ (16.05 и 21.05)',
                   'Сравнение параметров: 16.05.2026 (18 мин) vs 21.05.2026 (9.5 мин)')
    add_footer(slide)

    d84a = DATA.get('84a')
    d84b = DATA.get('84b')

    fig = plt.figure(figsize=(14, 8))
    fig.patch.set_facecolor(C_BG)
    gs = gridspec.GridSpec(3, 3, figure=fig, hspace=0.55, wspace=0.38)

    sessions = [
        ('84a', '16.05 (18 мин)', d84a, C_BLUE),
        ('84b', '21.05 (9.5 мин)', d84b, C_GREEN),
    ]

    param_info = [
        ('egt_avg',  'Ср. ЕГТ (°C)',  500,  C_ORANGE),
        ('rpm',      'Об/мин',         None, C_BLUE),
        ('oil_temp', 'Т масла (°C)',   105,  C_YELLOW),
    ]

    for col, (uid, sess_label, dd, col_color) in enumerate(sessions):
        for row, (param, param_label, limit, pcolor) in enumerate(param_info):
            ax = fig.add_subplot(gs[row, col])
            ax.set_facecolor(C_AXES)

            if dd and dd.get(param) is not None:
                sec = np.array(dd['sec'], dtype=float)
                vals = np.array(dd[param], dtype=float)
                ax.plot(sec/60, vals, color=pcolor, linewidth=1.5, alpha=0.9)
                ax.fill_between(sec/60, vals, alpha=0.12, color=pcolor)
                if limit:
                    ax.axhline(limit, color=C_RED, linewidth=1.1, linestyle='--',
                               alpha=0.8, label=f'Лим. {limit}')
                valid = vals[~np.isnan(vals)]
                if len(valid) > 0:
                    title_str = (f'{sess_label}\n{param_label}: '
                                 f'ср={np.nanmean(valid):.0f}, '
                                 f'макс={np.nanmax(valid):.0f}')
                else:
                    title_str = f'{sess_label}\n{param_label}'
                ax.set_title(title_str, color=C_WHITE, fontsize=7.5)
            else:
                ax.set_title(f'{sess_label}\n{param_label}: н/д', color=C_GRAY, fontsize=7.5)
                ax.text(0.5, 0.5, 'Нет данных', ha='center', va='center',
                        transform=ax.transAxes, color=C_GRAY, fontsize=10)

            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.grid(alpha=0.22)
            ax.set_xlabel('Время (мин)', color=C_GRAY, fontsize=6.5)
            if limit:
                ax.legend(fontsize=6.5)

    # Right column: bar comparison
    ax_cmp = fig.add_subplot(gs[:, 2])
    ax_cmp.set_facecolor(C_AXES)
    labels_cmp = ['ЕГТ ср.°C', 'ЕГТ макс.°C', 'Т масло°C', 'Нагрузка %']
    vals_a = [KNOWN_STATS['84a']['egt_avg'], KNOWN_STATS['84a']['egt_max'],
              KNOWN_STATS['84a']['oil_max'],  KNOWN_STATS['84a']['load']]
    vals_b = [KNOWN_STATS['84b']['egt_avg'], KNOWN_STATS['84b']['egt_max'],
              KNOWN_STATS['84b']['oil_max'],  KNOWN_STATS['84b']['load']]

    y_pos = np.arange(len(labels_cmp))
    w = 0.36
    ax_cmp.barh(y_pos - w/2, vals_a, w, color=C_BLUE,  label='16.05', alpha=0.9)
    ax_cmp.barh(y_pos + w/2, vals_b, w, color=C_GREEN, label='21.05', alpha=0.9)

    for i, (va, vb) in enumerate(zip(vals_a, vals_b)):
        mx = max(va, vb)
        ax_cmp.text(mx + mx*0.01, i, f'{va} / {vb}', va='center',
                    fontsize=7.5, color=C_WHITE)

    ax_cmp.set_yticks(y_pos)
    ax_cmp.set_yticklabels(labels_cmp, fontsize=9)
    ax_cmp.set_title('Сравнение\nсессий', color=C_WHITE, fontsize=10, fontweight='bold')
    ax_cmp.legend(fontsize=9)
    ax_cmp.spines['top'].set_visible(False)
    ax_cmp.spines['right'].set_visible(False)
    ax_cmp.grid(axis='x', alpha=0.25)

    fig.suptitle('Агрегат №84 — Прогрессия между сессиями (16.05 → 21.05)',
                 color=C_BLUE, fontsize=12, fontweight='bold')

    img_path = fig_to_tmp(fig)
    add_image(slide, img_path, 0.5, 2.4, 32.8, 15.8)

# ─── Slide 10: EGT Heatmap All Units ─────────────────────────────────────────

def make_slide10(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header_bar(slide, 'ТЕПЛОВАЯ КАРТА ЕГТ — ВСЕ АГРЕГАТЫ × 16 ЦИЛИНДРОВ',
                   'Средняя ЕГТ по цилиндрам | Ряд A (нечётные) vs Ряд B (чётные) | N/A = неисправность датчика')
    add_footer(slide)

    units_order = ['50', '43', '85', '56', '84a', '84b', '52', '87', '45']
    unit_labels_list = [UNIT_LABELS[u] for u in units_order]
    n_cyls  = 16
    n_units = len(units_order)

    matrix = np.full((n_units, n_cyls), np.nan)

    for row_i, uid in enumerate(units_order):
        d = DATA.get(uid)
        if d and d.get('cyl_egt'):
            for n in range(1, n_cyls + 1):
                if n in d['cyl_egt']:
                    vals = np.array(d['cyl_egt'][n], dtype=float)
                    valid = vals[~np.isnan(vals)]
                    if len(valid) > 0:
                        matrix[row_i, n - 1] = float(np.nanmean(valid))

    # Custom colormap
    colors_cmap = ['#1e2529', '#2040a0', C_GREEN, C_YELLOW, C_ORANGE, C_RED, '#880000']
    cmap = LinearSegmentedColormap.from_list('egt_heat', colors_cmap, N=256)

    fig, ax = plt.subplots(figsize=(14, 6.5))
    fig.patch.set_facecolor(C_BG)
    ax.set_facecolor(C_AXES)

    masked = np.ma.array(matrix, mask=np.isnan(matrix))
    cmap.set_bad(color='#333c42')

    im = ax.imshow(masked, cmap=cmap, aspect='auto',
                   vmin=380, vmax=570, interpolation='nearest')

    cbar = plt.colorbar(im, ax=ax, shrink=0.8, pad=0.02)
    cbar.set_label('ЕГТ средн. (°C)', color=C_WHITE, fontsize=9)
    cbar.ax.yaxis.set_tick_params(color=C_WHITE)
    plt.setp(cbar.ax.yaxis.get_ticklabels(), color=C_WHITE)

    ax.set_xticks(range(n_cyls))
    ax.set_xticklabels([f'Ц{i+1}' for i in range(n_cyls)], fontsize=8.5)
    ax.set_yticks(range(n_units))
    ax.set_yticklabels(unit_labels_list, fontsize=9)

    for tick, uid in zip(ax.get_yticklabels(), units_order):
        tick.set_color(RISK_COLOR[KNOWN_STATS[uid]['risk']])

    for i in range(n_units):
        for j in range(n_cyls):
            val = matrix[i, j]
            if not np.isnan(val):
                text_color = C_BG if val > 480 else C_WHITE
                ax.text(j, i, f'{val:.0f}', ha='center', va='center',
                        fontsize=6.5, color=text_color, fontweight='bold')
            else:
                ax.text(j, i, 'N/A', ha='center', va='center',
                        fontsize=6, color=C_GRAY)

    # Bank separator line
    ax.axvline(7.5, color=C_WHITE, linewidth=2, alpha=0.5)

    # Bank labels (above the heatmap)
    ax.text(3.5, -0.85, 'РЯД A (нечётные: 1,3,5...15)',
            ha='center', fontsize=8, color=C_GREEN, transform=ax.transData)
    ax.text(11.5, -0.85, 'РЯД B (чётные: 2,4,6...16)',
            ha='center', fontsize=8, color=C_BLUE, transform=ax.transData)

    ax.set_title('Матрица ЕГТ: агрегаты × цилиндры (средняя температура °C)',
                 color=C_WHITE, fontsize=11, fontweight='bold', pad=18)

    # Grid lines
    for x in np.arange(-0.5, n_cyls, 1):
        ax.axvline(x, color='#3a4550', linewidth=0.5, alpha=0.5)
    for y in np.arange(-0.5, n_units, 1):
        ax.axhline(y, color='#3a4550', linewidth=0.5, alpha=0.5)

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    img_path = fig_to_tmp(fig)
    add_image(slide, img_path, 0.5, 2.4, 32.8, 15.8)

# ─── Slide 11: Anomaly Detection Summary ─────────────────────────────────────

def make_slide11(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header_bar(slide, 'СВОДКА АНОМАЛИЙ — ДЕТЕКЦИЯ ОТКЛОНЕНИЙ',
                   'ЕГТ >500°C | Т масла >105°C | Пики RPM | Неисправности датчиков')
    add_footer(slide)

    units_order = ['50', '43', '85', '56', '84a', '84b', '52', '87', '45']
    anomaly_data = []

    for uid in units_order:
        d = DATA.get(uid)
        st = KNOWN_STATS[uid]

        # Minutes with EGT > 500
        egt_over500_min = 0.0
        if d and d.get('egt_avg') is not None:
            egt = np.array(d['egt_avg'], dtype=float)
            over500_frac = np.sum(egt > 500) / len(egt) if len(egt) > 0 else 0
            egt_over500_min = round(over500_frac * st['dur'], 1)

        # Oil temp > 105 minutes
        oil_over105_min = 0.0
        if d and d.get('oil_temp') is not None:
            oil = np.array(d['oil_temp'], dtype=float)
            oil_valid = oil[~np.isnan(oil)]
            if len(oil_valid) > 0:
                over105_frac = np.sum(oil_valid > 105) / len(oil_valid)
                oil_over105_min = round(over105_frac * st['dur'], 1)

        # Faulty sensors
        fault_sensors = 0
        if d and d.get('cyl_egt'):
            for n in range(1, 17):
                if n not in d['cyl_egt']:
                    fault_sensors += 1
                else:
                    v = np.array(d['cyl_egt'][n], dtype=float)
                    nan_pct = np.sum(np.isnan(v)) / len(v) if len(v) > 0 else 0
                    if nan_pct > 0.3:
                        fault_sensors += 1

        anomaly_data.append({
            'uid':        uid,
            'label':      UNIT_LABELS[uid],
            'risk':       st['risk'],
            'egt_avg':    st['egt_avg'],
            'egt_max':    st['egt_max'],
            'egt_over500': egt_over500_min,
            'oil_max':    st['oil_max'],
            'oil_over105': oil_over105_min,
            'rpm_max':    st['rpm_max'],
            'fault_sensors': fault_sensors,
        })

    fig = plt.figure(figsize=(14, 8))
    fig.patch.set_facecolor(C_BG)
    gs = gridspec.GridSpec(2, 2, figure=fig, hspace=0.4, wspace=0.35)

    labels = [d['label'] for d in anomaly_data]

    # Panel 1: Minutes with EGT > 500
    ax1 = fig.add_subplot(gs[0, 0])
    ax1.set_facecolor(C_AXES)
    over500_vals = [d['egt_over500'] for d in anomaly_data]
    bar_c = [RISK_COLOR[d['risk']] for d in anomaly_data]
    bars = ax1.bar(labels, over500_vals, color=bar_c, alpha=0.9)
    for bar, v in zip(bars, over500_vals):
        if v > 0:
            ax1.text(bar.get_x() + bar.get_width()/2, v + 0.05, f'{v}',
                     ha='center', fontsize=7.5, color=C_WHITE)
    ax1.set_title('Время с ЕГТ >500°C (мин)', color=C_WHITE, fontsize=9, fontweight='bold')
    ax1.set_ylabel('Минут', color=C_WHITE, fontsize=8)
    ax1.tick_params(axis='x', rotation=30, labelsize=7.5)
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.grid(axis='y', alpha=0.25)

    # Panel 2: Oil temp > 105 minutes
    ax2 = fig.add_subplot(gs[0, 1])
    ax2.set_facecolor(C_AXES)
    oil_over_vals = [d['oil_over105'] for d in anomaly_data]
    bar_c2 = [C_RED if d['oil_max'] >= 110 else (C_ORANGE if d['oil_max'] >= 105 else C_GREEN)
               for d in anomaly_data]
    bars2 = ax2.bar(labels, oil_over_vals, color=bar_c2, alpha=0.9)
    for bar, v in zip(bars2, oil_over_vals):
        if v > 0:
            ax2.text(bar.get_x() + bar.get_width()/2, v + 0.05, f'{v}',
                     ha='center', fontsize=7.5, color=C_WHITE)
    ax2.set_title('Время с T масла >105°C (мин)', color=C_WHITE, fontsize=9, fontweight='bold')
    ax2.set_ylabel('Минут', color=C_WHITE, fontsize=8)
    ax2.tick_params(axis='x', rotation=30, labelsize=7.5)
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    ax2.grid(axis='y', alpha=0.25)

    # Panel 3: Max EGT traffic light
    ax3 = fig.add_subplot(gs[1, 0])
    ax3.set_facecolor(C_AXES)
    egt_max_vals = [d['egt_max'] for d in anomaly_data]
    bar_c3 = [C_RED if v >= 530 else (C_ORANGE if v >= 500 else C_GREEN) for v in egt_max_vals]
    bars3 = ax3.bar(labels, egt_max_vals, color=bar_c3, alpha=0.9)
    ax3.axhline(530, color=C_RED, linewidth=1.2, linestyle='--', alpha=0.8, label='530°C крит.')
    ax3.axhline(500, color=C_ORANGE, linewidth=1.0, linestyle=':', alpha=0.8, label='500°C лимит')
    for bar, v in zip(bars3, egt_max_vals):
        ax3.text(bar.get_x() + bar.get_width()/2, v + 1, f'{v}',
                 ha='center', fontsize=7, color=C_WHITE)
    ax3.set_title('Макс. ЕГТ (светофор)', color=C_WHITE, fontsize=9, fontweight='bold')
    ax3.legend(fontsize=7.5)
    ax3.set_ylabel('°C', color=C_WHITE, fontsize=8)
    ax3.tick_params(axis='x', rotation=30, labelsize=7.5)
    ax3.spines['top'].set_visible(False)
    ax3.spines['right'].set_visible(False)
    ax3.grid(axis='y', alpha=0.25)
    ax3.set_ylim(440, 610)

    # Panel 4: Sensor faults
    ax4 = fig.add_subplot(gs[1, 1])
    ax4.set_facecolor(C_AXES)
    fault_vals = [d['fault_sensors'] for d in anomaly_data]
    fault_colors = [C_RED if v >= 3 else (C_ORANGE if v >= 1 else C_GREEN) for v in fault_vals]
    bars4 = ax4.bar(labels, fault_vals, color=fault_colors, alpha=0.9)
    for bar, v in zip(bars4, fault_vals):
        ax4.text(bar.get_x() + bar.get_width()/2, v + 0.05, str(v),
                 ha='center', fontsize=8, color=C_WHITE)
    ax4.set_title('Неисправности датчиков ЕГТ (кол-во)', color=C_WHITE,
                  fontsize=9, fontweight='bold')
    ax4.set_ylabel('Кол-во', color=C_WHITE, fontsize=8)
    ax4.tick_params(axis='x', rotation=30, labelsize=7.5)
    ax4.spines['top'].set_visible(False)
    ax4.spines['right'].set_visible(False)
    ax4.grid(axis='y', alpha=0.25)

    fig.suptitle('Сводка аномалий — Светофорная система оценки рисков',
                 color=C_WHITE, fontsize=12, fontweight='bold')

    img_path = fig_to_tmp(fig)
    add_image(slide, img_path, 0.5, 2.4, 32.8, 15.8)

# ─── Slide 12: Conclusions and Recommendations ───────────────────────────────

def make_slide12(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_header_bar(slide, 'ВЫВОДЫ И РЕКОМЕНДАЦИИ',
                   'Анализ DML лог-файлов | Парк NTE200 | Cummins QSK50 MCRS | 02.06.2026')
    add_footer(slide)

    # Left column: Risks
    add_rect(slide, 0.4, 2.5, 15.8, 0.6, fill_hex=C_RED)
    add_text(slide, 'ВЫЯВЛЕННЫЕ РИСКИ', 0.6, 2.55, 15, 0.5,
             font_size=11, bold=True, color=C_WHITE)

    risks = [
        ('КРИТИЧНО', C_RED,
         '#50 — ЕГТ 507°C ср., 528°C макс. | Т масла 110°C (лимит QSK50)'),
        ('КРИТИЧНО', C_RED,
         '#85 — Пиковый ЕГТ 566°C (рекорд парка) | Вероятна неисправность форсунки/клапана'),
        ('ВНИМАНИЕ', C_ORANGE,
         '#43 — История замены клапана 21.01.2026 | ЕГТ ср. 495°C (граница нормы)'),
        ('ВНИМАНИЕ', C_ORANGE,
         '#56 — ЕГТ макс. 526°C | Т масла 106°C | Разброс ЕГТ по цилиндрам'),
        ('ВНИМАНИЕ', C_ORANGE,
         '#84 (16.05) — ЕГТ ср. 473°C, Т масла 106°C | Динамика требует контроля'),
        ('АНАЛИЗ', C_YELLOW,
         '#45, #52 — Кратковременные тесты (2-2.5 мин): недостаточно для диагностики'),
        ('СПРАВКА', C_BLUE,
         'Агрегат #82 (тест 02.06.2026): макс. ЕГТ 567°C — сопоставим с #85 (566°C)'),
    ]

    for i, (badge, badge_col, text) in enumerate(risks):
        ry = 3.3 + i * 1.42
        add_rect(slide, 0.4, ry, 1.4, 1.1, fill_hex=badge_col)
        add_text(slide, badge, 0.4, ry + 0.2, 1.4, 0.75,
                 font_size=6.5, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        add_rect(slide, 1.85, ry, 14.3, 1.1, fill_hex=C_AXES,
                 line_hex=badge_col, line_width=0.5)
        add_text(slide, text, 2.0, ry + 0.15, 14.0, 0.9, font_size=8.5, color=C_WHITE)

    # Right column: Recommendations
    add_rect(slide, 17.1, 2.5, 16.3, 0.6, fill_hex=C_GREEN)
    add_text(slide, 'РЕКОМЕНДАЦИИ', 17.3, 2.55, 15.5, 0.5,
             font_size=11, bold=True, color=C_BG)

    recs = [
        ('СРОЧНО', C_RED,
         '#50: Немедленная диагностика | Проверить форсунки цилиндров с макс. ЕГТ | '
         'Контроль уровня и качества масла | Ограничить нагрузку до устранения причины'),
        ('СРОЧНО', C_RED,
         '#85: Диагностика причины пика 566°C | Анализ цилиндра с аномальным ЕГТ | '
         'Проверить турбину VGT, охладитель наддувочного воздуха'),
        ('ПЛАНОВОЕ', C_ORANGE,
         '#43: Расширенная диагностика через 500 м.ч. после ремонта клапана | '
         'Сравнить ЕГТ цилиндров 11, 13 с базовой линией'),
        ('ПЛАНОВОЕ', C_ORANGE,
         '#56: Мониторинг ЕГТ > 500°C | Проверить состояние клапанного механизма | '
         'Анализ разброса ЕГТ левый/правый ряд'),
        ('ПЛАНОВОЕ', C_ORANGE,
         '#84: Повторный тест на нагруженном режиме (>15 мин) | '
         'Сравнительный анализ с сессией 16.05'),
        ('МОНИТОР.', C_YELLOW,
         '#45, #52: Повторные тесты длительностью >15 мин | '
         'Тест при нагрузке для объективной оценки состояния'),
        ('СИСТЕМНО', C_BLUE,
         'Все агрегаты: ежеквартальный DML-мониторинг | Ведение базы данных ЕГТ | '
         'Корреляция с данными ECM и историей технического обслуживания'),
    ]

    for i, (badge, badge_col, text) in enumerate(recs):
        ry = 3.3 + i * 1.42
        add_rect(slide, 17.1, ry, 1.5, 1.1, fill_hex=badge_col)
        add_text(slide, badge, 17.1, ry + 0.2, 1.5, 0.75,
                 font_size=6.5, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        add_rect(slide, 18.65, ry, 14.7, 1.1, fill_hex=C_AXES,
                 line_hex=badge_col, line_width=0.5)
        add_text(slide, text, 18.8, ry + 0.05, 14.5, 1.0, font_size=7.5, color=C_WHITE)

    # Bottom key message bar
    add_rect(slide, 0, 13.45, 33.86, 5.1, fill_hex=C_DARK2)
    add_rect(slide, 0, 13.45, 33.86, 0.08, fill_hex=C_GREEN)
    add_text(slide, 'КЛЮЧЕВОЕ СООБЩЕНИЕ', 0.5, 13.65, 32, 0.7,
             font_size=10, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    key_msg = (
        'Агрегаты №50 и №85 требуют немедленного технического вмешательства. '
        'ЕГТ №50 (507°C ср.) и температура масла (110°C = лимит QSK50) указывают на критическое '
        'тепловое состояние двигателя. Пиковый ЕГТ 566°C у №85 — наивысший в парке (сравним с '
        '#82: 567°C от 02.06.2026). Агрегат №43, несмотря на замену клапанов 21.01.2026, '
        'демонстрирует ЕГТ на границе нормы (495°C ср.) — требует расширенного контроля. '
        'Рекомендуется переход на ежеквартальный DML-мониторинг всего парка NTE200 '
        'с созданием базы данных исторических ЕГТ для раннего выявления деградации.'
    )
    add_text(slide, key_msg, 0.8, 14.45, 32.2, 4.0,
             font_size=9, color=C_WHITE)

# ─── Main build ───────────────────────────────────────────────────────────────

def build_presentation():
    out_path = '/home/user/NTE200/DML_Анализ_LogFiles_NTE200.pptx'

    print("Building presentation...")
    prs = new_prs()

    print("  Slide 1: Title...")
    make_slide1(prs)

    print("  Slide 2: Overview table...")
    make_slide2(prs)

    print("  Slide 3: EGT comparison chart...")
    make_slide3(prs)

    print("  Slide 4: Key parameters comparison...")
    make_slide4(prs)

    print("  Slide 5: Unit 50 deep dive...")
    make_slide5(prs)

    print("  Slide 6: Unit 50 cylinder EGT...")
    make_slide6(prs)

    print("  Slide 7: Unit 85 deep dive...")
    make_slide7(prs)

    print("  Slide 8: Unit 43 deep dive...")
    make_slide8(prs)

    print("  Slide 9: Unit 84 two sessions...")
    make_slide9(prs)

    print("  Slide 10: EGT heatmap...")
    make_slide10(prs)

    print("  Slide 11: Anomaly detection...")
    make_slide11(prs)

    print("  Slide 12: Conclusions...")
    make_slide12(prs)

    print(f"Saving to {out_path}...")
    prs.save(out_path)
    print("Done!")

    cleanup_tmp()
    print(f"\nPresentation saved: {out_path}")
    print(f"Total slides: {len(prs.slides)}")
    return out_path


if __name__ == '__main__':
    build_presentation()
