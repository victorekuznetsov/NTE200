#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Add cylinder EGT bank asymmetry slides (with delta visualization)
to DML_Анализ_LogFiles_NTE200.pptx.
Replaces/appends slides after existing 12 analytical slides.
"""

import os
import tempfile
import warnings
import numpy as np
import pandas as pd

warnings.filterwarnings('ignore')

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_BG    = '#293136'
C_AXES  = '#1e2529'
C_GREEN = '#3EF0AF'
C_WHITE = '#FFFFFF'
C_RED   = '#FF4444'
C_ORANGE= '#FF8C42'
C_GRAY  = '#667788'
C_DARK2 = '#20272b'
C_BLUE  = '#4EC9F0'

def hex2rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def rgb_color(h):
    r, g, b = hex2rgb(h)
    return RGBColor(r, g, b)

plt.rcParams.update({
    'figure.facecolor': C_BG, 'axes.facecolor': C_AXES,
    'axes.edgecolor': C_GRAY, 'axes.labelcolor': C_WHITE,
    'xtick.color': C_GRAY, 'ytick.color': C_GRAY,
    'text.color': C_WHITE, 'grid.color': C_GRAY, 'grid.alpha': 0.3,
    'legend.facecolor': C_AXES, 'legend.edgecolor': C_GRAY,
    'font.family': 'DejaVu Sans', 'font.size': 9,
})

LOG_DIR  = '/home/user/NTE200/log_files_nte200'
PPTX_IN  = '/home/user/NTE200/DML_Анализ_LogFiles_NTE200.pptx'
PPTX_OUT = '/home/user/NTE200/DML_Анализ_LogFiles_NTE200.pptx'

FILE_CFG = {
    '43':  (f'{LOG_DIR}/NTE200 №43 16.05.2026.csv',  'cp1251', 'en', '16.05.2026'),
    '45':  (f'{LOG_DIR}/NTE200 №45 18.05.2026.csv',  'cp1251', 'ru', '18.05.2026'),
    '50':  (f'{LOG_DIR}/NTE200 №50 19.05.2026.csv',  'cp1251', 'ru', '19.05.2026'),
    '52':  (f'{LOG_DIR}/NTE200 №52 18.05.2026.csv',  'cp1251', 'ru', '18.05.2026'),
    '56':  (f'{LOG_DIR}/NTE200 №56 19.05.2026.csv',  'cp1251', 'ru', '19.05.2026'),
    '84a': (f'{LOG_DIR}/NTE200 №84 16.05.2026.csv',  'cp1251', 'en', '16.05.2026'),
    '84b': (f'{LOG_DIR}/NTE200 №84 21.05.2026.csv',  'utf-8',  'ru', '21.05.2026'),
    '85':  (f'{LOG_DIR}/NTE200 №85 19.05.2026.csv',  'cp1251', 'ru', '19.05.2026'),
    '87':  (f'{LOG_DIR}/NTE200 №87 21.05.2026.csv',  'utf-8',  'ru', '21.05.2026'),
}

UNIT_LABELS = {
    '43': '#43 (клапан)', '45': '#45', '50': '#50', '52': '#52',
    '56': '#56', '84a': '#84 (16.05)', '84b': '#84 (21.05)',
    '85': '#85', '87': '#87',
}

# ─── Data loading ─────────────────────────────────────────────────────────────

def get_col(df, *keywords, exclude=None):
    for col in df.columns:
        cl = col.lower()
        if all(k.lower() in cl for k in keywords):
            if exclude and any(e.lower() in cl for e in exclude):
                continue
            return col
    return None

def parse_time(series):
    def f(s):
        try:
            p = str(s).split(':')
            return float(p[0])*3600 + float(p[1])*60 + float(p[2])
        except: return np.nan
    t = series.apply(f)
    return (t - t.dropna().iloc[0]).values

def load_unit(uid):
    path, enc, lang, date = FILE_CFG[uid]
    df = pd.read_csv(path, skiprows=26, encoding=enc, decimal=',',
                     on_bad_lines='skip', low_memory=False)
    time_col = 'Time' if lang == 'en' else 'Время'
    if time_col not in df.columns: time_col = df.columns[1]
    sec = parse_time(df[time_col])

    cyl_egt = {}
    for n in range(1, 17):
        if lang == 'en':
            sfx = 'Identifier144' if n % 2 == 1 else 'Identifier1'
            col = get_col(df, f'Exhaust Temperature Sensor Cylinder {n} (', sfx, exclude=['Voltage'])
        else:
            sfx = 'Идентификатор144' if n % 2 == 1 else 'Идентификатор1'
            col = get_col(df, f'цилиндра {n} (°C)', sfx, exclude=['Напряжение', 'Voltage'])
        if col:
            v = df[col].values
            if hasattr(v[0], '__str__') and isinstance(v[0], (str,)):
                v = pd.to_numeric(pd.Series(v).str.replace(',', '.'), errors='coerce').values
            v = np.array(v, dtype=float)
            v = np.where((v < 50) | (v > 800), np.nan, v)
            cyl_egt[n] = v
    return sec, cyl_egt, date

# ─── Delta chart ──────────────────────────────────────────────────────────────

def make_delta_chart(uid, sec, cyl_egt, date):
    """
    3-row layout:
      Row 0 (left):  A-bank time series + intra-bank spread shading
      Row 0 (right): B-bank time series + intra-bank spread shading
      Row 1 (full):  Per-cylinder Δ from bank average (bar chart)
      Row 2 (full):  Inter-bank delta over time (B_avg - A_avg)
    """
    bank_a = [1, 3, 5, 7, 9, 11, 13, 15]
    bank_b = [2, 4, 6, 8, 10, 12, 14, 16]
    colors_a = plt.cm.YlOrRd(np.linspace(0.35, 0.95, 8))
    colors_b = plt.cm.cool(np.linspace(0.25, 0.85, 8))

    unit_lbl = UNIT_LABELS[uid]

    fig = plt.figure(figsize=(17, 11))
    fig.patch.set_facecolor(C_BG)
    fig.suptitle(
        f'NTE200 {unit_lbl}  |  {date}  |  Температурная асимметрия банков  |  Δ между цилиндрами',
        color=C_WHITE, fontsize=11, fontweight='bold', y=1.0)

    gs = gridspec.GridSpec(3, 2, figure=fig,
                           height_ratios=[3.2, 1.6, 1.6],
                           hspace=0.55, wspace=0.28)
    ax_a  = fig.add_subplot(gs[0, 0])   # A-bank time series
    ax_b  = fig.add_subplot(gs[0, 1])   # B-bank time series
    ax_db = fig.add_subplot(gs[1, :])   # Per-cyl Δ bar (full width)
    ax_dt = fig.add_subplot(gs[2, :])   # Inter-bank Δ over time (full width)

    valid_sec = sec[np.isfinite(sec)]
    t_max = valid_sec[-1] if len(valid_sec) > 0 else 1

    # ── Helper: build bank avg time series ────────────────────────────────────
    def bank_ts(bank):
        stacks = []
        for n in bank:
            if n in cyl_egt and np.any(np.isfinite(cyl_egt[n])):
                stacks.append(cyl_egt[n])
        if not stacks: return np.full_like(sec, np.nan, dtype=float)
        return np.nanmean(np.vstack(stacks), axis=0)

    def bank_max_ts(bank):
        stacks = [cyl_egt[n] for n in bank if n in cyl_egt and np.any(np.isfinite(cyl_egt[n]))]
        if not stacks: return np.full_like(sec, np.nan, dtype=float)
        return np.nanmax(np.vstack(stacks), axis=0)

    def bank_min_ts(bank):
        stacks = [cyl_egt[n] for n in bank if n in cyl_egt and np.any(np.isfinite(cyl_egt[n]))]
        if not stacks: return np.full_like(sec, np.nan, dtype=float)
        return np.nanmin(np.vstack(stacks), axis=0)

    def smooth(v, w=5):
        return pd.Series(v).rolling(w, center=True, min_periods=1).mean().values

    a_avg_ts = smooth(bank_ts(bank_a))
    b_avg_ts = smooth(bank_ts(bank_b))
    a_max_ts = smooth(bank_max_ts(bank_a))
    a_min_ts = smooth(bank_min_ts(bank_a))
    b_max_ts = smooth(bank_max_ts(bank_b))
    b_min_ts = smooth(bank_min_ts(bank_b))

    # ── A-bank and B-bank time series panels ──────────────────────────────────
    bank_stats = {}
    for ax_p, bank, colors, bname, ecm in [
        (ax_a, bank_a, colors_a, 'А-банк (нечётные)', 'ECM144'),
        (ax_b, bank_b, colors_b, 'В-банк (чётные)',   'ECM1'),
    ]:
        ax_p.set_facecolor(C_AXES)

        # Spread band (intra-bank max-min)
        mx_ts = smooth(bank_max_ts(bank))
        mn_ts = smooth(bank_min_ts(bank))
        ax_p.fill_between(sec, mn_ts, mx_ts, alpha=0.12, color=C_GRAY,
                          label='Разброс Δ')

        # Individual cylinder lines
        for cyl, color in zip(bank, colors):
            if cyl not in cyl_egt: continue
            v_s = smooth(cyl_egt[cyl])
            lw = 2.0 if cyl in [1, 2, 6, 9, 11, 15] else 1.3
            ax_p.plot(sec, v_s, color=color, lw=lw, label=f'Ц{cyl}', alpha=0.9)

        # Bank average line
        avg_ts = smooth(bank_ts(bank))
        ax_p.plot(sec, avg_ts, color=C_WHITE, lw=1.8, ls='--', alpha=0.6,
                  label='Avg банка', zorder=5)

        # Reference lines
        ax_p.axhline(500, color=C_RED,    lw=1.4, ls='--', alpha=0.85, label='500°C')
        ax_p.axhline(550, color='#FF0000',lw=1.0, ls='-',  alpha=0.45, label='550°C')

        ax_p.set_ylim(150, 640)
        ax_p.set_xlim(0, t_max)
        ax_p.set_title(f'{bname}  |  {ecm}', color=C_WHITE, fontsize=10, fontweight='bold')
        ax_p.set_ylabel('ЕГТ (°C)', color=C_WHITE, fontsize=8.5)
        ax_p.set_xlabel('Время (сек)', color=C_WHITE, fontsize=8.5)
        ax_p.legend(fontsize=6.5, ncol=3, loc='upper left', framealpha=0.35,
                    facecolor=C_AXES, edgecolor='#445055', labelcolor=C_WHITE)
        ax_p.grid(True, alpha=0.3)
        ax_p.tick_params(colors=C_WHITE, labelsize=7.5)
        for sp in ax_p.spines.values(): sp.set_edgecolor('#445055')

        # Stats + intra-bank Δ
        avgs, maxs = [], []
        for c in bank:
            if c in cyl_egt:
                finite = cyl_egt[c][np.isfinite(cyl_egt[c])]
                if len(finite) > 5:
                    avgs.append(np.mean(finite)); maxs.append(np.max(finite))
        if avgs:
            b_avg_v = np.mean(avgs); b_max_v = np.max(maxs)
            intra_delta = b_max_v - np.min([np.min(cyl_egt[c][np.isfinite(cyl_egt[c])]) for c in bank if c in cyl_egt and np.any(np.isfinite(cyl_egt[c]))])
        else:
            b_avg_v = b_max_v = intra_delta = 0.0

        box_c = C_RED if b_max_v >= 550 else (C_ORANGE if b_max_v >= 500 else C_GREEN)
        ax_p.text(0.98, 0.04,
                  f'Avg: {b_avg_v:.0f}°C  Max: {b_max_v:.0f}°C\nΔ внутри банка: {intra_delta:.0f}°C',
                  transform=ax_p.transAxes, ha='right', va='bottom',
                  color=C_WHITE, fontsize=8.5, fontweight='bold',
                  bbox=dict(boxstyle='round,pad=0.4', facecolor=C_AXES,
                            edgecolor=box_c, alpha=0.95, linewidth=1.5))
        bank_stats[bname] = (b_avg_v, b_max_v, intra_delta)

    # ── Per-cylinder Δ bar chart (deviation from bank average) ────────────────
    ax_db.set_facecolor(C_AXES)
    ax_db.set_title('Δ каждого цилиндра от среднего по своему банку', color=C_WHITE,
                    fontsize=9.5, fontweight='bold')

    cyl_delta_val = {}
    for n in range(1, 17):
        if n not in cyl_egt: continue
        finite = cyl_egt[n][np.isfinite(cyl_egt[n])]
        if len(finite) < 5: continue
        bank = bank_a if n % 2 == 1 else bank_b
        # Compute this bank's average (only from cylinders with data)
        b_vals = [cyl_egt[c][np.isfinite(cyl_egt[c])].mean()
                  for c in bank if c in cyl_egt and np.any(np.isfinite(cyl_egt[c]))]
        if not b_vals: continue
        bank_mean = np.mean(b_vals)
        cyl_delta_val[n] = np.mean(finite) - bank_mean

    if cyl_delta_val:
        cyls_sorted = sorted(cyl_delta_val.keys())
        deltas = [cyl_delta_val[n] for n in cyls_sorted]
        bar_c = []
        for n, dv in zip(cyls_sorted, deltas):
            if dv > 30: bar_c.append(C_RED)
            elif dv > 10: bar_c.append(C_ORANGE)
            elif dv < -15: bar_c.append(C_GREEN)
            else: bar_c.append('#4EC9F0' if n % 2 == 1 else '#7090B0')

        x_pos = np.arange(len(cyls_sorted))
        bars = ax_db.bar(x_pos, deltas, color=bar_c, alpha=0.88, width=0.7, zorder=3)
        ax_db.axhline(0, color=C_WHITE, lw=1.2, alpha=0.6, zorder=4)
        ax_db.axhline(20,  color=C_ORANGE, lw=0.8, ls=':', alpha=0.6)
        ax_db.axhline(-20, color=C_GREEN,  lw=0.8, ls=':', alpha=0.6)
        ax_db.set_xticks(x_pos)
        ax_db.set_xticklabels([f'Ц{n}' for n in cyls_sorted], fontsize=8.5, color=C_WHITE)
        ax_db.set_ylabel('Δ°C от avg', color=C_WHITE, fontsize=8.5)
        ax_db.grid(True, alpha=0.3, axis='y', zorder=1)
        ax_db.tick_params(colors=C_WHITE, labelsize=8)
        for sp in ax_db.spines.values(): sp.set_edgecolor('#445055')

        # Annotate each bar with delta value
        for bar, d, n in zip(bars, deltas, cyls_sorted):
            offset = 1.5 if d >= 0 else -4.0
            va = 'bottom' if d >= 0 else 'top'
            fc = C_RED if abs(d) > 25 else C_WHITE
            ax_db.text(bar.get_x() + bar.get_width()/2, d + offset,
                       f'{d:+.0f}°', ha='center', va=va,
                       fontsize=7.5, color=fc, fontweight='bold')

        # Bank separator
        ax_db.axvline(7.5, color=C_GRAY, lw=1.2, ls='--', alpha=0.5)
        ax_db.text(3.5, ax_db.get_ylim()[1]*0.88, 'А-банк (нечётные)',
                   ha='center', color=colors_a[4], fontsize=8, fontweight='bold', alpha=0.9)
        ax_db.text(11.5, ax_db.get_ylim()[1]*0.88, 'В-банк (чётные)',
                   ha='center', color=colors_b[4], fontsize=8, fontweight='bold', alpha=0.9)
    else:
        ax_db.text(0.5, 0.5, 'Нет данных по цилиндрам', transform=ax_db.transAxes,
                   ha='center', color=C_GRAY, fontsize=11)

    # ── Inter-bank delta over time ─────────────────────────────────────────────
    ax_dt.set_facecolor(C_AXES)
    ax_dt.set_title('Межбанковая дельта: В-банк avg − А-банк avg (по времени)',
                    color=C_WHITE, fontsize=9.5, fontweight='bold')

    inter_delta = b_avg_ts - a_avg_ts
    inter_delta_s = smooth(inter_delta, w=9)

    ax_dt.fill_between(sec, inter_delta_s, 0,
                       where=inter_delta_s > 0, alpha=0.35, color=C_RED,
                       label='В-банк горячее')
    ax_dt.fill_between(sec, inter_delta_s, 0,
                       where=inter_delta_s < 0, alpha=0.35, color=C_GREEN,
                       label='А-банк горячее')
    ax_dt.plot(sec, inter_delta_s, color=C_ORANGE, lw=2.0, alpha=0.9,
               label='Δ (В−А)')
    ax_dt.axhline(0,   color=C_WHITE, lw=1.2, alpha=0.5, ls='--')
    ax_dt.axhline(20,  color=C_ORANGE, lw=0.8, ls=':',  alpha=0.6, label='+20°C')
    ax_dt.axhline(-20, color=C_GREEN,  lw=0.8, ls=':',  alpha=0.6, label='-20°C')
    ax_dt.set_xlim(0, t_max)
    ax_dt.set_ylabel('Δ (°C)', color=C_WHITE, fontsize=8.5)
    ax_dt.set_xlabel('Время (сек)', color=C_WHITE, fontsize=8.5)
    ax_dt.legend(fontsize=7.5, ncol=4, loc='upper right', framealpha=0.35,
                 facecolor=C_AXES, edgecolor='#445055', labelcolor=C_WHITE)
    ax_dt.grid(True, alpha=0.3)
    ax_dt.tick_params(colors=C_WHITE, labelsize=7.5)
    for sp in ax_dt.spines.values(): sp.set_edgecolor('#445055')

    # Annotate mean inter-bank delta
    finite_d = inter_delta_s[np.isfinite(inter_delta_s)]
    if len(finite_d) > 0:
        mean_d = np.mean(finite_d)
        ax_dt.text(0.99, 0.85, f'Ср. Δ(В−А) = {mean_d:+.0f}°C',
                   transform=ax_dt.transAxes, ha='right', va='top',
                   color=C_RED if mean_d > 15 else (C_GREEN if mean_d < -15 else C_WHITE),
                   fontsize=9, fontweight='bold',
                   bbox=dict(boxstyle='round,pad=0.35', facecolor=C_AXES,
                             edgecolor=C_ORANGE, alpha=0.9))

    fig.tight_layout(rect=[0, 0, 1, 0.97])

    fd, tmppath = tempfile.mkstemp(suffix='.png')
    os.close(fd)
    fig.savefig(tmppath, dpi=150, bbox_inches='tight', facecolor=C_BG, edgecolor='none')
    plt.close(fig)
    return tmppath, bank_stats

# ─── PPTX helpers ─────────────────────────────────────────────────────────────

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color(C_BG)

def add_rect(slide, x, y, w, h, fill=None, line=None, lw=0):
    sh = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = rgb_color(fill)
    else: sh.fill.background()
    if line and lw > 0:
        sh.line.color.rgb = rgb_color(line); sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    return sh

def add_text(slide, text, x, y, w, h, fs=12, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(fs); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = rgb_color(color)
    return tb

def build_conclusion(bank_stats, cyl_delta_val):
    a_n, b_n = 'А-банк (нечётные)', 'В-банк (чётные)'
    a_avg, a_max, a_intra = bank_stats.get(a_n, (0, 0, 0))
    b_avg, b_max, b_intra = bank_stats.get(b_n, (0, 0, 0))
    delta_ab = b_avg - a_avg

    # Hottest cylinder
    if cyl_delta_val:
        hot_cyl = max(cyl_delta_val, key=lambda n: cyl_delta_val[n])
        hot_d   = cyl_delta_val[hot_cyl]
    else:
        hot_cyl, hot_d = '?', 0

    if abs(delta_ab) > 10:
        asym = f'В-банк горячее А-банка на {abs(delta_ab):.0f}°C' if delta_ab > 0 \
               else f'А-банк горячее В-банка на {abs(delta_ab):.0f}°C'
    else:
        asym = 'Банки симметричны'

    intra = f'  |  Разброс: А={a_intra:.0f}°C  В={b_intra:.0f}°C'
    if cyl_delta_val:
        cyl_note = f'  |  Цил.{hot_cyl} отклонение +{hot_d:.0f}°C от avg банка.'
    else:
        cyl_note = ''

    return f'ВЫВОД: {asym}.{intra}.{cyl_note}'

def add_bank_slide(prs, uid, sec, cyl_egt, date):
    tmppath, bank_stats = make_delta_chart(uid, sec, cyl_egt, date)

    # Compute per-cyl deltas again for conclusion
    bank_a = [1, 3, 5, 7, 9, 11, 13, 15]
    bank_b = [2, 4, 6, 8, 10, 12, 14, 16]
    cyl_delta_val = {}
    for n in range(1, 17):
        if n not in cyl_egt: continue
        finite = cyl_egt[n][np.isfinite(cyl_egt[n])]
        if len(finite) < 5: continue
        bank = bank_a if n % 2 == 1 else bank_b
        b_vals = [cyl_egt[c][np.isfinite(cyl_egt[c])].mean()
                  for c in bank if c in cyl_egt and np.any(np.isfinite(cyl_egt[c]))]
        if b_vals:
            cyl_delta_val[n] = np.mean(finite) - np.mean(b_vals)

    conclusion = build_conclusion(bank_stats, cyl_delta_val)

    sl = blank_slide(prs)
    set_bg(sl)
    add_rect(sl, 0, 0, 33.86, 2.0, fill=C_DARK2)
    add_rect(sl, 0, 1.85, 33.86, 0.12, fill=C_GREEN)
    add_text(sl, f'ТЕМПЕРАТУРА ЦИЛИНДРОВ — NTE200 {UNIT_LABELS[uid]}  |  {date}',
             0.5, 0.1, 32, 1.0, fs=15, bold=True)
    add_text(sl, 'Δ внутри банков  |  Δ каждого цилиндра от среднего  |  Межбанковая асимметрия по времени',
             0.5, 1.15, 32, 0.7, fs=9, color=C_GRAY)

    sl.shapes.add_picture(tmppath, Cm(0.2), Cm(2.05), Cm(33.46), Cm(15.55))

    add_rect(sl, 0, 17.68, 33.86, 1.37, fill=C_DARK2, line=C_ORANGE, lw=1.2)
    add_text(sl, conclusion, 0.4, 17.73, 33.0, 1.27, fs=9.5, bold=True, color=C_ORANGE)

    os.remove(tmppath)
    print(f'  Слайд добавлен: {uid}')

# ─── Slide trimming (remove existing bank slides 13+) ─────────────────────────

def trim_slides(prs, keep=12):
    """Remove all slides after index keep-1 (0-based)."""
    from pptx.oxml.ns import qn
    from lxml import etree
    sl_list = prs.slides._sldIdLst
    all_ids = list(sl_list)
    for sld_id in all_ids[keep:]:
        sl_list.remove(sld_id)
    # Also remove from _slides
    # python-pptx internal: re-sync slide list
    # Workaround: rebuild from scratch is safer; instead just rebuild
    pass  # We'll handle by rebuilding the file from scratch

# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    # Rebuild from the 12-slide presentation by running build_log_analysis_presentation.py
    # first if needed. Here we open the existing file and rebuild only the bank slides
    # by creating a fresh presentation for slides 1-12, then appending bank slides.
    # Strategy: open existing, delete slides >12, append new bank slides.

    import subprocess
    print("Rebuilding base 12-slide NTE200 presentation...")
    result = subprocess.run(['python3', 'build_log_analysis_presentation.py'],
                            capture_output=True, text=True)
    if result.returncode != 0:
        print("ERROR rebuilding base presentation:")
        print(result.stderr[-2000:])
        return
    print("  Base presentation rebuilt (12 slides).")

    print(f"Opening {PPTX_IN}...")
    prs = Presentation(PPTX_IN)
    print(f"  Slides: {len(prs.slides)}")

    # Section divider
    sl = blank_slide(prs)
    set_bg(sl)
    add_rect(sl, 0, 6.5, 33.86, 6.0, fill=C_DARK2)
    add_rect(sl, 0, 6.5, 33.86, 0.3, fill=C_GREEN)
    add_text(sl, '09', 1.5, 7.5, 5, 3, fs=72, bold=True, color=C_GREEN)
    add_text(sl, 'ТЕМПЕРАТУРА', 7, 7.5, 25, 2.5, fs=44, bold=True)
    add_text(sl, 'ПО ЦИЛИНДРАМ', 7, 9.7, 25, 2.0, fs=44, bold=True, color=C_GREEN)
    add_text(sl, 'А-банк (ECM144) vs В-банк (ECM1)  |  Δ цилиндров  |  9 агрегатов  |  май 2026',
             7, 12.2, 25, 1.0, fs=12, color=C_GRAY)
    print("  Section divider added.")

    for uid in ['43', '45', '50', '52', '56', '84a', '84b', '85', '87']:
        print(f"  Processing unit {uid}...")
        try:
            sec, cyl_egt, date = load_unit(uid)
            if not cyl_egt:
                print(f"    WARN: no cyl data for {uid}")
                continue
            add_bank_slide(prs, uid, sec, cyl_egt, date)
        except Exception as e:
            import traceback
            print(f"    ERROR {uid}: {e}")
            traceback.print_exc()

    print(f"\nTotal slides: {len(prs.slides)}")
    prs.save(PPTX_OUT)
    print(f"Saved: {PPTX_OUT}")

if __name__ == '__main__':
    main()
